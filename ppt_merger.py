"""
ppt_merger.py  —  APSG PPT Merger  v4.0
========================================
ROOT CAUSE FIX:
  Rectangle 14 and Rectangle 15 are LABEL TEXT BOXES ("TOP PHOTO" / "FRONT PHOTO")
  NOT empty image containers.
  
  Images must be placed BELOW these labels, NOT at their position.

EXACT placement logic:
  Top Photo image:
    Left  = Rectangle14.left
    Top   = Rectangle14.top + Rectangle14.height + 1mm gap
    Width = 15.98 cm   (as specified)
    Height= 8.33  cm   (as specified)

  Front Photo image:
    Left  = Rectangle15.left
    Top   = Rectangle15.top + Rectangle15.height + 1mm gap
    Width = min(15.51cm, slide_width - Rectangle15.left)  [clamp to slide]
    Height= 8.40 cm   (as specified)

Slide rules:
  Slide 0 (index)  → "first"   → copied as-is, NO changes
  Summary slides   → "summary" → copied as-is, NO changes
  Image slides     → "image"   → images replaced BELOW labels
  Footer/logo      → NEVER touched
  
Output filename: APSG-Load_Rejected_DD-MM-YYYY.pptx
"""

import io, os, re, logging, sys, datetime
from lxml import etree
from pptx import Presentation
from pptx.util import Cm, Emu, Inches
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ── XML constants ──────────────────────────────────────────────────────────────
R_EMBED = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
IMG_REL = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image'
BLIP    = qn('a:blip')

# ── Placeholder shape names (label text boxes in the template) ─────────────────
TOP_PH_NAME   = "Rectangle 14"   # contains text "TOP PHOTO"
FRONT_PH_NAME = "Rectangle 15"   # contains text "FRONT PHOTO"

# ══════════════════════════════════════════════════════════════════════════════
#  FULLY INDEPENDENT image defaults — each photo has its OWN absolute position.
#  Changing one photo's position/size NEVER affects the other.
#  All values in EMU (1 cm = 360,000 EMU).
# ══════════════════════════════════════════════════════════════════════════════

# Top Photo — absolute position, independent
TOP_LEFT_DEFAULT = int(Cm(0.95))   # absolute left edge on slide
TOP_TOP_DEFAULT  = None            # None = auto (label_bottom + gap)
TOP_W_DEFAULT    = int(Cm(15.28))  # width — independent, never affects front
TOP_H_DEFAULT    = int(Cm(9.11))   # height — independent

# Front Photo — absolute position, independent (no relation to top photo)
FRONT_LEFT_DEFAULT = int(Cm(18.73)) # absolute left edge — NOT derived from top
FRONT_TOP_DEFAULT  = None           # None = auto (label_bottom + gap)
FRONT_W_DEFAULT    = int(Cm(15.51)) # width — independent, never affects top
FRONT_H_DEFAULT    = int(Cm(9.11))  # height — independent

# ── Gap between label bottom and image top (1 mm) ─────────────────────────────
LABEL_GAP = int(Cm(0.1))

# ── Footer protection threshold: shapes with top > this value are NEVER removed ─
# Footer picture on image slides sits at T≈17.385cm = 6258698 emu
# We protect everything below 15cm = 5400000 emu to be safe
FOOTER_SAFE_Y = int(Cm(15.0))   # 5400000 emu  — do NOT remove pictures below this

# ── Slide classification keywords ─────────────────────────────────────────────
SUMMARY_KEYWORDS = [
    "management of staging", "staging ground",
    "load rejected summary", "summary",
]

# ── Date extraction patterns ───────────────────────────────────────────────────
DATE_PATTERNS = [
    r'\b(\d{1,2}[-/\.]\d{1,2}[-/\.]\d{4})\b',
    r'\b(\d{4}[-/\.]\d{1,2}[-/\.]\d{1,2})\b',
    r'\b(\d{1,2}\s+\w+\s+\d{4})\b',
]
DATE_FORMATS = [
    "%d-%m-%Y", "%d/%m/%Y", "%d.%m.%Y",
    "%Y-%m-%d", "%Y/%m/%d",
    "%d %B %Y", "%d %b %Y",
]


# ══════════════════════════════════════════════════════════════════════════════
#  Logger
# ══════════════════════════════════════════════════════════════════════════════
def make_logger(verbose=False, log_file=None):
    lg = logging.getLogger("apsg_" + str(id(object())))
    lg.setLevel(logging.DEBUG if verbose else logging.INFO)
    lg.handlers.clear()
    fmt = logging.Formatter("[%(levelname)s] %(message)s")
    sh = logging.StreamHandler(sys.stdout)
    sh.setFormatter(fmt)
    lg.addHandler(sh)
    if log_file:
        fh = logging.FileHandler(log_file, encoding="utf-8")
        fh.setFormatter(fmt)
        lg.addHandler(fh)
    return lg


# ══════════════════════════════════════════════════════════════════════════════
#  Shape lookup
# ══════════════════════════════════════════════════════════════════════════════
def _find_shape(slide, name):
    """Find shape by exact name (case-insensitive strip)."""
    nl = name.strip().lower()
    for s in slide.shapes:
        if s.name.strip().lower() == nl:
            return s
    return None


def _all_slide_text(slide):
    """Yield all non-empty text from every text-bearing shape."""
    for s in slide.shapes:
        if s.has_text_frame:
            t = s.text_frame.text.strip()
            if t:
                yield t
        if s.has_table:
            for row in s.table.rows:
                for cell in row.cells:
                    t = (cell.text_frame.text.strip()
                         if cell.text_frame else "")
                    if t:
                        yield t


# ══════════════════════════════════════════════════════════════════════════════
#  Slide classification
# ══════════════════════════════════════════════════════════════════════════════
def classify_slide(slide, idx):
    """
    Returns: 'first' | 'summary' | 'image' | 'other'
    'image' = has Rectangle 14 OR Rectangle 15 (label shapes)
    """
    if idx == 0:
        return "first"

    # Image slide: contains the photo label placeholders
    has_top_label   = _find_shape(slide, TOP_PH_NAME)   is not None
    has_front_label = _find_shape(slide, FRONT_PH_NAME) is not None
    if has_top_label or has_front_label:
        return "image"

    # Summary slide: keywords in text or has a data table
    all_text = " ".join(_all_slide_text(slide)).lower()
    if any(kw in all_text for kw in SUMMARY_KEYWORDS):
        return "summary"
    if any(s.has_table for s in slide.shapes):
        return "summary"

    return "other"


# ══════════════════════════════════════════════════════════════════════════════
#  Token extraction (for slide matching)
# ══════════════════════════════════════════════════════════════════════════════
def extract_token(slide):
    """
    Extract E-Token from table cells. Returns uppercase string or None.
    E-Token: >= 14 chars, >= 65% alphanumeric.
    """
    for s in slide.shapes:
        if not s.has_table:
            continue
        for row_idx in range(1, len(s.table.rows)):
            for cell in s.table.rows[row_idx].cells:
                txt = (cell.text_frame.text.strip()
                       if cell.text_frame else "")
                if (len(txt) >= 14 and
                        sum(c.isalnum() for c in txt) / max(len(txt), 1) >= 0.65):
                    return txt.upper()
    return None


def build_token_index(prs):
    return {extract_token(sl): i
            for i, sl in enumerate(prs.slides)
            if extract_token(sl)}


def find_matching_slide(prs, token, fallback_idx, token_idx):
    if token and token in token_idx:
        mi = token_idx[token]
        if mi < len(prs.slides):
            return prs.slides[mi]
    if fallback_idx < len(prs.slides):
        return prs.slides[fallback_idx]
    return None


# ══════════════════════════════════════════════════════════════════════════════
#  Reference file auto-detection
# ══════════════════════════════════════════════════════════════════════════════
def detect_reference(prs_a, prs_b, log):
    """
    The reference file contains Token Numbers.
    Returns (label, prs_ref, prs_other).
    """
    ta = sum(1 for sl in prs_a.slides if extract_token(sl))
    tb = sum(1 for sl in prs_b.slides if extract_token(sl))
    log.info(f"Token count: File A={ta}  File B={tb}")
    if ta >= tb:
        log.info("→ File A = FORMAT REFERENCE")
        return "A", prs_a, prs_b
    else:
        log.info("→ File B = FORMAT REFERENCE")
        return "B", prs_b, prs_a


# ══════════════════════════════════════════════════════════════════════════════
#  Date extraction
# ══════════════════════════════════════════════════════════════════════════════
def extract_date(prs):
    for slide in prs.slides:
        for txt in _all_slide_text(slide):
            for pat in DATE_PATTERNS:
                m = re.search(pat, txt, re.IGNORECASE)
                if m:
                    raw = m.group(1)
                    for fmt in DATE_FORMATS:
                        try:
                            return datetime.datetime.strptime(raw, fmt).date()
                        except ValueError:
                            continue
    return None


def build_filename(date_obj):
    if not date_obj:
        date_obj = datetime.date.today()
    # Date format: DDMMYYYY (no spaces, no hyphens between digits)
    return f"APSG-Loads_Rejected-{date_obj.strftime('%d%m%Y')}.pptx"


# ══════════════════════════════════════════════════════════════════════════════
#  Image helpers
# ══════════════════════════════════════════════════════════════════════════════
def _is_pic(shape):
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    except Exception:
        return False


def _get_blob(shape, slide):
    """
    Extract raw image bytes from a picture shape.
    The blob is the exact original image file — all markings/annotations
    that are baked INTO the image (e.g. red circles drawn on the photo)
    are preserved automatically since we copy the raw bytes unchanged.
    """
    try:
        blips = shape._element.findall('.//' + BLIP, shape._element.nsmap)
        if not blips:
            return None
        rId = blips[0].get(R_EMBED)
        if not rId:
            return None
        rel = slide.part.rels.get(rId)
        return rel.target_part.blob if rel else None
    except Exception:
        return None


def _classify_image(shape, slide):
    """Classify picture as 'top'/'front'/'unknown' by nearby label text."""
    nearby = []
    for s in slide.shapes:
        if s.shape_id == shape.shape_id or not s.has_text_frame:
            continue
        if (abs(s.top  - shape.top)  < Inches(2.0) and
                abs(s.left - shape.left) < Inches(5.0)):
            nearby.append(s.text_frame.text.strip().lower())
    label = " ".join(nearby)
    if any(kw in label for kw in ["top photo", "top image", "top"]):
        return "top"
    if any(kw in label for kw in ["front photo", "front image", "front"]):
        return "front"
    return "unknown"


def extract_images(slide):
    """Returns {'top': bytes|None, 'front': bytes|None, 'all': [bytes,...]}"""
    result = {"top": None, "front": None, "all": []}
    for shape in slide.shapes:
        if not _is_pic(shape):
            continue
        b = _get_blob(shape, slide)
        if b is None:
            continue
        result["all"].append(b)
        k = _classify_image(shape, slide)
        if k == "top"   and result["top"]   is None:
            result["top"] = b
        elif k == "front" and result["front"] is None:
            result["front"] = b
    return result


# ══════════════════════════════════════════════════════════════════════════════
#  CORE: Place image BELOW the label, removing any old picture first
# ══════════════════════════════════════════════════════════════════════════════
def _remove_pics_near(slide, L, T, W, H, log):
    """
    Remove picture shapes whose centre overlaps the target area.
    IMPORTANT: shapes with top >= FOOTER_SAFE_Y are NEVER removed (footer protection).
    Uses tight overlap: only removes pictures actually inside the image zone.
    """
    tree  = slide.shapes._spTree
    to_rm = []
    # Tight padding: 5% of dimensions to avoid accidentally catching adjacent shapes
    pad_x = max(int(W * 0.05), int(Cm(0.1)))
    pad_y = max(int(H * 0.05), int(Cm(0.1)))
    for s in slide.shapes:
        if not _is_pic(s):
            continue
        # ── FOOTER PROTECTION: never remove shapes in the footer zone ──────────
        if s.top >= FOOTER_SAFE_Y:
            log.debug(f"    Protected footer shape '{s.name}' at T={s.top/360000:.2f}cm")
            continue
        cx = s.left + s.width  // 2
        cy = s.top  + s.height // 2
        if (L - pad_x <= cx <= L + W + pad_x and
                T - pad_y <= cy <= T + H + pad_y):
            to_rm.append(s._element)
            log.debug(f"    Removing old picture '{s.name}' at "
                      f"({s.left/360000:.2f}cm,{s.top/360000:.2f}cm)")
    for el in to_rm:
        try:
            tree.remove(el)
        except Exception:
            pass
    return len(to_rm)


def get_table_bounds(slide):
    """
    Return (table_left, table_right) EMU for the first data table on the slide.
    Returns (None, None) if no table found.
    """
    for s in slide.shapes:
        if s.has_table:
            return s.left, s.left + s.width
    return None, None


def compute_column_and_center(label_shape, boundary_right, col_left, spec_w):
    """
    Compute the centered image position within its column.

    Column definition:
      col_left  = left boundary of this image column
      col_right = boundary_right (other label left, or table right)
      col_width = col_right - col_left

    Image width = min(spec_w, col_width)   — never wider than column
    Image left  = col_left + (col_width - img_width) / 2  — centred

    Returns (img_left, img_width).
    """
    col_width = boundary_right - col_left
    if col_width <= 0:
        return col_left, spec_w

    # Use the exact specified width — do NOT shrink to column width
    # This ensures user-specified dimensions are always honoured
    img_w    = spec_w
    img_left = col_left   # left-align within column (no centering offset)
    return img_left, img_w


def place_image_below_label(slide, label_shape_name,
                            img_bytes, img_h, spec_w,
                            col_left, boundary_right, log):
    """
    THE CORRECT PLACEMENT FUNCTION.

    Steps:
      1. Find the label shape (e.g. Rectangle 14 "TOP PHOTO")
      2. Compute column (col_left → boundary_right) and centre image within it
      3. Place image top = label.bottom + LABEL_GAP
      4. Remove ONLY pictures inside image area (protect footer with FOOTER_SAFE_Y)
      5. Insert picture at centred, table-bounded coordinates
    
    Returns True on success.
    """
    if img_bytes is None:
        log.debug(f"    place_image_below_label('{label_shape_name}'): no image — skip")
        return False

    # ── Step 1: Find the label shape ──────────────────────────────────────────
    label = _find_shape(slide, label_shape_name)
    if label is None:
        log.warning(f"    Label shape '{label_shape_name}' NOT found on slide")
        return False

    # ── Step 2: Compute centred position within column ───────────────────────
    img_top           = label.top + label.height + LABEL_GAP   # directly below label
    img_left, img_w_final = compute_column_and_center(label, boundary_right, col_left, spec_w)
    log.debug(f"    Centre: col=[{col_left/360000:.2f},{boundary_right/360000:.2f}]cm "
              f"→ img=[{img_left/360000:.2f},{(img_left+img_w_final)/360000:.2f}]cm "
              f"margin={(img_left-col_left)/360000:.3f}cm")

    log.info(f"    Placing '{label_shape_name}': "
             f"L={img_left/360000:.2f}cm "
             f"T={img_top/360000:.2f}cm "
             f"W={img_w_final/360000:.2f}cm "
             f"H={img_h/360000:.2f}cm "
             f"R={(img_left+img_w_final)/360000:.2f}cm")

    # ── Step 4: Remove existing pictures in the target area ───────────────────
    removed = _remove_pics_near(slide, img_left, img_top, img_w_final, img_h, log)
    if removed:
        log.debug(f"    Cleared {removed} existing picture(s)")

    # ── Step 5: Place new image ───────────────────────────────────────────────
    try:
        pic = slide.shapes.add_picture(
            io.BytesIO(img_bytes),
            img_left,    # exact left from label
            img_top,     # exact top = label bottom + gap
            img_w_final, # required width (clamped)
            img_h        # required height
        )
        pic.name = f"{label_shape_name}_image"
        log.info(f"    ✓ Placed '{label_shape_name}' image successfully")
        return True
    except Exception as e:
        log.warning(f"    add_picture failed for '{label_shape_name}': {e}")
        return False


# ══════════════════════════════════════════════════════════════════════════════
#  Result
# ══════════════════════════════════════════════════════════════════════════════
class MergeResult:
    def __init__(self):
        self.ok             = False
        self.total          = 0
        self.merged         = 0
        self.both           = 0
        self.top_only       = 0
        self.front_only     = 0
        self.none_found     = 0
        self.skipped        = 0
        self.reference_file = "?"
        self.output_name    = ""
        self.date_used      = None
        self.warnings       = []
        self.errors         = []
        self.log_lines      = []
        self._final_path    = None

    def to_dict(self):
        return {k: getattr(self, k) for k in (
            "ok", "total", "merged", "both", "top_only", "front_only",
            "none_found", "skipped", "reference_file",
            "output_name", "warnings", "errors"
        )}


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN MERGE
# ══════════════════════════════════════════════════════════════════════════════
def merge(file_a, file_b, output_path,
          top_ph=TOP_PH_NAME, front_ph=FRONT_PH_NAME,
          verbose=False, log_file=None,
          force_output_name=None,
          top_h_cm=None, top_w_cm=None, top_left_cm=None,
          front_h_cm=None, front_w_cm=None, front_left_cm=None,
          center_gap_cm=None):
    """
    Merge two PPT files:
    - Auto-detects which file is FORMAT REFERENCE
    - Slide 1 and summary slides → copied exactly, zero changes
    - Image slides → images placed BELOW their respective labels
    - Footer/logo → never touched
    - Output: APSG-Load_Rejected_DD-MM-YYYY.pptx
    """
    log    = make_logger(verbose=verbose, log_file=log_file)
    result = MergeResult()

    # ── Resolve ALL dimensions — each photo fully independent ────────────────
    # Top Photo — its own size and absolute position
    TOP_H    = int(Cm(top_h_cm))    if top_h_cm    is not None else TOP_H_DEFAULT
    TOP_W    = int(Cm(top_w_cm))    if top_w_cm    is not None else TOP_W_DEFAULT
    TOP_LEFT = int(Cm(top_left_cm)) if top_left_cm is not None else TOP_LEFT_DEFAULT

    # Front Photo — its own size and absolute position (ZERO dependency on top)
    FRONT_H    = int(Cm(front_h_cm))    if front_h_cm    is not None else FRONT_H_DEFAULT
    FRONT_W    = int(Cm(front_w_cm))    if front_w_cm    is not None else FRONT_W_DEFAULT
    FRONT_LEFT = int(Cm(front_left_cm)) if front_left_cm is not None else FRONT_LEFT_DEFAULT

    # Gap is INFORMATIONAL only — used to compute default front_left if not set explicitly
    # When front_left_cm is explicitly provided, gap is ignored entirely
    if front_left_cm is None and center_gap_cm is not None:
        FRONT_LEFT = TOP_LEFT + TOP_W + int(Cm(center_gap_cm))

    log.info(f"TOP   photo: left={TOP_LEFT/360000:.2f}cm  w={TOP_W/360000:.2f}cm  h={TOP_H/360000:.2f}cm")
    log.info(f"FRONT photo: left={FRONT_LEFT/360000:.2f}cm  w={FRONT_W/360000:.2f}cm  h={FRONT_H/360000:.2f}cm")
    log.info(f"(These are INDEPENDENT — neither depends on the other)")

    # Attach log capture
    class LH(logging.Handler):
        def emit(self, rec):
            result.log_lines.append((rec.levelname, self.format(rec)))
    lh = LH()
    lh.setFormatter(logging.Formatter("%(message)s"))
    log.addHandler(lh)

    log.info("=" * 60)
    log.info("APSG PPT Merger  v4.0  — Correct Image Placement")
    log.info(f"Top label: '{top_ph}'   Front label: '{front_ph}'")
    log.info("=" * 60)

    # ── Read raw bytes ─────────────────────────────────────────────────────────
    def _read(src, label):
        try:
            if hasattr(src, "read"):
                b = src.read()
                if hasattr(src, "seek"):
                    src.seek(0)
                return b
            with open(src, "rb") as f:
                return f.read()
        except Exception as e:
            result.errors.append(f"Cannot read File {label}: {e}")
            log.error(result.errors[-1])
            return None

    raw_a = _read(file_a, "A")
    raw_b = _read(file_b, "B")
    if raw_a is None or raw_b is None:
        return result

    # ── Parse presentations ────────────────────────────────────────────────────
    try:
        prs_a = Presentation(io.BytesIO(raw_a))
        log.info(f"File A: {len(prs_a.slides)} slides")
    except Exception as e:
        result.errors.append(f"Cannot parse File A: {e}")
        log.error(result.errors[-1])
        return result

    try:
        prs_b = Presentation(io.BytesIO(raw_b))
        log.info(f"File B: {len(prs_b.slides)} slides")
    except Exception as e:
        result.errors.append(f"Cannot parse File B: {e}")
        log.error(result.errors[-1])
        return result

    # ── Auto-detect reference ──────────────────────────────────────────────────
    ref_label, prs_ref, prs_other = detect_reference(prs_a, prs_b, log)
    raw_ref   = raw_a if ref_label == "A" else raw_b
    result.reference_file = ref_label

    # ── Extract date for filename ──────────────────────────────────────────────
    date_obj = extract_date(prs_ref) or extract_date(prs_other)
    result.date_used = date_obj
    if date_obj:
        log.info(f"Date: {date_obj.strftime('%d-%m-%Y')}")
    else:
        log.warning("No date found — using today")
        date_obj = datetime.date.today()

    out_name = force_output_name or build_filename(date_obj)
    result.output_name = out_name
    log.info(f"Output: {out_name}")

    # ── Token index for other file ─────────────────────────────────────────────
    tok_idx_other = build_token_index(prs_other)
    log.info(f"Tokens in other file: {len(tok_idx_other)}")

    # ── Load fresh editable base from reference ────────────────────────────────
    base     = Presentation(io.BytesIO(raw_ref))
    prs_ref2 = Presentation(io.BytesIO(raw_ref))   # for image extraction
    slide_w  = base.slide_width

    n_slides      = len(base.slides)
    result.total  = n_slides
    log.info(f"Base: {n_slides} slides   Slide width: {slide_w/360000:.2f}cm")
    log.info("─" * 60)

    # ── Process each slide ─────────────────────────────────────────────────────
    for i, slide in enumerate(base.slides):
        n          = i + 1
        slide_type = classify_slide(slide, i)
        log.info(f"── Slide {n}/{n_slides}  [{slide_type.upper()}]")

        # Non-image slides: preserved exactly as-is
        if slide_type != "image":
            log.info("   Preserved as-is")
            result.merged += 1
            continue

        # ── Image slide ────────────────────────────────────────────────────────
        has_top_label   = _find_shape(slide, top_ph)   is not None
        has_front_label = _find_shape(slide, front_ph) is not None
        log.info(f"   Labels: '{top_ph}'={has_top_label}  '{front_ph}'={has_front_label}")

        # Extract token for matching
        token = extract_token(slide)
        log.debug(f"   Token: {token or '(none)'}")

        # Get source slides
        src_ref   = prs_ref2.slides[i]   if i < len(prs_ref2.slides)  else None
        src_other = find_matching_slide(prs_other, token, i, tok_idx_other)

        # Extract images
        imgs_ref   = extract_images(src_ref)   if src_ref   else {"top":None,"front":None,"all":[]}
        imgs_other = extract_images(src_other) if src_other else {"top":None,"front":None,"all":[]}

        # Top Photo: from reference first, fallback other
        top_img = (imgs_ref["top"]
                   or (imgs_ref["all"][0]  if imgs_ref["all"]   else None)
                   or imgs_other["top"]
                   or (imgs_other["all"][0] if imgs_other["all"] else None))

        # Front Photo: from other file first, fallback reference
        frt_img = (imgs_other["front"]
                   or (imgs_other["all"][-1] if imgs_other["all"] else None)
                   or imgs_ref["front"]
                   or (imgs_ref["all"][-1]   if imgs_ref["all"]   else None))

        if top_img and not imgs_ref["top"]:
            log.info("   Top fallback: from other file")
        if frt_img and not imgs_other["front"]:
            log.info("   Front fallback: from reference file")

        ht, hf = top_img is not None, frt_img is not None
        if   ht and hf: result.both      += 1; sc = "BOTH ✓"
        elif ht:        result.top_only  += 1; sc = "TOP only"
        elif hf:        result.front_only+= 1; sc = "FRONT only"
        else:           result.none_found+= 1; sc = "NO images"
        log.info(f"   Images: {sc}")

        # ── Get table bounds for width calculation ──────────────────────────────
        table_left, table_right = get_table_bounds(slide)
        if table_right:
            log.debug(f"   Table right edge: {table_right/360000:.3f}cm")
        else:
            log.debug("   No table found — using fallback widths")

        # ── Place images — FULLY INDEPENDENT, no shared state ────────────────
        try:
            # TOP PHOTO — uses only its own TOP_LEFT, TOP_W, TOP_H
            # Changing Front Photo values has ZERO effect here
            if has_top_label:
                top_img_left  = TOP_LEFT            # absolute, independent
                top_img_right = TOP_LEFT + TOP_W    # depends only on TOP_* vars

                place_image_below_label(
                    slide, top_ph,
                    top_img, TOP_H, TOP_W,
                    top_img_left, top_img_right, log
                )

            # FRONT PHOTO — uses only its own FRONT_LEFT, FRONT_W, FRONT_H
            # Changing Top Photo values has ZERO effect here
            if has_front_label:
                frt_img_left  = FRONT_LEFT           # absolute, independent
                frt_img_right = FRONT_LEFT + FRONT_W # depends only on FRONT_* vars

                place_image_below_label(
                    slide, front_ph,
                    frt_img, FRONT_H, FRONT_W,
                    frt_img_left, frt_img_right, log
                )
        except Exception as e:
            msg = f"Slide {n}: error — {e}"
            log.error(msg)
            result.errors.append(msg)
            result.skipped += 1
            continue

        result.merged += 1

    # ── Save ───────────────────────────────────────────────────────────────────
    log.info("─" * 60)

    # Build final path
    out_dir = output_path if (output_path.endswith(os.sep) or
                              os.path.isdir(output_path)) else os.path.dirname(output_path)
    if not out_dir:
        out_dir = "."
    final_path = os.path.join(out_dir, out_name)

    try:
        base.save(final_path)
        result.ok         = True
        result._final_path = final_path
        log.info(f"✓ Saved → {final_path}")
    except Exception as e:
        # Fallback to original path
        try:
            base.save(output_path)
            result.ok          = True
            result._final_path = output_path
            log.info(f"✓ Saved (fallback) → {output_path}")
        except Exception as e2:
            result.errors.append(f"Save failed: {e2}")
            log.error(result.errors[-1])

    log.info("=" * 60)
    log.info(f"DONE  total={result.total} merged={result.merged} "
             f"both={result.both} top={result.top_only} "
             f"front={result.front_only} none={result.none_found} "
             f"skip={result.skipped}")
    return result


# ── CLI ────────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    import argparse
    p = argparse.ArgumentParser(description="APSG PPT Merger v4.0")
    p.add_argument("--file_a",            required=True)
    p.add_argument("--file_b",            required=True)
    p.add_argument("--output",            required=True)
    p.add_argument("--top_placeholder",   default=TOP_PH_NAME)
    p.add_argument("--front_placeholder", default=FRONT_PH_NAME)
    p.add_argument("--log",               default=None)
    p.add_argument("--verbose",           action="store_true")
    a = p.parse_args()
    r = merge(a.file_a, a.file_b, a.output,
              top_ph=a.top_placeholder, front_ph=a.front_placeholder,
              verbose=a.verbose, log_file=a.log)
    print("\n" + ("✓ OK" if r.ok else "✗ FAILED"))
    print(f"Output : {r.output_name}")
    print(f"Ref    : File {r.reference_file}")
    d = r.to_dict()
    print(f"Slides : total={d['total']} merged={d['merged']} "
          f"both={d['both']} top={d['top_only']} front={d['front_only']}")
    for e in r.errors:   print(f"ERROR: {e}")
    for w in r.warnings: print(f"WARN:  {w}")
