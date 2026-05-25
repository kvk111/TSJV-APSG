"""
APSG (Staging Ground) Report — Unified Application
Combined: Daily Report | Excel Rejection | PPT Alignment | Monthly Rejection Filter
"""

import gc
import matplotlib
matplotlib.use("Agg")   # set before any other matplotlib import — must be first
import os, io, uuid, threading, time, json, copy, re, traceback, hashlib, logging
import random
import pytz
import datetime as _dt
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from datetime import datetime, date
from collections import OrderedDict
from functools import wraps

from flask import (Flask, request, jsonify, send_file, make_response,
                   session, redirect, url_for, render_template_string)
from werkzeug.utils import secure_filename

# ── Python-PPTX (used by PPT Rejection + Photo Merge) ──────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

# ── Photo Merge engine ──────────────────────────────────────────────────────
import ppt_merger

# ── Daily Report engines ──────────────────────────────────────────────────────
from report_engine import (
    load_and_validate_file, generate_report, get_date_range,
    filter_preview, validate_and_flag, validate_net_weights, validate_etoken,
)
from wb_engine import (
    load_wb_file, wb_get_date_range, wb_filter_by_date,
    wb_apply_unified_logic, wb_build_pivot, wb_pivot_to_excel,
    wb_pivot_copy_text, wb_find_incomplete_rows, wb_apply_row_decisions,
    wb_net_weight_validation, wb_validate_etoken_match,
    _normalize_accepted as _wb_norm_acc,
)
from rectification_report import (
    build_rr_docx, fetch_row_by_token, fetch_row_from_series,
    apply_user_updates, fetch_and_generate, generate_table_image,
    _parse_dt, _MONTH_B,
    REASON_A, REASON_B, REASON_C,
)
import pandas as pd
import math

# ── Cycle Time Report — integrated module (Blueprint) ─────────────────────────
_CT_AVAILABLE = False
_ct_err_detail = ''
try:
    from ct_module import ct as ct_blueprint
    _CT_AVAILABLE = True
    print("✓ CT module imported successfully")
except Exception as _ct_err:
    import traceback as _ct_tb
    _ct_err_detail = _ct_tb.format_exc()
    # Print full traceback so it appears in Render logs
    print(f"[CT Module] IMPORT FAILED — blueprint will NOT be registered")
    print(_ct_err_detail)

# ── DB / state ──────────────────────────────────────────────────────────────
import sqlite3

# ── JSON-safe numeric helpers ────────────────────────────────────────────────
def safe_float(v, default=0.0):
    """Convert v to float, returning default for NaN/None/inf/NA."""
    try:
        result = float(v)
        if math.isnan(result) or math.isinf(result):
            return default
        return result
    except (TypeError, ValueError):
        return default

def safe_int(v, default=0):
    """Convert v to int, returning default for NaN/None/NA."""
    try:
        f = float(v)
        if math.isnan(f) or math.isinf(f):
            return default
        return int(f)
    except (TypeError, ValueError):
        return default

def safe_col_sum(df, col, default=0.0):
    """Sum a DataFrame column safely, returning default if col missing or all NaN."""
    if col not in df.columns:
        return default
    return safe_float(df[col].dropna().sum(), default)

app = Flask(__name__, static_folder='static', static_url_path='/static')
app.secret_key = os.environ.get("SECRET_KEY", "apsg-report-secret-2024")
# Security: session cookie flags
app.config.update(
    SESSION_COOKIE_HTTPONLY  = True,
    SESSION_COOKIE_SAMESITE  = 'Lax',
    SESSION_COOKIE_SECURE    = os.environ.get('RENDER') == 'true',
    PERMANENT_SESSION_LIFETIME = __import__('datetime').timedelta(hours=8),
)

# ── Global error handlers — always return JSON, never HTML ────────────────────
@app.errorhandler(400)
def err_400(e): return jsonify({"ok":False,"error":str(e)}), 400
@app.errorhandler(401)
def err_401(e): return jsonify({"ok":False,"error":"Not authenticated"}), 401
@app.errorhandler(403)
def err_403(e): return jsonify({"ok":False,"error":"Forbidden"}), 403
@app.errorhandler(404)
def err_404(e): return jsonify({"ok":False,"error":"Not found"}), 404
@app.errorhandler(405)
def err_405(e): return jsonify({"ok":False,"error":"Method not allowed"}), 405
@app.errorhandler(500)
def err_500(e): return jsonify({"ok":False,"error":"Server error","detail":str(e)}), 500
@app.errorhandler(Exception)
def err_any(e):
    import traceback as _tb
    print(f"[Unhandled] {e}\n{_tb.format_exc()}")
    return jsonify({"ok":False,"error":str(e)}), 500

# ── Register Cycle Time Blueprint ──────────────────────────────────────────────
if _CT_AVAILABLE:
    app.register_blueprint(ct_blueprint)
    print("✓ Cycle Time module registered at /ct/")

# ── NaN-safe JSON serialization ──────────────────────────────────────────────
# Flask's default encoder passes NaN to json.dumps which produces invalid JSON.
# Override to replace NaN/Inf with null so browsers never see a parse error.
import json as _json
class _NanSafeProvider(app.json_provider_class):
    def dumps(self, obj, **kw):
        return _json.dumps(obj, **kw, allow_nan=False,
                           default=lambda o: None if (isinstance(o, float) and (math.isnan(o) or math.isinf(o))) else str(o))
    def loads(self, s, **kw):
        return _json.loads(s, **kw)
app.json_provider_class = _NanSafeProvider
app.json = _NanSafeProvider(app)

_BASE = os.path.dirname(os.path.abspath(__file__))
# Fly.io: use /app/data (persistent volume mount) if available, else local
_DATA = os.environ.get("DATA_DIR", os.path.join(_BASE, "data") if os.path.isdir(os.path.join(_BASE, "data")) else _BASE)
DB_PATH = os.environ.get("DB_PATH", os.path.join(_DATA, "apsg_report.db"))
UPLOAD_DIR = os.environ.get("UPLOAD_DIR", os.path.join(_DATA, "uploads"))
OUTPUT_DIR = os.environ.get("OUTPUT_DIR", os.path.join(_DATA, "outputs"))
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(OUTPUT_DIR, exist_ok=True)

def _cleanup_old_tempfiles():
    """Remove temp output files older than 24 hours to prevent disk accumulation."""
    import glob, time as _time
    cutoff = _time.time() - 86400   # 24 hours
    for pattern in ['*.pkl', '*.pptx', '*.xlsx']:
        for f in glob.glob(os.path.join(OUTPUT_DIR, pattern)):
            try:
                if os.path.getmtime(f) < cutoff:
                    os.remove(f)
            except OSError:
                pass  # already deleted or permission issue — ignore

# Run cleanup in background thread on startup (non-blocking)
import threading as _threading
_threading.Thread(target=_cleanup_old_tempfiles, daemon=True).start()

# ── PPT Rejection in-memory DB ──────────────────────────────────────────────
PPT_DB: dict = {"records": [], "slide_map": {}, "presentations": []}
jobs: dict = {}
jobs_lock = threading.Lock()

# ═══════════════════════════════════════════════════════════════════════════════
#  DATABASE SETUP
# ═══════════════════════════════════════════════════════════════════════════════

def get_db():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn

def init_db():
    """Create tables and seed default admin. Safe to call multiple times."""
    conn = get_db()
    c = conn.cursor()
    c.execute("""CREATE TABLE IF NOT EXISTS users (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        username TEXT UNIQUE NOT NULL,
        name TEXT NOT NULL,
        password_hash TEXT NOT NULL,
        plaintext_pw TEXT DEFAULT '',
        role TEXT DEFAULT 'user',
        email TEXT DEFAULT '',
        email_verified INTEGER DEFAULT 0,
        mobile TEXT DEFAULT '',
        otp_code TEXT DEFAULT '',
        otp_expires TEXT DEFAULT '',
        created_at TEXT DEFAULT (datetime('now'))
    )""")
    # Add new columns for upgrading from older schema
    for col_def in [
        "ALTER TABLE users ADD COLUMN plaintext_pw TEXT DEFAULT ''",
        "ALTER TABLE users ADD COLUMN email TEXT DEFAULT ''",
        "ALTER TABLE users ADD COLUMN email_verified INTEGER DEFAULT 0",
        "ALTER TABLE users ADD COLUMN mobile TEXT DEFAULT ''",
        "ALTER TABLE users ADD COLUMN otp_code TEXT DEFAULT ''",
        "ALTER TABLE users ADD COLUMN otp_expires TEXT DEFAULT ''",
        "ALTER TABLE users ADD COLUMN last_login TEXT DEFAULT ''",
    ]:
        try:
            c.execute(col_def)
        except Exception:
            pass  # Column already exists
    c.execute("""CREATE TABLE IF NOT EXISTS activity_log (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        username TEXT NOT NULL,
        action TEXT NOT NULL,
        detail TEXT,
        ip TEXT,
        ts TEXT DEFAULT (datetime('now'))
    )""")
    # ── Indexes for performance ────────────────────────────────────────────────
    c.execute("CREATE INDEX IF NOT EXISTS idx_activity_username ON activity_log(username)")
    c.execute("CREATE INDEX IF NOT EXISTS idx_activity_ts ON activity_log(ts DESC)")
    c.execute("CREATE INDEX IF NOT EXISTS idx_activity_action ON activity_log(action)")
    c.execute("CREATE INDEX IF NOT EXISTS idx_users_username ON users(username)")
    # Always ensure admin exists (re-seeds after Render restarts)
    c.execute("""INSERT OR IGNORE INTO users (username,name,password_hash,role)
                 VALUES (?,?,?,?)""",
              ("admin", "Administrator", hash_pw("Admin@1234"), "admin"))
    # Also ensure karthi account exists
    c.execute("""INSERT OR IGNORE INTO users (username,name,password_hash,role)
                 VALUES (?,?,?,?)""",
              ("karthi", "Karthi", hash_pw("karthi123"), "admin"))
    conn.commit()
    conn.close()

# (DB init is called after all functions are defined — see bottom of file)

def hash_pw(pw: str) -> str:
    return hashlib.sha256(pw.encode()).hexdigest()

def verify_pw(pw: str, h: str) -> bool:
    return hash_pw(pw) == h

@app.before_request
def ensure_db_and_auth():
    """1. Re-init DB if missing. 2. Enforce login on every route."""
    # DB check
    try:
        conn = get_db()
        conn.execute("SELECT 1 FROM users LIMIT 1")
        conn.close()
    except Exception:
        try:
            init_db()
        except Exception as e:
            print(f"DB re-init error: {e}")
    # Auth enforcement — public endpoints that don't need login
    public = {'/login', '/register', '/api/health', '/api/quote',
              '/api/send_otp', '/api/verify_otp', '/api/create_account',
              '/api/reg_send_otp', '/api/reg_verify_otp',
              '/api/forgot_password', '/api/reset_password', '/static'}
    path = request.path
    if any(path.startswith(p) for p in public):
        return  # allow through
    if 'username' not in session:
        if request.is_json or path.startswith('/api/'):
            return jsonify({'error': 'Not authenticated', 'redirect': '/login'}), 401
        return redirect(url_for('login_page'))

ADMIN_EMAIL = os.environ.get("ADMIN_EMAIL", "")
SMTP_HOST   = os.environ.get("SMTP_HOST", "smtp.gmail.com")
SMTP_PORT   = int(os.environ.get("SMTP_PORT", "587"))
SMTP_USER   = os.environ.get("SMTP_USER", "")
SMTP_PASS   = os.environ.get("SMTP_PASS", "")

_notify_lock = threading.Lock()

def _send_email_notification(username: str, action: str, detail: str, ip: str, ts: str):
    """Send admin email notification in background thread. Silent if not configured."""
    if not all([ADMIN_EMAIL, SMTP_USER, SMTP_PASS]):
        return
    def _send():
        try:
            subject = f"[APSG] Activity: {action} by {username}"
            body = (
                f"<h3>APSG System Activity Notification</h3>"
                f"<table style='border-collapse:collapse;font-family:monospace;font-size:14px;'>"
                f"<tr><td style='padding:4px 12px;font-weight:bold;'>User</td><td>{username}</td></tr>"
                f"<tr><td style='padding:4px 12px;font-weight:bold;'>Action</td><td>{action}</td></tr>"
                f"<tr><td style='padding:4px 12px;font-weight:bold;'>Detail</td><td>{detail or chr(8212)}</td></tr>"
                f"<tr><td style='padding:4px 12px;font-weight:bold;'>IP</td><td>{ip}</td></tr>"
                f"<tr><td style='padding:4px 12px;font-weight:bold;'>Timestamp</td><td>{ts}</td></tr>"
                f"</table>"
                f"<p style='color:#888;font-size:12px;'>APSG Report System &middot; Developed by Karthi</p>"
            )
            msg = MIMEMultipart("alternative")
            msg["Subject"] = subject
            msg["From"]    = SMTP_USER
            msg["To"]      = ADMIN_EMAIL
            msg.attach(MIMEText(body, "html"))
            with smtplib.SMTP(SMTP_HOST, SMTP_PORT, timeout=10) as srv:
                srv.ehlo(); srv.starttls(); srv.login(SMTP_USER, SMTP_PASS)
                srv.sendmail(SMTP_USER, ADMIN_EMAIL, msg.as_string())
        except Exception as e:
            print(f"[Email] Notification failed: {e}")
    with _notify_lock:
        threading.Thread(target=_send, daemon=True).start()


def send_activity_email(username: str, action: str, detail: str, ts: str, ip: str):
    """Send email notification for key actions (non-blocking, best-effort)."""
    import smtplib, os
    from email.mime.text import MIMEText
    notify_email = os.environ.get("NOTIFY_EMAIL", "")
    smtp_host    = os.environ.get("SMTP_HOST", "smtp.gmail.com")
    smtp_port    = int(os.environ.get("SMTP_PORT", "587"))
    smtp_user    = os.environ.get("SMTP_USER", "")
    smtp_pass    = os.environ.get("SMTP_PASS", "")
    if not (notify_email and smtp_user and smtp_pass):
        return  # Email not configured — skip silently
    try:
        subject = "[APSG] " + action + " - " + username
        sep = "-" * 40
        body = "\n".join([
            "APSG Activity Notification",
            sep,
            "User      : " + username,
            "Action    : " + action,
            "Detail    : " + (detail or "-"),
            "Timestamp : " + ts,
            "IP        : " + (ip or "unknown"),
        ])
        msg = MIMEText(body)
        msg["Subject"] = subject
        msg["From"] = smtp_user
        msg["To"] = notify_email
        with smtplib.SMTP(smtp_host, smtp_port) as srv:
            srv.ehlo(); srv.starttls(); srv.ehlo()
            srv.login(smtp_user, smtp_pass)
            srv.sendmail(smtp_user, [notify_email], msg.as_string())
    except Exception:
        pass  # Never let email errors break the main flow

_EMAIL_ACTIONS = {"LOGIN", "OPEN_APP", "DAILY_REPORT", "RECTIFICATION_REPORT",
                  "EXCEL_REJECTION", "PPT_GENERATE"}

def log_activity(username: str, action: str, detail: str = ""):
    ts = datetime.utcnow().strftime("%Y-%m-%d %H:%M:%S UTC")
    ip = ""
    try:
        ip = request.remote_addr or ""
    except Exception:
        pass
    # Send email notification for key actions (in background thread)
    if action in _EMAIL_ACTIONS:
        import threading
        threading.Thread(
            target=send_activity_email,
            args=(username, action, detail, ts, ip),
            daemon=True
        ).start()
    try:
        conn = get_db()
        conn.execute("INSERT INTO activity_log (username,action,detail,ip) VALUES (?,?,?,?)",
                     (username, action, detail, ip))
        conn.commit()
        conn.close()
    except Exception:
        pass
    _send_email_notification(username, action, detail, ip, ts)

# ─── auth helpers ────────────────────────────────────────────────────────────
def login_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if "username" not in session:
            # Return JSON for API/fetch calls so the frontend gets a readable error
            if request.is_json or request.path.startswith('/api/'):
                return jsonify({"ok": False, "error": "Session expired. Please log in again.",
                                "redirect": "/login"}), 401
            return redirect(url_for("login_page"))
        return f(*args, **kwargs)
    return decorated

def admin_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if session.get("role") != "admin":
            return jsonify({"error": "Admin only"}), 403
        return f(*args, **kwargs)
    return decorated

# ═══════════════════════════════════════════════════════════════════════════════
#  PPT REJECTION LOGIC  (unchanged from original)
# ═══════════════════════════════════════════════════════════════════════════════

SLIDE_W  = 12192000; SLIDE_H  = 6858000
TBL_LEFT = 597159;   TBL_TOP  = 1189793; TBL_W = 11306716
COL_WIDTHS = [289286,1425361,739037,826948,1091559,928816,774237,2169617,626565,2435290]
ROW_H_HEADER = 681522; ROW_H_DATA = 734233
TITLE_LEFT=0; TITLE_TOP=288093; TITLE_W=12283440; TITLE_H=901700
HEADER_COLOR='00B050'; BORDER_W=6350

COMPANY_ABBREV = {
    'KKL':'KOH KOCK LEONG ENTERPRISE PTE LTD','KTC':'KTC Civil Engineering & Construction Pte Ltd',
    'UNC':'Unity Contractors Pte Ltd','CCE':'Chan & Chan Engineering Pte Ltd',
    'CLM':'Chuan Lim Construction Pte Ltd','HHT':'Hong Aik Engineering Pte Ltd',
    'JME':'JME E&C PTE LTD','KTE':'KOK TONG EARTHWORKS & ENGINEERING PTE LTD',
    'QQC':'QUEK & QUEK CIVIL ENGINEERING PTE LTD','BKH':'Backho (S) Pte Ltd',
    'OTP':'OKT TRANSPORT PTE LTD','RCE':'RECLAIMS ENTERPRISE PTE LTD',
    'SIN':'Sin Heng Transport Pte Ltd','HSE':'Hanshika Engineering & Construction Pte Ltd',
    'JFT':'JIN FENG TRANSPORT PTE LTD','CJC':'Chye Joo Construction Pte. Ltd.',
    'GGC':'Guan Gi Construction Pte Ltd','HTC':'HUATIONG CONTRACTOR PTE LTD',
    'KHT':'KENG HO TRADING AND TRANSPORT PTE LTD','MSH':'Megastone Holdings Pte Ltd',
    'MTC':'Metrocon Pte Ltd','TGT':'Tengah Transportation & Construction Pte Ltd',
    'SAS':'SASAN CONSTRUCTION PTE LTD','SHC':'SIN HUA CIVIL ENGINEERING & CONSTRUCTION PTE LTD',
    'WCF':'Wang Cheng Foundation Pte Ltd','YEC':'YONGSHENG E & C PTE LTD',
    'EEH':'EE HUP CONSTRUCTION PTE LTD','SLE':'Sam Lain Equipment Services Pte Ltd',
}

_BANNER_BLOB_CACHE = [None]

def _get_banner_blob():
    if _BANNER_BLOB_CACHE[0]: return _BANNER_BLOB_CACHE[0]
    for prs in PPT_DB.get('presentations',[]):
        for slide in prs.slides:
            for sh in slide.shapes:
                if sh.shape_type==13 and hasattr(sh,'width') and sh.width>9000000 and sh.top>5500000:
                    xml_s=etree.tostring(sh._element).decode()
                    embeds=re.findall(r'r:embed="(rId\d+)"',xml_s)
                    for rId in embeds:
                        rel=slide.part.rels.get(rId)
                        if rel and 'image' in rel.reltype:
                            _BANNER_BLOB_CACHE[0]=rel.target_part.blob
                            return _BANNER_BLOB_CACHE[0]
    return None

def _add_footer_banner(dst_slide):
    blob=_get_banner_blob()
    if not blob: return
    BANNER_H=557487; BANNER_Y=SLIDE_H-BANNER_H
    img_part,new_rId=dst_slide.part.get_or_add_image_part(io.BytesIO(blob))
    NS_P='xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
    NS_A='xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    NS_R=f'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"'
    pic_xml=(f'<p:pic {NS_P} {NS_A} {NS_R}><p:nvPicPr><p:cNvPr id="999" name="FooterBanner"/>'
             f'<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr><p:nvPr/></p:nvPicPr>'
             f'<p:blipFill><a:blip r:embed="{new_rId}" cstate="email"><a:extLst/></a:blip>'
             f'<a:stretch><a:fillRect/></a:stretch></p:blipFill>'
             f'<p:spPr><a:xfrm><a:off x="0" y="{BANNER_Y}"/><a:ext cx="{SLIDE_W}" cy="{BANNER_H}"/>'
             f'</a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr></p:pic>')
    dst_slide.shapes._spTree.append(etree.fromstring(pic_xml.encode()))

def _borders():
    fill=('<a:solidFill><a:srgbClr val="000000"/></a:solidFill>'
          '<a:prstDash val="solid"/><a:round/>'
          '<a:headEnd type="none" w="med" len="med"/>'
          '<a:tailEnd type="none" w="med" len="med"/>')
    return ''.join(f'<a:{s} w="{BORDER_W}" cap="flat" cmpd="sng" algn="ctr">{fill}</a:{s}>'
                   for s in ['lnL','lnR','lnT','lnB'])

def parse_date_ppt(s):
    try: return datetime.strptime(s.strip(),'%d-%m-%Y')
    except: return None

def cell_xml(text,header,col):
    ns='xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    safe=(str(text or '')).replace('&','&amp;').replace('<','&lt;').replace('>','&gt;').replace('"','&quot;')
    if header:
        algn,fa='ctr','b'; sz,bold,color,face=1200,1,HEADER_COLOR,'Calibri'
        mar='marL="8912" marR="8912" marT="8912" marB="0"'
    else:
        algn,fa='ctr','ctr'; sz,bold,color,face=1100,0,'000000','Calibri'
        mar='marL="9525" marR="9525" marT="9525" marB="0"'
    return (f'<a:tc {ns}><a:txBody><a:bodyPr/><a:lstStyle/>'
            f'<a:p><a:pPr algn="{algn}" fontAlgn="{fa}"/>'
            f'<a:r><a:rPr lang="en-SG" sz="{sz}" b="{bold}" i="0" u="none" '
            f'strike="noStrike" dirty="0">'
            f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
            f'<a:effectLst/><a:latin typeface="{face}"/>'
            f'</a:rPr><a:t>{safe}</a:t></a:r></a:p>'
            f'</a:txBody>'
            f'<a:tcPr {mar} anchor="ctr">{_borders()}<a:noFill/></a:tcPr>'
            f'</a:tc>')

def is_summary_slide(slide):
    return any(sh.has_table and len(sh.table.columns)==10 for sh in slide.shapes)

def is_detail_slide(slide):
    return any(sh.has_table and len(sh.table.columns)==8 for sh in slide.shapes)

def extract_from_prs(prs, prs_idx):
    records,slide_map=[],{}
    for si,slide in enumerate(prs.slides):
        if is_summary_slide(slide):
            for sh in slide.shapes:
                if sh.has_table and len(sh.table.columns)==10:
                    tbl=sh.table
                    for ri in range(1,len(tbl.rows)):
                        cells=[c.text.strip() for c in tbl.rows[ri].cells]
                        if len(cells)>=10 and cells[1]:
                            records.append({'sn':cells[0],'ticket_no':cells[1],'veh_no':cells[2],
                                'material':cells[3],'source_site':cells[4],'date':cells[5],
                                'time':cells[6],'e_token':cells[7],'accepted':cells[8],
                                'reject_reason':cells[9],'_prs':prs_idx})
        elif is_detail_slide(slide):
            for sh in slide.shapes:
                if sh.has_table and len(sh.table.columns)==8:
                    tbl=sh.table
                    if len(tbl.rows)>1:
                        tn=tbl.rows[1].cells[0].text.strip()
                        if tn: slide_map[tn]=(prs_idx,si)
                    break
    return records,slide_map

def _safe_layout(prs, preferred=6):
    """Return a slide layout safely — falls back to last or first available."""
    layouts = prs.slide_layouts
    if len(layouts) > preferred:
        return layouts[preferred]
    if len(layouts) > 0:
        # prefer blank (index 6) or last layout
        return layouts[-1]
    raise ValueError("Presentation has no slide layouts")

def add_summary_slide(out_prs,chunk,slide_num,sn_offset=1):
    slide=out_prs.slides.add_slide(_safe_layout(out_prs, 6))
    tx=slide.shapes.add_textbox(Emu(TITLE_LEFT),Emu(TITLE_TOP),Emu(TITLE_W),Emu(TITLE_H))
    tf=tx.text_frame; tf.word_wrap=True
    r0=chunk[0] if chunk else {}
    sites=list(dict.fromkeys(r.get('source_site','') for r in chunk if r.get('source_site')))
    site_lbl=(sites[0] if len(sites)==1 else
              (f'{sites[0]}, {sites[1]}' if len(sites)==2 else f'{sites[0]} et al.'))
    p1=tf.paragraphs[0]; p1.alignment=PP_ALIGN.CENTER
    r1=p1.add_run(); r1.text='Rejected Loads Summary'
    r1.font.name='Calibri'; r1.font.size=Pt(28); r1.font.bold=True
    r1.font.underline=True; r1.font.color.rgb=RGBColor(0,0,0)
    p2=tf.add_paragraph(); p2.alignment=PP_ALIGN.CENTER
    r2=p2.add_run(); r2.text=site_lbl
    r2.font.name='Calibri'; r2.font.size=Pt(20); r2.font.bold=True
    r2.font.color.rgb=RGBColor(0x22,0x22,0x22)
    n_data=len(chunk)
    tbl_frm=slide.shapes.add_table(1+n_data,10,Emu(TBL_LEFT),Emu(TBL_TOP),
                                   Emu(TBL_W),Emu(ROW_H_HEADER+ROW_H_DATA*n_data))
    tbl=tbl_frm.table
    for ci,w in enumerate(COL_WIDTHS): tbl.columns[ci].width=w
    tbl.rows[0].height=ROW_H_HEADER
    for ri in range(1,1+n_data): tbl.rows[ri].height=ROW_H_DATA
    HEADERS=['S/N','Ticket No','Veh No','Material','Source Site','Date In','Time In','E-Token','Accepted','Reject Reason']
    for ci,h in enumerate(HEADERS):
        tc=tbl.rows[0].cells[ci]._tc
        tc.getparent().replace(tc,etree.fromstring(cell_xml(h,True,ci)))
    for ri,rec in enumerate(chunk):
        vals=[str(sn_offset+ri),rec.get('ticket_no',''),rec.get('veh_no',''),
              rec.get('material',''),rec.get('source_site',''),rec.get('date',''),
              rec.get('time',''),rec.get('e_token',''),rec.get('accepted','NO'),rec.get('reject_reason','')]
        for ci,v in enumerate(vals):
            tc=tbl.rows[ri+1].cells[ci]._tc
            tc.getparent().replace(tc,etree.fromstring(cell_xml(v,False,ci)))
    nb=slide.shapes.add_textbox(Emu(8610600),Emu(6356350),Emu(2743200),Emu(365125))
    np_=nb.text_frame.paragraphs[0]; np_.alignment=PP_ALIGN.RIGHT
    nr=np_.add_run(); nr.text=str(slide_num); nr.font.size=Pt(12); nr.font.color.rgb=RGBColor(0,0,0)
    _add_footer_banner(slide)
    return slide

def _replace_rids_single_pass(xml_str,rId_map):
    pattern=re.compile(r'(r:embed|r:link)="(rId\d+)"')
    def replacer(m): return f'{m.group(1)}="{rId_map.get(m.group(2),m.group(2))}"'
    return pattern.sub(replacer,xml_str)

def clone_slide(src_prs,src_idx,dst_prs,remove_slide_number=True,sn=None):
    src=src_prs.slides[src_idx]
    dst=dst_prs.slides.add_slide(_safe_layout(dst_prs, 6))
    src_tree=src.shapes._spTree; dst_tree=dst.shapes._spTree
    for ch in list(dst_tree)[2:]: dst_tree.remove(ch)
    for ch in list(src_tree)[2:]: dst_tree.append(copy.deepcopy(ch))
    NS_P='http://schemas.openxmlformats.org/presentationml/2006/main'
    NS_A='http://schemas.openxmlformats.org/drawingml/2006/main'
    if remove_slide_number:
        for sp in list(dst_tree.iter(f'{{{NS_P}}}sp')):
            ph=sp.find(f'.//{{{NS_P}}}ph')
            if ph is not None and ph.get('type')=='sldNum':
                sp.getparent().remove(sp); break
    rId_map={}
    for rel in src.part.rels.values():
        if 'image' not in rel.reltype: continue
        try:
            _,new_rId=dst.part.get_or_add_image_part(io.BytesIO(rel.target_part.blob))
            rId_map[rel.rId]=new_rId
        except Exception as e:
            app.logger.warning(f'Image copy rId={rel.rId}: {e}')
    if rId_map:
        xml_str=etree.tostring(dst_tree).decode()
        xml_fixed=_replace_rids_single_pass(xml_str,rId_map)
        new_tree=etree.fromstring(xml_fixed.encode())
        dst_tree.getparent().replace(dst_tree,new_tree)
    return dst

def get_group_key(source_id):
    if len(source_id)<6: return source_id.upper()
    substr=source_id[5:8]
    alpha=re.sub(r'[^A-Za-z]','',substr).upper()
    return alpha if alpha else substr.upper()

def build_report(records):
    groups=OrderedDict()
    for rec in sorted(records,key=lambda r:(r.get('source_site',''),r.get('e_token',''))):
        site=rec.get('source_site','Unknown')
        groups.setdefault(site,[]).append(rec)
    out=Presentation(); out.slide_width=Emu(SLIDE_W); out.slide_height=Emu(SLIDE_H)
    slide_num=1
    for site,site_recs in groups.items():
        sn_offset=1
        for chunk in [site_recs[i:i+5] for i in range(0,len(site_recs),5)]:
            add_summary_slide(out,chunk,slide_num,sn_offset=sn_offset)
            sn_offset+=len(chunk); slide_num+=1
        for sn_idx,rec in enumerate(site_recs,start=1):
            tn=rec.get('ticket_no','')
            if tn in PPT_DB['slide_map']:
                pi,si=PPT_DB['slide_map'][tn]
                try: clone_slide(PPT_DB['presentations'][pi],si,out,remove_slide_number=True,sn=sn_idx)
                except Exception as e: app.logger.warning(f'Clone failed {tn}: {e}')
            slide_num+=1
    buf=io.BytesIO(); out.save(buf); buf.seek(0); return buf

def build_excel_ppt(records):
    wb=openpyxl.Workbook(); ws=wb.active; ws.title='Rejection Report'
    hdr_font=Font(name='Calibri',bold=True,color='FFFFFF',size=11)
    hdr_fill=PatternFill('solid',fgColor='00B050')
    hdr_align=Alignment(horizontal='center',vertical='center',wrap_text=True)
    hdr_border=Border(left=Side(style='thin'),right=Side(style='thin'),top=Side(style='thin'),bottom=Side(style='thin'))
    data_font=Font(name='Calibri',size=10)
    data_align=Alignment(horizontal='center',vertical='center',wrap_text=True)
    data_border=Border(left=Side(style='thin'),right=Side(style='thin'),top=Side(style='thin'),bottom=Side(style='thin'))
    HEADERS=['S/N','Ticket No','Veh No','Material','Source Site','Date In','Time In','E-Token','Accepted','Reject Reason']
    COL_W_XLSX=[6,18,12,12,18,14,10,32,10,26]
    for ci,(h,w) in enumerate(zip(HEADERS,COL_W_XLSX),start=1):
        cell=ws.cell(row=1,column=ci,value=h)
        cell.font=hdr_font; cell.fill=hdr_fill; cell.alignment=hdr_align; cell.border=hdr_border
        ws.column_dimensions[get_column_letter(ci)].width=w
    ws.row_dimensions[1].height=28
    sorted_recs=sorted(records,key=lambda r:(r.get('source_site',''),r.get('e_token','')))
    current_site=None; row=2; sn=1
    for rec in sorted_recs:
        site=rec.get('source_site','')
        if site!=current_site:
            if current_site is not None:
                ws.row_dimensions[row].height=6; row+=1
            ws.merge_cells(start_row=row,start_column=1,end_row=row,end_column=10)
            gc=ws.cell(row=row,column=1,value=f'▶  Source Site: {site}')
            gc.font=Font(name='Calibri',bold=True,size=11,color='1B5E20')
            gc.fill=PatternFill('solid',fgColor='C8E6C9')
            gc.alignment=Alignment(horizontal='left',vertical='center')
            gc.border=hdr_border; ws.row_dimensions[row].height=22; row+=1; current_site=site; sn=1
        vals=[sn,rec.get('ticket_no',''),rec.get('veh_no',''),rec.get('material',''),
              rec.get('source_site',''),rec.get('date',''),rec.get('time',''),
              rec.get('e_token',''),rec.get('accepted','NO'),rec.get('reject_reason','')]
        is_alt=(sn%2==0)
        row_fill=PatternFill('solid',fgColor='F0FAF4') if is_alt else None
        for ci,v in enumerate(vals,start=1):
            cell=ws.cell(row=row,column=ci,value=v)
            cell.font=data_font; cell.border=data_border; cell.alignment=data_align
            if row_fill: cell.fill=row_fill
        ws.row_dimensions[row].height=20; row+=1; sn+=1
    ws.freeze_panes='A2'
    ws.auto_filter.ref=f'A1:{get_column_letter(len(HEADERS))}1'
    buf=io.BytesIO(); wb.save(buf); buf.seek(0); return buf

def build_zip_ppt(records):
    import zipfile as _zipfile
    company_groups=OrderedDict()
    for rec in sorted(records,key=lambda r:(get_group_key(r.get('source_site','')),r.get('source_site',''),r.get('e_token',''))):
        gkey=get_group_key(rec.get('source_site',''))
        company_groups.setdefault(gkey,[]).append(rec)
    zip_buf=io.BytesIO()
    with _zipfile.ZipFile(zip_buf,'w',_zipfile.ZIP_DEFLATED) as zf:
        for gkey,grp_recs in company_groups.items():
            pptx_buf=build_report(grp_recs)
            safe_key=re.sub(r'[\/:*?"<>|]','_',gkey)
            zf.writestr(f'{safe_key}_Rejected_Reports/{safe_key}_Rejected_Report.pptx',pptx_buf.read())
    zip_buf.seek(0); return zip_buf

# ═══════════════════════════════════════════════════════════════════════════════
#  PHOTO MERGE WORKER
# ═══════════════════════════════════════════════════════════════════════════════

def _save_upload(file_obj, tag=""):
    name=secure_filename(file_obj.filename or "upload.pptx")
    uid=uuid.uuid4().hex[:8]
    path=os.path.join(UPLOAD_DIR,f"{uid}_{tag}_{name}")
    file_obj.save(path); return path

def _run_merge(job_id,path_a,path_b,top_ph,front_ph,verbose,
               top_h_cm,top_w_cm,top_left_cm,front_h_cm,front_w_cm,front_left_cm,center_gap_cm):
    def upd(pct,msg):
        with jobs_lock:
            jobs[job_id]["progress"]=pct; jobs[job_id]["status"]=msg
    with jobs_lock:
        jobs[job_id]={"status":"Starting…","progress":0,"log_lines":[],"result_file":None,
                      "output_name":"","reference_file":"?","error":None,"stats":None}
    try:
        upd(5,"Loading files…"); time.sleep(0.1)
        upd(25,"Extracting tokens and images…")
        tmp_out=os.path.join(OUTPUT_DIR,f"{job_id}_output.pptx")
        result=ppt_merger.merge(path_a,path_b,tmp_out,
            top_ph=top_ph,front_ph=front_ph,verbose=verbose,
            top_h_cm=top_h_cm,top_w_cm=top_w_cm,top_left_cm=top_left_cm,
            front_h_cm=front_h_cm,front_w_cm=front_w_cm,front_left_cm=front_left_cm,
            center_gap_cm=center_gap_cm)
        upd(90,"Finalising…"); time.sleep(0.15)
        final_path=getattr(result,"_final_path",None)
        if final_path and not os.path.exists(final_path): final_path=tmp_out
        with jobs_lock:
            jobs[job_id]["log_lines"]=result.log_lines
            jobs[job_id]["stats"]=result.to_dict()
            jobs[job_id]["output_name"]=result.output_name
            jobs[job_id]["reference_file"]=result.reference_file
            jobs[job_id]["progress"]=100
            if result.ok and final_path and os.path.exists(final_path):
                jobs[job_id]["status"]="complete"
                jobs[job_id]["result_file"]=final_path
                jobs[job_id]["output_name"]=result.output_name
            else:
                err="; ".join(result.errors) if result.errors else "Merge failed"
                jobs[job_id]["status"]="error"; jobs[job_id]["error"]=err
    except Exception as e:
        with jobs_lock:
            jobs[job_id]["status"]="error"; jobs[job_id]["error"]=str(e); jobs[job_id]["progress"]=100
    finally:
        for p in [path_a,path_b]:
            if p and os.path.exists(p):
                try: os.remove(p)
                except: pass

# ═══════════════════════════════════════════════════════════════════════════════
#  ROUTES — Auth
# ═══════════════════════════════════════════════════════════════════════════════


@app.route('/static/bg.jpg')
def serve_bg():
    import os
    path = os.path.join(os.path.dirname(__file__), 'static', 'bg.jpg')
    return send_file(path, mimetype='image/jpeg', max_age=86400)


@app.route("/login", methods=["GET","POST"])
def login_page():
    if request.method=="POST":
        data     = request.get_json(force=True, silent=True) or request.form
        username = data.get("username","").strip().lower()
        password = data.get("password","")
        client_ip = request.remote_addr or "unknown"

        # Rate-limit check
        allowed, rate_err = _check_rate_limit(client_ip)
        if not allowed:
            if request.is_json: return jsonify({"ok": False, "error": rate_err}), 429
            return redirect(url_for("login_page"))

        conn = get_db()
        user = conn.execute("SELECT * FROM users WHERE username=?", (username,)).fetchone()
        if user and verify_pw(password, user["password_hash"]):
            conn.execute("UPDATE users SET last_login=? WHERE username=?",
                         (_dt.datetime.utcnow().strftime('%Y-%m-%d %H:%M'), username))
            conn.commit()
            conn.close()
            _clear_login_attempts(client_ip)
            session.permanent = True   # respects PERMANENT_SESSION_LIFETIME (8 hours)
            session["username"] = user["username"]
            session["name"]     = user["name"]
            session["role"]     = user["role"]
            log_activity(username, "LOGIN", "Successful login")
            if request.is_json: return jsonify({"ok": True, "name": user["name"], "role": user["role"]})
            return redirect(url_for("dashboard"))
        conn.close()
        _record_failed_login(client_ip)
        log_activity(username or "unknown", "FAILED_LOGIN", f"Bad password from {client_ip}")
        if request.is_json: return jsonify({"ok": False, "error": "Invalid username or password"}), 401
        return redirect(url_for("login_page"))
    return render_template_string(AUTH_HTML, page="login")


# ── Simple in-memory login rate limiter ──────────────────────────────────────
# Tracks failed attempts per IP. Resets after RATE_WINDOW seconds.
_login_attempts: dict = {}   # ip → {'count': int, 'first': datetime, 'locked_until': datetime|None}
_RATE_MAX     = 10           # max failed attempts before lockout
_RATE_WINDOW  = 300          # seconds before attempt counter resets (5 min)
_RATE_LOCKOUT = 600          # lockout duration in seconds (10 min)

def _check_rate_limit(ip: str) -> tuple[bool, str]:
    """Returns (allowed, error_message). Cleans up stale entries automatically."""
    now = _dt.datetime.utcnow()
    entry = _login_attempts.get(ip)
    if entry:
        # Clear expired lockout
        if entry.get('locked_until') and now >= entry['locked_until']:
            _login_attempts.pop(ip, None)
            entry = None
        # Clear stale window
        elif not entry.get('locked_until') and (now - entry['first']).seconds > _RATE_WINDOW:
            _login_attempts.pop(ip, None)
            entry = None

    if entry and entry.get('locked_until'):
        remaining = int((entry['locked_until'] - now).total_seconds())
        return False, f"Too many failed attempts. Try again in {remaining // 60 + 1} minute(s)."
    return True, ""

def _record_failed_login(ip: str):
    now = _dt.datetime.utcnow()
    entry = _login_attempts.setdefault(ip, {'count': 0, 'first': now, 'locked_until': None})
    entry['count'] += 1
    if entry['count'] >= _RATE_MAX:
        entry['locked_until'] = now + _dt.timedelta(seconds=_RATE_LOCKOUT)

def _clear_login_attempts(ip: str):
    _login_attempts.pop(ip, None)

def _send_otp_email(to_email: str, otp: str, name: str) -> bool:
    """Send OTP to user email. Returns True on success."""
    smtp_host = os.environ.get("SMTP_HOST", "smtp.gmail.com")
    smtp_port = int(os.environ.get("SMTP_PORT", "587"))
    smtp_user = os.environ.get("SMTP_USER", "")
    smtp_pass = os.environ.get("SMTP_PASS", "")
    if not (smtp_user and smtp_pass):
        return False  # Email not configured
    try:
        msg = MIMEMultipart("alternative")
        msg["Subject"] = "APSG Report — Your Verification Code"
        msg["From"]    = smtp_user
        msg["To"]      = to_email
        html_body = f"""
        <div style="font-family:Inter,sans-serif;max-width:420px;margin:0 auto;padding:32px 24px;
                    background:#0D1120;border-radius:16px;border:1px solid rgba(99,102,241,.3);">
          <div style="text-align:center;margin-bottom:24px;">
            <div style="font-size:2rem;margin-bottom:8px;">🔐</div>
            <h2 style="color:#A5B4FC;margin:0;font-size:1.1rem;">Email Verification</h2>
            <p style="color:#64748B;font-size:.82rem;margin-top:6px;">APSG (Staging Ground) Report</p>
          </div>
          <p style="color:#CBD5E1;font-size:.9rem;">Hi <strong style="color:#fff">{name}</strong>,</p>
          <p style="color:#94A3B8;font-size:.85rem;margin:12px 0;">
            Your one-time verification code is:
          </p>
          <div style="text-align:center;margin:20px 0;">
            <span style="font-size:2.4rem;font-weight:900;letter-spacing:.3em;
                         color:#818CF8;background:rgba(99,102,241,.12);
                         padding:14px 28px;border-radius:12px;
                         border:1px solid rgba(99,102,241,.35);display:inline-block;">
              {otp}
            </span>
          </div>
          <p style="color:#64748B;font-size:.78rem;text-align:center;">
            This code expires in <strong style="color:#F59E0B">10 minutes</strong>.<br>
            Do not share this code with anyone.
          </p>
        </div>"""
        msg.attach(MIMEText(html_body, "html"))
        with smtplib.SMTP(smtp_host, smtp_port, timeout=10) as srv:
            srv.starttls()
            srv.login(smtp_user, smtp_pass)
            srv.sendmail(smtp_user, [to_email], msg.as_string())
        return True
    except Exception as e:
        print(f"[OTP Email] Failed: {e}")
        return False

# Temporary OTP store (in-memory, keyed by email) — cleared after use or expiry
_otp_store: dict = {}  # email → {otp, expires, name, username, password, confirmed_pw}


def _norm_mobile(m: str) -> str:
    """Strip non-digits for uniqueness comparison."""
    import re as _re
    return _re.sub(r'\D', '', str(m))


# ── Registration flow: collect all fields → Email OTP → create account ──

@app.route("/api/reg_send_otp", methods=["POST", "OPTIONS"])
def api_reg_send_otp():
    """
    Registration Step 1: validate all fields, check for duplicate username/email/mobile,
    generate OTP, send to email address.
    """
    if request.method == "OPTIONS":
        return '', 200
    data    = request.get_json(force=True, silent=True) or {}
    name     = str(data.get("name", "")).strip()
    username = str(data.get("username", "")).strip().lower()
    email    = str(data.get("email", "")).strip().lower()
    mobile   = str(data.get("mobile", "")).strip()
    password = str(data.get("password", ""))
    confirm  = str(data.get("confirm", ""))

    # ── Basic field validation ─────────────────────────────────────────────
    if not all([name, username, email, mobile, password, confirm]):
        return jsonify({"ok": False, "error": "All fields are required"}), 400
    if "@" not in email or "." not in email.split("@")[-1]:
        return jsonify({"ok": False, "error": "Please enter a valid email address"}), 400
    if _norm_mobile(mobile) and len(_norm_mobile(mobile)) < 7:
        return jsonify({"ok": False, "error": "Please enter a valid mobile number"}), 400
    if password != confirm:
        return jsonify({"ok": False, "error": "Passwords do not match"}), 400
    if len(password) < 6:
        return jsonify({"ok": False, "error": "Password must be at least 6 characters"}), 400

    # ── Duplicate checks (return specific messages per field) ──────────────
    conn = get_db()
    if conn.execute("SELECT 1 FROM users WHERE LOWER(username)=?", (username,)).fetchone():
        conn.close()
        return jsonify({"ok": False, "field": "username",
                        "error": "Username already registered"}), 409
    if conn.execute("SELECT 1 FROM users WHERE LOWER(email)=?", (email,)).fetchone():
        conn.close()
        return jsonify({"ok": False, "field": "email",
                        "error": "Email ID already registered"}), 409
    mob_norm = _norm_mobile(mobile)
    if mob_norm:
        dup_mob = conn.execute(
            "SELECT 1 FROM users WHERE replace(replace(replace(replace(replace(mobile,'+',' '),' ',''),'-',''),'(',''),')','') LIKE ?",
            (f"%{mob_norm[-9:]}%",)
        ).fetchone()
        if dup_mob:
            conn.close()
            return jsonify({"ok": False, "field": "mobile",
                            "error": "Mobile number already registered"}), 409
    conn.close()

    # ── Generate OTP, store registration data, send OTP to email ──────────
    otp     = str(random.randint(100000, 999999))
    expires = (_dt.datetime.utcnow() + _dt.timedelta(minutes=10)).isoformat()
    _otp_store[f"reg:{email}"] = {
        "otp": otp, "expires": expires, "type": "reg",
        "name": name, "username": username, "email": email,
        "mobile": mobile, "password": password,
    }

    sent = _send_otp_email(email, otp, name)
    if not sent:
        return jsonify({"ok": True, "dev_otp": otp,
                        "message": "Email not configured — dev mode OTP shown"})
    return jsonify({"ok": True, "message": f"OTP sent to {email}"})


@app.route("/api/reg_verify_otp", methods=["POST", "OPTIONS"])
def api_reg_verify_otp():
    """
    Registration Step 2: verify mobile OTP → create account.
    """
    if request.method == "OPTIONS":
        return '', 200
    data  = request.get_json(force=True, silent=True) or {}
    email = str(data.get("email", "")).strip().lower()
    code  = str(data.get("otp", "")).strip()

    key     = f"reg:{email}"
    pending = _otp_store.get(key)
    if not pending or pending.get("type") != "reg":
        return jsonify({"ok": False,
                        "error": "No pending registration. Please restart."}), 400
    if _dt.datetime.utcnow().isoformat() > pending["expires"]:
        _otp_store.pop(key, None)
        return jsonify({"ok": False, "error": "OTP expired. Please request a new code."}), 400
    if code != pending["otp"]:
        return jsonify({"ok": False, "error": "Incorrect code. Please try again."}), 400

    # OTP correct — create the account
    _otp_store.pop(key, None)
    name     = pending["name"]
    username = pending["username"]
    mobile   = pending["mobile"]
    password = pending["password"]
    try:
        conn = get_db()
        # Re-check duplicates at insertion time (race condition guard)
        if conn.execute("SELECT 1 FROM users WHERE LOWER(username)=?", (username,)).fetchone():
            conn.close()
            return jsonify({"ok": False, "error": "Username already registered"}), 409
        if conn.execute("SELECT 1 FROM users WHERE LOWER(email)=?", (email,)).fetchone():
            conn.close()
            return jsonify({"ok": False, "error": "Email ID already registered"}), 409
        conn.execute(
            "INSERT INTO users (username,name,password_hash,plaintext_pw,role,email,email_verified,mobile)"
            " VALUES (?,?,?,?,'user',?,1,?)",
            (username, name, hash_pw(password), password, email, mobile)
        )
        conn.commit(); conn.close()
        log_activity(username, "REGISTER", f"Account created, mobile {mobile}, email {email}")
        session["username"] = username
        session["name"]     = name
        session["role"]     = "user"
        return jsonify({"ok": True, "name": name, "role": "user"})
    except sqlite3.IntegrityError:
        return jsonify({"ok": False, "error": "Username already registered"}), 409



def api_send_otp():
    """Phase A of registration: validate email, generate & send OTP."""
    if request.method == "OPTIONS":
        return '', 200
    data = request.get_json(force=True, silent=True) or {}
    email = str(data.get("email","")).strip().lower()
    name  = str(data.get("name","")).strip() or (email.split("@")[0] if "@" in email else "User")

    if not email:
        return jsonify({"ok": False, "error": "Email address is required"}), 400
    if "@" not in email or "." not in email.split("@")[-1]:
        return jsonify({"ok": False, "error": "Please enter a valid email address"}), 400

    otp     = str(random.randint(100000, 999999))
    expires = (_dt.datetime.utcnow() + _dt.timedelta(minutes=10)).isoformat()
    _otp_store[email] = {"otp": otp, "expires": expires, "name": name}

    sent = _send_otp_email(email, otp, name)
    if not sent:
        return jsonify({"ok": True, "dev_otp": otp,
                        "message": "Email not configured — dev mode OTP shown above"})
    return jsonify({"ok": True, "message": f"OTP sent to {email}"})

@app.route("/api/verify_otp", methods=["POST","OPTIONS"])
def api_verify_otp():
    """Verify OTP."""
    if request.method == "OPTIONS":
        return '', 200
    data     = request.get_json(force=True, silent=True) or {}
    email    = str(data.get("email","")).strip().lower()
    code     = str(data.get("otp","")).strip()
    otp_only = data.get("otp_only", False)

    pending = _otp_store.get(email)
    if not pending:
        return jsonify({"ok": False, "email_ok": False,
                        "error": "No pending verification. Please restart registration."}), 400
    if _dt.datetime.utcnow().isoformat() > pending["expires"]:
        _otp_store.pop(email, None)
        return jsonify({"ok": False, "email_ok": False,
                        "error": "Code expired. Please request a new OTP."}), 400
    if code != pending["otp"]:
        return jsonify({"ok": False, "email_ok": False,
                        "error": "Incorrect code. Please try again."}), 400

    if otp_only:
        # Just confirm email is verified — mark it in store, don't create account yet
        _otp_store[email]["email_verified"] = True
        return jsonify({"ok": True, "email_ok": True})

    # Full flow (legacy): OTP correct — create the account now
    _otp_store.pop(email, None)
    name     = pending["name"]
    username = pending["username"]
    password = pending["password"]
    if name == "__otp__" or username == "__otp__":
        return jsonify({"ok": False, "email_ok": True,
                        "error": "Use /api/create_account to finalize registration."}), 400
    try:
        conn = get_db()
        conn.execute(
            "INSERT INTO users (username,name,password_hash,plaintext_pw,role,email,email_verified) VALUES (?,?,?,?,'user',?,1)",
            (username, name, hash_pw(password), password, email)
        )
        conn.commit(); conn.close()
        log_activity(username, "REGISTER", f"Verified via {email}")
        session["username"] = username
        session["name"]     = name
        session["role"]     = "user"
        return jsonify({"ok": True, "email_ok": True, "name": name, "role": "user"})
    except sqlite3.IntegrityError:
        return jsonify({"ok": False, "error": "Username already exists"}), 409

@app.route("/api/create_account", methods=["POST","OPTIONS"])
def api_create_account():
    """Final step: create account after email is OTP-verified."""
    if request.method == 'OPTIONS':
        return '', 200
    data     = request.get_json(force=True, silent=True) or {}
    name     = str(data.get("name","")).strip()
    username = str(data.get("username","")).strip().lower()
    password = str(data.get("password",""))
    confirm  = str(data.get("confirm",""))
    email    = str(data.get("email","")).strip().lower()

    if not all([name, username, password, email]):
        return jsonify({"ok": False, "error": "All fields required"}), 400
    if password != confirm:
        return jsonify({"ok": False, "error": "Passwords do not match"}), 400
    if len(password) < 6:
        return jsonify({"ok": False, "error": "Password must be at least 6 characters"}), 400

    # Confirm email was OTP-verified in this session
    pending = _otp_store.get(email)
    if not pending or not pending.get("email_verified"):
        return jsonify({"ok": False, "error": "Email not verified. Please verify OTP first."}), 400

    _otp_store.pop(email, None)
    try:
        conn = get_db()
        conn.execute(
            "INSERT INTO users (username,name,password_hash,plaintext_pw,role,email,email_verified) VALUES (?,?,?,?,'user',?,1)",
            (username, name, hash_pw(password), password, email)
        )
        conn.commit(); conn.close()
        log_activity(username, "REGISTER", f"Account created with verified email {email}")
        session["username"] = username
        session["name"]     = name
        session["role"]     = "user"
        return jsonify({"ok": True, "name": name, "role": "user"})
    except sqlite3.IntegrityError:
        return jsonify({"ok": False, "error": "Username already taken. Please choose another."}), 409


# ══════════════════════════════════════════════════════════════════════════════
#  FORGOT PASSWORD — OTP-based self-service password reset
# ══════════════════════════════════════════════════════════════════════════════

def _send_reset_otp_email(to_email: str, otp: str, name: str) -> bool:
    """Send a password-reset OTP email. Returns True on success."""
    smtp_host = os.environ.get("SMTP_HOST", "smtp.gmail.com")
    smtp_port = int(os.environ.get("SMTP_PORT", "587"))
    smtp_user = os.environ.get("SMTP_USER", "")
    smtp_pass = os.environ.get("SMTP_PASS", "")
    if not (smtp_user and smtp_pass):
        return False
    try:
        msg = MIMEMultipart("alternative")
        msg["Subject"] = "APSG Report — Password Reset Code"
        msg["From"]    = smtp_user
        msg["To"]      = to_email
        html_body = (
            '<div style="font-family:Inter,sans-serif;max-width:420px;margin:0 auto;padding:32px 24px;' +
            'background:#0D1120;border-radius:16px;border:1px solid rgba(239,68,68,.3);">' +
            '<div style="text-align:center;margin-bottom:24px;">' +
            '<div style="font-size:2rem;margin-bottom:8px;">🔑</div>' +
            '<h2 style="color:#F87171;margin:0;font-size:1.1rem;">Password Reset</h2>' +
            '<p style="color:#64748B;font-size:.82rem;margin-top:6px;">APSG (Staging Ground) Report</p>' +
            '</div><p style="color:#CBD5E1;font-size:.9rem;">Hi <strong style="color:#fff">' +
            name + '</strong>,</p>' +
            '<p style="color:#94A3B8;font-size:.85rem;margin:12px 0;">Your password reset code is:</p>' +
            '<div style="text-align:center;margin:20px 0;">' +
            '<span style="font-size:2.4rem;font-weight:900;letter-spacing:.3em;color:#F87171;' +
            'background:rgba(239,68,68,.12);padding:14px 28px;border-radius:12px;' +
            'border:1px solid rgba(239,68,68,.35);display:inline-block;">' +
            otp +
            '</span></div>' +
            '<p style="color:#64748B;font-size:.78rem;text-align:center;">' +
            'Expires in <strong style="color:#F59E0B">10 minutes</strong>.<br>' +
            'If you did not request this, please ignore this email.</p></div>'
        )
        msg.attach(MIMEText(html_body, "html"))
        with smtplib.SMTP(smtp_host, smtp_port, timeout=10) as srv:
            srv.starttls()
            srv.login(smtp_user, smtp_pass)
            srv.sendmail(smtp_user, [to_email], msg.as_string())
        return True
    except Exception as e:
        print(f"[Reset OTP Email] Failed: {e}")
        return False


@app.route("/api/forgot_password", methods=["POST", "OPTIONS"])
def api_forgot_password():
    if request.method == "OPTIONS":
        return '', 200
    data  = request.get_json(force=True, silent=True) or {}
    email = str(data.get("email", "")).strip().lower()
    if not email or "@" not in email or "." not in email.split("@")[-1]:
        return jsonify({"ok": False, "error": "Please enter a valid email address"}), 400
    conn = get_db()
    user = conn.execute(
        "SELECT username, name FROM users WHERE LOWER(email)=?", (email,)
    ).fetchone()
    conn.close()
    if not user:
        return jsonify({"ok": False, "error": "No account found with this email address."}), 404
    otp     = str(random.randint(100000, 999999))
    expires = (_dt.datetime.utcnow() + _dt.timedelta(minutes=10)).isoformat()
    _otp_store[f"reset:{email}"] = {
        "otp": otp, "expires": expires,
        "name": user["name"], "username": user["username"], "type": "reset"
    }
    sent = _send_reset_otp_email(email, otp, user["name"])
    if not sent:
        return jsonify({"ok": True, "dev_otp": otp,
                        "message": "Dev mode — OTP shown (email not configured)"})
    return jsonify({"ok": True, "message": f"Reset code sent to {email}"})


@app.route("/api/reset_password", methods=["POST", "OPTIONS"])
def api_reset_password_self():
    """Verify OTP + current password, then set new password."""
    if request.method == "OPTIONS":
        return '', 200
    data       = request.get_json(force=True, silent=True) or {}
    email      = str(data.get("email", "")).strip().lower()
    otp        = str(data.get("otp", "")).strip()
    current_pw = str(data.get("current_password", ""))
    new_pw     = str(data.get("password", ""))
    confirm    = str(data.get("confirm", ""))

    if not all([email, otp, current_pw, new_pw, confirm]):
        return jsonify({"ok": False, "error": "All fields are required"}), 400
    if new_pw != confirm:
        return jsonify({"ok": False, "error": "Passwords do not match"}), 400
    if len(new_pw) < 6:
        return jsonify({"ok": False, "error": "Password must be at least 6 characters"}), 400

    # Verify OTP
    key     = f"reset:{email}"
    pending = _otp_store.get(key)
    if not pending or pending.get("type") != "reset":
        return jsonify({"ok": False, "error": "No pending reset. Please request a new OTP."}), 400
    if _dt.datetime.utcnow().isoformat() > pending["expires"]:
        _otp_store.pop(key, None)
        return jsonify({"ok": False, "error": "Code expired. Please request a new OTP."}), 400
    if otp != pending["otp"]:
        return jsonify({"ok": False, "error": "Incorrect code. Please try again."}), 400

    # Verify current password
    conn = get_db()
    user = conn.execute("SELECT password_hash FROM users WHERE LOWER(email)=?", (email,)).fetchone()
    if not user:
        conn.close()
        return jsonify({"ok": False, "error": "Account not found."}), 404
    if not verify_pw(current_pw, user["password_hash"]):
        conn.close()
        return jsonify({"ok": False, "field": "current_password",
                        "error": "Current password is incorrect."}), 400

    # Both OTP and current password verified — update password
    _otp_store.pop(key, None)
    conn.execute("UPDATE users SET password_hash=?, plaintext_pw=? WHERE LOWER(email)=?",
                 (hash_pw(new_pw), new_pw, email))
    conn.commit(); conn.close()
    log_activity(pending["username"], "PASSWORD_RESET", f"Self-service reset via {email}")
    return jsonify({"ok": True, "message": "Password reset successfully! You can now sign in."})

@app.route("/register", methods=["POST"])
def register():
    data=request.get_json(force=True, silent=True) or {}
    name=data.get("name","").strip()
    username=data.get("username","").strip().lower()
    password=data.get("password","")
    confirm=data.get("confirm","")
    if not all([name,username,password]):
        return jsonify({"ok":False,"error":"All fields required"}),400
    if password!=confirm:
        return jsonify({"ok":False,"error":"Passwords do not match"}),400
    if len(password)<6:
        return jsonify({"ok":False,"error":"Password must be at least 6 characters"}),400
    try:
        conn=get_db()
        conn.execute("INSERT INTO users (username,name,password_hash,plaintext_pw,role) VALUES (?,?,?,?,'user')",
                     (username,name,hash_pw(password),password))
        conn.commit(); conn.close()
        log_activity(username,"REGISTER","New user registered")
        # Auto-login after registration
        session["username"] = username
        session["name"] = name
        session["role"] = "user"
        return jsonify({"ok":True,"name":name,"role":"user"})
    except sqlite3.IntegrityError:
        return jsonify({"ok":False,"error":"Username already exists"}),409

@app.route("/api/user/change_password", methods=["POST"])
@login_required
def user_change_password():
    """Self-service: logged-in user changes their own password."""
    data       = request.get_json(force=True, silent=True) or {}
    current_pw = data.get("current_password","").strip()
    new_pw     = data.get("new_password","").strip()
    confirm    = data.get("confirm","").strip()
    username   = session.get("username")

    if not all([current_pw, new_pw, confirm]):
        return jsonify({"ok": False, "error": "All fields required"}), 400
    if len(new_pw) < 6:
        return jsonify({"ok": False, "error": "New password must be at least 6 characters"}), 400
    if new_pw != confirm:
        return jsonify({"ok": False, "error": "Passwords do not match"}), 400

    conn = get_db()
    try:
        user = conn.execute("SELECT password_hash FROM users WHERE username=?", (username,)).fetchone()
        if not user or not verify_pw(current_pw, user["password_hash"]):
            return jsonify({"ok": False, "error": "Current password is incorrect"}), 401
        conn.execute("UPDATE users SET password_hash=?,plaintext_pw=? WHERE username=?",
                     (hash_pw(new_pw), new_pw, username))
        conn.commit()
        log_activity(username, "CHANGE_PW", "User changed their own password")
        return jsonify({"ok": True})
    finally:
        conn.close()

@app.route("/logout")
def logout():
    username=session.get("username","")
    if username: log_activity(username,"LOGOUT","")
    session.clear()
    return redirect(url_for("login_page"))

@app.route("/")
@login_required
def dashboard():
    conn = get_db()
    row  = conn.execute("SELECT last_login FROM users WHERE username=?",
                        (session.get("username"),)).fetchone()
    conn.close()
    last_login = (row["last_login"] or "First login") if row else ""
    return render_template_string(DASHBOARD_HTML,
        name=session.get("name",""), username=session.get("username",""),
        role=session.get("role",""), last_login=last_login)

# ═══════════════════════════════════════════════════════════════════════════════
#  ROUTES — App pages (served as embedded SPAs)
# ═══════════════════════════════════════════════════════════════════════════════

@app.route("/app/ppt-rejection")
@login_required
def ppt_rejection_page():
    log_activity(session["username"],"OPEN_APP","PPT Alignment / Rejection Filter")
    return render_template_string(PPT_REJECTION_HTML)

@app.route("/app/photo-merge")
@login_required
def photo_merge_page():
    log_activity(session["username"],"OPEN_APP","Front and Top Photo Merge")
    return render_template_string(PHOTO_MERGE_HTML)

@app.route("/app/excel-rejection")
@login_required
def excel_rejection_page():
    log_activity(session["username"],"OPEN_APP","Excel Rejection Report")
    return render_template_string(EXCEL_REJECTION_HTML)

@app.route("/app/daily-report")
@login_required
def daily_report_page():
    log_activity(session["username"],"OPEN_APP","Daily Report")
    return render_template_string(DAILY_REPORT_HTML)

@app.route("/app/cga-report")
@login_required
def cga_report_page():
    log_activity(session["username"],"OPEN_APP","Daily Report / Slack CGA Report")
    return render_template_string(CGA_REPORT_HTML)

@app.route("/app/hourly-report")
@login_required
def hourly_report_page():
    log_activity(session["username"],"OPEN_APP","Hourly Report")
    return render_template_string(HOURLY_REPORT_HTML)

# ═══════════════════════════════════════════════════════════════════════════════
#  ROUTES — CGA Report API
# ═══════════════════════════════════════════════════════════════════════════════

_CGA_VALID_SITES = ["TMSGPKKL154", "TMSGPKTC173", "TMSGPPKJ001"]

_CGA_HOUR_BUCKETS = []
for _h in range(7, 24):
    _CGA_HOUR_BUCKETS.append((_h, f"{_h:02d}:00 - {(_h+1)%24:02d}:00"))
for _h in range(0, 5):
    _CGA_HOUR_BUCKETS.append((_h, f"{_h:02d}:00 - {_h+1:02d}:00"))

_CGA_HOUR_TO_LABEL = {h: lbl for h, lbl in _CGA_HOUR_BUCKETS}
# Hour 5 (5:01-5:59 AM) is mapped to the 04:00-05:00 bucket so late arrivals
# are still counted rather than excluded from the operational window.
_CGA_HOUR_TO_LABEL[5] = "04:00 - 05:00"
_CGA_BUCKET_ORDER  = [lbl for _, lbl in _CGA_HOUR_BUCKETS]


def _cga_process(df, selected_date, selected_sites=None):
    """Filter, validate and bucket the uploaded dataframe for CGA pivot.
    selected_sites: list of site codes to filter by. Defaults to _CGA_VALID_SITES.
    """
    from datetime import timedelta as _td
    df = df.copy()
    df.columns = [c.strip().upper() for c in df.columns]

    for col in ("SITE CODE", "ACCEPTED", "WB IN TIME"):
        if col not in df.columns:
            return None, f"Required column '{col}' not found."

    # Use provided sites or fall back to defaults
    sites_to_use = selected_sites if selected_sites else _CGA_VALID_SITES
    df = df[df["SITE CODE"].isin(sites_to_use)].copy()
    df = df[df["ACCEPTED"].astype(str).str.strip().str.upper() == "YES"].copy()
    if df.empty:
        return None, "No rows remain after site/accepted filtering."

    # Parse WB IN TIME — try multiple formats
    wbin = None
    for fmt in ("%d-%m-%Y %H:%M:%S", "%d/%m/%Y %H:%M:%S",
                "%d-%m-%Y %H:%M",    "%d/%m/%Y %H:%M",
                "%Y-%m-%d %H:%M:%S", "%Y-%m-%d %H:%M"):
        try:
            wbin = pd.to_datetime(df["WB IN TIME"], format=fmt, errors="raise")
            break
        except Exception:
            pass
    if wbin is None:
        wbin = pd.to_datetime(df["WB IN TIME"], errors="coerce")
    df["_wbin"] = wbin
    df = df.dropna(subset=["_wbin"]).copy()

    op_start = datetime.combine(selected_date, datetime.min.time()).replace(hour=7)
    # Extend window to 06:00 so entries arriving at 5:01-5:59 AM are captured
    # and placed into the "04:00-05:00" bucket via _CGA_HOUR_TO_LABEL[5].
    op_end   = datetime.combine(selected_date + _td(days=1),
                                datetime.min.time()).replace(hour=6)
    df = df[(df["_wbin"] >= op_start) & (df["_wbin"] < op_end)].copy()
    if df.empty:
        return None, (f"No data in window "
                      f"{op_start.strftime('%d-%b-%Y %H:%M')} → "
                      f"{op_end.strftime('%d-%b-%Y %H:%M')}.")

    df["_hour"]  = df["_wbin"].dt.hour
    df["Bucket"] = df["_hour"].map(_CGA_HOUR_TO_LABEL)
    return df, None


def _cga_build_pivot(df, selected_sites=None):
    """Build the hourly trucks pivot from the processed dataframe.
    selected_sites: ordered list of site codes to include. Defaults to _CGA_VALID_SITES.
    """
    sites_to_use = selected_sites if selected_sites else _CGA_VALID_SITES
    pivot = (df.groupby(["Bucket", "SITE CODE"])
               .size()
               .unstack(fill_value=0))
    for s in sites_to_use:
        if s not in pivot.columns:
            pivot[s] = 0
    available_sites = [s for s in sites_to_use if s in pivot.columns]
    pivot = pivot[available_sites]
    pivot["Grand Total"] = pivot.sum(axis=1)

    present = [b for b in _CGA_BUCKET_ORDER if b in pivot.index]
    pivot   = pivot.loc[present]

    total = pivot.sum(axis=0).rename("Grand Total")
    pivot = pd.concat([pivot, total.to_frame().T])
    pivot.index.name = "Row Labels"
    return pivot


def _cga_render_png(pivot, report_title="CGA"):
    """Render pivot table to PNG bytes using matplotlib — no ax.table()."""
    # matplotlib backend is already set to Agg at import time via ct_module;
    # calling matplotlib.use() again after first import silently fails on some
    # versions — so we import directly without re-setting the backend.
    import matplotlib.pyplot as _plt
    import matplotlib.patches as _mp

    col_labels = list(pivot.columns)
    row_labels  = list(pivot.index)
    n_cols      = len(col_labels) + 1

    col_widths  = [2.30] + [1.65] * len(col_labels)
    row_height  = 0.38
    total_w     = sum(col_widths)
    total_rows  = 1 + 1 + len(row_labels)    # title + header + data
    total_h     = row_height * total_rows
    fig_w       = total_w + 0.30
    fig_h       = total_h + 0.20

    fig, ax = _plt.subplots(figsize=(fig_w, fig_h))
    ax.set_xlim(0, total_w)
    ax.set_ylim(0, total_h)
    ax.axis("off")
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")

    col_x = [sum(col_widths[:i]) for i in range(n_cols)]
    BLK   = (0, 0, 0)
    WHT   = (1, 1, 1)
    BRD   = 1.2

    def _cell(x, y, w, h, bg, txt, fg, bold=False, fs=14.78):
        ax.add_patch(_mp.Rectangle((x, y), w, h,
                     linewidth=BRD, edgecolor=BLK, facecolor=bg, zorder=2))
        ax.text(x + w/2, y + h/2, txt, ha="center", va="center",
                fontsize=fs, fontweight="bold" if bold else "normal",
                color=fg, zorder=3, clip_on=False)

    # Title row
    ty = total_h - row_height
    _cell(0, ty, total_w, row_height, WHT, report_title, BLK, bold=True, fs=15.78)

    # Header row
    hy = ty - row_height
    _cell(col_x[0], hy, col_widths[0], row_height, WHT, "Row Labels", BLK, bold=True)
    for ci, lbl in enumerate(col_labels):
        _cell(col_x[ci+1], hy, col_widths[ci+1], row_height, WHT, lbl, BLK, bold=True)

    # Data rows
    for ri, rlbl in enumerate(row_labels):
        ry     = hy - (ri+1) * row_height
        is_gt  = (str(rlbl) == "Grand Total")
        bg     = WHT
        fg     = BLK
        _cell(col_x[0], ry, col_widths[0], row_height, bg, str(rlbl), fg, bold=is_gt)
        for ci, clbl in enumerate(col_labels):
            val = pivot.loc[rlbl, clbl]
            _cell(col_x[ci+1], ry, col_widths[ci+1], row_height,
                  bg, str(int(val)), fg, bold=is_gt)

    # Outer border
    ax.add_patch(_mp.Rectangle((0, 0), total_w, total_h,
                 linewidth=BRD*1.5, edgecolor=BLK, facecolor="none", zorder=4))

    _plt.subplots_adjust(left=0.01, right=0.99, top=0.99, bottom=0.01)
    buf = io.BytesIO()
    _plt.savefig(buf, format="png", dpi=180, bbox_inches="tight",
                 pad_inches=0.08, facecolor="white", edgecolor="none")
    _plt.close(fig)
    gc.collect()
    buf.seek(0)
    return buf.read()


@app.route("/api/cga/sources", methods=["POST"])
@login_required
def cga_get_sources():
    """Extract all available Source IDs and auto-detect operational date from the uploaded file."""
    try:
        f = request.files.get("file")
        if not f:
            return jsonify({"error": "No file uploaded"}), 400
        raw = io.BytesIO(f.read())
        try:
            df_raw = pd.read_excel(raw)
        except Exception:
            for enc in ("utf-8", "latin1", "cp1252"):
                try:
                    raw.seek(0); df_raw = pd.read_csv(raw, encoding=enc); break
                except Exception:
                    continue
            else:
                return jsonify({"error": "Could not read file"}), 400
        df_raw.columns = [c.strip().upper() for c in df_raw.columns]
        if "SITE CODE" not in df_raw.columns:
            return jsonify({"error": "'SITE CODE' column not found"}), 400
        all_sites = sorted(df_raw["SITE CODE"].dropna().astype(str).str.strip().unique().tolist())

        # Auto-detect operational date from WB IN TIME column
        auto_date = None
        if "WB IN TIME" in df_raw.columns:
            try:
                wbin = pd.to_datetime(df_raw["WB IN TIME"], errors="coerce")
                valid = wbin.dropna()
                if not valid.empty:
                    # Use the most common date (mode)
                    date_counts = valid.dt.date.value_counts()
                    if not date_counts.empty:
                        auto_date = str(date_counts.index[0])  # YYYY-MM-DD
            except Exception:
                pass

        return jsonify({
            "ok": True,
            "all_sources": all_sites,
            "default_sources": _CGA_VALID_SITES,
            "auto_date": auto_date,
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/cga/process", methods=["POST"])
@login_required
def cga_process():
    """Parse uploaded file, return pivot table JSON."""
    try:
        f = request.files.get("file")
        if not f:
            return jsonify({"error": "No file uploaded"}), 400
        date_str = request.form.get("op_date", "")
        selected_sources_raw = request.form.get("selected_sources", "")
        try:
            sel_date = datetime.strptime(date_str, "%Y-%m-%d").date()
        except ValueError:
            return jsonify({"error": "Invalid date format (expected YYYY-MM-DD)"}), 400

        selected_sites = None
        if selected_sources_raw.strip():
            import json as _json_mod
            try:
                parsed = _json_mod.loads(selected_sources_raw)
                if isinstance(parsed, list) and parsed:
                    selected_sites = [str(s).strip() for s in parsed if str(s).strip()]
            except Exception:
                pass

        raw = io.BytesIO(f.read())
        try:
            df_raw = pd.read_excel(raw)
        except Exception:
            for enc in ("utf-8", "latin1", "cp1252"):
                try:
                    raw.seek(0); df_raw = pd.read_csv(raw, encoding=enc); break
                except Exception:
                    continue
            else:
                return jsonify({"error": "Could not read file"}), 400

        # ── Pivot for selected sites (unchanged — drives table + PNG) ──────────
        df, err = _cga_process(df_raw, sel_date, selected_sites=selected_sites)
        if err:
            return jsonify({"error": err}), 400

        pivot = _cga_build_pivot(df, selected_sites=selected_sites)

        total_loads = int(pivot.loc["Grand Total", "Grand Total"])
        mat_col = next((c for c in df.columns if "MATERIAL" in c.upper()), None)
        mat_counts = df[mat_col].astype(str).str.strip().str.upper().value_counts() if mat_col else pd.Series(dtype=int)
        soft       = int(mat_counts.get("SOFT CLAY",  mat_counts.get("SOFT", 0)))
        hard       = int(mat_counts.get("HARD CLAY",  mat_counts.get("HARD", 0)))
        good_earth = int(mat_counts.get("GOOD EARTH", mat_counts.get("GOODEARTH",
                     mat_counts.get("GOOD EARTH CLAY", 0))))
        mat_breakdown = {k: int(v) for k, v in mat_counts.items()} if mat_col else {}
        piv_json = pivot.reset_index().to_dict(orient="records")

        # ── All-sites totals for the side panel ────────────────────────────────
        # Re-run with EVERY site in the uploaded file so the panel shows
        # the full operational picture, not just the 3 selected pivot sites.
        try:
            df_raw_col = [c.strip().upper() for c in df_raw.columns]
            site_col = next((c for c in df_raw.columns if c.strip().upper() == "SITE CODE"), None)
            all_sites_in_file = (
                df_raw[site_col].dropna().astype(str).str.strip().unique().tolist()
                if site_col else []
            )
            df_all, _err_all = _cga_process(
                df_raw.copy(), sel_date,
                selected_sites=all_sites_in_file if all_sites_in_file else None
            )
            if df_all is not None and not df_all.empty:
                mc_all = next((c for c in df_all.columns if "MATERIAL" in c.upper()), None)
                mc_counts = (df_all[mc_all].astype(str).str.strip().str.upper().value_counts()
                             if mc_all else pd.Series(dtype=int))
                panel_total = int(len(df_all))
                panel_soft  = int(mc_counts.get("SOFT CLAY",  mc_counts.get("SOFT", 0)))
                panel_ge    = int(mc_counts.get("GOOD EARTH", mc_counts.get("GOODEARTH",
                              mc_counts.get("GOOD EARTH CLAY", 0))))
                panel_bd    = {k: int(v) for k, v in mc_counts.items()}
            else:
                panel_total = total_loads
                panel_soft  = soft
                panel_ge    = good_earth
                panel_bd    = mat_breakdown
        except Exception:
            panel_total = total_loads
            panel_soft  = soft
            panel_ge    = good_earth
            panel_bd    = mat_breakdown

        log_activity(session["username"], "CGA_PROCESS",
                     f"date={sel_date} rows={len(df)} total={total_loads} panel_total={panel_total}")
        df_raw = None; df = None; gc.collect()

        return jsonify({
            "ok":             True,
            "total_loads":    total_loads,
            "all_sites_total":total_loads,
            "soft_loads":     soft,
            "hard_loads":     hard,
            "good_earth":     good_earth,
            "total_qty":      total_loads,
            "mat_breakdown":  mat_breakdown,
            # Panel fields: ALL sites in file, full accepted window
            "panel_total":    panel_total,
            "panel_soft":     panel_soft,
            "panel_ge":       panel_ge,
            "panel_breakdown":panel_bd,
            "pivot":          piv_json,
            "columns":        list(pivot.columns),
            "op_date":        date_str,
        })
    except Exception as e:
        return jsonify({"error": str(e), "trace": traceback.format_exc()}), 500


@app.route("/api/cga/download", methods=["POST"])
@login_required
def cga_download():
    """Re-render pivot as PNG and stream for download."""
    try:
        data   = request.get_json(force=True, silent=True) or {}
        pivot_records = data.get("pivot", [])
        columns       = data.get("columns", [])
        op_date       = data.get("op_date", "")
        if not pivot_records or not columns:
            return jsonify({"error": "No pivot data"}), 400

        pivot = pd.DataFrame(pivot_records).set_index("Row Labels")
        pivot = pivot[columns]
        for c in pivot.columns:
            pivot[c] = pd.to_numeric(pivot[c], errors="coerce").fillna(0).astype(int)

        png = _cga_render_png(pivot, report_title="CGA")
        pivot = None; gc.collect()

        try:
            date_label = datetime.strptime(op_date, "%Y-%m-%d").strftime("%d%b%Y")
        except Exception:
            date_label = op_date.replace("-", "")
        fname = f"CGA_Pivot_{date_label}.png"

        resp = make_response(png)
        png  = None
        resp.headers["Content-Type"]        = "image/png"
        resp.headers["Content-Disposition"] = f'attachment; filename="{fname}"'
        return resp
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ── CGA Report HTML ──────────────────────────────────────────────────────────
CGA_REPORT_HTML = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>CGA Report — APSG</title>
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800&family=Poppins:wght@700;800&display=swap" rel="stylesheet">
<style>
*{box-sizing:border-box;margin:0;padding:0;}
:root{--bg:#0A0F1E;--s1:#111827;--s2:#1C2333;--s3:#232D42;
  --bdr:rgba(255,255,255,.08);--t1:#F1F5F9;--t2:#94A3B8;--mu:#64748B;
  --acc:#2DD4BF;--ind:#6366F1;--ind-l:#818CF8;--grn:#10B981;--r:14px;}
html,body{background:url('/static/bg.jpg') center/cover fixed #08101E;}
body::before{content:'';position:fixed;inset:0;z-index:0;background:rgba(3,7,18,.52);pointer-events:none;}
body{font-family:'Inter',sans-serif;color:var(--t1);min-height:100vh;padding-bottom:3rem;}
body>*{position:relative;z-index:1;}

/* top-bar */
.topbar{position:sticky;top:0;z-index:200;height:44px;display:flex;align-items:center;
  padding:0 1.2rem;gap:.65rem;background:rgba(6,10,28,.65);backdrop-filter:blur(18px);
  border-bottom:1px solid var(--bdr);}
.topbar .brand{font-family:'Poppins',sans-serif;font-size:.8rem;font-weight:800;color:var(--ind-l);}
.topbar .sep{width:1px;height:18px;background:var(--bdr);}
.topbar .lbl{font-size:.76rem;font-weight:600;color:var(--t2);}
.topbar .sp{flex:1;}
.back{background:rgba(99,102,241,.18);border:1px solid rgba(99,102,241,.35);border-radius:7px;
  padding:.28rem .8rem;font-size:.7rem;font-weight:600;color:#C5D5FF;text-decoration:none;transition:.2s;}
.back:hover{background:rgba(99,102,241,.32);}

/* layout */
.wrap{max-width:820px;margin:0 auto;padding:1.4rem 1rem 4rem;}
.panel{background:rgba(8,14,38,.75);border:1px solid var(--bdr);border-radius:var(--r);
  padding:1.4rem;margin-bottom:1.1rem;backdrop-filter:blur(14px);}
.ptitle{font-size:.7rem;font-weight:700;color:var(--acc);text-transform:uppercase;
  letter-spacing:.1em;margin-bottom:.85rem;}

/* upload zone */
.dz{border:2px dashed rgba(45,212,191,.35);border-radius:11px;padding:1.6rem 1rem;
  text-align:center;cursor:pointer;transition:.22s;position:relative;
  background:rgba(45,212,191,.03);user-select:none;}
.dz:hover,.dz.over{border-color:var(--acc);background:rgba(45,212,191,.09);}
.dz.has{border-color:var(--acc);border-style:solid;background:rgba(45,212,191,.06);}
.dz .ico{font-size:2rem;display:block;margin-bottom:.3rem;}
.dz .hint{font-size:.72rem;color:var(--mu);margin-top:.25rem;}
.dz .fn{font-size:.78rem;color:var(--acc);font-weight:700;margin-top:.35rem;}

/* grid row */
.row2{display:grid;grid-template-columns:1fr 185px;gap:1rem;align-items:start;}
@media(max-width:540px){.row2{grid-template-columns:1fr;}}
.fgl label{display:block;font-size:.72rem;font-weight:600;color:var(--t2);margin-bottom:.38rem;}

/* date input */
.di{width:100%;padding:.62rem .85rem;background:rgba(5,9,28,.85);border:1.5px solid rgba(99,102,241,.4);
  border-radius:9px;color:var(--t1);font-size:.87rem;font-family:'Inter',sans-serif;color-scheme:dark;}
.di:focus{outline:none;border-color:var(--acc);}

/* source selector */
.ss-wrap{margin-top:.9rem;}
.ss-search{width:100%;padding:.55rem .8rem;background:rgba(5,9,28,.85);
  border:1.5px solid rgba(99,102,241,.4);border-radius:8px;color:var(--t1);
  font-size:.83rem;font-family:'Inter',sans-serif;margin-bottom:.4rem;}
.ss-search:focus{outline:none;border-color:var(--acc);}
.ss-list{max-height:140px;overflow-y:auto;border:1px solid var(--bdr);border-radius:8px;
  background:rgba(5,9,28,.7);scrollbar-width:thin;}
.ss-item{display:flex;align-items:center;gap:.5rem;padding:.38rem .7rem;cursor:pointer;
  border-bottom:1px solid rgba(255,255,255,.04);transition:.15s;}
.ss-item:hover{background:rgba(45,212,191,.07);}
.ss-item:last-child{border-bottom:none;}
.ss-item input[type=checkbox]{accent-color:var(--acc);width:14px;height:14px;cursor:pointer;}
.ss-item .sn{font-size:.74rem;color:var(--t1);font-family:monospace;}
.ss-item .def{font-size:.62rem;color:var(--acc);background:rgba(45,212,191,.12);
  border:1px solid rgba(45,212,191,.2);border-radius:4px;padding:.05rem .3rem;margin-left:auto;}
.ss-actions{display:flex;gap:.4rem;margin-top:.4rem;}
.ss-btn{font-size:.68rem;padding:.22rem .65rem;background:rgba(5,9,28,.7);
  border:1px solid var(--bdr);border-radius:5px;color:var(--acc);cursor:pointer;transition:.15s;}
.ss-btn:hover{background:rgba(45,212,191,.12);}
.ss-btn.red{color:#F87171;}
.ss-count{font-size:.7rem;color:var(--mu);margin-left:auto;align-self:center;}
.no-file-chips{display:flex;flex-wrap:wrap;gap:.4rem;}
.chip{background:rgba(5,9,28,.7);border:1px solid rgba(45,212,191,.25);border-radius:6px;
  padding:.22rem .6rem;font-size:.72rem;font-weight:600;color:var(--acc);}

/* generate btn */
.btn-gen{width:100%;padding:.9rem;background:linear-gradient(135deg,#2DD4BF,#06B6D4);
  border:none;border-radius:10px;color:#0A0F1E;font-size:1rem;font-weight:800;
  cursor:pointer;transition:.2s;margin-top:.65rem;}
.btn-gen:hover{opacity:.9;transform:translateY(-1px);}
.btn-gen:disabled{opacity:.38;cursor:not-allowed;transform:none;}

/* err */
.err{background:rgba(239,68,68,.09);border:1px solid rgba(239,68,68,.3);border-radius:8px;
  padding:.65rem 1rem;color:#FCA5A5;font-size:.82rem;display:none;margin-top:.55rem;}

/* spinner */
.loading{display:none;text-align:center;padding:2rem;color:var(--t2);}
.spin{width:34px;height:34px;border:3px solid rgba(45,212,191,.2);border-top-color:var(--acc);
  border-radius:50%;animation:sp .7s linear infinite;margin:0 auto .7rem;}
@keyframes sp{to{transform:rotate(360deg);}}

/* metrics */
.metrics{display:grid;grid-template-columns:repeat(3,1fr);gap:.75rem;margin-bottom:1.1rem;}
@media(max-width:480px){.metrics{grid-template-columns:1fr 1fr;}}
.met{background:rgba(8,14,38,.75);border:1px solid var(--bdr);border-radius:10px;
  padding:.85rem 1rem;text-align:center;}
.met-v{font-size:1.65rem;font-weight:800;color:var(--acc);}
.met-l{font-size:.68rem;color:var(--mu);font-weight:600;margin-top:.15rem;}

/* pivot */
.pscroll{overflow-x:auto;margin-top:.65rem;}
table.piv{border-collapse:collapse;width:100%;font-size:.77rem;}
table.piv th,table.piv td{border:1px solid rgba(255,255,255,.09);padding:.42rem .7rem;
  text-align:center;white-space:nowrap;}
table.piv th{background:var(--s3);font-weight:700;color:var(--t1);}
table.piv td{color:var(--t2);}
table.piv td:first-child{text-align:left;color:var(--t1);}
table.piv tr.grand td{background:var(--s2);color:var(--acc);font-weight:700;}

/* download btn */
.btn-dl{width:100%;padding:.82rem;background:rgba(5,9,28,.8);border:1px solid rgba(45,212,191,.35);
  border-radius:10px;color:var(--acc);font-size:.9rem;font-weight:700;cursor:pointer;
  transition:.2s;margin-top:.75rem;display:none;}
.btn-dl:hover{background:rgba(45,212,191,.1);}

/* preview */
.prev-wrap{display:none;}
.prev-wrap img{width:100%;border-radius:8px;border:1px solid var(--bdr);}

/* date auto badge */
.abadge{display:inline-flex;align-items:center;gap:.35rem;padding:.22rem .65rem;
  background:rgba(45,212,191,.12);border:1px solid rgba(45,212,191,.28);border-radius:6px;
  font-size:.72rem;color:var(--acc);font-weight:600;margin-top:.4rem;}

/* material dashboard panel — large, high-visibility */
#matPanel{display:none;position:fixed;top:54px;right:0;z-index:900;
  width:310px;background:rgba(7,12,32,.97);
  border:1.5px solid rgba(45,212,191,.4);border-left:3px solid var(--acc);
  border-radius:0 0 0 16px;
  box-shadow:-6px 6px 40px rgba(0,0,0,.65);overflow:hidden;}
.mp-head{background:linear-gradient(135deg,rgba(45,212,191,.2),rgba(6,182,212,.1));
  padding:.85rem 1.1rem .75rem;border-bottom:1px solid rgba(45,212,191,.18);}
.mp-label{font-size:.68rem;font-weight:800;color:var(--acc);letter-spacing:.12em;
  text-transform:uppercase;margin-bottom:.55rem;display:flex;align-items:center;gap:.4rem;}
.mp-sub{font-size:.62rem;color:rgba(148,163,184,.7);font-weight:400;text-transform:none;
  letter-spacing:0;margin-left:auto;}
.mp-stats{display:grid;grid-template-columns:1fr 1fr 1fr;gap:.45rem;}
.mp-stat{background:rgba(45,212,191,.08);border:1px solid rgba(45,212,191,.22);border-radius:9px;
  padding:.5rem .55rem;text-align:center;}
.mp-stat-n{font-size:1.45rem;font-weight:900;color:var(--acc);line-height:1;}
.mp-stat-l{font-size:.6rem;color:var(--mu);margin-top:.18rem;font-weight:600;}
.mp-stat.ge{background:rgba(245,158,11,.1);border-color:rgba(245,158,11,.25);}
.mp-stat.ge .mp-stat-n{color:#FBBF24;}
.mp-stat.hard{background:rgba(129,140,248,.1);border-color:rgba(129,140,248,.25);}
.mp-stat.hard .mp-stat-n{color:#A5B4FC;}
.mp-divider{height:1px;background:rgba(45,212,191,.12);margin:.1rem 0;}
.mp-rows{padding:.55rem .95rem;max-height:200px;overflow-y:auto;scrollbar-width:thin;}
.mp-row{display:flex;align-items:center;gap:.45rem;margin-bottom:.32rem;}
.mp-dot{width:9px;height:9px;border-radius:50%;flex-shrink:0;}
.mp-name{font-size:.73rem;color:#E2E8F0;flex:1;white-space:nowrap;overflow:hidden;text-overflow:ellipsis;}
.mp-cnt{font-size:.73rem;font-weight:800;min-width:30px;text-align:right;}
.mp-pct{font-size:.65rem;color:var(--mu);min-width:32px;text-align:right;}
.mp-foot{padding:.3rem .95rem .5rem;display:flex;justify-content:space-between;align-items:center;
  border-top:1px solid rgba(255,255,255,.05);}
.mp-total-lbl{font-size:.68rem;color:var(--t2);font-weight:600;}
.mp-total-n{font-size:1rem;font-weight:900;color:#fff;}
.mp-close{font-size:.68rem;color:var(--mu);background:none;border:none;cursor:pointer;
  padding:.15rem .4rem;border-radius:4px;transition:.15s;}
.mp-close:hover{background:rgba(255,255,255,.08);color:var(--t2);}
.mp-copy-btn{font-size:.7rem;color:var(--acc);background:rgba(45,212,191,.1);
  border:1px solid rgba(45,212,191,.3);cursor:pointer;padding:.25rem .75rem;
  border-radius:6px;font-weight:700;transition:.18s;white-space:nowrap;}
.mp-copy-btn:hover{background:rgba(45,212,191,.22);}
.mp-copy-ok{font-size:.66rem;font-weight:700;color:#34D399;display:none;
  padding:.15rem .5rem;background:rgba(16,185,129,.12);border-radius:4px;
  border:1px solid rgba(16,185,129,.25);}
</style>
</head>
<body>

<div class="topbar">
  <span class="brand">APSG</span>
  <div class="sep"></div>
  <span class="lbl">🚛 CGA Report</span>
  <div class="sp"></div>
  <span style="font-size:.65rem;color:var(--ind-l);font-weight:700;margin-right:.35rem;">✦ Karthi</span>
  <a href="/" class="back">← Dashboard</a>
</div>

<div class="wrap">

  <div class="panel">
    <div class="ptitle">📁 Upload Source File</div>
    <div style="background:rgba(45,212,191,.07);border-left:3px solid var(--acc);
      border-radius:8px;padding:.65rem 1rem;font-size:.77rem;color:var(--t2);margin-bottom:.9rem;">
      Window: <strong>Selected date 07:00 → next day 05:00</strong> &nbsp;|&nbsp;
      Timestamp column: <strong>WB IN TIME</strong> &nbsp;|&nbsp; Only <strong>Accepted</strong> records counted
    </div>

    <div class="row2">
      <div class="fgl">
        <label>Source File (Excel or CSV)</label>
        <!-- Upload zone – click OR drag-drop, both work -->
        <div class="dz" id="dz">
          <input type="file" id="fi" accept=".xlsx,.xls,.csv" style="display:none">
          <span class="ico">📄</span>
          <strong style="font-size:.84rem;">Click to browse or drag &amp; drop</strong>
          <div class="hint">.xlsx · .xls · .csv</div>
          <div class="fn" id="fn"></div>
          <div id="autoBadge"></div>
        </div>
      </div>
      <div class="fgl">
        <label>Operational Date</label>
        <input type="date" class="di" id="opDate">
      </div>
    </div>

    <!-- Source site selector -->
    <div class="ss-wrap">
      <div style="display:flex;align-items:center;gap:.5rem;margin-bottom:.38rem;">
        <span style="font-size:.72rem;font-weight:600;color:var(--mu);">Source Site Selection</span>
        <span class="ss-count" id="ssCount"></span>
      </div>
      <!-- shown before file upload -->
      <div id="defaultChips" class="no-file-chips">
        <span class="chip">TMSGPKKL154</span>
        <span class="chip">TMSGPKTC173</span>
        <span class="chip">TMSGPPKJ001</span>
      </div>
      <!-- shown after file upload -->
      <div id="ssSelector" style="display:none;">
        <input type="text" class="ss-search" id="ssSearch" placeholder="🔍  Search source sites…"
               oninput="filterSites()">
        <div class="ss-list" id="ssList"></div>
        <div class="ss-actions">
          <button class="ss-btn" onclick="selAll()">Select All</button>
          <button class="ss-btn" onclick="selDefaults()">★ Defaults</button>
          <button class="ss-btn red" onclick="clearSites()">Clear</button>
        </div>
      </div>
      <div id="ssLoading" style="display:none;font-size:.72rem;color:var(--mu);padding:.4rem 0;">
        ⏳ Loading source IDs from file…
      </div>
    </div>

    <button class="btn-gen" id="btnGen" onclick="generate()" disabled>🚀 Generate Pivot Table</button>
    <div class="err" id="errBox"></div>
  </div>

  <!-- Loading -->
  <div class="loading" id="loadDiv">
    <div class="spin"></div>
    Processing file and building pivot table…
  </div>

  <!-- Metrics -->
  <div class="metrics" id="metricsDiv" style="display:none;">
    <div class="met"><div class="met-v" id="mTotal">—</div><div class="met-l">Total Loads</div></div>
    <div class="met"><div class="met-v" id="mSoft">—</div><div class="met-l">Soft Clay</div></div>
    <div class="met"><div class="met-v" id="mHard">—</div><div class="met-l">Hard Clay</div></div>
  </div>

  <!-- Pivot table -->
  <div class="panel" id="pivPanel" style="display:none;">
    <div class="ptitle">📊 Pivot Table Preview</div>
    <div class="pscroll"><table class="piv" id="pivTable"></table></div>
    <button class="btn-dl" id="btnDl" onclick="dlPng()">⬇ Download as PNG Image</button>
  </div>

  <!-- PNG preview -->
  <div class="panel prev-wrap" id="prevPanel">
    <div class="ptitle">🖼 Image Preview</div>
    <img id="prevImg" src="" alt="pivot">
  </div>

</div><!-- /wrap -->

<!-- Material Dashboard Panel — large side panel -->
<div id="matPanel">
  <div class="mp-head">
    <div class="mp-label">
      📊 Accepted Load Summary
      <span class="mp-sub">7:00 AM → 5:00 AM</span>
    </div>
    <div class="mp-stats">
      <div class="mp-stat">
        <div class="mp-stat-n" id="mpTotal">—</div>
        <div class="mp-stat-l">Total Accepted</div>
      </div>
      <div class="mp-stat">
        <div class="mp-stat-n" id="mpSoft">—</div>
        <div class="mp-stat-l">Soft Clay</div>
      </div>
      <div class="mp-stat ge">
        <div class="mp-stat-n" id="mpGE">—</div>
        <div class="mp-stat-l">GoodEarth</div>
      </div>
    </div>
    <div style="margin-top:.45rem;font-size:.64rem;color:rgba(148,163,184,.6);">
      ✦ All selected source sites combined &nbsp;·&nbsp; Rejected entries excluded
    </div>
  </div>
  <div class="mp-divider"></div>
  <div class="mp-rows" id="mpRows"></div>
  <div class="mp-foot">
    <span class="mp-total-lbl">Grand Total &nbsp;<span class="mp-total-n" id="mpGrandTotal">—</span> loads</span>
    <div style="display:flex;align-items:center;gap:.4rem;">
      <button class="mp-close" onclick="document.getElementById('matPanel').style.display='none'">✕ hide</button>
    </div>
  </div>
</div>

<script>
// ── State ────────────────────────────────────────────────────────────────────
var _file=null, _pivData=null, _pivCols=null, _opDate=null, _blob=null;
var _allSites=[], _defSites=[], _selSites=[];

// ── Set today's date ─────────────────────────────────────────────────────────
document.getElementById('opDate').value = new Date().toISOString().slice(0,10);

// ── Upload zone: click ────────────────────────────────────────────────────────
var dz = document.getElementById('dz');
var fi = document.getElementById('fi');

dz.addEventListener('click', function(){ fi.click(); });
fi.addEventListener('change', function(){ if(fi.files[0]) applyFile(fi.files[0]); });

// ── Upload zone: drag & drop ──────────────────────────────────────────────────
dz.addEventListener('dragenter', function(e){ e.preventDefault(); e.stopPropagation(); dz.classList.add('over'); });
dz.addEventListener('dragover',  function(e){ e.preventDefault(); e.stopPropagation(); dz.classList.add('over'); });
dz.addEventListener('dragleave', function(e){
  e.stopPropagation();
  if(!dz.contains(e.relatedTarget)){ dz.classList.remove('over'); }
});
dz.addEventListener('drop', function(e){
  e.preventDefault(); e.stopPropagation(); dz.classList.remove('over');
  var f = e.dataTransfer && e.dataTransfer.files && e.dataTransfer.files[0];
  if(f) applyFile(f);
});

// ── Apply chosen file ─────────────────────────────────────────────────────────
function applyFile(f){
  _file = f;
  document.getElementById('fn').textContent = '✓ ' + f.name;
  dz.classList.add('has');
  checkReady();
  loadSourceIds(f);
}

// ── Load source IDs ───────────────────────────────────────────────────────────
async function loadSourceIds(f){
  document.getElementById('ssLoading').style.display='block';
  document.getElementById('defaultChips').style.display='none';
  document.getElementById('ssSelector').style.display='none';
  document.getElementById('autoBadge').innerHTML='';

  var fd=new FormData(); fd.append('file',f);
  try{
    var r=await fetch('/api/cga/sources',{method:'POST',body:fd});
    var j=await r.json();
    document.getElementById('ssLoading').style.display='none';
    if(!j.ok){ document.getElementById('defaultChips').style.display='flex'; return; }
    _allSites=j.all_sources||[];
    _defSites=j.default_sources||[];
    _selSites=_defSites.slice();

    if(j.auto_date){
      document.getElementById('opDate').value=j.auto_date;
      checkReady();
      document.getElementById('autoBadge').innerHTML=
        '<span class="abadge">📅 Date auto-detected: '+escH(j.auto_date)+'</span>';
    }
    renderSiteList('');
    document.getElementById('ssSelector').style.display='block';
    updateCount();
  }catch(ex){
    document.getElementById('ssLoading').style.display='none';
    document.getElementById('defaultChips').style.display='flex';
  }
}

function renderSiteList(q){
  var list=document.getElementById('ssList');
  list.innerHTML='';
  var lq=q.toLowerCase();
  var shown=0;
  _allSites.forEach(function(s){
    if(lq && s.toLowerCase().indexOf(lq)<0) return;
    shown++;
    var chk=_selSites.indexOf(s)>=0;
    var isDef=_defSites.indexOf(s)>=0;
    var d=document.createElement('div'); d.className='ss-item';
    var c=document.createElement('input'); c.type='checkbox'; c.value=s; c.checked=chk;
    c.addEventListener('change',function(){
      if(this.checked){ if(_selSites.indexOf(s)<0) _selSites.push(s); }
      else { _selSites=_selSites.filter(function(x){return x!==s;}); }
      updateCount();
    });
    var nm=document.createElement('span'); nm.className='sn'; nm.textContent=s;
    d.appendChild(c); d.appendChild(nm);
    if(isDef){ var b=document.createElement('span'); b.className='def'; b.textContent='★'; d.appendChild(b); }
    d.addEventListener('click',function(e){ if(e.target!==c) c.click(); });
    list.appendChild(d);
  });
  if(shown===0){
    list.innerHTML='<div style="padding:.6rem .9rem;font-size:.72rem;color:var(--mu);">No sites match search</div>';
  }
}

function filterSites(){ renderSiteList(document.getElementById('ssSearch').value); }
function selAll(){ _selSites=_allSites.slice(); renderSiteList(document.getElementById('ssSearch').value); updateCount(); }
function selDefaults(){ _selSites=_defSites.slice(); renderSiteList(document.getElementById('ssSearch').value); updateCount(); }
function clearSites(){ _selSites=[]; renderSiteList(document.getElementById('ssSearch').value); updateCount(); }
function updateCount(){
  document.getElementById('ssCount').textContent = _selSites.length>0 ? '('+_selSites.length+' selected)' : '';
}

function checkReady(){
  document.getElementById('btnGen').disabled=!(_file && document.getElementById('opDate').value);
}
document.getElementById('opDate').addEventListener('change',checkReady);

// ── Generate ─────────────────────────────────────────────────────────────────
function showErr(m){ var b=document.getElementById('errBox'); b.textContent='❌ '+m; b.style.display='block'; }
function clearErr(){ document.getElementById('errBox').style.display='none'; }

async function generate(){
  clearErr();
  document.getElementById('metricsDiv').style.display='none';
  document.getElementById('pivPanel').style.display='none';
  document.getElementById('prevPanel').style.display='none';
  document.getElementById('matPanel').style.display='none';
  document.getElementById('loadDiv').style.display='block';
  document.getElementById('btnGen').disabled=true;

  var fd=new FormData();
  fd.append('file',_file);
  fd.append('op_date',document.getElementById('opDate').value);
  var sites = _selSites.length>0 ? _selSites : _defSites;
  fd.append('selected_sources',JSON.stringify(sites));

  try{
    var r=await fetch('/api/cga/process',{method:'POST',body:fd});
    var j=await r.json();
    document.getElementById('loadDiv').style.display='none';
    document.getElementById('btnGen').disabled=false;
    if(!r.ok||j.error){ showErr(j.error||'Server error'); return; }

    _pivData=j.pivot; _pivCols=j.columns; _opDate=j.op_date;

    document.getElementById('mTotal').textContent=(j.total_loads||0).toLocaleString();
    document.getElementById('mSoft').textContent=(j.soft_loads||0).toLocaleString();
    document.getElementById('mHard').textContent=(j.hard_loads||0).toLocaleString();
    document.getElementById('metricsDiv').style.display='grid';

    buildTable(j.pivot,j.columns);
    document.getElementById('pivPanel').style.display='block';
    document.getElementById('btnDl').style.display='block';

    showMatPanel(j);
    await previewPng();
  }catch(ex){
    document.getElementById('loadDiv').style.display='none';
    document.getElementById('btnGen').disabled=false;
    showErr('Network error: '+ex.message);
  }
}

function buildTable(rows,cols){
  var t=document.getElementById('pivTable'); t.innerHTML='';
  var h=t.insertRow(); var th0=document.createElement('th'); th0.textContent='Row Labels'; h.appendChild(th0);
  cols.forEach(function(c){ var th=document.createElement('th'); th.textContent=c; h.appendChild(th); });
  rows.forEach(function(row){
    var tr=t.insertRow();
    if(row['Row Labels']==='Grand Total') tr.className='grand';
    var td0=document.createElement('td'); td0.textContent=row['Row Labels']; tr.appendChild(td0);
    cols.forEach(function(c){ var td=document.createElement('td'); td.textContent=row[c]!==undefined?row[c]:''; tr.appendChild(td); });
  });
}

async function previewPng(){
  try{
    var r=await fetch('/api/cga/download',{method:'POST',
      headers:{'Content-Type':'application/json'},
      body:JSON.stringify({pivot:_pivData,columns:_pivCols,op_date:_opDate})});
    if(!r.ok){ var j=await r.json().catch(function(){return{};}); showErr(j.error||'PNG failed'); return; }
    var blob=await r.blob(); _blob=blob;
    document.getElementById('prevImg').src=URL.createObjectURL(blob);
    document.getElementById('prevPanel').style.display='block';
  }catch(ex){ showErr('PNG error: '+ex.message); }
}

async function dlPng(){
  if(_blob){
    var a=document.createElement('a');
    try{
      var d=new Date(_opDate);
      var mo=['Jan','Feb','Mar','Apr','May','Jun','Jul','Aug','Sep','Oct','Nov','Dec'];
      a.download='CGA_Pivot_'+String(d.getDate()).padStart(2,'0')+mo[d.getMonth()]+d.getFullYear()+'.png';
    }catch(e){ a.download='CGA_Pivot.png'; }
    a.href=URL.createObjectURL(_blob); a.click(); return;
  }
  var r=await fetch('/api/cga/download',{method:'POST',headers:{'Content-Type':'application/json'},
    body:JSON.stringify({pivot:_pivData,columns:_pivCols,op_date:_opDate})});
  if(!r.ok){ showErr('Download failed'); return; }
  var bl=await r.blob();
  var a=document.createElement('a'); a.href=URL.createObjectURL(bl); a.download='CGA_Pivot.png'; a.click();
}

// ── Material dashboard panel (large, all-sites combined) ─────────────────────
function showMatPanel(j){
  var bd     = j.mat_breakdown || {};
  // Use panel_* (all sites in file) for reference display
  var total  = j.panel_total || j.all_sites_total || j.total_loads || 0;
  var soft   = j.panel_soft  || j.soft_loads  || 0;
  var ge     = j.panel_ge    || j.good_earth  || 0;
  var bd     = j.panel_breakdown || j.mat_breakdown || {};

  document.getElementById('mpTotal').textContent      = total.toLocaleString();
  document.getElementById('mpSoft').textContent       = soft.toLocaleString();
  document.getElementById('mpGE').textContent         = ge.toLocaleString();
  document.getElementById('mpGrandTotal').textContent = total.toLocaleString();

  var sorted = Object.entries(bd).sort(function(a,b){return b[1]-a[1];});
  var html = sorted.map(function(entry){
    var m=entry[0], cnt=entry[1];
    var pct = total>0 ? Math.round(cnt/total*100) : 0;
    var up  = m.toUpperCase();
    var isSoft = up.indexOf('SOFT')>=0;
    var isGE   = up.replace(/\s+/g,'').indexOf('GOODEARTH')>=0 || up.indexOf('GOOD EARTH')>=0;
    var isHard = up.indexOf('HARD')>=0;
    var col = isGE?'#FBBF24' : isSoft?'#2DD4BF' : isHard?'#A5B4FC' : '#94A3B8';
    return '<div class="mp-row">'
      +'<div class="mp-dot" style="background:'+col+'"></div>'
      +'<div class="mp-name" title="'+escH(m)+'">'+escH(m)+'</div>'
      +'<div class="mp-cnt" style="color:'+col+'">'+cnt.toLocaleString()+'</div>'
      +'<div class="mp-pct">'+pct+'%</div>'
      +'</div>';
  }).join('');
  document.getElementById('mpRows').innerHTML =
    html || '<div style="color:var(--mu);font-size:.73rem;padding:.3rem 0;">No material breakdown available</div>';

  // Always show — it's the main dashboard widget
  document.getElementById('matPanel').style.display = 'block';
}

function escH(s){ return String(s||'').replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;'); }
</script>
<div style="position:fixed;bottom:0;left:0;right:0;z-index:9999;text-align:center;padding:.28rem 1rem;
  background:rgba(4,7,20,.72);backdrop-filter:blur(8px);border-top:1px solid rgba(255,255,255,.07);
  font-size:11px;font-weight:600;color:rgba(200,220,255,.7);letter-spacing:.06em;
  pointer-events:none;user-select:none;">
  ✦ Internal Reporting Platform — APSG Staging Ground &nbsp;·&nbsp; Developed by Karthik
</div>
</body>
</html>"""


# ═══════════════════════════════════════════════════════════════════════════════
#  ROUTES — Hourly Report (WB Operations Dashboard)
# ═══════════════════════════════════════════════════════════════════════════════

def _hr_hour_label(h):
    def fmt(x):
        s = "AM" if x < 12 else "PM"
        d = x % 12 or 12
        return f"{d:02d} {s}"
    return fmt(h) + " → " + fmt((h + 1) % 24)

def _hr_fmt_short(h):
    s = "AM" if h < 12 else "PM"
    d = h % 12 or 12
    return f"{d}{s}"

def _hr_op_sort_key(h):
    return h + 24 if h < 7 else h

def _hr_load_and_process(f_bytes, fname):
    import io
    try:
        if fname.lower().endswith((".xlsx", ".xls")):
            raw = pd.read_excel(io.BytesIO(f_bytes), engine="openpyxl")
        else:
            try:
                raw = pd.read_excel(io.BytesIO(f_bytes), engine="openpyxl")
            except Exception:
                raw = pd.read_csv(io.BytesIO(f_bytes), on_bad_lines="skip")
    except Exception as e:
        return None, {}, str(e)

    raw.columns = [str(c).strip().upper() for c in raw.columns]
    df = (raw[raw["ACCEPTED"].astype(str).str.strip().str.upper() == "YES"].copy()
          if "ACCEPTED" in raw.columns else raw.copy())

    for col in ["WB IN TIME", "WB OUT TIME"]:
        if col in df.columns:
            df[col] = pd.to_datetime(df[col], errors="coerce", dayfirst=True)

    if "MATERIAL" in df.columns:
        df["MATERIAL"] = df["MATERIAL"].astype(str).str.strip().str.upper()

    if "WB IN TIME" not in df.columns or df["WB IN TIME"].isna().all():
        return None, {}, "No valid WB IN TIME column found."

    valid_rows = df[df["WB IN TIME"].notna()]
    rows_7am   = valid_rows[valid_rows["WB IN TIME"].dt.hour == 7]

    if not rows_7am.empty:
        first_7am  = rows_7am["WB IN TIME"].min()
        op_start   = first_7am.normalize() + pd.Timedelta(hours=7)
        op_end_cap = op_start + pd.Timedelta(hours=22)
        df_op = df[
            df["WB IN TIME"].notna() &
            (df["WB IN TIME"] >= op_start) &
            (df["WB IN TIME"] <  op_end_cap)
        ].copy()
    else:
        df_op    = df[df["WB IN TIME"].notna()].copy()
        op_start = df_op["WB IN TIME"].min()

    if df_op.empty:
        return None, {}, "No data found in the operational window (07:00 → next day 05:00)."

    df_op["IN_HOUR"] = df_op["WB IN TIME"].dt.hour
    if "WB OUT TIME" in df_op.columns:
        df_op["OUT_HOUR"] = df_op["WB OUT TIME"].dt.hour

    latest_wbin = df_op["WB IN TIME"].max()
    last_vh     = int(latest_wbin.hour)
    now_sgt     = datetime.now(pytz.timezone("Asia/Singapore"))
    curr_h      = now_sgt.hour
    hours_ahead = (curr_h - last_vh) % 24
    latest_hour = last_vh  # always use last vehicle's hour as latest

    start_date  = op_start.date()
    end_date    = latest_wbin.date()
    end_h_str   = _hr_fmt_short((latest_hour + 1) % 24)

    if start_date == end_date:
        date_range_copy = op_start.strftime("%d/%m/%Y") + " (7AM ~ " + end_h_str + ")"
    else:
        date_range_copy = (op_start.strftime("%d/%m/%Y") + " ~ " +
                           latest_wbin.strftime("%d/%m/%Y") + " (7AM ~ " + end_h_str + ")")

    meta = {
        "op_start":        op_start.isoformat(),
        "latest_wbin":     latest_wbin.isoformat(),
        "latest_hour":     latest_hour,
        "win_start_lbl":   op_start.strftime("%d %b %Y  %I:%M %p").upper(),
        "win_end_lbl":     latest_wbin.strftime("%d %b %Y  %I:%M %p").upper(),
        "report_date_lbl": latest_wbin.strftime("%d %b %Y"),
        "date_range_copy": date_range_copy,
    }
    return df_op, meta, None

def _hr_calc_queue_wait(df):
    if "WB IN TIME" not in df.columns:
        return 0, 0
    has_in  = df["WB IN TIME"].notna()
    has_out = df["WB OUT TIME"].notna() if "WB OUT TIME" in df.columns else pd.Series(False, index=df.index)
    open_df = df[has_in & ~has_out]
    queue   = int(len(open_df))
    if queue == 0 or open_df["WB IN TIME"].isna().all():
        return queue, 0
    wait_min = max(0, int(
        (df["WB IN TIME"].max() - open_df["WB IN TIME"].min()).total_seconds() // 60
    ))
    return queue, wait_min

@app.route("/api/hourly/process", methods=["POST"])
@login_required
def hourly_process():
    f = request.files.get("file")
    if not f:
        return jsonify({"error": "No file uploaded"}), 400
    fname    = f.filename or ""
    f_bytes  = f.read()
    df, meta, err = _hr_load_and_process(f_bytes, fname)
    if err:
        return jsonify({"error": err}), 400
    if df is None or df.empty:
        return jsonify({"error": "No data found after applying operational window filter."}), 400

    queue, wait_min = _hr_calc_queue_wait(df)

    in_rows = df[df["WB IN TIME"].notna()].copy()
    in_rows["_GE"] = (in_rows["MATERIAL"] == "GOOD EARTH").astype(int) if "MATERIAL" in in_rows.columns else 0
    in_rows["_SC"] = (in_rows["MATERIAL"] == "SOFT CLAY").astype(int)  if "MATERIAL" in in_rows.columns else 0

    grp_in = in_rows.groupby("IN_HOUR", as_index=False).agg(
        total=("WB IN TIME", "count"), ge=("_GE", "sum"), sc=("_SC", "sum")
    )
    grp_in["_sk"] = grp_in["IN_HOUR"].apply(_hr_op_sort_key)
    grp_in = grp_in.sort_values("_sk", ascending=False).drop(columns=["_sk"]).reset_index(drop=True)

    if "WB OUT TIME" in df.columns and df["WB OUT TIME"].notna().any():
        out_rows = df[df["WB OUT TIME"].notna()].copy()
        grp_out  = (
            out_rows.groupby("OUT_HOUR", as_index=False)
            .agg(out_total=("WB OUT TIME", "count"))
            .rename(columns={"OUT_HOUR": "IN_HOUR"})
        )
    else:
        grp_out = pd.DataFrame(columns=["IN_HOUR", "out_total"])

    merged = grp_in.merge(grp_out, on="IN_HOUR", how="left")
    merged["out_total"] = merged["out_total"].fillna(0).astype(int)

    latest_hour = int(merged.iloc[0]["IN_HOUR"])
    latest_out  = int(merged.iloc[0]["out_total"])
    sc_total    = int(in_rows["_SC"].sum())
    ge_total    = int(in_rows["_GE"].sum())
    truck_total = int(len(in_rows))

    cards = []
    for _, row in merged.iterrows():
        h   = int(row["IN_HOUR"])
        sc  = int(row["sc"])
        ge  = int(row["ge"])
        tot = int(row["total"])
        out = int(row["out_total"])
        cards.append({
            "hour":      h,
            "label":     _hr_hour_label(h),
            "sc":        sc,
            "ge":        ge,
            "total":     tot,
            "out_total": out,
            "is_latest": h == latest_hour,
        })

    summary_copy = (
        meta["date_range_copy"] + "\n"
        "HDB = " + str(sc_total) + "\n"
        "HDB GE = " + str(ge_total) + "\n"
        "TOTAL = " + str(truck_total)
    )
    card_copy = (
        "Date/Period of reporting: " + meta["report_date_lbl"] + "\n"
        + _hr_fmt_short(latest_hour) + " ~" + _hr_fmt_short((latest_hour + 1) % 24) + "\n"
        "Waiting time: " + str(wait_min) + " min\n"
        "Queue: " + str(queue) + " Trucks\n"
        "Exit: " + str(latest_out) + " Trucks"
    )

    log_activity(session["username"], "HOURLY_PROCESS", fname)
    return jsonify({
        "ok":           True,
        "meta":         meta,
        "sc_total":     sc_total,
        "ge_total":     ge_total,
        "truck_total":  truck_total,
        "queue":        queue,
        "wait_min":     wait_min,
        "cards":        cards,
        "summary_copy": summary_copy,
        "card_copy":    card_copy,
    })


# ── Hourly Report HTML ─────────────────────────────────────────────────────────
HOURLY_REPORT_HTML = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>Hourly Report — APSG</title>
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800&family=Poppins:wght@700;800&family=JetBrains+Mono:wght@500;600&display=swap" rel="stylesheet">
<style>
*{box-sizing:border-box;margin:0;padding:0;}
:root{--bg:#0A0F1E;--s1:#111827;--s2:#1C2333;--s3:#232D42;
  --bdr:rgba(255,255,255,.08);--t1:#F1F5F9;--t2:#94A3B8;--mu:#64748B;
  --acc:#818CF8;--acc2:#6366F1;--acc-l:#A5B4FC;--grn:#10B981;--r:14px;}
html,body{background:url('/static/bg.jpg') center/cover fixed #08101E;}
body::before{content:'';position:fixed;inset:0;z-index:0;background:rgba(3,7,18,.52);pointer-events:none;}
body{font-family:'Inter',sans-serif;color:var(--t1);min-height:100vh;padding-bottom:4rem;}
body>*{position:relative;z-index:1;}

/* topbar */
.topbar{position:sticky;top:0;z-index:200;height:44px;display:flex;align-items:center;
  padding:0 1.2rem;gap:.65rem;background:rgba(6,10,28,.65);backdrop-filter:blur(18px);
  border-bottom:1px solid var(--bdr);}
.topbar .brand{font-family:'Poppins',sans-serif;font-size:.8rem;font-weight:800;color:var(--acc-l);}
.topbar .sep{width:1px;height:18px;background:var(--bdr);}
.topbar .lbl{font-size:.76rem;font-weight:600;color:var(--t2);}
.topbar .sp{flex:1;}
.back{background:rgba(99,102,241,.18);border:1px solid rgba(99,102,241,.35);border-radius:7px;
  padding:.28rem .8rem;font-size:.7rem;font-weight:600;color:#C5D5FF;text-decoration:none;transition:.2s;}
.back:hover{background:rgba(99,102,241,.32);}

/* layout */
.wrap{max-width:860px;margin:0 auto;padding:1.4rem 1rem 4rem;}
.panel{background:rgba(8,14,38,.75);border:1px solid var(--bdr);border-radius:var(--r);
  padding:1.4rem;margin-bottom:1.1rem;backdrop-filter:blur(14px);}
.ptitle{font-size:.7rem;font-weight:700;color:var(--acc);text-transform:uppercase;
  letter-spacing:.1em;margin-bottom:.85rem;}

/* upload zone */
.dz{border:2px dashed rgba(129,140,248,.35);border-radius:11px;padding:1.6rem 1rem;
  text-align:center;cursor:pointer;transition:.22s;position:relative;
  background:rgba(129,140,248,.03);user-select:none;}
.dz:hover,.dz.over{border-color:var(--acc);background:rgba(129,140,248,.09);}
.dz.has{border-color:var(--acc);border-style:solid;background:rgba(129,140,248,.06);}
.dz .ico{font-size:2rem;display:block;margin-bottom:.3rem;}
.dz .hint{font-size:.72rem;color:var(--mu);margin-top:.25rem;}
.dz .fn{font-size:.78rem;color:var(--acc);font-weight:700;margin-top:.35rem;}

/* generate btn */
.btn-gen{width:100%;padding:.9rem;background:linear-gradient(135deg,#6366F1,#818CF8);
  border:none;border-radius:10px;color:#fff;font-size:1rem;font-weight:800;
  cursor:pointer;transition:.2s;margin-top:.65rem;}
.btn-gen:hover{opacity:.9;transform:translateY(-1px);}
.btn-gen:disabled{opacity:.38;cursor:not-allowed;transform:none;}

/* error */
.err{background:rgba(239,68,68,.09);border:1px solid rgba(239,68,68,.3);border-radius:8px;
  padding:.65rem 1rem;color:#FCA5A5;font-size:.82rem;display:none;margin-top:.55rem;}

/* spinner */
.loading{display:none;text-align:center;padding:2rem;color:var(--t2);}
.spin{width:34px;height:34px;border:3px solid rgba(129,140,248,.2);border-top-color:var(--acc);
  border-radius:50%;animation:sp .7s linear infinite;margin:0 auto .7rem;}
@keyframes sp{to{transform:rotate(360deg);}}

/* summary bar */
.sum-bar{display:grid;grid-template-columns:1fr 1fr 1fr auto;gap:.75rem;align-items:center;
  margin-bottom:1rem;}
@media(max-width:560px){.sum-bar{grid-template-columns:1fr 1fr;}}
.sum-box{background:rgba(8,14,38,.88);border:1.5px solid rgba(129,140,248,.28);border-radius:12px;
  padding:.9rem 1.1rem;text-align:center;box-shadow:0 4px 18px rgba(0,0,0,.35);}
.sum-v{font-family:'JetBrains Mono',monospace;font-size:1.8rem;font-weight:800;color:#D0E0FF;line-height:1;text-shadow:0 0 18px rgba(129,140,248,.5);}
.sum-l{font-size:.74rem;font-weight:800;letter-spacing:.12em;text-transform:uppercase;
  color:#A5B4FC;margin-top:.3rem;}
.copy-sum-btn{padding:.72rem 1.4rem;border:1.5px solid rgba(129,140,248,.65);border-radius:9px;
  background:linear-gradient(135deg,rgba(99,102,241,.32),rgba(129,140,248,.2));color:#D0E0FF;font-size:.8rem;font-weight:800;
  cursor:pointer;transition:.2s;white-space:nowrap;letter-spacing:.06em;box-shadow:0 3px 14px rgba(99,102,241,.3);}
.copy-sum-btn:hover{background:linear-gradient(135deg,rgba(99,102,241,.5),rgba(129,140,248,.35));border-color:rgba(165,180,252,.9);box-shadow:0 5px 20px rgba(99,102,241,.45);transform:translateY(-1px);}
.copy-sum-btn.done{background:rgba(16,185,129,.18);border-color:rgba(16,185,129,.4);color:#4ADE80;}

/* op window bar */
.opbar{display:flex;align-items:center;gap:.7rem;padding:.55rem 1rem;
  background:rgba(8,14,38,.65);border:1px solid var(--bdr);border-radius:9px;
  margin-bottom:1rem;font-size:.72rem;flex-wrap:wrap;}
.opbar-lbl{font-weight:800;color:#94A3B8;text-transform:uppercase;letter-spacing:.09em;font-size:.74rem;}
.opbar-val{font-family:'JetBrains Mono',monospace;font-weight:700;color:#E2ECFF;font-size:.78rem;}
.opbar-sep{color:var(--mu);}

/* hourly cards */
.cards-wrap{display:none;}
.sec-lbl{font-size:.74rem;font-weight:800;letter-spacing:.13em;text-transform:uppercase;
  color:#A5B4FC;margin-bottom:.7rem;text-shadow:0 0 10px rgba(129,140,248,.3);}
.hscroll{display:flex;gap:8px;overflow-x:auto;padding-bottom:6px;
  scrollbar-width:thin;scrollbar-color:rgba(129,140,248,.3) transparent;}
.hscroll::-webkit-scrollbar{height:4px;}
.hscroll::-webkit-scrollbar-thumb{background:rgba(129,140,248,.3);border-radius:2px;}

/* latest card (highlighted) */
.cl{flex:0 0 auto;width:182px;background:rgba(10,17,44,.9);
  border:2px solid var(--acc);border-radius:11px;padding:.75rem .9rem;position:relative;
  box-shadow:0 0 20px rgba(99,102,241,.18);}
/* normal cards */
.co{flex:0 0 auto;width:152px;background:rgba(8,14,38,.75);
  border:1px solid var(--bdr);border-radius:11px;padding:.75rem .9rem;transition:border-color .15s;}
.co:hover{border-color:rgba(129,140,248,.35);}
.latest-tag{position:absolute;top:7px;right:8px;font-size:.52rem;font-weight:700;
  letter-spacing:.1em;color:var(--acc-l);text-transform:uppercase;
  background:rgba(99,102,241,.2);border:1px solid rgba(129,140,248,.3);
  border-radius:4px;padding:.05rem .32rem;}
.ch{font-family:'JetBrains Mono',monospace;font-size:.78rem;font-weight:700;color:var(--t1);
  padding-bottom:.4rem;margin-bottom:.4rem;border-bottom:1px solid var(--bdr);white-space:nowrap;}
.cr{display:flex;justify-content:space-between;align-items:center;padding:2.5px 0;}
.cr-l{font-size:.68rem;font-weight:800;color:#94A3B8;letter-spacing:.07em;text-transform:uppercase;}
.cr-v{font-family:'JetBrains Mono',monospace;font-size:1.02rem;font-weight:700;color:#E8F0FF;}
.cdiv{border:none;border-top:1px solid var(--bdr);margin:.4rem 0;}
/* copy button inside latest card */
.card-copy-btn{display:flex;align-items:center;justify-content:center;gap:5px;
  margin-top:.6rem;padding:.45rem .7rem;
  background:linear-gradient(135deg,rgba(99,102,241,.28),rgba(129,140,248,.18));border:1.5px solid rgba(129,140,248,.55);border-radius:7px;
  cursor:pointer;width:100%;font-size:.68rem;font-weight:800;letter-spacing:.07em;
  text-transform:uppercase;color:#C7D9FF;transition:.2s;box-shadow:0 2px 10px rgba(99,102,241,.22);}
.card-copy-btn:hover{background:linear-gradient(135deg,rgba(99,102,241,.45),rgba(129,140,248,.32));border-color:rgba(165,180,252,.8);box-shadow:0 4px 14px rgba(99,102,241,.38);transform:translateY(-1px);}
.card-copy-btn.done{background:rgba(16,185,129,.18);border-color:rgba(16,185,129,.4);color:#4ADE80;}

/* footer bar */
.footer-bar{position:fixed;bottom:0;left:0;right:0;z-index:9999;text-align:center;
  padding:.28rem 1rem;background:rgba(4,7,20,.72);backdrop-filter:blur(8px);
  border-top:1px solid rgba(255,255,255,.07);font-size:11px;font-weight:600;
  color:rgba(200,220,255,.7);letter-spacing:.06em;pointer-events:none;user-select:none;}
</style>
</head>
<body>

<div class="topbar">
  <span class="brand">APSG</span>
  <div class="sep"></div>
  <span class="lbl">⏱ Hourly Report</span>
  <div class="sp"></div>
  <span style="font-size:.65rem;color:var(--acc-l);font-weight:700;margin-right:.35rem;">✦ Karthi</span>
  <a href="/" class="back">← Dashboard</a>
</div>

<div class="wrap">

  <!-- Upload Panel -->
  <div class="panel">
    <div class="ptitle">📁 Upload WB Data File</div>
    <div style="background:rgba(129,140,248,.07);border-left:3px solid var(--acc);
      border-radius:8px;padding:.65rem 1rem;font-size:.77rem;color:var(--t2);margin-bottom:.9rem;">
      Op window: <strong>First 07:00 in data → +22 hrs</strong> &nbsp;|&nbsp;
      Only <strong>Accepted</strong> records counted &nbsp;|&nbsp;
      Columns: <strong>WB IN TIME · MATERIAL</strong>
    </div>
    <div class="dz" id="dz">
      <input type="file" id="fi" accept=".xlsx,.xls,.csv" style="display:none">
      <span class="ico">📄</span>
      <strong style="font-size:.84rem;">Click to browse or drag &amp; drop</strong>
      <div class="hint">.xlsx · .xls · .csv</div>
      <div class="fn" id="fn"></div>
    </div>
    <button class="btn-gen" id="btnGen" onclick="generate()" disabled>⚡ Generate Hourly Report</button>
    <div class="err" id="errBox"></div>
  </div>

  <!-- Spinner -->
  <div class="loading" id="loadDiv">
    <div class="spin"></div>
    Processing weighbridge data…
  </div>

  <!-- Results -->
  <div id="resultsDiv" style="display:none;">

    <!-- Summary bar -->
    <div class="sum-bar" id="sumBar">
      <div class="sum-box">
        <div class="sum-v" id="scTot">—</div>
        <div class="sum-l">SC Total (HDB)</div>
      </div>
      <div class="sum-box">
        <div class="sum-v" id="geTot">—</div>
        <div class="sum-l">GE Total</div>
      </div>
      <div class="sum-box">
        <div class="sum-v" id="trTot">—</div>
        <div class="sum-l">Truck Total</div>
      </div>
      <div>
        <button class="copy-sum-btn" id="sumCopyBtn" onclick="copySummary()">⎘ &nbsp;COPY</button>
      </div>
    </div>

    <!-- Op window bar -->
    <div class="opbar">
      <span class="opbar-lbl">Op Window</span>
      <span class="opbar-val" id="opWin">—</span>
    </div>

    <!-- Cards -->
    <div class="cards-wrap" id="cardsWrap">
      <div class="sec-lbl">Hourly Summary — Latest First</div>
      <div class="hscroll" id="hscroll"></div>
    </div>

  </div>

</div><!-- /wrap -->

<div class="footer-bar">
  ✦ Internal Reporting Platform — APSG Staging Ground &nbsp;·&nbsp; Developed by Karthik
</div>

<script>
// ── State ─────────────────────────────────────────────────────────────────────
var _file=null, _summCopy='', _cardCopy='';

// ── Drop zone ─────────────────────────────────────────────────────────────────
var dz=document.getElementById('dz');
var fi=document.getElementById('fi');
dz.addEventListener('click',function(){ fi.click(); });
fi.addEventListener('change',function(){ if(fi.files[0]) applyFile(fi.files[0]); });
dz.addEventListener('dragenter',function(e){ e.preventDefault(); e.stopPropagation(); dz.classList.add('over'); });
dz.addEventListener('dragover', function(e){ e.preventDefault(); e.stopPropagation(); dz.classList.add('over'); });
dz.addEventListener('dragleave',function(e){
  e.stopPropagation();
  if(!dz.contains(e.relatedTarget)) dz.classList.remove('over');
});
dz.addEventListener('drop',function(e){
  e.preventDefault(); e.stopPropagation(); dz.classList.remove('over');
  var f=e.dataTransfer&&e.dataTransfer.files&&e.dataTransfer.files[0];
  if(f) applyFile(f);
});

function applyFile(f){
  _file=f;
  document.getElementById('fn').textContent='✓ '+f.name;
  dz.classList.add('has');
  document.getElementById('btnGen').disabled=false;
}

function showErr(m){ var b=document.getElementById('errBox'); b.textContent='❌ '+m; b.style.display='block'; }
function clearErr(){ document.getElementById('errBox').style.display='none'; }

// ── Generate ──────────────────────────────────────────────────────────────────
async function generate(){
  clearErr();
  document.getElementById('resultsDiv').style.display='none';
  document.getElementById('cardsWrap').style.display='none';
  document.getElementById('loadDiv').style.display='block';
  document.getElementById('btnGen').disabled=true;

  var fd=new FormData();
  fd.append('file',_file);

  try{
    var r=await fetch('/api/hourly/process',{method:'POST',body:fd});
    var j=await r.json();
    document.getElementById('loadDiv').style.display='none';
    document.getElementById('btnGen').disabled=false;
    if(!r.ok||j.error){ showErr(j.error||'Server error'); return; }

    _summCopy = j.summary_copy||'';
    _cardCopy = j.card_copy||'';

    document.getElementById('scTot').textContent = (j.sc_total||0).toLocaleString();
    document.getElementById('geTot').textContent = (j.ge_total||0).toLocaleString();
    document.getElementById('trTot').textContent = (j.truck_total||0).toLocaleString();

    var m=j.meta||{};
    var opWin = (m.win_start_lbl||'') + ' → ' + (m.win_end_lbl||'');
    document.getElementById('opWin').textContent = opWin;

    buildCards(j.cards||[], j.queue||0, j.wait_min||0, j.card_copy||'');

    document.getElementById('resultsDiv').style.display='block';
    document.getElementById('cardsWrap').style.display='block';

  }catch(ex){
    document.getElementById('loadDiv').style.display='none';
    document.getElementById('btnGen').disabled=false;
    showErr('Network error: '+ex.message);
  }
}

function buildCards(cards, queue, wait_min, cardCopy){
  var scr=document.getElementById('hscroll');
  scr.innerHTML='';
  cards.forEach(function(c){
    var div=document.createElement('div');
    if(c.is_latest){
      div.className='cl';
      div.innerHTML=
        '<span class="latest-tag">LATEST</span>'+
        '<div class="ch">'+escH(c.label)+'</div>'+
        '<div class="cr"><span class="cr-l">SC</span><span class="cr-v">'+c.sc+'</span></div>'+
        '<div class="cr"><span class="cr-l">GE</span><span class="cr-v">'+c.ge+'</span></div>'+
        '<div class="cr"><span class="cr-l">TOTAL</span><span class="cr-v">'+c.total+'</span></div>'+
        '<div class="cr"><span class="cr-l">OUT</span><span class="cr-v">'+c.out_total+'</span></div>'+
        '<hr class="cdiv">'+
        '<div class="cr"><span class="cr-l">QUEUE</span><span class="cr-v">'+queue+'</span></div>'+
        '<div class="cr"><span class="cr-l">WAIT</span><span class="cr-v">'+wait_min+'m</span></div>'+
        '<hr class="cdiv">'+
        '<button class="card-copy-btn" id="ccb" onclick="copyCard()">⎘ &nbsp;COPY SUMMARY</button>';
    } else {
      div.className='co';
      div.innerHTML=
        '<div class="ch">'+escH(c.label)+'</div>'+
        '<div class="cr"><span class="cr-l">SC</span><span class="cr-v">'+c.sc+'</span></div>'+
        '<div class="cr"><span class="cr-l">GE</span><span class="cr-v">'+c.ge+'</span></div>'+
        '<div class="cr"><span class="cr-l">TOTAL</span><span class="cr-v">'+c.total+'</span></div>'+
        '<div class="cr"><span class="cr-l">OUT</span><span class="cr-v">'+c.out_total+'</span></div>';
    }
    scr.appendChild(div);
  });
}

function copySummary(){
  doCopy(_summCopy, 'sumCopyBtn', '⎘ &nbsp;COPY');
}
function copyCard(){
  doCopy(_cardCopy, 'ccb', '⎘ &nbsp;COPY SUMMARY');
}

function doCopy(text, btnId, resetLabel){
  var btn=document.getElementById(btnId);
  function onOk(){
    btn.innerHTML='✓ &nbsp;COPIED';
    btn.classList.add('done');
    setTimeout(function(){ btn.innerHTML=resetLabel; btn.classList.remove('done'); }, 2200);
  }
  if(navigator.clipboard&&navigator.clipboard.writeText){
    navigator.clipboard.writeText(text).then(onOk).catch(function(){ fallbackCopy(text,onOk); });
  } else { fallbackCopy(text,onOk); }
}
function fallbackCopy(text,cb){
  var ta=document.createElement('textarea'); ta.value=text;
  ta.style.cssText='position:fixed;top:0;left:0;opacity:0;';
  document.body.appendChild(ta); ta.focus(); ta.select();
  try{ document.execCommand('copy'); cb(); }catch(e){}
  document.body.removeChild(ta);
}
function escH(s){ return String(s||'').replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;'); }
</script>
</body>
</html>"""


# ═══════════════════════════════════════════════════════════════════════════════
#  ROUTES — PPT Rejection API
# ═══════════════════════════════════════════════════════════════════════════════

@app.route("/api/ppt/upload", methods=["POST"])
@login_required
def ppt_upload():
    try:
        if request.args.get("reset","true").lower()=="true":
            PPT_DB["records"]=[]; PPT_DB["slide_map"]={}; PPT_DB["presentations"]=[]
        loaded,errors=0,[]
        for f in request.files.getlist("files"):
            if not f.filename: continue
            name=f.filename
            try:
                data=f.read()
                if not data: errors.append(f'{name}: empty'); continue
                prs=Presentation(io.BytesIO(data))
                prs_idx=len(PPT_DB["presentations"])
                PPT_DB["presentations"].append(prs)
                recs,smap=extract_from_prs(prs,prs_idx)
                PPT_DB["records"].extend(recs); PPT_DB["slide_map"].update(smap); loaded+=1
            except Exception as e: errors.append(f'{name}: {e}')
        out_recs=[{k:v for k,v in r.items() if k!='_prs'} for r in PPT_DB["records"]]
        result={"records":out_recs,"files_loaded":loaded,"errors":errors}
        log_activity(session["username"],"PPT_UPLOAD",f"{loaded} files, {len(out_recs)} records")
        if loaded==0 and errors:
            result["error"]="All files failed: "+"; ".join(errors[:3])
            return jsonify(result),400
        return jsonify(result)
    except Exception as e:
        return jsonify({"error":str(e),"records":[],"files_loaded":0,"errors":[]}),500

@app.route("/api/ppt/generate", methods=["POST"])
@login_required
def ppt_generate():
    try:
        records=( request.get_json(force=True, silent=True) or {} ).get("records",[])
        if not records: return jsonify({"error":"No records"}),400
        buf=build_report(records)
        log_activity(session["username"],"PPT_GENERATE",f"{len(records)} records")
        resp = send_file(buf,mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                         as_attachment=True,download_name="Rejection_Report.pptx")
        # Free PPT presentation objects from memory after build — Render free tier
        PPT_DB["presentations"].clear()
        gc.collect()
        return resp
    except Exception as e:
        return jsonify({"error":str(e)}),500

@app.route("/api/ppt/export_excel", methods=["POST"])
@login_required
def ppt_export_excel():
    try:
        records=( request.get_json(force=True, silent=True) or {} ).get("records",[])
        if not records: return jsonify({"error":"No records"}),400
        buf=build_excel_ppt(records)
        fname=f"Rejection_Report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        log_activity(session["username"],"PPT_EXCEL",f"{len(records)} records")
        return send_file(buf,mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                         as_attachment=True,download_name=fname)
    except Exception as e:
        return jsonify({"error":str(e)}),500

@app.route("/api/ppt/generate_zip", methods=["POST"])
@login_required
def ppt_generate_zip():
    try:
        records=( request.get_json(force=True, silent=True) or {} ).get("records",[])
        if not records: return jsonify({"error":"No records"}),400
        buf=build_zip_ppt(records)
        fname=f"Rejection_Report_ZIP_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip"
        log_activity(session["username"],"PPT_ZIP",f"{len(records)} records")
        return send_file(buf,mimetype="application/zip",as_attachment=True,download_name=fname)
    except Exception as e:
        return jsonify({"error":str(e)}),500

# ═══════════════════════════════════════════════════════════════════════════════
#  ROUTES — Photo Merge API
# ═══════════════════════════════════════════════════════════════════════════════

@app.route("/api/merge", methods=["POST"])
@login_required
def api_merge():
    file_a=request.files.get("file_a"); file_b=request.files.get("file_b")
    if not file_a: return jsonify({"error":"File A is required"}),400
    if not file_b: return jsonify({"error":"File B is required"}),400
    path_a=_save_upload(file_a,"A"); path_b=_save_upload(file_b,"B")
    job_id=uuid.uuid4().hex
    top_ph=request.form.get("top_placeholder",ppt_merger.TOP_PH_NAME)
    front_ph=request.form.get("front_placeholder",ppt_merger.FRONT_PH_NAME)
    verbose=request.form.get("verbose","0")=="1"
    top_h_cm=float(request.form.get("top_h_cm",9.11)); top_w_cm=float(request.form.get("top_w_cm",15.28))
    top_left_cm=float(request.form.get("top_left_cm",0.95)); front_h_cm=float(request.form.get("front_h_cm",9.11))
    front_w_cm=float(request.form.get("front_w_cm",15.51)); front_left_cm=float(request.form.get("front_left_cm",18.73))
    center_gap_cm=float(request.form.get("center_gap_cm",2.5))
    t=threading.Thread(target=_run_merge,args=(job_id,path_a,path_b,top_ph,front_ph,verbose,
        top_h_cm,top_w_cm,top_left_cm,front_h_cm,front_w_cm,front_left_cm,center_gap_cm),daemon=True)
    t.start()
    log_activity(session["username"],"PHOTO_MERGE_START",f"job={job_id}")
    return jsonify({"job_id":job_id})

@app.route("/api/job/<job_id>")
@login_required
def api_job_status(job_id):
    with jobs_lock:
        job=jobs.get(job_id)
    if not job: return jsonify({"error":"Job not found"}),404
    return jsonify({"status":job["status"],"progress":job["progress"],"log_lines":job["log_lines"],
                    "stats":job["stats"],"error":job["error"],"output_name":job.get("output_name",""),
                    "reference_file":job.get("reference_file","?"),"has_result":job["result_file"] is not None})

@app.route("/api/download/<job_id>")
@login_required
def api_download_merge(job_id):
    with jobs_lock:
        job=jobs.get(job_id)
    if not job or not job.get("result_file"): return jsonify({"error":"No result"}),404
    path=job["result_file"]
    if not os.path.exists(path): return jsonify({"error":"File not found"}),404
    display_name=job.get("output_name") or os.path.basename(path)
    log_activity(session["username"],"PHOTO_MERGE_DOWNLOAD",display_name)
    return send_file(path,as_attachment=True,download_name=display_name,
                     mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")

# ═══════════════════════════════════════════════════════════════════════════════
#  ROUTES — Excel Rejection API (Streamlit-like logic via REST)
# ═══════════════════════════════════════════════════════════════════════════════

@app.route("/api/excel/dates", methods=["POST"])
@login_required
def excel_get_dates():
    try:
        import sys; sys.path.insert(0,os.path.dirname(__file__))
        from generate_ppt_excel import get_available_dates, format_date
        f=request.files.get("file")
        if not f: return jsonify({"error":"No file"}),400
        dates=get_available_dates(f)
        return jsonify({"dates":[{"label":format_date(d),"value":str(d)} for d in dates]})
    except Exception as e:
        return jsonify({"error":str(e)}),500

@app.route("/api/excel/preview", methods=["POST"])
@login_required
def excel_preview():
    """Return stats (rejections, sites, est_slides) for the selected date — matches standalone."""
    try:
        import sys; sys.path.insert(0, os.path.dirname(__file__))
        from generate_ppt_excel import load_and_filter, format_date
        import datetime as dt
        f = request.files.get("file")
        date_str = request.form.get("date", "")
        if not f:
            return jsonify({"error": "No file"}), 400
        selected_date = dt.date.fromisoformat(date_str) if date_str else None
        groups, data_date = load_and_filter(f, filter_date=selected_date)
        total_rej   = sum(len(v) for v in groups.values())
        total_sites = len(groups)
        est_slides  = 1 + total_sites * 2  # cover + (summary+detail) per site
        eff_date    = selected_date or data_date
        badges = [{"site": site, "count": len(rows)} for site, rows in groups.items()]
        return jsonify({
            "ok": True,
            "rejections": total_rej,
            "sites": total_sites,
            "est_slides": est_slides,
            "date_label": format_date(eff_date) if eff_date else "",
            "badges": badges,
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/api/excel/generate", methods=["POST"])
@login_required
def excel_generate():
    try:
        import sys; sys.path.insert(0,os.path.dirname(__file__))
        from generate_ppt_excel import generate_ppt, get_available_dates, load_and_filter
        import datetime as dt
        f=request.files.get("file")
        date_str=request.form.get("date","")
        secondary_ppt_f = request.files.get("secondary_ppt")   # optional — image swap
        if not f: return jsonify({"error":"No file"}),400
        selected_date=dt.date.fromisoformat(date_str) if date_str else None

        # Resolve the report date: filter date first, then date from data, never today
        f.seek(0)
        _, data_date = load_and_filter(f, filter_date=selected_date)
        report_date = selected_date or data_date or dt.date.today()

        f.seek(0)
        ppt_bytes=generate_ppt(f,report_date_obj=None,template_path=None,
                                photo_folder=None,filter_date=selected_date)

        # ── Optional: replace matching image slides from secondary PPT ────────
        replaced_count = 0
        if secondary_ppt_f:
            secondary_bytes = secondary_ppt_f.read()
            if secondary_bytes:
                ppt_bytes, replaced_count = _apply_secondary_slide_replacement(
                    ppt_bytes, secondary_bytes
                )

        fname=f"APSG-Loads Rejected {report_date.strftime('%d%m%Y')}.pptx"
        log_activity(session["username"],"EXCEL_REJECTION",
                     fname + (f" [img_swapped={replaced_count}]" if replaced_count else ""))
        resp = make_response(send_file(
            io.BytesIO(ppt_bytes),
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            as_attachment=True, download_name=fname
        ))
        if replaced_count:
            resp.headers["X-Replaced-Count"] = str(replaced_count)
        return resp
    except Exception as e:
        return jsonify({"error":str(e)}),500

# ═══════════════════════════════════════════════════════════════════════════════
#  ROUTES — Excel Rejection PPT Merge (Compare with Previous PPT)
# ═══════════════════════════════════════════════════════════════════════════════


def _move_slide(prs, from_idx, to_idx):
    """Reorder slides: move slide at from_idx to to_idx position (0-based)."""
    from pptx.oxml.ns import qn as _qn
    # prs.presentation does not exist in all python-pptx versions;
    # prs.slides._sldIdLst is the correct cross-version accessor.
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst.findall(_qn('p:sldId')))
    n = len(ids)
    if n < 2 or from_idx == to_idx: return
    from_idx = max(0, min(from_idx, n-1))
    to_idx   = max(0, min(to_idx,   n-1))
    if from_idx == to_idx: return
    entry = ids.pop(from_idx)
    ids.insert(to_idx, entry)
    # Rebuild XML
    for el in list(sldIdLst.findall(_qn('p:sldId'))):
        sldIdLst.remove(el)
    for el in ids:
        sldIdLst.append(el)


def _copy_slide_into(src_slide, dst_prs):
    """Append a copy of src_slide to dst_prs. Returns the new slide index."""
    dst = dst_prs.slides.add_slide(_safe_layout(dst_prs, 6))
    src_tree = src_slide.shapes._spTree
    dst_tree = dst.shapes._spTree
    for ch in list(dst_tree)[2:]: dst_tree.remove(ch)
    for ch in list(src_tree)[2:]: dst_tree.append(copy.deepcopy(ch))
    rId_map = {}
    for rel in src_slide.part.rels.values():
        if 'image' not in rel.reltype: continue
        try:
            _, new_rId = dst.part.get_or_add_image_part(io.BytesIO(rel.target_part.blob))
            rId_map[rel.rId] = new_rId
        except Exception: pass
    if rId_map:
        xml_str = etree.tostring(dst_tree).decode()
        def _fix(m): return f'{m.group(1)}="{rId_map.get(m.group(2), m.group(2))}"'
        xml_fixed = re.sub(r'(r:embed|r:link)="(rId\d+)"', _fix, xml_str)
        new_tree = etree.fromstring(xml_fixed.encode())
        dst_tree.getparent().replace(dst_tree, new_tree)
    return len(dst_prs.slides) - 1


# ══════════════════════════════════════════════════════════════════════════════
#  IMAGE SLIDE REPLACEMENT HELPERS  (secondary PPT swap feature)
# ══════════════════════════════════════════════════════════════════════════════

def _extract_detail_slide_key(slide):
    """
    Extract (ticket_no, e_token) from a detail slide's 8-column table.
    Both values are returned upper-cased.  Returns None when the slide
    has no 8-col table or the key fields are empty.
    Summary slides (10-col table) are deliberately excluded.
    """
    # Reject summary slides
    if any(sh.has_table and len(sh.table.columns) == 10 for sh in slide.shapes):
        return None
    for sh in slide.shapes:
        if not sh.has_table or len(sh.table.columns) != 8:
            continue
        tbl = sh.table
        if len(tbl.rows) < 2:
            continue
        cells = [c.text.strip() for c in tbl.rows[1].cells]
        ticket = cells[0].upper() if len(cells) > 0 else ''
        etoken = cells[6].upper() if len(cells) > 6 else ''
        if ticket or etoken:
            return (ticket, etoken)
    return None


def _get_slide_heading_text(slide):
    """Return the full heading string (e.g. '2. Rejection Due To Stone…') or None."""
    for sh in slide.shapes:
        if sh.has_text_frame:
            txt = sh.text_frame.text.strip()
            if re.match(r'^\d+[\.\s]', txt):
                return txt
    return None


def _replace_slide_in_place(dst_slide, src_slide):
    """
    Overwrite dst_slide's content with src_slide's content — images included.
    The slide OBJECT stays in its original position in the presentation;
    only the inner XML tree is swapped.  No formatting is applied or altered.
    """
    dst_tree = dst_slide.shapes._spTree
    src_tree = src_slide.shapes._spTree

    # Clear existing content shapes (keep the first two XML children which are
    # the spTree namespace node and the group shape node)
    for ch in list(dst_tree)[2:]:
        dst_tree.remove(ch)

    # Clone every source shape into the destination
    for ch in list(src_tree)[2:]:
        dst_tree.append(copy.deepcopy(ch))

    # Transfer image parts and remap relationship IDs
    rId_map = {}
    for rel in src_slide.part.rels.values():
        if 'image' not in rel.reltype:
            continue
        try:
            _, new_rId = dst_slide.part.get_or_add_image_part(
                io.BytesIO(rel.target_part.blob)
            )
            rId_map[rel.rId] = new_rId
        except Exception:
            pass

    if rId_map:
        xml_str  = etree.tostring(dst_tree).decode()
        def _fix_r(m): return f'{m.group(1)}="{rId_map.get(m.group(2), m.group(2))}"'
        xml_fixed = re.sub(r'(r:embed|r:link)="(rId\d+)"', _fix_r, xml_str)
        new_tree  = etree.fromstring(xml_fixed.encode())
        dst_tree.getparent().replace(dst_tree, new_tree)


def _restore_slide_heading(slide, heading_text):
    """
    Write heading_text into the heading shape of *slide*.
    The heading shape is identified as the text frame whose combined text
    starts with a digit followed by '.' (e.g. "2. Rejection Due To…").
    Only the <a:t> run text is changed — all other XML attributes, fonts,
    colours, and formatting are left exactly as the secondary PPT produced.
    """
    NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text.strip()
        if not re.match(r'^\d+[\.\s]', txt):
            continue
        tf_xml    = sh.text_frame._txBody
        first_para = tf_xml.find(f'{{{NS_A}}}p')
        if first_para is None:
            continue
        t_elems = first_para.findall(f'.//{{{NS_A}}}t')
        if not t_elems:
            continue
        # Put the full heading text in the first run, clear subsequent runs
        t_elems[0].text = heading_text
        for t in t_elems[1:]:
            t.text = ''
        break


def _apply_secondary_slide_replacement(primary_ppt_bytes, secondary_ppt_bytes):
    """
    For every detail/image slide in the primary PPT that has a matching slide
    in the secondary PPT (matched on Ticket No. + E-Token), replace the
    primary slide's content with the secondary slide's content.

    Rules:
    • Matching: BOTH Ticket No. AND E-Token must match (case-insensitive).
      Fallback: Ticket No. alone if E-Token is absent in either slide.
    • Summary slides (10-column table) are completely ignored on both sides.
    • Primary slide's heading (serial number) is preserved after replacement
      so serial-number continuity is never broken.
    • Slides with no match in the secondary PPT are kept exactly as-is.
    • No slide reordering; no duplicates introduced.

    Returns (modified_ppt_bytes: bytes, replaced_count: int).
    """
    primary_prs   = Presentation(io.BytesIO(primary_ppt_bytes))
    secondary_prs = Presentation(io.BytesIO(secondary_ppt_bytes))

    # ── Build secondary slide index ────────────────────────────────────────
    # Full key: (ticket, etoken)  — preferred
    # Ticket-only key: ticket     — fallback when etoken is missing on either side
    sec_full   = {}   # (ticket, etoken) → slide
    sec_ticket = {}   # ticket           → slide  (first occurrence wins)

    for sec_slide in secondary_prs.slides:
        key = _extract_detail_slide_key(sec_slide)
        if not key:
            continue   # summary slide or no key
        ticket, etoken = key
        if ticket and etoken:
            sec_full.setdefault((ticket, etoken), sec_slide)
        if ticket:
            sec_ticket.setdefault(ticket, sec_slide)

    if not sec_full and not sec_ticket:
        buf = io.BytesIO()
        primary_prs.save(buf)
        return buf.getvalue(), 0

    # ── Replace matching primary slides ────────────────────────────────────
    replaced = 0
    for pri_slide in primary_prs.slides:
        key = _extract_detail_slide_key(pri_slide)
        if not key:
            continue   # summary or no key — skip entirely

        pri_ticket, pri_etoken = key

        # Look up secondary match (full match preferred, ticket-only fallback)
        sec_slide = sec_full.get((pri_ticket, pri_etoken))
        if not sec_slide and pri_ticket:
            sec_slide = sec_ticket.get(pri_ticket)

        if not sec_slide:
            continue   # no match — leave primary slide unchanged

        # Save primary heading so we can restore serial continuity
        saved_heading = _get_slide_heading_text(pri_slide)

        # Swap slide content in-place (images, table, layout from secondary)
        _replace_slide_in_place(pri_slide, sec_slide)

        # Restore the primary slide's heading (serial number continuity)
        if saved_heading:
            _restore_slide_heading(pri_slide, saved_heading)

        replaced += 1

    buf = io.BytesIO()
    primary_prs.save(buf)
    return buf.getvalue(), replaced


def _norm_row(row, site):
    """Robustly extract all rejection fields from a row dict returned by
    load_and_filter, regardless of which column-name variant is used."""
    def _g(*keys):
        for k in keys:
            v = row.get(k, '')
            sv = str(v).strip() if v is not None else ''
            if sv and sv.upper() not in ('NAN', 'NONE', 'N/A', 'NAT', '-', ''):
                return sv
        return ''

    # ─ Ticket number ─
    ticket = _g('Ticket No', 'ticket_no', 'TICKET NO', 'Ticket Number',
                'ticket_number', 'Ticket', 'ticket', 'TKT', 'tkt')

    # ─ Vehicle number ─
    veh = _g('Vehicle No', 'veh_no', 'VEHICLE NO', 'Vehicle Number',
             'vehicle_number', 'Vehicle', 'vehicle', 'Veh No', 'Plate No',
             'plate_no', 'Lorry No', 'lorry_no')

    # ─ Material ─
    mat = _g('Material', 'material', 'MATERIAL', 'Material Type',
             'material_type', 'Mat Type', 'mat_type', 'Type', 'type')

    # ─ Date ─
    date = _g('Date In', 'date_in', 'DATE IN', 'Date', 'date', 'DATE',
              'Entry Date', 'entry_date', 'Weighbridge Date', 'WB Date')

    # ─ Time ─
    time_ = _g('Time In', 'time_in', 'TIME IN', 'Time', 'time', 'TIME',
               'Entry Time', 'entry_time', 'Weighbridge Time', 'WB Time')

    # If date/time are combined in WB IN TIME column, split them
    if (not date or not time_):
        wbin = _g('WB IN TIME', 'wb_in_time', 'WB In Time', 'WbInTime')
        if wbin:
            parts = wbin.replace('T', ' ').split(' ')
            if len(parts) >= 2:
                if not date:  date  = parts[0]
                if not time_: time_ = parts[1][:5]   # HH:MM

    # ─ E-Token ─
    etoken = _g('E-Token', 'e_token', 'EToken', 'etoken', 'E TOKEN',
                'E_TOKEN', 'E-token', 'Etoken', 'Token', 'token',
                'E-Tok', 'ETok')

    # ─ Accepted ─
    accepted = _g('Accepted', 'accepted', 'ACCEPTED', 'Status', 'status',
                  'Accept', 'accept')

    # ─ Reject Reason ─
    reason = _g('Reject Reason', 'reject_reason', 'REJECT REASON',
                'Rejection Reason', 'rejection_reason', 'Reason',
                'reason', 'REASON', 'Remarks', 'remarks')

    return {
        'ticket_no':     ticket,
        'veh_no':        veh,
        'material':      mat,
        'source_site':   site,
        'date':          date,
        'time':          time_,
        'e_token':       etoken,
        'accepted':      accepted if accepted else 'YES',
        'reject_reason': reason,
    }


MAX_ROWS_PER_SLIDE = 5   # maximum rejection rows per summary slide


def _renumber_detail_slide(slide, old_sn, new_sn):
    """
    Replace the leading serial-number prefix in a detail slide title.
    Uses direct XML <a:t> element manipulation — works for ANY run structure.

    Handles all cases from generate_ppt_excel:
      ["1. Rejection..."]              → single run
      ["1", ". Rejection..."]         → SN split across runs
      ["1", ".", " Rejection..."]    → SN split into 3 parts

    Algorithm:
      1. Find the title shape (full text starts with old_sn).
      2. Collect all <a:t> elements from the FIRST paragraph.
      3. Accumulate their text until "old_sn." appears in the buffer.
      4. Zero-out the <a:t> elements consumed before the match point.
      5. Set the matched <a:t> to: prefix + new_sn + "." + remainder.
         (The <a:t> elements AFTER the match are left untouched — they hold
          the rest of the title and must not be modified.)
    """
    if old_sn == new_sn:
        return

    old_str = str(old_sn) + '.'
    new_str = str(new_sn) + '.'
    NS_A    = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue

        # Identify the title shape: its combined text must start with old_sn
        full = sh.text_frame.text.strip()
        if not (full.startswith(old_str) or full.startswith(str(old_sn))):
            continue

        # Work on the first <a:p> paragraph only (title is always first para)
        txBody     = sh.text_frame._txBody
        first_para = txBody.find(f'{{{NS_A}}}p')
        if first_para is None:
            continue

        # Collect ALL <a:t> leaf elements within this paragraph
        t_elems = first_para.findall(f'.//{{{NS_A}}}t')
        if not t_elems:
            continue

        # Accumulate text from t_elems until old_str appears
        buf       = ''
        found_i   = -1
        for i, t_el in enumerate(t_elems):
            buf += (t_el.text or '')
            if old_str in buf:
                found_i = i
                break

        if found_i < 0:
            # old_str not found — nothing to do
            continue

        sn_pos = buf.find(old_str)
        prefix = buf[:sn_pos]                    # text before "1." (usually empty)
        suffix = buf[sn_pos + len(old_str):]     # text inside this t_el after "1."
        # Note: t_elems AFTER found_i still hold the rest of the title — do not touch them

        # Step 4: zero-out every t_el that was fully consumed before found_i
        for j in range(found_i):
            t_elems[j].text = ''

        # Step 5: overwrite the t_el where the match completed
        t_elems[found_i].text = prefix + new_str + suffix
        break   # only one title shape per slide


def _build_detail_slides_from_rows(rows, site, prs, start_sn=1):
    """
    Build rejection detail slides that EXACTLY match the standard rejection
    report template (same as generate_ppt_excel produces):

      • Heading: "[N]. Rejection — [Reject Reason]"
      • 8-col table  (header row + 1 data row):
          Ticket No. | Veh No. | Materials | Source Site |
          Date In | Time In | E-Token | Reject Reasons
      • LEFT half of slide:  bold "TOP PHOTO" placeholder text
      • RIGHT half of slide: bold "FRONT PHOTO" placeholder text
      • Footer banner

    Slides are appended to prs. Returns list of new slide objects.
    """
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # ── Layout constants (widescreen 13.33"×7.5") ─────────────────────────────
    # TBL_LEFT, TBL_W, TITLE_LEFT, TITLE_TOP, TITLE_W, TITLE_H,
    # ROW_H_HEADER, ROW_H_DATA  are all global constants defined earlier.

    # 8-col widths that sum to TBL_W (11306716 EMU)
    _C8 = [
        1800000,   # Ticket No.    1.97"
          900000,  # Veh No.       0.98"
        1000000,   # Materials     1.09"
        1300000,   # Source Site   1.42"
          900000,  # Date In       0.98"
          800000,  # Time In       0.87"
        2600000,   # E-Token       2.84"
        2006716,   # Reject Reason 2.19"
    ]
    _HDR_NAMES = [
        'Ticket No.', 'Veh No.', 'Materials', 'Source Site',
        'Date In', 'Time In', 'E-Token', 'Reject Reasons',
    ]

    HDR_BG = RGBColor(0x1F, 0x38, 0x64)  # dark navy — same as standard tables
    HDR_FG = RGBColor(0xFF, 0xFF, 0xFF)  # white text
    DAT_FG = RGBColor(0x00, 0x00, 0x00)  # black text
    PHO_C  = RGBColor(0x60, 0x60, 0x60)  # grey for placeholder label

    # Table bottom = TBL_TOP + ROW_H_HEADER + ROW_H_DATA
    _tbl_bot = TBL_TOP + ROW_H_HEADER + ROW_H_DATA
    # Photo area starts just below the table, ends just above footer
    _photo_top = _tbl_bot + 200000            # small gap after table
    _photo_h   = 6858000 - _photo_top - 550000  # leave ~0.6" for footer
    _half_w    = TBL_W // 2                   # split slide in two equal halves

    new_slides = []
    for i, row in enumerate(rows):
        n  = _norm_row(row, site)
        sn = start_sn + i

        slide = prs.slides.add_slide(_safe_layout(prs, 6))

        # Remove layout placeholders ("Click to add title" etc.)
        for sp in list(slide.shapes):
            try:
                if sp.is_placeholder:
                    sp.element.getparent().remove(sp.element)
            except Exception:
                pass

        # ── Heading ───────────────────────────────────────────────────────────
        reason_txt = n['reject_reason'] or 'Unknown Reason'
        heading    = f"{sn}. Rejection — {reason_txt}"

        tx = slide.shapes.add_textbox(
            Emu(TITLE_LEFT), Emu(TITLE_TOP), Emu(TITLE_W), Emu(TITLE_H)
        )
        tf = tx.text_frame; tf.word_wrap = True
        p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r  = p.add_run()
        r.text            = heading
        r.font.name       = 'Calibri'
        r.font.size       = Pt(24)
        r.font.bold       = True
        r.font.color.rgb  = RGBColor(0, 0, 0)

        # ── 8-column table ────────────────────────────────────────────────────
        tbl_frm = slide.shapes.add_table(
            2, 8,
            Emu(TBL_LEFT), Emu(TBL_TOP),
            Emu(TBL_W), Emu(ROW_H_HEADER + ROW_H_DATA)
        )
        tbl = tbl_frm.table
        for ci, cw in enumerate(_C8):
            tbl.columns[ci].width = Emu(cw)
        tbl.rows[0].height = Emu(ROW_H_HEADER)
        tbl.rows[1].height = Emu(ROW_H_DATA)

        # Header cells
        for ci, h in enumerate(_HDR_NAMES):
            cell = tbl.rows[0].cells[ci]
            cell.text = h
            # No fill — white/transparent background to match standard template
            cell.fill.background()
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.bold      = True
                    run.font.size      = Pt(10)
                    run.font.name      = 'Calibri'
                    run.font.color.rgb = RGBColor(0, 0, 0)  # black text on white background

        # Data cells
        vals = [
            n['ticket_no'], n['veh_no'], n['material'], n['source_site'],
            n['date'], n['time'], n['e_token'], n['reject_reason'],
        ]
        for ci, v in enumerate(vals):
            cell = tbl.rows[1].cells[ci]
            cell.text = v
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.size      = Pt(10)
                    run.font.name      = 'Calibri'
                    run.font.color.rgb = DAT_FG

        # ── TOP PHOTO placeholder (left half) ─────────────────────────────────
        tp = slide.shapes.add_textbox(
            Emu(TBL_LEFT), Emu(_photo_top), Emu(_half_w), Emu(_photo_h)
        )
        tp_tf = tp.text_frame; tp_tf.word_wrap = False
        tp_p  = tp_tf.paragraphs[0]; tp_p.alignment = PP_ALIGN.CENTER
        tp_r  = tp_p.add_run()
        tp_r.text            = 'TOP PHOTO'
        tp_r.font.bold       = True
        tp_r.font.size       = Pt(20)
        tp_r.font.name       = 'Calibri'
        tp_r.font.color.rgb  = PHO_C

        # ── FRONT PHOTO placeholder (right half) ──────────────────────────────
        fp = slide.shapes.add_textbox(
            Emu(TBL_LEFT + _half_w), Emu(_photo_top), Emu(_half_w), Emu(_photo_h)
        )
        fp_tf = fp.text_frame; fp_tf.word_wrap = False
        fp_p  = fp_tf.paragraphs[0]; fp_p.alignment = PP_ALIGN.CENTER
        fp_r  = fp_p.add_run()
        fp_r.text            = 'FRONT PHOTO'
        fp_r.font.bold       = True
        fp_r.font.size       = Pt(20)
        fp_r.font.name       = 'Calibri'
        fp_r.font.color.rgb  = PHO_C

        # ── Footer ────────────────────────────────────────────────────────────
        try:
            _add_footer_banner(slide)
        except Exception:
            pass

        new_slides.append(slide)
    return new_slides


def _insert_rows_paginated(base_prs, summary_slide_idx, missing_rows, site,
                           existing_row_count, last_slide_idx, insert_offset_ref,
                           report_date=None, sn_start=None):
    """
    Insert missing_rows into summary slides starting at summary_slide_idx.
    existing_row_count : rows already in the LAST summary slide (for capacity checks).
    sn_start           : total existing rows across ALL summary slides (for SN continuity).
                         If None, defaults to existing_row_count.
    """
    slides_added   = 0
    current_s_idx  = summary_slide_idx   # index of the current summary slide
    current_count  = existing_row_count  # rows in LAST summary slide (capacity tracking)
    # sn_start is the TOTAL existing rows so SN continues correctly across all slides
    sn             = (sn_start if sn_start is not None else existing_row_count) + 1

    for row in missing_rows:
        norm = _norm_row(row, site)

        if current_count >= MAX_ROWS_PER_SLIDE:
            # Current slide is full — create a continuation summary slide
            # and insert it immediately after current summary slide's group
            new_summary = base_prs.slides.add_slide(_safe_layout(base_prs, 6))
            _init_summary_slide_header(new_summary, site, report_date=report_date)
            end_idx = len(base_prs.slides) - 1   # currently at end
            # Insert after the LAST slide of the current section
            target = last_slide_idx + slides_added + 1
            _move_slide(base_prs, end_idx, target)
            slides_added  += 1
            current_s_idx  = target
            current_count  = 0

        # Add row to current summary slide (guard against bad index)
        n_now = len(base_prs.slides)
        if current_s_idx >= n_now:
            current_s_idx = n_now - 1
        slide = base_prs.slides[current_s_idx]
        for sh in slide.shapes:
            if sh.has_table and len(sh.table.columns) == 10:
                _append_row_to_table(sh, norm, sn)
                break
        current_count += 1
        sn += 1

    return slides_added, last_slide_idx + slides_added


def _ordinal_suffix(n):
    """Return e.g. '13th', '1st', '2nd', '3rd'."""
    if 11 <= (n % 100) <= 13:
        return f"{n}th"
    return f"{n}{['th','st','nd','rd','th','th','th','th','th','th'][n % 10]}"


def _init_summary_slide_header(slide, site, report_date=None):
    """Add heading + empty 10-col table to a blank continuation slide.
    Heading matches the generate_ppt_excel format exactly:
    'Source Site [SITE] Load Rejected Summary List Dated [dd][th] [Month] [YYYY] ([Weekday])'
    """
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    if report_date:
        day_s    = _ordinal_suffix(report_date.day)
        month_yr = report_date.strftime('%B %Y')
        weekday  = report_date.strftime('%A')
        heading  = (f"Source Site {site} Load Rejected Summary List "
                    f"Dated {day_s} {month_yr} ({weekday})")
    else:
        heading = f"Source Site {site} Load Rejected Summary List"

    tx = slide.shapes.add_textbox(Emu(TITLE_LEFT), Emu(TITLE_TOP), Emu(TITLE_W), Emu(TITLE_H))
    tf = tx.text_frame; tf.word_wrap = True
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = heading
    r1.font.name = 'Calibri'; r1.font.size = Pt(22)
    r1.font.bold = True; r1.font.underline = True
    r1.font.color.rgb = RGBColor(0, 0, 0)

    # Add empty 1-row (header only) table
    tbl_frm = slide.shapes.add_table(
        1, 10, Emu(TBL_LEFT), Emu(TBL_TOP),
        Emu(TBL_W), Emu(ROW_H_HEADER)
    )
    tbl = tbl_frm.table
    for ci, w in enumerate(COL_WIDTHS): tbl.columns[ci].width = w
    tbl.rows[0].height = ROW_H_HEADER
    HEADERS = ['S/N', 'Ticket No', 'Veh No', 'Material', 'Source Site',
               'Date In', 'Time In', 'E-Token', 'Accepted', 'Reject Reason']
    for ci, h in enumerate(HEADERS):
        tc = tbl.rows[0].cells[ci]._tc
        tc.getparent().replace(tc, etree.fromstring(cell_xml(h, True, ci)))
    _add_footer_banner(slide)

# ── E-Token extraction helpers ───────────────────────────────────────────────

_ETOKEN_KEY_NORMS = frozenset([
    'ETOKEN','TOKEN','ETOKENNUMBER','ETOKENNO','ETOKID',
])

def _etoken_col_idx(header_cells):
    """Return the column index of the E-Token column from a header row, or -1."""
    for ci, h in enumerate(header_cells):
        norm = h.strip().upper().replace('-','').replace('_','').replace(' ','')
        if norm in _ETOKEN_KEY_NORMS:
            return ci
    return -1


def _etokens_from_table(tbl):
    """Extract all E-Token values from a pptx table.
    Uses the header row to locate the correct column — works for any layout."""
    result = set()
    if len(tbl.rows) < 1:
        return result
    # Row 0 should be the header
    header = [c.text.strip() for c in tbl.rows[0].cells]
    e_idx  = _etoken_col_idx(header)
    if e_idx == -1:
        # Fallback: use assumed positions based on column count
        nc = len(tbl.columns)
        if nc == 10: e_idx = 7   # S/N, Ticket, Veh, Mat, Site, Date, Time, EToken, Accepted, Reason
        elif nc == 8: e_idx = 6  # Ticket, Veh, Mat, Site, Date, Time, EToken, Reason
        else: return result
    for ri in range(1, len(tbl.rows)):
        row = tbl.rows[ri]
        if e_idx >= len(row.cells):
            continue
        val = row.cells[e_idx].text.strip().upper()
        # Exclude header text that might have leaked into data rows
        if val and val not in ('', 'E-TOKEN', 'ETOKEN', 'E TOKEN', 'NAN', 'NONE', 'N/A'):
            result.add(val)
    return result


def _etoken_from_row(row, site=''):
    """Extract E-Token from a row dict (from load_and_filter) using all field-name variants.
    Falls back to _norm_row for maximum compatibility."""
    # Fast path: try the most common keys first
    for k in ('E-Token', 'e_token', 'EToken', 'E TOKEN', 'E_TOKEN',
              'E-token', 'Etoken', 'ETOKEN', 'Token', 'token',
              'E-Tok', 'ETok', 'E_Token'):
        v = row.get(k)
        if v is not None:
            sv = str(v).strip().upper()
            if sv and sv not in ('', 'NAN', 'NONE', 'N/A', 'NAT'):
                return sv
    # Slow path: case-insensitive search over all keys
    for k, v in row.items():
        if v is None: continue
        k_norm = k.strip().upper().replace('-','').replace('_','').replace(' ','')
        if k_norm in _ETOKEN_KEY_NORMS:
            sv = str(v).strip().upper()
            if sv and sv not in ('', 'NAN', 'NONE', 'N/A', 'NAT'):
                return sv
    # Last resort: use _norm_row which handles every variant we know of
    try:
        return _norm_row(row, site).get('e_token', '').strip().upper()
    except Exception:
        return ''



@app.route("/api/excel/merge_ppt", methods=["POST"])
@login_required

def excel_merge_ppt():
    """
    Compare & Merge mode.

    Strategy: generate a FRESH mini PPT for the missing records using the
    EXACT SAME generate_ppt_excel.generate_ppt call as normal mode — so every
    slide looks byte-for-byte identical to a normally-generated report.

    • Summary rows   → added to existing summary slides via _insert_rows_paginated
    • Detail slides  → copied directly from the mini PPT (correct format guaranteed)
    • Filename       → same format as normal generation (no "Merged" suffix)
    """
    try:
        import sys, datetime as dt
        sys.path.insert(0, os.path.dirname(__file__))
        from generate_ppt_excel import load_and_filter, generate_ppt as _gpe_gen

        excel_file    = request.files.get("file")
        prev_ppt_file = request.files.get("prev_ppt")
        date_str      = request.form.get("date", "")

        if not excel_file:    return jsonify({"error": "No Excel file uploaded"}), 400
        if not prev_ppt_file: return jsonify({"error": "No previous PPT file uploaded"}), 400

        # Read ALL bytes upfront — we'll need the Excel file twice
        excel_bytes = excel_file.read()
        prev_bytes  = prev_ppt_file.read()
        selected_date = dt.date.fromisoformat(date_str) if date_str else None

        # ── STEP 1: Parse existing PPT ─────────────────────────────────────────
        base_prs      = Presentation(io.BytesIO(prev_bytes))
        existing_keys = set()   # set of E-Token strings from primary PPT
        site_map      = {}

        # Use _etokens_from_table (header-driven) for robust extraction
        # regardless of which column position E-Token occupies.
        for si, slide in enumerate(base_prs.slides):
            for sh in slide.shapes:
                if not sh.has_table: continue
                tbl = sh.table
                nc  = len(tbl.columns)

                # Extract ALL E-Tokens from this table (any column count)
                for et in _etokens_from_table(tbl):
                    existing_keys.add(et)

                # For 10-col summary slides: also build site_map
                if nc == 10 and len(tbl.rows) > 1:
                    # Identify site code column (header row, col 4 in standard layout)
                    header = [c.text.strip() for c in tbl.rows[0].cells]
                    site_ci = 4  # default
                    for ci2, h in enumerate(header):
                        hn = h.upper().replace(' ','').replace('-','').replace('_','')
                        if hn in ('SOURCESITE','SITE','SITECODE'):
                            site_ci = ci2; break

                    s_site = None; rows_here = 0
                    for ri in range(1, len(tbl.rows)):
                        cells = [c.text.strip() for c in tbl.rows[ri].cells]
                        if not any(cells): continue
                        if site_ci < len(cells) and cells[site_ci]:
                            s_site = cells[site_ci]; rows_here += 1
                    if s_site:
                        if s_site not in site_map:
                            site_map[s_site] = {
                                "summary_idx": si, "last_summary_idx": si,
                                "last_summary_row_count": rows_here,
                                "total_row_count": rows_here,
                                "last_idx": si
                            }
                        else:
                            site_map[s_site]["last_summary_idx"]      = si
                            site_map[s_site]["last_summary_row_count"] = rows_here
                            site_map[s_site]["total_row_count"]       += rows_here

        cur_site = None
        for si, slide in enumerate(base_prs.slides):
            has10 = any(sh.has_table and len(sh.table.columns)==10 for sh in slide.shapes)
            if has10:
                for sh in slide.shapes:
                    if not sh.has_table or len(sh.table.columns) != 10: continue
                    for ri in range(1, len(sh.table.rows)):
                        cells = [c.text.strip() for c in sh.table.rows[ri].cells]
                        if len(cells) >= 5 and cells[4]: cur_site = cells[4]; break
                    break
            if cur_site and cur_site in site_map:
                site_map[cur_site]["last_idx"] = si

        app.logger.info(f"merge: PPT has {len(existing_keys)} unique E-Tokens across {len(site_map)} sites")

        # ── STEP 2: Find missing records from Excel ────────────────────────────
        excel_io = io.BytesIO(excel_bytes)
        groups_all, data_date = load_and_filter(excel_io, filter_date=selected_date)

        # Comparison key = E-Token ONLY.
        # _etoken_from_row() handles every possible field-name variant that
        # load_and_filter may return (E-TOKEN, EToken, e_token, TOKEN, etc.)
        missing_groups = {}; missing_etokens = set()
        total_new = total_missing = 0
        for site, rows in groups_all.items():
            miss = []; seen_batch = set()
            for row in rows:
                total_new += 1
                etoken = _etoken_from_row(row, site)  # robust extraction
                if etoken and etoken not in existing_keys and etoken not in seen_batch:
                    miss.append(row); total_missing += 1
                    missing_etokens.add(etoken)
                    seen_batch.add(etoken)
                    existing_keys.add(etoken)
            if miss: missing_groups[site] = miss

        app.logger.info(
            f"merge: Excel has {total_new} records | PPT has "
            f"{len(existing_keys)-total_missing} | missing={total_missing}"
        )

        if not missing_groups:
            return jsonify({
                "ok": True, "merged": False, "total_new": total_new, "total_missing": 0,
                "message": (
                    f"No missing entries detected. "
                    f"The existing PPT already contains all {total_new} rejection records."
                )
            })

        # ── STEP 3: Generate mini PPT for missing records using ORIGINAL Excel ─
        # Filter original Excel to only the missing rows (by ticket+etoken+vehicle key).
        # Then call generate_ppt EXACTLY as normal mode does — guaranteed identical format.
        raw_buf = io.BytesIO(excel_bytes)
        try:
            df_orig = pd.read_excel(raw_buf)
        except Exception:
            for enc in ("utf-8", "latin1", "cp1252"):
                try:
                    raw_buf.seek(0); df_orig = pd.read_csv(raw_buf, encoding=enc); break
                except Exception: continue
            else:
                return jsonify({"error": "Could not re-read Excel file"}), 400

        # Find the ticket / vehicle / etoken columns in the original file
        cols_up = {c.strip().upper(): c for c in df_orig.columns}
        t_col = next((cols_up[k] for k in cols_up if 'TICKET' in k), None)
        v_col = next((cols_up[k] for k in cols_up
                      if k in ('VEHICLE NO', 'VEH NO', 'VEHICLE NUMBER', 'VEH')), None)
        e_col = next((cols_up[k] for k in cols_up
                      if 'TOKEN' in k or k in ('E-TOKEN', 'ETOKEN')), None)

        if e_col:
            # Filter by E-Token only
            def _keep(r):
                return str(r[e_col]).strip().upper() in missing_etokens
            df_miss = df_orig[df_orig.apply(_keep, axis=1)]
            df_miss = df_miss.drop_duplicates(subset=[e_col])
        elif missing_etokens:
            # e_col not found: try all columns for E-Token values
            def _keep_any(r):
                for v in r.values:
                    if str(v).strip().upper() in missing_etokens:
                        return True
                return False
            df_miss = df_orig[df_orig.apply(_keep_any, axis=1)]
        else:
            df_miss = pd.DataFrame()  # nothing missing — safe empty fallback

        # Write filtered DataFrame to in-memory Excel and call generate_ppt
        mini_excel = io.BytesIO()
        df_miss.to_excel(mini_excel, index=False)
        mini_excel.seek(0)

        mini_ppt_bytes = _gpe_gen(
            mini_excel,
            report_date_obj=None,
            template_path=None,
            photo_folder=None,   # same as normal excel_generate route
            filter_date=None,    # already filtered above
        )
        mini_prs = Presentation(io.BytesIO(mini_ppt_bytes))

        # Parse mini_prs: collect detail slides per site
        # (summary slides → skip, we add rows via _insert_rows_paginated instead)
        mini_detail = {}   # site → [slide_obj, ...]
        cur_mini    = None
        after_sum   = False
        for ms in mini_prs.slides:
            has10 = any(sh.has_table and len(sh.table.columns)==10 for sh in ms.shapes)
            if has10:
                after_sum = True
                for sh in ms.shapes:
                    if not sh.has_table or len(sh.table.columns) != 10: continue
                    for ri in range(1, len(sh.table.rows)):
                        cells = [c.text.strip() for c in sh.table.rows[ri].cells]
                        if len(cells) >= 5 and cells[4]:
                            cur_mini = cells[4]
                            mini_detail.setdefault(cur_mini, [])
                            break
                    break
                # Do NOT copy summary slide — rows are added via _insert_rows_paginated
            elif after_sum and cur_mini:
                mini_detail.setdefault(cur_mini, []).append(ms)
            # else: cover/heading slide before any summary → skip

        # ── STEP 4: Insert missing content into base_prs ──────────────────────
        existing_order      = sorted(site_map.items(), key=lambda x: x[1]["summary_idx"])
        existing_site_names = [s for s, _ in existing_order]
        new_sites           = [s for s in missing_groups if s not in site_map]

        insert_offset = 0

        for site in (existing_site_names + new_sites):
            if site not in missing_groups: continue
            missing_rows = missing_groups[site]

            if site in site_map:
                info     = site_map[site]
                n_slides = len(base_prs.slides)

                last_s_raw   = info.get("last_summary_idx", info["summary_idx"])
                last_all_raw = info["last_idx"]
                last_s_rows  = info.get("last_summary_row_count", info.get("row_count", 0))

                adj_last_s   = last_s_raw   + insert_offset
                adj_last_all = last_all_raw + insert_offset

                if adj_last_s >= n_slides:
                    app.logger.warning(f"merge: skip {site} — adj_last_s={adj_last_s} >= {n_slides}")
                    continue
                adj_last_all = min(adj_last_all, n_slides - 1)

                # total_rows = cumulative count across ALL summary slides for this site.
                # This ensures new serial numbers continue correctly from the last existing SN.
                # e.g. Site has slides 1-5 and 6-8 → total_rows=8 → next SN=9
                total_rows = info.get("total_row_count",
                             info.get("last_summary_row_count",
                             info.get("row_count", 0)))
                next_sn    = total_rows + 1   # first serial number for missing entries

                # PHASE A — summary rows  (capacity = last_s_rows; SN = next_sn)
                pag_added, _ = _insert_rows_paginated(
                    base_prs, adj_last_s,
                    missing_rows, site, last_s_rows,
                    adj_last_s, insert_offset,
                    report_date=(selected_date or data_date),
                    sn_start=total_rows          # ← SN continues from total, not last-slide count
                )
                insert_offset += pag_added

                # PHASE B — detail slides  (serial numbers must match summary: next_sn, next_sn+1, …)
                det_slides = mini_detail.get(site, [])

                if det_slides:
                    # ── Guarantee ascending serial-number order ───────────────
                    # Sort the detail slides by their heading serial number so
                    # slides always appear in ascending order under each Summary
                    # section, regardless of how mini_prs happened to order them.
                    def _slide_heading_sn(sl):
                        for _sh in sl.shapes:
                            if _sh.has_text_frame:
                                _txt = _sh.text_frame.text.strip()
                                _m = re.match(r'^(\d+)[\.\s]', _txt)
                                if _m:
                                    return int(_m.group(1))
                        return 999999   # unknown → sort last
                    det_slides = sorted(det_slides, key=_slide_heading_sn)
                    # ─────────────────────────────────────────────────────────

                    # Copy from mini_prs and renumber to continue primary PPT's
                    # serial sequence (next_sn, next_sn+1, …) without resetting.
                    insert_after = adj_last_all + pag_added + 1
                    for k, src_slide in enumerate(det_slides):
                        new_idx = _copy_slide_into(src_slide, base_prs)
                        copied  = base_prs.slides[new_idx]
                        # old_sn = the actual heading number in this slide;
                        # derive it from the slide itself for robustness rather
                        # than assuming it equals k+1 (handles edge cases where
                        # mini_prs numbering differs from strict 1,2,3... sequence)
                        old_sn_actual = _slide_heading_sn(src_slide)
                        if old_sn_actual == 999999:
                            old_sn_actual = k + 1   # fallback
                        _renumber_detail_slide(copied, old_sn_actual, next_sn + k)
                        target  = insert_after + k
                        if new_idx != target and target < len(base_prs.slides):
                            _move_slide(base_prs, new_idx, target)
                    det_count = len(det_slides)
                else:
                    # Fallback — build detail slides with correct starting SN
                    fallback = _build_detail_slides_from_rows(
                        missing_rows, site, base_prs, start_sn=next_sn
                    )
                    det_count = len(fallback)
                    insert_after = adj_last_all + pag_added + 1
                    for k in range(det_count):
                        end_now = len(base_prs.slides) - (det_count - k)
                        target  = insert_after + k
                        if end_now != target and target < len(base_prs.slides):
                            _move_slide(base_prs, end_now, target)

                insert_offset += det_count

            else:
                # Brand-new site — copy ALL slides from mini_prs for this site
                # Order: summary slide(s) first, then detail slides
                site_sum_slides  = []
                site_det_slides  = []
                in_this_site     = False
                for ms in mini_prs.slides:
                    has10 = any(sh.has_table and len(sh.table.columns)==10 for sh in ms.shapes)
                    if has10:
                        s_in_slide = None
                        for sh in ms.shapes:
                            if not sh.has_table or len(sh.table.columns) != 10: continue
                            for ri in range(1, len(sh.table.rows)):
                                cells = [c.text.strip() for c in sh.table.rows[ri].cells]
                                if len(cells) >= 5 and cells[4]: s_in_slide = cells[4]; break
                            break
                        if s_in_slide == site:
                            in_this_site = True
                            site_sum_slides.append(ms)
                        elif in_this_site:
                            in_this_site = False   # moved to next site's summary
                    elif in_this_site:
                        site_det_slides.append(ms)

                for ms in site_sum_slides: _copy_slide_into(ms, base_prs)
                if site_det_slides:
                    for ms in site_det_slides: _copy_slide_into(ms, base_prs)
                else:
                    _build_detail_slides_from_rows(missing_rows, site, base_prs, start_sn=1)

                insert_offset += len(site_sum_slides) + max(len(site_det_slides), len(missing_rows))

        # ── STEP 5: Save + respond ─────────────────────────────────────────────
        buf = io.BytesIO(); base_prs.save(buf); buf.seek(0)
        eff   = selected_date or data_date or dt.date.today()
        fname = f"APSG-Loads Rejected {eff.strftime('%d%m%Y')}.pptx"
        log_activity(session["username"], "EXCEL_REJECTION_MERGE",
                     f"missing={total_missing} total={total_new} file={fname}")
        resp = make_response(send_file(buf,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            as_attachment=True, download_name=fname))
        resp.headers["X-Missing-Count"]  = str(total_missing)
        resp.headers["X-Total-Count"]    = str(total_new)
        resp.headers["X-Merge-Filename"] = fname
        return resp

    except Exception as e:
        return jsonify({"error": str(e), "trace": traceback.format_exc()}), 500




def _generate_ppt_from_groups(groups, report_date=None):
    """Generate PPT bytes from a groups dict (site -> list of row dicts)."""
    import sys, io as _io, datetime as dt
    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

    # Import the internal builder from generate_ppt_excel
    import generate_ppt_excel as _gpe

    # Create a minimal DataFrame from groups and call generate_ppt
    rows = []
    for site, site_rows in groups.items():
        for row in site_rows:
            rows.append(row)

    if not rows:
        # Return empty presentation
        prs = Presentation()
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    import pandas as pd
    df = pd.DataFrame(rows)

    # Remap column names to what generate_ppt_excel expects if needed
    col_map = {
        "ticket_no": "Ticket No", "veh_no": "Vehicle No",
        "material": "Material", "source_site": "Source Site",
        "date": "Date In", "time": "Time In", "e_token": "E-Token",
        "accepted": "Accepted", "reject_reason": "Reject Reason"
    }
    df = df.rename(columns={k: v for k, v in col_map.items() if k in df.columns})

    # Write to a temp Excel file and call generate_ppt
    excel_buf = io.BytesIO()
    df.to_excel(excel_buf, index=False)
    excel_buf.seek(0)

    ppt_bytes = _gpe.generate_ppt(
        excel_buf,
        report_date_obj=report_date,
        template_path=None,
        photo_folder=None,
        filter_date=None  # Data already filtered
    )
    return ppt_bytes


def _merge_ppts(base_ppt_bytes, new_ppt_bytes):
    """Append all slides from new_ppt into base_ppt and return merged bytes."""
    base_prs = Presentation(io.BytesIO(base_ppt_bytes))
    new_prs = Presentation(io.BytesIO(new_ppt_bytes))

    # Clone each slide from new_prs and append to base_prs
    for si in range(len(new_prs.slides)):
        dst = base_prs.slides.add_slide(_safe_layout(base_prs, 6))
        src = new_prs.slides[si]
        src_tree = src.shapes._spTree
        dst_tree = dst.shapes._spTree
        # Remove placeholder elements from blank layout
        for ch in list(dst_tree)[2:]:
            dst_tree.remove(ch)
        # Copy all source shapes
        for ch in list(src_tree)[2:]:
            dst_tree.append(copy.deepcopy(ch))
        # Copy image relationships
        rId_map = {}
        for rel in src.part.rels.values():
            if "image" not in rel.reltype:
                continue
            try:
                _, new_rId = dst.part.get_or_add_image_part(io.BytesIO(rel.target_part.blob))
                rId_map[rel.rId] = new_rId
            except Exception:
                pass
        if rId_map:
            import re as _re
            xml_str = etree.tostring(dst_tree).decode()
            def _fix(m):
                return f'{m.group(1)}="{rId_map.get(m.group(2), m.group(2))}"'
            xml_fixed = _re.sub(r'(r:embed|r:link)="(rId\d+)"', _fix, xml_str)
            new_tree = etree.fromstring(xml_fixed.encode())
            dst_tree.getparent().replace(dst_tree, new_tree)

    buf = io.BytesIO()
    base_prs.save(buf)
    return buf.getvalue()


def _detect_table_data_sz(sh):
    """Read the sz (font size, in hundredths of pt) from first existing data row.
    Returns 1100 as fallback. Ensures new rows match existing row font size."""
    try:
        tbl = sh.table
        for ri in range(1, len(tbl.rows)):
            row_str = etree.tostring(tbl.rows[ri]._tr).decode()
            m = re.search(r'<a:rPr[^>]*sz="(\d+)"', row_str)
            if m:
                return int(m.group(1))
    except Exception:
        pass
    return 1100   # default = 11pt


def _cell_xml_sz(text, header, col, data_sz=1100):
    """Like cell_xml but with configurable data-row font size."""
    ns = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    safe = (str(text or '')).replace('&','&amp;').replace('<','&lt;').replace('>','&gt;').replace('"','&quot;')
    if header:
        algn, fa = 'ctr', 'b'
        sz, bold, color, face = 1200, 1, HEADER_COLOR, 'Calibri'
        mar = 'marL="8912" marR="8912" marT="8912" marB="0"'
    else:
        algn, fa = 'ctr', 'ctr'
        sz, bold, color, face = data_sz, 0, '000000', 'Calibri'
        mar = 'marL="9525" marR="9525" marT="9525" marB="0"'
    return (f'<a:tc {ns}><a:txBody><a:bodyPr/><a:lstStyle/>'
            f'<a:p><a:pPr algn="{algn}" fontAlgn="{fa}"/>'
            f'<a:r><a:rPr lang="en-SG" sz="{sz}" b="{bold}" i="0" u="none" '
            f'strike="noStrike" dirty="0">'
            f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
            f'<a:effectLst/><a:latin typeface="{face}"/>'
            f'</a:rPr><a:t>{safe}</a:t></a:r></a:p>'
            f'</a:txBody>'
            f'<a:tcPr {mar} anchor="ctr">{_borders()}<a:noFill/></a:tcPr>'
            f'</a:tc>')


def _append_row_to_table(sh, rec, sn):
    """Append one data row to a 10-column summary table shape and expand the frame height.
    Font size is auto-detected from existing rows to ensure visual consistency."""
    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

    # Match font size of existing data rows — prevents mixed-font-size tables
    data_sz = _detect_table_data_sz(sh)

    vals = [
        str(sn),
        rec.get('ticket_no', ''), rec.get('veh_no', ''), rec.get('material', ''),
        rec.get('source_site', ''), rec.get('date', ''), rec.get('time', ''),
        rec.get('e_token', ''), rec.get('accepted', 'NO'), rec.get('reject_reason', ''),
    ]
    row_xml = f'<a:tr xmlns:a="{NS_A}" h="{ROW_H_DATA}">'
    for ci, v in enumerate(vals):
        row_xml += _cell_xml_sz(v, False, ci, data_sz)
    row_xml += '</a:tr>'
    try:
        new_row_el = etree.fromstring(row_xml.encode())
        sh.table._tbl.append(new_row_el)
        # Expand the graphicFrame height
        gf = sh.element
        xfrm = gf.find(f'{{{NS_P}}}xfrm')
        if xfrm is not None:
            ext = xfrm.find(f'{{{NS_A}}}ext')
            if ext is not None:
                ext.set('cy', str(int(ext.get('cy', 0)) + ROW_H_DATA))
    except Exception as _e:
        app.logger.warning(f'_append_row_to_table: {_e}')


def _smart_merge_ppts(base_ppt_bytes, new_ppt_bytes):
    """
    Smart merge that avoids duplicate heading/cover slides and section headers.

    Rules:
    - Slides with no 10-col or 8-col tables that appear BEFORE the first summary
      slide are treated as heading/cover slides and are SKIPPED entirely.
    - For each summary slide (10-col table):
        * If the site already exists in base PPT → add new rows to the existing
          summary table instead of creating a duplicate summary slide.
        * If the site is new → copy the summary slide as normal.
    - All detail slides (8-col tables OR image slides that appear after a
      summary slide) are always copied so photo evidence is preserved.
    """
    base_prs = Presentation(io.BytesIO(base_ppt_bytes))
    new_prs  = Presentation(io.BytesIO(new_ppt_bytes))

    # Build site → first summary slide index in base PPT
    base_site_summary_idx = {}
    for si, slide in enumerate(base_prs.slides):
        for sh in slide.shapes:
            if not sh.has_table or len(sh.table.columns) != 10:
                continue
            for ri in range(1, len(sh.table.rows)):
                cells = [c.text.strip() for c in sh.table.rows[ri].cells]
                if len(cells) >= 5 and cells[4]:
                    base_site_summary_idx.setdefault(cells[4], si)
                    break
            break

    def _copy_slide_to_base(src_slide):
        dst = base_prs.slides.add_slide(_safe_layout(base_prs, 6))
        src_tree = src_slide.shapes._spTree
        dst_tree = dst.shapes._spTree
        for ch in list(dst_tree)[2:]:
            dst_tree.remove(ch)
        for ch in list(src_tree)[2:]:
            dst_tree.append(copy.deepcopy(ch))
        rId_map = {}
        for rel in src_slide.part.rels.values():
            if 'image' not in rel.reltype:
                continue
            try:
                _, new_rId = dst.part.get_or_add_image_part(io.BytesIO(rel.target_part.blob))
                rId_map[rel.rId] = new_rId
            except Exception:
                pass
        if rId_map:
            xml_str = etree.tostring(dst_tree).decode()
            def _fix(m):
                return f'{m.group(1)}="{rId_map.get(m.group(2), m.group(2))}"'
            xml_fixed = re.sub(r'(r:embed|r:link)="(rId\d+)"', _fix, xml_str)
            new_tree = etree.fromstring(xml_fixed.encode())
            dst_tree.getparent().replace(dst_tree, new_tree)

    seen_first_summary = False   # once True, all slides are data slides
    current_site_exists = False  # whether the current section's site already exists

    for si in range(len(new_prs.slides)):
        src = new_prs.slides[si]
        has_10col = any(sh.has_table and len(sh.table.columns) == 10 for sh in src.shapes)
        has_8col  = any(sh.has_table and len(sh.table.columns) == 8  for sh in src.shapes)

        if has_10col:
            # ── Summary slide ────────────────────────────────────────────────
            seen_first_summary = True
            # Determine the site name from the table data rows
            current_site = None
            for sh in src.shapes:
                if not sh.has_table or len(sh.table.columns) != 10:
                    continue
                for ri in range(1, len(sh.table.rows)):
                    cells = [c.text.strip() for c in sh.table.rows[ri].cells]
                    if len(cells) >= 5 and cells[4]:
                        current_site = cells[4]
                        break
                break

            current_site_exists = current_site in base_site_summary_idx

            if current_site_exists:
                # Insert new rows into the EXISTING summary table — no duplicate slide
                base_slide = base_prs.slides[base_site_summary_idx[current_site]]
                for sh_src in src.shapes:
                    if not sh_src.has_table or len(sh_src.table.columns) != 10:
                        continue
                    new_tbl = sh_src.table
                    for sh_base in base_slide.shapes:
                        if not sh_base.has_table or len(sh_base.table.columns) != 10:
                            continue
                        base_row_count = len(sh_base.table.rows) - 1  # exclude header
                        for ri in range(1, len(new_tbl.rows)):
                            cells = [c.text.strip() for c in new_tbl.rows[ri].cells]
                            if not any(cells):
                                continue
                            rec = {
                                'ticket_no':    cells[1]  if len(cells) > 1 else '',
                                'veh_no':       cells[2]  if len(cells) > 2 else '',
                                'material':     cells[3]  if len(cells) > 3 else '',
                                'source_site':  cells[4]  if len(cells) > 4 else '',
                                'date':         cells[5]  if len(cells) > 5 else '',
                                'time':         cells[6]  if len(cells) > 6 else '',
                                'e_token':      cells[7]  if len(cells) > 7 else '',
                                'accepted':     cells[8]  if len(cells) > 8 else 'NO',
                                'reject_reason':cells[9]  if len(cells) > 9 else '',
                            }
                            _append_row_to_table(sh_base, rec, base_row_count + ri)
                        break  # only one summary table per slide
                    break
                # Do NOT copy this slide — rows already merged above
            else:
                # Brand-new site — copy the summary slide normally
                _copy_slide_to_base(src)

        elif not seen_first_summary:
            # Before the first summary = heading / cover slide → SKIP (prevents duplicate titles)
            continue

        else:
            # Detail slide (8-col table, photo slide, etc.) — always copy
            _copy_slide_to_base(src)

    buf = io.BytesIO()
    base_prs.save(buf)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════════════
#  ROUTES — Admin Panel
# ═══════════════════════════════════════════════════════════════════════════════

@app.route("/admin")
@login_required
def admin_panel():
    if session.get("role")!="admin":
        return redirect(url_for("dashboard"))
    return render_template_string(ADMIN_HTML, name=session.get("name",""))

@app.route("/api/admin/users")
@login_required
@admin_required
def admin_users():
    conn=get_db()
    users=conn.execute(
        "SELECT id,username,name,role,email,email_verified,created_at FROM users ORDER BY created_at DESC"
    ).fetchall()
    conn.close()
    return jsonify([dict(u) for u in users])

@app.route("/api/admin/activity")
@login_required
@admin_required
def admin_activity():
    limit  = int(request.args.get("limit", 200))
    action = request.args.get("action", "").strip().upper()
    conn   = get_db()
    if action:
        rows = conn.execute(
            "SELECT * FROM activity_log WHERE action=? ORDER BY ts DESC LIMIT ?",
            (action, limit)
        ).fetchall()
    else:
        rows = conn.execute(
            "SELECT * FROM activity_log ORDER BY ts DESC LIMIT ?", (limit,)
        ).fetchall()
    conn.close()
    return jsonify([dict(r) for r in rows])

@app.route("/api/admin/activity/download")
@login_required
@admin_required
def admin_activity_download():
    conn=get_db()
    rows=conn.execute("SELECT * FROM activity_log ORDER BY ts DESC").fetchall()
    conn.close()
    wb=openpyxl.Workbook(); ws=wb.active; ws.title="Activity Log"
    ws.append(["ID","Username","Action","Detail","IP","Timestamp"])
    for r in rows: ws.append(list(r))
    buf=io.BytesIO(); wb.save(buf); buf.seek(0)
    return send_file(buf,mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                     as_attachment=True,download_name=f"activity_log_{datetime.now().strftime('%Y%m%d')}.xlsx")

@app.route("/api/health")
def health():
    """Lightweight health check — also reports RAM usage and CT module status."""
    import gc as _gc
    _gc.collect()
    mem_info = {"ram_mb": None}
    try:
        import resource as _res
        rss = _res.getrusage(_res.RUSAGE_SELF).ru_maxrss
        mem_info["ram_mb"] = round(rss / 1024 if rss > 100000 else rss / 1024 / 1024, 1)
    except Exception:
        pass
    return jsonify({
        "status": "ok",
        "version": "2.0.0",
        "app": "APSG (Staging Ground) Report",
        "ct_module": "registered" if _CT_AVAILABLE else f"FAILED: {_ct_err_detail[:200]}",
        "memory": mem_info,
    })

@app.teardown_request
def _gc_after_request(exc=None):
    """
    Run GC after every request on heavy routes to keep RAM low on free tier.
    Only collects generation 0 (fast) on normal routes; full collect on CT/WB builds.
    """
    import gc as _gc
    path = request.path if request else ''
    if any(seg in path for seg in ('/build/', '/process', '/generate', '/phase2')):
        _gc.collect()   # full collect after heavy operations
    else:
        _gc.collect(0)  # generation-0 only — fast, < 1 ms



@app.route("/api/quote")
def api_quote():
    """
    Returns a random 'Interesting Thing' — rotates across categories:
    productivity tips, tech facts, history snippets, nature facts, life insights.
    Public endpoint, no auth needed.
    """
    import random as _rand
    _cards = [
        # ── Productivity Tips ─────────────────────────────────────────────────
        {"cat": "💡 Productivity Tip", "text": "The 2-minute rule: if a task takes less than 2 minutes, do it now instead of scheduling it later."},
        {"cat": "💡 Productivity Tip", "text": "Writing down 3 priorities the night before reduces morning decision fatigue by 40%."},
        {"cat": "💡 Productivity Tip", "text": "A 5-minute break every 52 minutes keeps your focus sharper than working 2 hours straight."},
        {"cat": "💡 Productivity Tip", "text": "Closing unused browser tabs before starting work reduces cognitive load and speeds up thinking."},
        {"cat": "💡 Productivity Tip", "text": "The best time to check email is twice a day — not 40 times. It triples deep work time."},
        {"cat": "💡 Productivity Tip", "text": "Done imperfectly is worth infinitely more than perfect but never started."},
        # ── Tech Facts ────────────────────────────────────────────────────────
        {"cat": "🔧 Tech Fact", "text": "The first computer bug was a real insect — a moth found stuck in a relay of the Harvard Mark II in 1947."},
        {"cat": "🔧 Tech Fact", "text": "Excel's maximum row count — 1,048,576 — is exactly 2 to the power of 20."},
        {"cat": "🔧 Tech Fact", "text": "The Python programming language is named after Monty Python's Flying Circus, not the snake."},
        {"cat": "🔧 Tech Fact", "text": "WiFi was accidentally invented while trying to detect evaporating black holes. Science is unpredictable."},
        {"cat": "🔧 Tech Fact", "text": "Google's search index contains over 100 million gigabytes of data — and it doubles every few years."},
        {"cat": "🔧 Tech Fact", "text": "A single modern smartphone has more computing power than all of NASA had during the 1969 moon landing."},
        # ── History Snippets ─────────────────────────────────────────────────
        {"cat": "📜 Did You Know", "text": "Singapore became an independent nation in 1965 — and within 30 years became one of the world's top 10 economies."},
        {"cat": "📜 Did You Know", "text": "The Post-it Note was invented accidentally in 1968 when a scientist tried to create a super-strong adhesive."},
        {"cat": "📜 Did You Know", "text": "The Eiffel Tower grows about 15 cm taller in summer because heat expands the iron structure."},
        {"cat": "📜 Did You Know", "text": "Construction helmets (hard hats) were first used in 1919 during the Hoover Dam project — saving hundreds of lives."},
        {"cat": "📜 Did You Know", "text": "The world's first spreadsheet program, VisiCalc (1979), sold 700,000 copies in 2 years — the first 'killer app'."},
        # ── Nature & Science ─────────────────────────────────────────────────
        {"cat": "🌿 Nature Fact", "text": "A tree planted today will remove approximately 1 tonne of CO₂ over its lifetime. Small actions have long impact."},
        {"cat": "🌿 Nature Fact", "text": "Octopuses have three hearts, blue blood, and can solve complex puzzles — they are genuinely alien-level intelligent."},
        {"cat": "🌿 Nature Fact", "text": "The Amazon rainforest produces 20% of the world's oxygen — often called 'the lungs of the Earth'."},
        {"cat": "🌿 Nature Fact", "text": "Ants have been farming fungi for over 50 million years — long before humans discovered agriculture."},
        # ── Life Insights ────────────────────────────────────────────────────
        {"cat": "✨ Insight", "text": "Compound interest is called the 8th wonder of the world — it applies equally to money, skills, and habits."},
        {"cat": "✨ Insight", "text": "You don't rise to the level of your goals — you fall to the level of your systems. Build better systems."},
        {"cat": "✨ Insight", "text": "The people who read 20 minutes a day end up reading 1.8 million words a year more than those who don't."},
        {"cat": "✨ Insight", "text": "Most overnight successes took 10 years. The work was invisible; only the result was sudden."},
        {"cat": "✨ Insight", "text": "Asking 'what went well today?' before sleeping rewires the brain toward optimism over time. Simple but powerful."},
        # ── Work & Data ──────────────────────────────────────────────────────
        {"cat": "📊 Work & Data", "text": "Bad data costs businesses an average of $12.9 million per year. Clean data is not a detail — it's the foundation."},
        {"cat": "📊 Work & Data", "text": "A report that takes 2 hours to build manually but saves 1 hour per week pays back in 2 weeks. Automate early."},
        {"cat": "📊 Work & Data", "text": "The first step of any analysis is not calculation — it's understanding what question you're actually trying to answer."},
        {"cat": "📊 Work & Data", "text": "Visualising data is not decoration — it activates a different part of the brain and reveals patterns numbers hide."},
    ]
    card = _rand.choice(_cards)
    return jsonify({"quote": card["text"], "category": card["cat"]})

@app.route("/api/admin/user/<username>/password", methods=["POST"])
@login_required
@admin_required
def admin_change_password(username):
    """Admin changes any user's password."""
    data = request.get_json(force=True, silent=True) or {}
    new_pw = data.get("password","").strip()
    if len(new_pw) < 6:
        return jsonify({"ok":False,"error":"Password must be at least 6 characters"}), 400
    conn = get_db()
    row = conn.execute("SELECT id FROM users WHERE username=?", (username,)).fetchone()
    if not row:
        conn.close()
        return jsonify({"ok":False,"error":"User not found"}), 404
    conn.execute("UPDATE users SET password_hash=?,plaintext_pw=? WHERE username=?", (hash_pw(new_pw), new_pw, username))
    conn.commit()
    conn.close()
    log_activity(session["username"], "ADMIN_CHANGE_PW", f"Changed password for user: {username}")
    return jsonify({"ok":True})

@app.route("/api/admin/user/<username>/view_password", methods=["GET"])
@login_required
@admin_required
def admin_view_password(username):
    """
    Return the stored plaintext password if available.
    If plaintext_pw is set (stored at registration/change), return it directly.
    Otherwise generate and set a new temp password (old behaviour).
    """
    conn = get_db()
    row = conn.execute(
        "SELECT id,name,plaintext_pw FROM users WHERE username=?", (username,)
    ).fetchone()
    if not row:
        conn.close()
        return jsonify({"ok":False,"error":"User not found"}), 404
    stored_pw = (row["plaintext_pw"] or "").strip()
    if stored_pw:
        # Return the plaintext password stored at registration/last change
        conn.close()
        log_activity(session["username"], "ADMIN_VIEW_PW", f"Viewed password for: {username}")
        return jsonify({"ok":True,"password":stored_pw,"username":username,"name":row["name"],"source":"stored"})
    # Fallback: generate temp password
    import secrets, string
    temp_pw = ''.join(secrets.choice(string.ascii_letters + string.digits) for _ in range(12))
    conn.execute("UPDATE users SET password_hash=?,plaintext_pw=? WHERE username=?",
                 (hash_pw(temp_pw), temp_pw, username))
    conn.commit()
    conn.close()
    log_activity(session["username"], "ADMIN_VIEW_PW", f"Generated temp password for: {username}")
    return jsonify({"ok":True,"password":temp_pw,"username":username,"name":row["name"],"source":"generated"})

@app.route("/api/admin/verify_master", methods=["POST"])
@login_required
@admin_required
def admin_verify_master():
    """Verify the admin's own password before allowing sensitive operations."""
    data    = request.get_json(force=True, silent=True) or {}
    pw      = str(data.get("password",""))
    me      = session.get("username")
    conn    = get_db()
    row     = conn.execute("SELECT password_hash FROM users WHERE username=?", (me,)).fetchone()
    conn.close()
    if row and verify_pw(pw, row["password_hash"]):
        # Issue a short-lived token stored in session so the UI can proceed
        token = str(random.randint(100000, 999999))
        session["master_token"] = token
        session["master_token_exp"] = (_dt.datetime.utcnow() + _dt.timedelta(minutes=5)).isoformat()
        return jsonify({"ok": True, "token": token})
    log_activity(me, "MASTER_PW_FAIL", "Wrong master password entered in admin panel")
    return jsonify({"ok": False, "error": "Incorrect admin password"}), 401

@app.route("/api/admin/user/<username>/reset_password", methods=["POST"])
@login_required
@admin_required
def admin_reset_password(username):
    """Admin resets any user's password — requires prior master password verification."""
    data  = request.get_json(force=True, silent=True) or {}
    token = str(data.get("master_token",""))
    new_pw = str(data.get("new_password","")).strip()
    confirm = str(data.get("confirm","")).strip()

    # Validate master token
    stored_token = session.get("master_token","")
    token_exp    = session.get("master_token_exp","")
    if not stored_token or token != stored_token:
        return jsonify({"ok": False, "error": "Master authentication required. Please verify your password first."}), 403
    if _dt.datetime.utcnow().isoformat() > token_exp:
        session.pop("master_token", None)
        session.pop("master_token_exp", None)
        return jsonify({"ok": False, "error": "Master token expired. Please re-authenticate."}), 403

    if len(new_pw) < 6:
        return jsonify({"ok": False, "error": "Password must be at least 6 characters"}), 400
    if new_pw != confirm:
        return jsonify({"ok": False, "error": "Passwords do not match"}), 400

    conn = get_db()
    try:
        row = conn.execute("SELECT id FROM users WHERE username=?", (username,)).fetchone()
        if not row:
            return jsonify({"ok": False, "error": "User not found"}), 404
        conn.execute("UPDATE users SET password_hash=?,plaintext_pw=? WHERE username=?",
                     (hash_pw(new_pw), new_pw, username))
        conn.commit()
        log_activity(session["username"], "ADMIN_RESET_PW", f"Reset password for: {username}")
        # Invalidate master token after use
        session.pop("master_token", None)
        session.pop("master_token_exp", None)
        return jsonify({"ok": True})
    finally:
        conn.close()

@app.route("/api/admin/user/<username>/delete", methods=["DELETE"])
@login_required
@admin_required
def admin_delete_user(username):
    """Admin deletes a user (cannot delete own account or built-in admin)."""
    if username == session.get("username"):
        return jsonify({"ok":False,"error":"Cannot delete your own account"}), 400
    if username == "admin":
        return jsonify({"ok":False,"error":"Cannot delete the built-in admin account"}), 400
    conn = get_db()
    conn.execute("DELETE FROM users WHERE username=?", (username,))
    conn.commit()
    conn.close()
    log_activity(session["username"], "ADMIN_DELETE_USER", f"Deleted user: {username}")
    return jsonify({"ok":True})

@app.route("/api/admin/analytics")
@login_required
@admin_required
def admin_analytics():
    conn = get_db()
    users = conn.execute("SELECT id,username,name,role,created_at FROM users ORDER BY created_at DESC").fetchall()
    # Login attempts
    logins = conn.execute("""SELECT username, COUNT(*) as cnt, MAX(ts) as last_login
        FROM activity_log WHERE action='LOGIN' GROUP BY username ORDER BY last_login DESC""").fetchall()
    # All activity summary per user
    activity_sum = conn.execute("""SELECT username, COUNT(*) as total_actions,
        MIN(ts) as first_seen, MAX(ts) as last_seen,
        GROUP_CONCAT(DISTINCT action) as actions_used
        FROM activity_log GROUP BY username ORDER BY last_seen DESC""").fetchall()
    # Login attempts (failed = no matching login row within same minute)
    attempts = conn.execute("""SELECT username, COUNT(*) as attempts, MAX(ts) as last_ts
        FROM activity_log WHERE action IN ('LOGIN','LOGIN_FAIL')
        GROUP BY username ORDER BY last_ts DESC""").fetchall()
    total_users = len(users)
    total_logins = conn.execute("SELECT COUNT(*) FROM activity_log WHERE action='LOGIN'").fetchone()[0]
    total_actions = conn.execute("SELECT COUNT(*) FROM activity_log").fetchone()[0]
    conn.close()
    return jsonify({
        'total_users': total_users,
        'total_logins': total_logins,
        'total_actions': total_actions,
        'users': [dict(u) for u in users],
        'logins': [dict(l) for l in logins],
        'activity_summary': [dict(a) for a in activity_sum],
    })

@app.route("/api/admin/analytics/download")
@login_required
@admin_required
def admin_analytics_download():
    conn = get_db()
    users = conn.execute("SELECT id,username,name,role,created_at FROM users ORDER BY created_at").fetchall()
    activity = conn.execute("""SELECT a.username, u.name, a.action, a.detail, a.ip, a.ts
        FROM activity_log a LEFT JOIN users u ON a.username=u.username
        ORDER BY a.ts DESC""").fetchall()
    conn.close()
    wb = openpyxl.Workbook()
    # Users sheet
    ws1 = wb.active; ws1.title = 'Users'
    ws1.append(['ID','Username','Name','Role','Registered'])
    hf = Font(bold=True,color='FFFFFF'); hfill = PatternFill('solid',fgColor='1E3A8A')
    for cell in ws1[1]: cell.font=hf; cell.fill=hfill; cell.alignment=Alignment(horizontal='center')
    for u in users: ws1.append(list(u))
    for col in ws1.columns:
        ws1.column_dimensions[col[0].column_letter].width = max(len(str(c.value or '')) for c in col)+4
    # Activity sheet
    ws2 = wb.create_sheet('Activity Log')
    ws2.append(['Username','Full Name','Action','Detail','IP','Timestamp'])
    for cell in ws2[1]: cell.font=hf; cell.fill=hfill; cell.alignment=Alignment(horizontal='center')
    for row in activity: ws2.append(list(row))
    for col in ws2.columns:
        ws2.column_dimensions[col[0].column_letter].width = max(len(str(c.value or '')) for c in col)+4
    # Summary sheet
    ws3 = wb.create_sheet('Usage Summary')
    ws3.append(['Username','Total Actions','First Seen','Last Seen'])
    for cell in ws3[1]: cell.font=hf; cell.fill=hfill; cell.alignment=Alignment(horizontal='center')
    conn2 = get_db()
    rows = conn2.execute("""SELECT username,COUNT(*) as total,MIN(ts),MAX(ts)
        FROM activity_log GROUP BY username ORDER BY total DESC""").fetchall()
    conn2.close()
    for r in rows: ws3.append(list(r))
    buf = io.BytesIO(); wb.save(buf); buf.seek(0)
    return send_file(buf,mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                     as_attachment=True,download_name=f'APSG_Analytics_{datetime.now().strftime("%Y%m%d")}.xlsx')

# ═══════════════════════════════════════════════════════════════════════════════
#  ROUTES — Daily Report API (engine-powered, no Streamlit needed)
# ═══════════════════════════════════════════════════════════════════════════════

@app.route("/api/daily/upload", methods=["POST"])
@login_required
def daily_upload():
    """Load and validate the Online export file, return metadata + date range."""
    try:
        f = request.files.get("file")
        if not f:
            return jsonify({"error": "No file uploaded"}), 400
        result = load_and_validate_file(f)
        if result["error"]:
            return jsonify({"error": result["error"]}), 400
        df = result["df"]
        min_date, max_date = get_date_range(df)
        # Store df in session via temp file
        import tempfile, pickle
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pkl",
                                          dir=OUTPUT_DIR)
        pickle.dump(df, tmp); tmp.close()
        return jsonify({
            "ok": True,
            "rows": len(df),
            "filename": f.filename,
            "tmp_path": tmp.name,
            "min_date": str(min_date),
            "max_date": str(max_date),
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/daily/generate", methods=["POST"])
@login_required
def daily_generate():
    """Generate the Online report Excel file."""
    try:
        import pickle
        data = request.get_json(force=True, silent=True) or {}
        tmp_path = data.get("tmp_path", "")
        filter_date_str = data.get("filter_date", "")
        corrections = data.get("corrections", {})

        if not tmp_path or not os.path.exists(tmp_path):
            return jsonify({"error": "Session expired — please re-upload the file"}), 400

        with open(tmp_path, "rb") as fh:
            df = pickle.load(fh)

        from datetime import date as dt_date
        filter_date = dt_date.fromisoformat(filter_date_str) if filter_date_str else None
        if not filter_date:
            return jsonify({"error": "Filter date required"}), 400

        # Build corrections dict with int keys
        corr_int = {int(k): v for k, v in corrections.items()}

        def _build_flag_map(corr_dict):
            fm = {}
            for idx, corr in corr_dict.items():
                decision = corr.get("Accepted", "Yes")
                red_cols = ({"Accepted", "Out Weight", "Net Weight"}
                            if str(decision).strip().lower() == "no"
                            else {"Out Weight", "Net Weight"})
                fm[idx] = {"flagged": True, "red_cols": red_cols}
            return fm

        flag_map = _build_flag_map(corr_int)
        excel_buf, fname, stats = generate_report(
            df, filter_date, filter_date,
            corrections=corr_int, flag_map=flag_map,
        )
        log_activity(session["username"], "DAILY_REPORT", fname)
        return send_file(
            excel_buf,
            mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            as_attachment=True,
            download_name=fname,
        )
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/daily/validate", methods=["POST"])
@login_required
def daily_validate():
    """Run all validations on filtered data and return results."""
    try:
        import pickle
        data = request.get_json(force=True, silent=True) or {}
        tmp_path = data.get("tmp_path", "")
        filter_date_str = data.get("filter_date", "")
        corrections = data.get("corrections", {})

        if not tmp_path or not os.path.exists(tmp_path):
            return jsonify({"error": "Session expired — please re-upload the file"}), 400

        with open(tmp_path, "rb") as fh:
            df = pickle.load(fh)

        from datetime import date as dt_date
        filter_date = dt_date.fromisoformat(filter_date_str) if filter_date_str else None
        if not filter_date:
            return jsonify({"error": "Filter date required"}), 400

        preview = filter_preview(df, filter_date, filter_date)
        if preview.empty:
            return jsonify({"ok": True, "rows": 0, "incomplete": [], "errors": [], "stats": {}})

        corr_int = {int(k): v for k, v in corrections.items()}
        val = validate_and_flag(preview)
        incomplete_list = []
        for idx in val["incomplete_accepted"]:
            row = preview.loc[idx]
            incomplete_list.append({
                "idx": int(idx),
                "token": str(row.get("Token", idx)),
                "in_weight": safe_float(row.get("In Weight")),
                "veh": str(row.get("Vehicle Number", "") or ""),
            })

        # Apply corrections to preview
        preview_corr = preview.copy()
        for idx, corr in corr_int.items():
            if idx in preview_corr.index:
                if "Accepted"   in corr: preview_corr.at[idx, "Accepted"]   = corr["Accepted"]
                if "Out Weight" in corr: preview_corr.at[idx, "Out Weight"] = float(corr["Out Weight"])
                if "Net Weight" in corr: preview_corr.at[idx, "Net Weight"] = float(corr["Net Weight"])

        nw_errors = validate_net_weights(preview_corr)
        errors = [f"E-Token [{e['token']}] – Ledger mismatch" for e in nw_errors]

        # Stats
        acc_mask = preview_corr["Accepted"].astype(str).str.strip().str.lower().isin(
            ("yes","1","true","accepted"))
        rej_mask = preview_corr["Accepted"].astype(str).str.strip().str.lower().isin(
            ("no","0","false","rejected","reject"))
        acc_rows = preview_corr[acc_mask]
        stats = {
            "accepted": int(acc_mask.sum()),
            "rejected": int(rej_mask.sum()),
            "wi": safe_col_sum(acc_rows, "In Weight"),
            "wo": safe_col_sum(acc_rows, "Out Weight"),
            "nw": safe_col_sum(acc_rows, "Net Weight"),
        }

        return jsonify({
            "ok": True,
            "rows": len(preview),
            "incomplete": incomplete_list,
            "errors": errors,
            "stats": stats,
            "all_resolved": len(incomplete_list) == 0 or all(
                item["idx"] in corr_int for item in incomplete_list),
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/daily/wb_upload", methods=["POST"])
@login_required
def daily_wb_upload():
    """Load WB file."""
    try:
        import pickle
        f = request.files.get("file")
        if not f:
            return jsonify({"error": "No file"}), 400
        result = load_wb_file(f)
        if result["error"]:
            return jsonify({"error": result["error"]}), 400
        df = result["df"]
        min_date, max_date = wb_get_date_range(df)
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pkl", dir=OUTPUT_DIR)
        pickle.dump(df, tmp); tmp.close()
        return jsonify({
            "ok": True, "rows": len(df), "filename": f.filename,
            "tmp_path": tmp.name,
            "min_date": str(min_date) if min_date else "",
            "max_date": str(max_date) if max_date else "",
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500

import tempfile

@app.route("/api/daily/wb_pivot", methods=["POST"])
@login_required
def daily_wb_pivot():
    """
    Build WB pivot using the exact same logic as the standalone Daily Report app.
    Runs:
      1. wb_filter_by_date          — date filter
      2. Normalize Accepted (0/1 → Yes/No)
      3. Filter accepted rows only (remove Accepted=No)
      4. wb_find_incomplete_rows    — detect missing Date Out / Time Out
      5. wb_apply_row_decisions     — apply user decisions for incomplete rows
      6. wb_apply_unified_logic     — sync Online rejections + resolve blank Out Weights via E-Token
      7. wb_net_weight_validation   — per-row In−Out=Net check
      8. wb_validate_etoken_match   — E-Token cross-comparison with Online
      9. wb_build_pivot             — build pivot table
      10. Compute wb_stats from wb_df_for_processing (INTENDED user decisions,
          same as standalone) so comparison table matches Online
    """
    try:
        import pickle
        data = request.get_json(force=True, silent=True) or {}
        wb_tmp          = data.get("wb_tmp_path", "")
        filter_date_str = data.get("filter_date", "")
        wb_decisions    = data.get("wb_decisions", {})
        online_tmp      = data.get("online_tmp_path", "")
        online_date_str = data.get("online_date", filter_date_str)
        corrections_raw = data.get("corrections", {})

        if not wb_tmp or not os.path.exists(wb_tmp):
            return jsonify({"error": "WB session expired — re-upload WB file"}), 400

        with open(wb_tmp, "rb") as fh:
            wb_df_raw = pickle.load(fh)

        from datetime import date as dt_date
        filter_date = dt_date.fromisoformat(filter_date_str) if filter_date_str else None
        if not filter_date:
            return jsonify({"error": "Filter date required"}), 400

        # ── Step 1: Date filter ──────────────────────────────────────────────
        wb_df_filtered = wb_filter_by_date(wb_df_raw, filter_date)
        if wb_df_filtered.empty:
            return jsonify({"error": f"No WB records for {filter_date_str}"}), 400

        # ── Step 2: Normalize Accepted (0/1 → Yes/No) ───────────────────────
        wb_df_filtered = wb_df_filtered.copy()
        wb_df_filtered["Accepted"] = wb_df_filtered["Accepted"].apply(_wb_norm_acc)

        wb_total  = len(wb_df_filtered)
        _acc_mask = wb_df_filtered["Accepted"].str.strip().str.lower() == "yes"
        wb_rej_raw = int((~_acc_mask).sum())

        # ── Step 3: Filter accepted rows only ───────────────────────────────
        wb_df_accepted = wb_df_filtered[_acc_mask].copy()
        if wb_df_accepted.empty:
            return jsonify({"error": "No accepted records remain after filtering Accepted=No rows."}), 400

        # ── Step 4: Find incomplete rows (missing Date Out / Time Out) ───────
        incomplete_rows = wb_find_incomplete_rows(wb_df_accepted)

        # ── Step 5: Apply user row decisions ────────────────────────────────
        dec_int = {int(k): v for k, v in wb_decisions.items()}

        # Check if all incomplete rows have decisions — if not, return them for UI
        if not incomplete_rows.empty:
            undecided = [
                {
                    "idx": int(idx),
                    "etoken": str(row.get("E-Token", f"Row {idx}")),
                    "in_weight": safe_float(row.get("In Weight", 0)),
                }
                for idx, row in incomplete_rows.iterrows()
                if int(idx) not in dec_int
            ]
            if undecided:
                return jsonify({"ok": True, "incomplete": undecided, "rows": [], "wb_stats": {}, "errors": []})

        _user_rejected_count = sum(1 for d in dec_int.values() if d.get("decision") == 0)
        wb_rej_raw_total = wb_rej_raw + _user_rejected_count

        wb_df_for_processing, _ = wb_apply_row_decisions(wb_df_accepted, dec_int)

        # ── Step 6: Load Online data for cross-comparison ───────────────────
        online_filtered = pd.DataFrame()
        corrections = {}

        if online_tmp and os.path.exists(online_tmp):
            try:
                with open(online_tmp, "rb") as fh:
                    online_df_raw = pickle.load(fh)
                corrections = {int(k): v for k, v in corrections_raw.items()}
                online_filtered = filter_preview(online_df_raw, filter_date, filter_date)
            except Exception:
                pass  # graceful degradation — proceed without Online comparison

        # ── Step 7: wb_apply_unified_logic ──────────────────────────────────
        # Matches standalone: syncs Online rejections + resolves blank Out Weights via E-Token
        wb_proc = wb_apply_unified_logic(wb_df_for_processing, online_filtered, corrections)
        wb_accepted_df = wb_proc["wb_accepted_df"]

        # ── Step 8: Collect validation errors (same as standalone) ──────────
        wb_all_errors = []

        # Net weight validation on wb_df_for_processing (pre-sync, matches standalone)
        nw_val_errors = wb_net_weight_validation(wb_df_for_processing)
        for e in nw_val_errors:
            wb_all_errors.append(f"E-Token [{e['etoken']}] – Ledger mismatch")

        # Also include errors from unified logic pass
        for e in wb_proc.get("nw_errors", []):
            msg = f"E-Token [{e['etoken']}] – Ledger mismatch"
            if not any(f"E-Token [{e['etoken']}]" in m and "Ledger mismatch" in m for m in wb_all_errors):
                wb_all_errors.append(msg)

        # E-Token cross-comparison with Online (only if Online data available)
        if not online_filtered.empty:
            etoken_result = wb_validate_etoken_match(
                wb_accepted_df, online_filtered, corrections,
                online_rejected_tokens=wb_proc.get("synced_tokens", set()),
            )
            for tok in etoken_result.get("wb_only", []):
                wb_all_errors.append(f"E-Token [{tok}] – Mismatch between Online & Weighbridge")
            for tok in etoken_result.get("online_only", []):
                wb_all_errors.append(f"E-Token [{tok}] – Mismatch between Online & Weighbridge")
            for tok in etoken_result.get("wb_dupes", []):
                wb_all_errors.append(f"E-Token [{tok}] – Duplicate E-Token in Weighbridge data")

        # ── Step 9: Build pivot ──────────────────────────────────────────────
        pivot_df = wb_build_pivot(wb_accepted_df)

        rows = []
        for _, r in pivot_df.iterrows():
            rows.append({
                "type":  r.get("_row_type", "data"),
                "label": str(r.get("Row Labels", "")),
                "loads": safe_int(r.get("Sum of Loads", 0)) if r.get("_row_type") != "mat_header" else 0,
                "wi":    safe_float(r.get("Sum of Weight In (T)", 0)),
                "wo":    safe_float(r.get("Sum of Weight Out (T)", 0)),
                "nw":    safe_float(r.get("Sum of Net Weight (T)", 0)),
            })

        # ── Step 10: Stats — use wb_df_for_processing (INTENDED accepted count) ─
        # Matches standalone: uses pre-sync df so stats align with what user decided,
        # not what was removed by Online-sync (which would cause count mismatch display).
        _wb_proc_acc_mask = (
            wb_df_for_processing["Accepted"].astype(str).str.strip().str.lower().isin(
                ("yes", "1", "true", "accepted")
            )
            if not wb_df_for_processing.empty and "Accepted" in wb_df_for_processing.columns
            else pd.Series(dtype=bool)
        )
        _wb_proc_accepted = (
            wb_df_for_processing[_wb_proc_acc_mask]
            if not wb_df_for_processing.empty and not _wb_proc_acc_mask.empty
            else pd.DataFrame()
        )
        wb_stats = {
            "accepted": int(_wb_proc_acc_mask.sum()) if not _wb_proc_acc_mask.empty else 0,
            "rejected": wb_rej_raw_total,
            "wi": safe_col_sum(_wb_proc_accepted, "In Weight"),
            "wo": safe_col_sum(_wb_proc_accepted, "Out Weight"),
            "nw": safe_col_sum(_wb_proc_accepted, "Net Weight"),
        }

        # ── Excel pivot download ─────────────────────────────────────────────
        dl_bytes, dl_fname = wb_pivot_to_excel(pivot_df, filter_date=filter_date)
        pivot_tmp = os.path.join(OUTPUT_DIR, f"pivot_{uuid.uuid4().hex}.xlsx")
        with open(pivot_tmp, "wb") as fh:
            fh.write(dl_bytes)

        return jsonify({
            "ok": True,
            "rows": rows,
            "wb_stats": wb_stats,
            "errors": wb_all_errors,
            "pivot_tmp": pivot_tmp,
            "pivot_fname": dl_fname,
        })
    except Exception as e:
        import traceback as _tb
        return jsonify({"error": str(e), "trace": _tb.format_exc()}), 500


@app.route("/api/daily/wb_pivot_download")
@login_required
def daily_wb_pivot_download():
    path = request.args.get("path","")
    fname = request.args.get("fname","pivot.xlsx")
    if not path or not os.path.exists(path):
        return jsonify({"error":"File not found"}),404
    return send_file(path, as_attachment=True, download_name=fname,
                     mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")


# ═══════════════════════════════════════════════════════════════════════════════
#  ROUTE — Rectification Report (Daily Report → Action Required)
# ═══════════════════════════════════════════════════════════════════════════════

@app.route("/api/daily/generate_rr", methods=["POST"])
@login_required
def daily_generate_rr():
    """
    Generate a Rectification Report .docx for a single Action Required row.

    Expected JSON payload
    ---------------------
    {
        "tmp_path"     : str,   # path to the pickled online DataFrame
        "token"        : str,   # E-Token of the affected row
        "rr_serial"    : str,   # 4-digit serial e.g. "0290"
        "accepted"     : str,   # "YES" or "NO"
        "out_weight"   : float, # operator-entered Out Weight (0 if NO)
        "net_weight"   : float, # computed Net Weight
        "weight_label" : str,   # optional page-5 heading override (YES only)
        "reason"       : str    # "Accepted Towing Vehicle" | "Rejected Towing Vehicle" | "Late Time / Breakdown"
    }
    """
    try:
        import pickle
        from datetime import timedelta

        data         = request.get_json(force=True, silent=True) or {}
        tmp_path     = data.get("tmp_path", "")
        token        = str(data.get("token", "")).strip()
        rr_serial    = str(data.get("rr_serial", "")).strip()
        accepted     = str(data.get("accepted", "")).upper().strip()
        out_weight   = data.get("out_weight", 0)
        net_weight   = data.get("net_weight", 0)
        weight_label = str(data.get("weight_label", "")).strip()
        reason       = str(data.get("reason", REASON_A)).strip() or REASON_A
        filter_date  = str(data.get("filter_date", "")).strip()  # YYYY-MM-DD from UI

        # Validate inputs
        if not tmp_path or not os.path.exists(tmp_path):
            return jsonify({"error": "Session expired — please re-upload the Online file"}), 400
        if not token:
            return jsonify({"error": "E-Token is required"}), 400
        if not rr_serial:
            return jsonify({"error": "RR Serial Number is required"}), 400
        if accepted not in ("YES", "NO"):
            return jsonify({"error": "Accepted must be YES or NO"}), 400

        # Load Online DataFrame
        with open(tmp_path, "rb") as fh:
            df = pickle.load(fh)

        # Build action_data dict
        # For NO: out_weight = in_weight (from JS corrections), net_weight = 0
        # For YES: user-entered out_weight and computed net_weight
        # Use explicit None/missing check (not truthiness) so 0 is preserved
        action_data = {
            "ACCEPTED":       accepted,
            "OUT WEIGHT":     str(out_weight) if out_weight not in (None, "", "None") else "",
            "NET WEIGHT":     str(net_weight) if net_weight not in (None, "", "None") else "0",
            "UNLADEN WEIGHT": str(out_weight) if (accepted == "YES" and out_weight) else "",
        }
        # Override for NO: ensure net_weight is always exactly 0
        if accepted == "NO":
            action_data["NET WEIGHT"] = "0"
            # Out weight for NO = In Weight (truck rejected, nothing unloaded)
            if not action_data["OUT WEIGHT"]:
                action_data["OUT WEIGHT"] = str(out_weight) if out_weight else ""

        # Fetch original row
        before_dict = fetch_row_by_token(df, token, source="online")
        after_dict  = apply_user_updates(before_dict, action_data)

        # Parse dates
        arr_dt = _parse_dt(before_dict.get("DATETIME ARRIVAL") or "")
        rpt_dt = arr_dt + timedelta(days=1)
        b_code = _MONTH_B.get(arr_dt.month, 44)
        rr_line = f"Rectification Report No. RR/B-{b_code}/{arr_dt.year}/{rr_serial}"

        # weight_label is now ALWAYS derived from reason (UI dropdown removed)
        # Option A / B → "Refer Unladen Weight"
        # Option C     → "Out Weight from Weighbridge Indicator: X T"
        from rectification_report import REASON_C
        out_w_str = str(out_weight).strip() if out_weight else ""
        if reason == REASON_C:
            weight_label = (
                f"Out Weight from Weighbridge Indicator: {out_w_str} T"
                if out_w_str else "Out Weight from Weighbridge Indicator"
            )
        else:
            weight_label = "Refer Unladen Weight"

        # Generate table images
        tbl1_jpg, _ = fetch_and_generate(
            df=df, token=token, source="online",
            is_table2=False, override_values=None, dpi=300,
            force_outnet_yellow=True,
        )
        tbl2_jpg, _ = fetch_and_generate(
            df=df, token=token, source="online",
            is_table2=True, override_values=action_data, dpi=300,
        )

        # Build .docx
        docx_bytes, filename = build_rr_docx(
            token                = token,
            rr_serial            = rr_serial,
            rr_line              = rr_line,
            before_dict          = before_dict,
            after_dict           = after_dict,
            action_data          = action_data,
            arr_dt               = arr_dt,
            rpt_dt               = rpt_dt,
            tbl1_jpg             = tbl1_jpg,
            tbl2_jpg             = tbl2_jpg,
            excel_screenshot_jpg = None,
            weight_label         = weight_label,
            reason               = reason,
            filter_date_str      = filter_date,
        )

        log_activity(session["username"], "RECTIFICATION_REPORT", filename)
        return send_file(
            io.BytesIO(docx_bytes),
            mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            as_attachment=True,
            download_name=filename,
        )

    except ValueError as ve:
        return jsonify({"error": str(ve)}), 400
    except Exception as e:
        return jsonify({"error": str(e), "trace": traceback.format_exc()}), 500


# ═══════════════════════════════════════════════════════════════════════════════
#  HTML TEMPLATES
# ═══════════════════════════════════════════════════════════════════════════════

AUTH_HTML = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>APSG (Staging Ground) Report — Sign In</title>
<link rel="preconnect" href="https://fonts.googleapis.com">
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&family=Poppins:wght@600;700;800&display=swap" rel="stylesheet">
<style>
*{box-sizing:border-box;margin:0;padding:0;}
:root{
  --bg:#080C1A;--card-bg:rgba(13,18,35,0.92);
  --indigo:#6366F1;--indigo-l:#818CF8;--cyan:#22D3EE;
  --purple:#A855F7;--green:#10B981;--red:#F87171;
  --text:#E8EEF8;--muted:#64748B;--border:rgba(99,102,241,0.18);
}
body{font-family:'Inter',system-ui,sans-serif;min-height:100vh;
  display:flex;align-items:flex-start;justify-content:center;
  padding-top:5vh;overflow-y:auto;}
.bg-mesh{position:fixed;inset:0;z-index:0;pointer-events:none;
  background:
    radial-gradient(ellipse 70% 60% at 15% 15%,rgba(99,102,241,.12) 0%,transparent 65%),
    radial-gradient(ellipse 60% 50% at 85% 85%,rgba(168,85,247,.10) 0%,transparent 65%),
    radial-gradient(ellipse 40% 35% at 50% 50%,rgba(34,211,238,.05) 0%,transparent 65%);
}
.grid-lines{position:fixed;inset:0;z-index:0;pointer-events:none;opacity:.025;
  background-image:linear-gradient(var(--indigo) 1px,transparent 1px),
    linear-gradient(90deg,var(--indigo) 1px,transparent 1px);
  background-size:60px 60px;}
.wrap{width:100%;max-width:440px;padding:1.5rem;position:relative;z-index:1;padding-top:.5rem;}
.brand{text-align:center;margin-bottom:1.8rem;}
.brand-logo{display:inline-flex;align-items:center;justify-content:center;
  width:64px;height:64px;border-radius:18px;margin-bottom:1rem;
  background:linear-gradient(135deg,rgba(99,102,241,.25),rgba(168,85,247,.15));
  border:1px solid rgba(99,102,241,.3);font-size:28px;
  box-shadow:0 0 40px rgba(99,102,241,.2);}
.brand-title{font-family:'Poppins',sans-serif;font-size:1.45rem;font-weight:800;
  color:var(--text);letter-spacing:-.03em;line-height:1.2;}
.brand-sub{font-size:.73rem;color:var(--muted);margin-top:.3rem;font-weight:400;letter-spacing:.03em;}
.card{background:rgba(8,14,38,0.72);border:1px solid rgba(255,255,255,.10);
  border-radius:20px;padding:2rem 2.2rem;
  backdrop-filter:blur(20px);-webkit-backdrop-filter:blur(20px);
  box-shadow:0 8px 48px rgba(0,0,0,.45);}
.card-header{font-size:.95rem;font-weight:700;color:var(--text);
  text-align:center;margin-bottom:1.6rem;display:flex;
  align-items:center;justify-content:center;gap:.5rem;letter-spacing:-.01em;}
.field{margin-bottom:1rem;}
.field label{display:block;font-size:.67rem;font-weight:700;color:var(--muted);
  letter-spacing:.12em;text-transform:uppercase;margin-bottom:.4rem;}
.field input{width:100%;padding:.78rem 1rem;
  background:rgba(6,10,24,.8);border:1.5px solid rgba(30,41,86,.8);
  border-radius:12px;color:var(--text);font-size:.9rem;
  font-family:'Inter',sans-serif;transition:all .2s;}
.field input:focus{outline:none;border-color:var(--indigo);
  box-shadow:0 0 0 3px rgba(99,102,241,.18);background:rgba(10,14,32,.9);}
.field input::placeholder{color:rgba(100,116,139,.5);font-weight:300;}
.btn{width:100%;padding:.85rem;border-radius:12px;border:none;cursor:pointer;
  font-size:.88rem;font-weight:700;font-family:'Inter',sans-serif;
  letter-spacing:.01em;transition:all .22s;position:relative;overflow:hidden;}
.btn-primary{
  background:linear-gradient(135deg,#4338CA 0%,#6366F1 50%,#818CF8 100%);
  color:#fff;box-shadow:0 4px 20px rgba(99,102,241,.4);margin-top:.3rem;}
.btn-primary::before{content:'';position:absolute;inset:0;
  background:linear-gradient(135deg,transparent,rgba(255,255,255,.1),transparent);
  transform:translateX(-100%);transition:transform .5s;}
.btn-primary:hover{transform:translateY(-2px);
  box-shadow:0 8px 32px rgba(99,102,241,.6);}
.btn-primary:hover::before{transform:translateX(100%);}
.btn-primary:active{transform:translateY(0) scale(.98);}
.btn-outline{background:transparent;border:1.5px solid rgba(99,102,241,.3);
  color:var(--indigo-l);margin-top:.65rem;}
.btn-outline:hover{background:rgba(99,102,241,.1);border-color:rgba(99,102,241,.55);}
.divider{display:flex;align-items:center;gap:.75rem;margin:.8rem 0;}
.divider::before,.divider::after{content:'';flex:1;height:1px;background:rgba(99,102,241,.12);}
.divider span{font-size:.68rem;color:rgba(100,116,139,.6);white-space:nowrap;}
.alert{border-radius:10px;padding:.65rem .9rem;font-size:.78rem;
  margin-bottom:.9rem;display:none;font-weight:500;}
.alert.show{display:block;}
.alert-error{background:rgba(239,68,68,.08);border:1px solid rgba(239,68,68,.2);color:#F87171;}
.alert-success{background:rgba(16,185,129,.08);border:1px solid rgba(16,185,129,.2);color:#34D399;}
.footer-links{text-align:center;margin-top:1rem;font-size:.72rem;color:var(--muted);}
.footer-links a{color:var(--indigo-l);text-decoration:none;font-weight:600;}
.footer-links a:hover{color:var(--cyan);}
.tag{display:inline-block;background:rgba(99,102,241,.1);
  border:1px solid rgba(99,102,241,.2);border-radius:6px;
  padding:.15rem .5rem;font-size:.62rem;font-weight:600;
  color:var(--indigo-l);letter-spacing:.05em;margin-bottom:1.4rem;}



/* Action cards */
.action-card { background: rgba(10,16,42,.85) !important; border-color: rgba(245,158,11,.5) !important; }
.rr-panel { background: rgba(6,12,38,.88) !important; border-color: rgba(99,102,241,.35) !important; }

/* Stats / text helpers */
.sec-hint, .upload-hint { color: #9BB8E0 !important; }
.chip-ok { color: #4ADE80 !important; }
.chip-info { color: #818CF8 !important; }
.chip-wait { color: #F59E0B !important; }

/* ── Developed by Karthik — fixed footer ── */
.dev-credit {
  position: fixed; bottom: 0; left: 0; right: 0; z-index: 9999;
  text-align: center; padding: .3rem 1rem;
  background: rgba(5, 8, 22, 0.75); backdrop-filter: blur(8px);
  border-top: 1px solid rgba(99,102,241,.2);
  font-size: 11px; font-weight: 600; color: rgba(160,180,220,.75);
  letter-spacing: .06em; font-family: 'Inter', system-ui, sans-serif;
  pointer-events: none; user-select: none;
}

/* ═══ GLOBAL BACKGROUND & TRANSPARENCY — v3 ═══════════════════════ */
html, body {
  background-image: url('/static/bg.jpg') !important;
  background-size: cover !important;
  background-position: center center !important;
  background-attachment: fixed !important;
  background-repeat: no-repeat !important;
  background-color: #08101E !important;
}
/* Single very-light overlay — image stays visible */
body::before {
  content: '' !important;
  position: fixed !important;
  inset: 0 !important;
  z-index: 0 !important;
  background: rgba(3, 7, 18, 0.45) !important;
  pointer-events: none !important;
}
body > * { position: relative; z-index: 1; }

/* ── Top-bar: fully transparent glass, no black border ── */
.top-bar {
  background: rgba(6, 10, 28, 0.55) !important;
  backdrop-filter: blur(18px) !important;
  -webkit-backdrop-filter: blur(18px) !important;
  border-bottom: 1px solid rgba(255,255,255,0.08) !important;
  position: sticky !important;
  top: 0 !important;
  z-index: 200 !important;
}

/* ── Cards / Panels — glass, no solid black fill ── */
.card, .section-card, .panel, .sec-card, .stats-box, .action-card {
  background: rgba(8, 14, 38, 0.70) !important;
  border: 1px solid rgba(255,255,255,0.10) !important;
  backdrop-filter: blur(16px) !important;
  -webkit-backdrop-filter: blur(16px) !important;
  box-shadow: 0 4px 32px rgba(0,0,0,0.35) !important;
}
.panel-head, .card-header {
  background: rgba(10, 18, 52, 0.72) !important;
  border-bottom: 1px solid rgba(255,255,255,0.07) !important;
}
.panel-body { background: rgba(5, 10, 30, 0.60) !important; }

/* ── Upload zones ── */
.upload-zone, .dz {
  background: rgba(6, 12, 34, 0.55) !important;
  border: 2px dashed rgba(99,102,241,0.55) !important;
}

/* ── Typography: all white/light ── */
body, h1, h2, h3, h4, p, span, div, td, th, label, a {
  color: #EEF3FF !important;
}
.hero-title, .brand-title, .card-title, .admin-hero-title {
  color: #FFFFFF !important;
  text-shadow: 0 2px 16px rgba(0,0,0,0.7) !important;
  font-weight: 800 !important;
}
.hero-sub, .brand-sub, .admin-hero-sub, .sec-hint {
  color: rgba(210, 225, 255, 0.80) !important;
}
.top-mini-brand, .top-page-label, .top-brand .brand-text {
  color: #FFFFFF !important;
  font-weight: 700 !important;
}
.back-btn {
  color: #C5D5FF !important;
  background: rgba(99,102,241,0.18) !important;
  border: 1px solid rgba(99,102,241,0.35) !important;
}
.back-btn:hover { background: rgba(99,102,241,0.32) !important; }

/* ── Inputs: legible on transparent backgrounds ── */
input[type=text], input[type=number], input[type=date],
input[type=password], select, textarea,
.input-num, .action-select, .rr-input, .rr-select, .date-input {
  background: rgba(5, 9, 28, 0.80) !important;
  border: 1.5px solid rgba(99,102,241,0.40) !important;
  color: #EEF3FF !important;
  font-size: 14px !important;
}
input::placeholder, textarea::placeholder {
  color: rgba(180, 200, 240, 0.50) !important;
}

/* ── Muted / secondary text ── */
.muted, .sec-label, [style*="color:#475569"],
[style*="color:#64748B"], [style*="color:#374167"] {
  color: rgba(190, 210, 255, 0.70) !important;
}

/* ── Action + RR cards ── */
.action-card {
  background: rgba(12, 18, 48, 0.78) !important;
  border-color: rgba(245,158,11,0.55) !important;
}
.rr-panel {
  background: rgba(8, 14, 42, 0.80) !important;
  border-color: rgba(99,102,241,0.40) !important;
}

/* ── Global font sizes ── */
body { font-size: 14px !important; }
h1 { font-size: 26px !important; }
h2 { font-size: 22px !important; }
h3, .hero-title { font-size: 22px !important; }
h4 { font-size: 17px !important; }

/* ── Dev-credit fixed footer ── */
.apsg-footer {
  position: fixed !important; bottom: 0 !important;
  left: 0 !important; right: 0 !important; z-index: 9999 !important;
  text-align: center !important; padding: .28rem 1rem !important;
  background: rgba(4, 7, 20, 0.70) !important;
  backdrop-filter: blur(8px) !important;
  border-top: 1px solid rgba(255,255,255,0.07) !important;
  font-size: 11px !important; font-weight: 600 !important;
  color: rgba(200, 220, 255, 0.70) !important;
  letter-spacing: .06em !important; pointer-events: none !important;
  user-select: none !important;
}
</style>
</head>
<body>
<div class="wrap">
  <div class="brand">
    <div class="brand-logo">📊</div>
    <div class="brand-title">APSG (Staging Ground) Report</div>
    <div class="brand-sub">Staging Ground Report System · Phase 3</div>
  </div>
  <!-- Login Card -->
  <div class="card" id="loginCard">
    <div class="card-header">🔐 Sign In to Continue</div>
    <div class="alert alert-error" id="loginError"></div>
    <div class="alert alert-success" id="loginSuccess"></div>
    <div class="field">
      <label>Username</label>
      <input type="text" id="loginUser" placeholder="Enter your username" autocomplete="username" autofocus>
    </div>
    <div class="field">
      <label>Password</label>
      <input type="password" id="loginPass" placeholder="Enter your password" autocomplete="current-password">
    </div>
    <div style="display:flex;align-items:center;gap:.6rem;margin:.5rem 0 .8rem;">
        <input type="checkbox" id="rememberMe" style="width:16px;height:16px;accent-color:#6366F1;cursor:pointer;margin:0;flex-shrink:0;">
        <label for="rememberMe" style="font-size:13px;color:#A0BAD8;cursor:pointer;font-weight:500;margin:0;user-select:none;">Remember me</label>
      </div>
      <button class="btn btn-primary" onclick="doLogin()">Sign In →</button>
    <div class="divider"><span>New to the system?</span></div>
    <button class="btn btn-outline" onclick="showRegister()">✨ Register for New User</button>
    <div class="footer-links" style="margin-top:.9rem;">
      <a href="#" onclick="showForgot();return false;" style="color:#F87171;">🔑 Forgot Password?</a>
    </div>
  </div>
  <!-- Register Card — 2-phase: All Details → Mobile OTP → Account Created -->
  <div class="card" id="registerCard" style="display:none">
    <div class="card-header">✨ Create New Account</div>
    <div class="alert alert-error"   id="regError"></div>
    <div class="alert alert-success" id="regSuccess"></div>

    <!-- Phase A: All registration fields -->
    <div id="regPhaseA">
      <div class="field">
        <label>Full Name</label>
        <input type="text" id="regName" placeholder="Enter your full name" autocomplete="name">
      </div>
      <div class="field">
        <label>Username</label>
        <input type="text" id="regUser" placeholder="Choose a unique username" autocomplete="username"
          oninput="clearFieldError('username')">
      </div>
      <div class="field">
        <label>Email Address</label>
        <input type="email" id="regEmail" placeholder="your@email.com" autocomplete="email"
          oninput="clearFieldError('email')">
      </div>
      <div class="field">
        <label>Mobile Number</label>
        <input type="tel" id="regMobile" placeholder="+65 XXXX XXXX" autocomplete="tel"
          oninput="clearFieldError('mobile')">
      </div>
      <div class="field">
        <label>Password</label>
        <input type="password" id="regPass" placeholder="Minimum 6 characters" autocomplete="new-password">
      </div>
      <div class="field">
        <label>Confirm Password</label>
        <input type="password" id="regPass2" placeholder="Re-enter password" autocomplete="new-password">
      </div>
      <button class="btn btn-primary" id="sendOtpBtn" onclick="doRegSendOtp()">
        📧 &nbsp;Send OTP to Email →
      </button>
    </div>

    <!-- Phase B: Email OTP verification -->
    <div id="regPhaseB" style="display:none">
      <div style="text-align:center;margin-bottom:1rem;padding:.7rem .9rem;
        background:rgba(16,185,129,.08);border:1px solid rgba(16,185,129,.2);border-radius:10px;">
        <div style="font-size:.82rem;color:#34D399;font-weight:600;">📧 OTP sent to email</div>
        <div id="otpMobileDisplay" style="font-size:.88rem;color:#A5B4FC;font-weight:700;margin-top:.2rem;"></div>
      </div>
      <div class="field">
        <label>Enter OTP Code</label>
        <input type="text" id="otpCode" placeholder="Enter the 6-digit code"
          maxlength="6" inputmode="numeric"
          style="text-align:center;font-size:1.4rem;font-weight:700;letter-spacing:.35em;">
      </div>
      <button class="btn btn-primary" id="verifyOtpBtn" onclick="doRegVerifyOtp()">✓ &nbsp;Verify OTP &amp; Create Account →</button>
      <button class="btn btn-outline" style="margin-top:.5rem;" onclick="doRegSendOtp(true)">↩ Resend OTP</button>
    </div>

    <div class="divider"><span>Already registered?</span></div>
    <button class="btn btn-outline" onclick="showLogin()">← Back to Sign In</button>
  </div>

  <!-- Forgot Password Card — 3-phase: Email → OTP → New Password -->
  <div class="card" id="forgotCard" style="display:none">
    <div class="card-header">🔑 Reset Your Password</div>
    <div class="alert alert-error"   id="fpError"></div>
    <div class="alert alert-success" id="fpSuccess"></div>

    <!-- Phase A: Enter email -->
    <div id="fpPhaseA">
      <div class="field">
        <label>Registered Email Address</label>
        <input type="email" id="fpEmail" placeholder="your@email.com" autocomplete="email">
      </div>
      <button class="btn btn-primary" id="fpSendBtn" onclick="doFpSendOtp()"
        style="background:linear-gradient(135deg,#991B1B,#DC2626,#F87171);">
        📧 &nbsp;Send Reset Code →
      </button>
    </div>

    <!-- Phase B: OTP verification -->
    <div id="fpPhaseB" style="display:none">
      <div style="text-align:center;margin-bottom:1rem;padding:.7rem .9rem;
        background:rgba(239,68,68,.08);border:1px solid rgba(239,68,68,.2);border-radius:10px;">
        <div style="font-size:.82rem;color:#F87171;font-weight:600;">🔑 Reset code sent to</div>
        <div id="fpEmailDisplay" style="font-size:.88rem;color:#A5B4FC;font-weight:700;margin-top:.2rem;"></div>
      </div>
      <div class="field">
        <label>Enter Reset Code</label>
        <input type="text" id="fpOtpCode" placeholder="Enter the 6-digit code"
          maxlength="6" inputmode="numeric"
          style="text-align:center;font-size:1.4rem;font-weight:700;letter-spacing:.35em;">
      </div>
      <button class="btn btn-primary" onclick="doFpVerifyOtp()"
        style="background:linear-gradient(135deg,#991B1B,#DC2626,#F87171);">
        ✓ &nbsp;Verify Code →
      </button>
      <button class="btn btn-outline" style="margin-top:.5rem;" onclick="doFpSendOtp(true)">↩ Resend Code</button>
    </div>

    <!-- Phase C: Set new password (requires current password as extra verification) -->
    <div id="fpPhaseC" style="display:none">
      <div style="text-align:center;margin-bottom:1rem;padding:.6rem .9rem;
        background:rgba(16,185,129,.08);border:1px solid rgba(16,185,129,.25);border-radius:10px;">
        <div style="font-size:.88rem;color:#34D399;font-weight:700;">✅ Identity Verified — Set New Password</div>
      </div>
      <div class="field">
        <label>Current Password</label>
        <input type="password" id="fpCurrPass" placeholder="Enter your current password">
      </div>
      <div class="field">
        <label>New Password</label>
        <input type="password" id="fpNewPass" placeholder="Minimum 6 characters">
      </div>
      <div class="field">
        <label>Re-enter New Password</label>
        <input type="password" id="fpNewPass2" placeholder="Re-enter new password">
      </div>
      <button class="btn btn-primary" id="fpResetBtn" onclick="doFpReset()"
        style="background:linear-gradient(135deg,#065F38,#059669,#10B981);">
        🔒 &nbsp;Reset Password →
      </button>
    </div>

    <div class="divider"><span>Remembered your password?</span></div>
    <button class="btn btn-outline" onclick="showLogin()">← Back to Sign In</button>
  </div>

  <div class="footer-links" style="margin-top:1rem;">
    Need help? <a href="mailto:karthickkv02@gmail.com">Contact Admin</a>
    &nbsp;·&nbsp; Developed by <strong>Karthi</strong>
  </div>
</div>
<script>
function show(id){const el=document.getElementById(id);if(el)el.style.display='block';}
function hide(id){const el=document.getElementById(id);if(el)el.style.display='none';}
function showAlert(id,msg){
  const el=document.getElementById(id);if(!el)return;
  el.textContent=msg;el.style.display='block';el.classList.add('show');
}
function clearAlerts(ids){
  ids.forEach(id=>{
    const el=document.getElementById(id);if(!el)return;
    el.classList.remove('show');el.textContent='';el.style.display='none';
  });
}
function showRegister(){
  hide('loginCard');hide('forgotCard');show('registerCard');
  show('regPhaseA');hide('regPhaseB');
  clearAlerts(['loginError','loginSuccess','regError','regSuccess','fpError','fpSuccess']);
  const el=document.getElementById('regName');if(el)setTimeout(()=>el.focus(),80);
}
function showLogin(){
  hide('registerCard');hide('forgotCard');show('loginCard');
  clearAlerts(['loginError','loginSuccess','regError','regSuccess','fpError','fpSuccess']);
}
function showForgot(){
  hide('loginCard');hide('registerCard');show('forgotCard');
  show('fpPhaseA');hide('fpPhaseB');hide('fpPhaseC');
  const em=document.getElementById('fpEmail');if(em)em.value='';
  clearAlerts(['loginError','loginSuccess','regError','regSuccess','fpError','fpSuccess']);
  setTimeout(()=>{if(em)em.focus();},80);
}

// Clear per-field error highlight on input
function clearFieldError(field){clearAlerts(['regError']);}

async function doLogin(){
  const rememberMe=document.getElementById('rememberMe')?.checked;
  if(rememberMe) localStorage.setItem('apsg_remember_user',document.getElementById('loginUser').value);
  else localStorage.removeItem('apsg_remember_user');
  clearAlerts(['loginError','loginSuccess']);
  const user=document.getElementById('loginUser').value.trim();
  const pass=document.getElementById('loginPass').value;
  if(!user||!pass){showAlert('loginError','Please enter username and password.');return;}
  try{
    const r=await fetch('/login',{method:'POST',headers:{'Content-Type':'application/json'},
      body:JSON.stringify({username:user,password:pass})});
    const d=await r.json();
    if(d.ok) window.location.href='/';
    else showAlert('loginError',d.error||'Invalid username or password.');
  }catch(e){showAlert('loginError','Connection error. Please try again.');}
}

let _pendingEmail='', _pendingMobile='';

// Registration Phase A: validate all fields + check duplicates + send mobile OTP
async function doRegSendOtp(resend=false){
  clearAlerts(['regError','regSuccess']);
  const name   = document.getElementById('regName').value.trim();
  const user   = document.getElementById('regUser').value.trim();
  const email  = document.getElementById('regEmail').value.trim();
  const mobile = document.getElementById('regMobile').value.trim();
  const pass   = document.getElementById('regPass').value;
  const pass2  = document.getElementById('regPass2').value;

  if(!name||!user||!email||!mobile||!pass||!pass2){
    showAlert('regError','All fields are required.');return;
  }
  if(!email.includes('@')||!email.includes('.')){
    showAlert('regError','Please enter a valid email address.');return;
  }
  if(pass!==pass2){showAlert('regError','Passwords do not match.');return;}
  if(pass.length<6){showAlert('regError','Password must be at least 6 characters.');return;}

  const btn=document.getElementById('sendOtpBtn');
  btn.disabled=true;btn.textContent='Checking & Sending…';
  try{
    const r=await fetch('/api/reg_send_otp',{method:'POST',
      headers:{'Content-Type':'application/json'},
      body:JSON.stringify({name,username:user,email,mobile,password:pass,confirm:pass2})});
    const ct=r.headers.get('Content-Type')||'';
    let d;
    if(ct.includes('application/json')){d=await r.json();}
    else{d={ok:false,error:'Server error ('+r.status+'). Please try again.'};}
    btn.disabled=false;btn.textContent='📧 Send OTP to Email →';
    if(!d.ok){
      showAlert('regError',d.error||'Failed to send OTP.');
      // Highlight the specific duplicate field
      if(d.field){
        const fieldMap={username:'regUser',email:'regEmail',mobile:'regMobile'};
        const el=document.getElementById(fieldMap[d.field]);
        if(el){el.style.borderColor='#F87171';el.focus();}
      }
      return;
    }
    _pendingEmail=email;_pendingMobile=mobile;
    document.getElementById('otpMobileDisplay').textContent=email;
    document.getElementById('otpCode').value='';
    hide('regPhaseA');show('regPhaseB');
    clearAlerts(['regError','regSuccess']);
    if(d.dev_otp){
      showAlert('regSuccess','🔑 Dev mode — your OTP is: '+d.dev_otp);
    }else{
      showAlert('regSuccess','✅ OTP sent to '+email+'. Check your inbox (and spam folder).');
    }
    setTimeout(()=>document.getElementById('otpCode').focus(),100);
  }catch(e){
    btn.disabled=false;btn.textContent='📱 Send OTP to Mobile →';
    showAlert('regError','Connection error: '+e.message);
  }
}

// Registration Phase B: verify mobile OTP → create account
async function doRegVerifyOtp(){
  clearAlerts(['regError','regSuccess']);
  const code=document.getElementById('otpCode').value.trim();
  if(code.length!==6){showAlert('regError','Please enter the complete 6-digit code.');return;}
  const btn=document.getElementById('verifyOtpBtn');
  btn.disabled=true;btn.textContent='Creating account…';
  try{
    const r=await fetch('/api/reg_verify_otp',{method:'POST',
      headers:{'Content-Type':'application/json'},
      body:JSON.stringify({email:_pendingEmail,otp:code})});
    const d=await r.json();
    if(d.ok){
      showAlert('regSuccess','✅ Account created! Redirecting…');
      setTimeout(()=>window.location.href='/',900);
    }else{
      showAlert('regError',d.error||'Verification failed.');
      btn.disabled=false;btn.textContent='✓ Verify OTP & Create Account →';
    }
  }catch(e){
    showAlert('regError','Connection error: '+e.message);
    btn.disabled=false;btn.textContent='✓ Verify OTP & Create Account →';
  }
}

// Pre-fill remembered username
const _rem=localStorage.getItem('apsg_remember_user');
if(_rem){const el=document.getElementById('loginUser');if(el){el.value=_rem;document.getElementById('rememberMe').checked=true;}}

// ── Forgot Password flow ──────────────────────────────────────────────────
let _fpEmail='';

async function doFpSendOtp(resend=false){
  clearAlerts(['fpError','fpSuccess']);
  const email=document.getElementById('fpEmail').value.trim();
  if(!email||!email.includes('@')||!email.includes('.')){
    showAlert('fpError','Please enter your registered email address.');return;
  }
  const btn=document.getElementById('fpSendBtn');
  btn.disabled=true;btn.textContent='Sending…';
  try{
    const r=await fetch('/api/forgot_password',{method:'POST',
      headers:{'Content-Type':'application/json'},
      body:JSON.stringify({email})});
    const ct=r.headers.get('Content-Type')||'';
    let d;
    if(ct.includes('application/json')){
      d=await r.json();
    } else {
      d={ok:false,error:'Server error ('+r.status+'). Please try again.'};
    }
    btn.disabled=false;btn.textContent='📧 Send Reset Code →';
    if(!d.ok){showAlert('fpError',d.error||'Failed to send reset code.');return;}
    _fpEmail=email;
    document.getElementById('fpEmailDisplay').textContent=email;
    document.getElementById('fpOtpCode').value='';
    hide('fpPhaseA');show('fpPhaseB');
    clearAlerts(['fpError','fpSuccess']);
    if(d.dev_otp){
      showAlert('fpSuccess','🔑 Dev mode — your reset code is: '+d.dev_otp);
    } else {
      showAlert('fpSuccess','✅ Reset code sent! Check your inbox (and spam folder).');
    }
    setTimeout(()=>document.getElementById('fpOtpCode').focus(),100);
  }catch(e){
    btn.disabled=false;btn.textContent='📧 Send Reset Code →';
    showAlert('fpError','Connection error: '+e.message);
  }
}

async function doFpVerifyOtp(){
  clearAlerts(['fpError','fpSuccess']);
  const code=document.getElementById('fpOtpCode').value.trim();
  if(code.length!==6){showAlert('fpError','Please enter the complete 6-digit code.');return;}
  // We verify the OTP by attempting the reset — but here we just move to Phase C
  // (actual OTP verification happens server-side at reset time)
  hide('fpPhaseB');show('fpPhaseC');
  clearAlerts(['fpError','fpSuccess']);
  setTimeout(()=>document.getElementById('fpNewPass').focus(),100);
}

async function doFpReset(){
  clearAlerts(['fpError','fpSuccess']);
  const otp      = document.getElementById('fpOtpCode').value.trim();
  const currPass = document.getElementById('fpCurrPass').value;
  const pass     = document.getElementById('fpNewPass').value;
  const pass2    = document.getElementById('fpNewPass2').value;
  if(!currPass){showAlert('fpError','Please enter your current password.');return;}
  if(!pass||!pass2){showAlert('fpError','Please enter and confirm your new password.');return;}
  if(pass!==pass2){showAlert('fpError','New passwords do not match.');return;}
  if(pass.length<6){showAlert('fpError','New password must be at least 6 characters.');return;}
  const btn=document.getElementById('fpResetBtn');
  btn.disabled=true;btn.textContent='Resetting…';
  try{
    const r=await fetch('/api/reset_password',{method:'POST',
      headers:{'Content-Type':'application/json'},
      body:JSON.stringify({email:_fpEmail,otp,current_password:currPass,password:pass,confirm:pass2})});
    const d=await r.json();
    if(d.ok){
      showAlert('fpSuccess','✅ '+d.message);
      btn.disabled=false;btn.textContent='🔒 Reset Password →';
      setTimeout(()=>showLogin(),2000);
    } else {
      showAlert('fpError',d.error||'Reset failed. Please try again.');
      btn.disabled=false;btn.textContent='🔒 Reset Password →';
      if(d.field==='current_password'){
        document.getElementById('fpCurrPass').style.borderColor='#F87171';
        document.getElementById('fpCurrPass').focus();
      }
      if(d.error&&(d.error.includes('Incorrect code')||d.error.includes('expired'))){
        show('fpPhaseB');hide('fpPhaseC');
      }
    }
  }catch(e){
    showAlert('fpError','Connection error: '+e.message);
    btn.disabled=false;btn.textContent='🔒 Reset Password →';
  }
}

// Enter key
document.addEventListener('keydown',e=>{
  if(e.key!=='Enter') return;
  const lc=document.getElementById('loginCard');
  const fp=document.getElementById('forgotCard');
  const rc=document.getElementById('registerCard');
  const fpB=document.getElementById('fpPhaseB');
  const fpC=document.getElementById('fpPhaseC');
  const rgB=document.getElementById('regPhaseB');
  if(lc?.style.display!=='none'){doLogin();return;}
  if(fp?.style.display!=='none'){
    if(fpB?.style.display!=='none') doFpVerifyOtp();
    else if(fpC?.style.display!=='none') doFpReset();
    else doFpSendOtp();
    return;
  }
  if(rc?.style.display!=='none'){
    if(rgB?.style.display!=='none') doRegVerifyOtp();
    else doRegSendOtp();
    return;
  }
});
</script>
<div class="apsg-footer">✦ Internal Reporting Platform — APSG Staging Ground &nbsp;·&nbsp; Developed by Karthik</div>
</body>
</html>"""


DASHBOARD_HTML = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>Dashboard - APSG (Staging Ground) Report</title>
<link rel="preconnect" href="https://fonts.googleapis.com">
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800&family=Poppins:wght@600;700;800&display=swap" rel="stylesheet">
<style>
*{box-sizing:border-box;margin:0;padding:0;}
:root{
  --bg:#060912;
  --surface:#0D1120;
  --surface2:#111827;
  --indigo:#6366F1;
  --indigo-d:#4F46E5;
  --indigo-l:#A5B4FC;
  --cyan:#22D3EE;
  --purple:#A855F7;
  --green:#10B981;
  --amber:#F59E0B;
  --red:#F87171;
  --text:#F0F4FF;
  --text2:#CBD5E1;
  --muted:#64748B;
  --border:rgba(148,163,184,0.08);
  --border-h:rgba(148,163,184,0.18);
}
body{font-family:"Inter",system-ui,sans-serif;min-height:100vh;
  background:var(--bg);color:var(--text);font-size:14px;overflow-x:hidden;}

/* ── Subtle aurora — reduced opacity, less distraction ── */
.aurora{position:fixed;inset:0;z-index:0;pointer-events:none;overflow:hidden;}
.aurora-blob{position:absolute;border-radius:50%;filter:blur(110px);opacity:.07;animation:blobDrift 18s ease-in-out infinite;}
.aurora-blob:nth-child(1){width:800px;height:800px;background:#6366F1;top:-300px;left:-200px;animation-delay:0s;}
.aurora-blob:nth-child(2){width:600px;height:600px;background:#A855F7;bottom:-200px;right:-150px;animation-delay:-8s;}
.aurora-blob:nth-child(3){width:400px;height:400px;background:#22D3EE;top:40%;left:40%;animation-delay:-14s;}
@keyframes blobDrift{0%,100%{transform:translate(0,0);}50%{transform:translate(30px,-20px);}}

/* ── TOP BAR ── */
.top-bar{position:sticky;top:0;z-index:300;height:56px;display:flex;
  align-items:center;padding:0 2rem;gap:.75rem;
  background: rgba(4,8,22,0.50);backdrop-filter:blur(20px);
  border-bottom: 1px solid rgba(255,255,255,0.08);}
.top-brand{display:flex;align-items:center;gap:.6rem;}
.brand-mark{width:28px;height:28px;border-radius:8px;flex-shrink:0;
  background:linear-gradient(135deg,#4F46E5,#7C3AED);
  display:flex;align-items:center;justify-content:center;font-size:14px;
  box-shadow:0 0 0 1px rgba(99,102,241,.4);}
.brand-text{font-family:"Poppins",sans-serif;font-size:.82rem;font-weight:700;
  color:var(--text);letter-spacing:-.01em;}
.brand-text span{color:var(--indigo-l);}
.top-spacer{flex:1;}
.top-right{display:flex;align-items:center;gap:.4rem;}
.top-pill{border-radius:6px;padding:.28rem .8rem;font-size:.68rem;font-weight:600;
  text-decoration:none;transition:all .18s;border:1px solid transparent;
  white-space:nowrap;letter-spacing:.01em;}
.pill-user{background:rgba(99,102,241,.09);border-color:rgba(99,102,241,.2);color:var(--indigo-l);}
.pill-admin{background:rgba(245,158,11,.08);border-color:rgba(245,158,11,.2);color:var(--amber);}
.pill-logout{background:rgba(248,113,113,.06);border-color:rgba(248,113,113,.15);color:var(--red);}
.pill-admin:hover,.pill-logout:hover{opacity:.85;transform:translateY(-1px);}

/* ── PAGE ── */
.page{padding:2.2rem 1.8rem 5rem;max-width:1160px;margin:0 auto;position:relative;z-index:1;}

/* ── HERO ── */
.hero{text-align:center;padding:2rem 1rem 2.6rem;}
.hero-eyebrow{display:inline-flex;align-items:center;gap:.5rem;padding:.28rem .9rem;
  background:rgba(99,102,241,.07);border:1px solid rgba(99,102,241,.18);
  border-radius:4px;font-size:.66rem;color:var(--indigo-l);font-weight:600;
  letter-spacing:.12em;text-transform:uppercase;margin-bottom:1.2rem;}
.live-dot{width:6px;height:6px;border-radius:50%;background:#10B981;
  animation:livePulse 2s ease-in-out infinite;flex-shrink:0;}
@keyframes livePulse{0%,100%{box-shadow:0 0 0 0 rgba(16,185,129,.5);}60%{box-shadow:0 0 0 5px rgba(16,185,129,0);}}
.hero-title{font-family:"Poppins",sans-serif;
  font-size:clamp(1.75rem,4.5vw,2.8rem);font-weight:800;
  color:var(--text);line-height:1.1;margin-bottom:.6rem;letter-spacing:-.03em;}
.hero-sub{font-size:.84rem;color:var(--muted);margin-bottom:.45rem;}
.hero-user{font-size:.8rem;color:var(--indigo-l);font-weight:600;}

/* ── SECTION LABEL ── */
.section-label{font-size:.65rem;font-weight:700;letter-spacing:.14em;
  text-transform:uppercase;color:var(--muted);margin-bottom:1rem;
  display:flex;align-items:center;gap:.6rem;}
.section-label::after{content:"";flex:1;height:1px;background:var(--border);}

/* ── CARDS GRID ── */
.cards-grid{display:grid;grid-template-columns:repeat(auto-fit,minmax(300px,1fr));gap:1rem;}

.app-card{
  position:relative;display:block;text-decoration:none;color:inherit;
  background: rgba(8,14,40,0.72);
  border: 1px solid rgba(255,255,255,0.10);
  border-radius:14px;padding:1.6rem 1.5rem;
  transition:transform .22s ease,box-shadow .22s ease,border-color .22s ease,background .22s ease;
  cursor:pointer;overflow:hidden;}

/* Left accent bar — colour identifier, always visible */
.app-card::before{
  content:"";position:absolute;left:0;top:16px;bottom:16px;width:3px;
  border-radius:0 3px 3px 0;
  background:var(--accent,#6366F1);
  opacity:.7;transition:opacity .22s,top .22s,bottom .22s;}

/* Hover state — clear, readable, not over-glowing */
.app-card:hover{
  background:var(--surface2);
  border-color:var(--border-h);
  transform:translateY(-4px);
  box-shadow:0 12px 40px rgba(0,0,0,.45),0 0 0 1px var(--accent-border,rgba(99,102,241,.25));}
.app-card:hover::before{opacity:1;top:12px;bottom:12px;}

/* Active / click feedback */
.app-card:active{transform:translateY(-2px) scale(.99);transition-duration:.08s;}

/* Hover — title becomes full white */
.app-card:hover .card-title{color:#fff;}
/* Hover — description becomes clearly readable */
.app-card:hover .card-desc{color:var(--text2);}
/* Hover — number becomes accent colour */
.app-card:hover .card-num{color:var(--accent,#6366F1);opacity:1;}
/* Hover — arrow moves and brightens */
.app-card:hover .card-arrow{transform:translateX(5px);opacity:1;}

/* Card colour themes — only sets accent, no background changes */
.card-blue{--accent:#60A5FA;--accent-border:rgba(96,165,250,.22);}
.card-green{--accent:#34D399;--accent-border:rgba(52,211,153,.22);}
.card-purple{--accent:#C084FC;--accent-border:rgba(192,132,252,.22);}
.card-amber{--accent:#FBBF24;--accent-border:rgba(251,191,36,.22);}
.card-gray{--accent:#6B7280;--accent-border:rgba(107,114,128,.15);}
.card-teal{--accent:#2DD4BF;--accent-border:rgba(45,212,191,.22);}
.card-indigo{--accent:#818CF8;--accent-border:rgba(129,140,248,.22);}

.app-card.disabled{cursor:not-allowed;opacity:.35;pointer-events:none;}

/* Content */
.card-num{font-size:.58rem;font-weight:700;color:var(--muted);
  letter-spacing:.16em;text-transform:uppercase;margin-bottom:.6rem;
  transition:color .22s;}
.card-icon-wrap{display:flex;align-items:flex-start;justify-content:space-between;margin-bottom:.9rem;}
.card-icon{width:44px;height:44px;border-radius:10px;display:flex;
  align-items:center;justify-content:center;font-size:20px;flex-shrink:0;
  background:rgba(255,255,255,.04);border:1px solid var(--border);
  transition:background .22s,transform .22s;}
.app-card:hover .card-icon{background:rgba(255,255,255,.07);transform:scale(1.06);}
.card-arrow{color:var(--muted);font-size:1rem;transition:transform .22s,opacity .22s;opacity:.5;}

.card-title{font-family:"Poppins",sans-serif;font-size:.93rem;font-weight:700;
  color:var(--text);margin-bottom:.4rem;line-height:1.35;
  transition:color .22s;}
.card-desc{font-size:.76rem;color:var(--muted);line-height:1.7;
  transition:color .22s;}
.card-badge{display:inline-flex;align-items:center;gap:.28rem;
  padding:.18rem .65rem;border-radius:4px;
  font-size:.61rem;font-weight:700;margin-top:.9rem;letter-spacing:.05em;}
.badge-active{background:rgba(16,185,129,.09);color:#34D399;border:1px solid rgba(16,185,129,.2);}
.badge-soon{background:rgba(107,114,128,.07);color:var(--muted);border:1px solid var(--border);}

/* Card entry animation */
.app-card{animation:slideIn .4s cubic-bezier(.16,1,.3,1) both;}
.app-card:nth-child(1){animation-delay:.04s;}
.app-card:nth-child(2){animation-delay:.09s;}
.app-card:nth-child(3){animation-delay:.14s;}
.app-card:nth-child(4){animation-delay:.19s;}
.app-card:nth-child(5){animation-delay:.24s;}
.app-card:nth-child(6){animation-delay:.29s;}
@keyframes slideIn{from{opacity:0;transform:translateY(20px);}to{opacity:1;transform:none;}}

/* ── FOOTER ── */
.page-footer{text-align:center;padding:2rem 0 .5rem;margin-top:2.5rem;
  border-top:1px solid var(--border);}

.dev-credit{
  position:relative;display:inline-flex;align-items:center;gap:.5rem;
  background:rgba(99,102,241,.07);border:1px solid rgba(99,102,241,.22);
  border-radius:30px;padding:.4rem 1.1rem .4rem .55rem;cursor:pointer;
  font-size:.73rem;font-weight:600;color:var(--text2);
  transition:all .22s;user-select:none;}
.dev-credit:hover{border-color:rgba(99,102,241,.45);background:rgba(99,102,241,.13);
  color:#A5B4FC;transform:translateY(-2px);
  box-shadow:0 4px 16px rgba(99,102,241,.18);}
.dev-credit:active{transform:translateY(0);}

/* Avatar */
.karthi-avatar{
  width:28px;height:28px;border-radius:50%;flex-shrink:0;
  background:linear-gradient(135deg,#4F46E5,#7C3AED,#DB2777);
  display:flex;align-items:center;justify-content:center;font-size:14px;
  animation:gentleWave 3s ease-in-out infinite;}
@keyframes gentleWave{0%,100%{transform:rotate(0deg);}35%{transform:rotate(-8deg);}65%{transform:rotate(8deg);}}

/* Motivational popup animation */
@keyframes popIn{0%{opacity:0;transform:scale(.88) translateY(16px);}100%{opacity:1;transform:scale(1) translateY(0);}}

.footer-meta{font-size:.61rem;color:rgba(100,116,139,.4);margin-top:.55rem;letter-spacing:.04em;}

/* ── WAKE SCREEN ── */
.wake{position:fixed;inset:0;z-index:9999;display:flex;flex-direction:column;
  align-items:center;justify-content:center;gap:1rem;background: rgba(4,8,22,0.85);}
.wake-spinner{width:36px;height:36px;border:2px solid rgba(99,102,241,.15);
  border-top-color:var(--indigo);border-radius:50%;animation:spin 1s linear infinite;}
@keyframes spin{to{transform:rotate(360deg);}}
.wake-label{font-family:"Poppins",sans-serif;font-size:.9rem;font-weight:700;color:var(--text2);}
.wake-sub{font-size:.72rem;color:var(--muted);}

/* ── RESPONSIVE ── */
@media(max-width:600px){
  .cards-grid{grid-template-columns:1fr;}
  .top-bar{padding:0 1rem;}
  .page{padding:1.5rem 1rem 4rem;}
}



/* Action cards */
.action-card { background: rgba(10,16,42,.85) !important; border-color: rgba(245,158,11,.5) !important; }
.rr-panel { background: rgba(6,12,38,.88) !important; border-color: rgba(99,102,241,.35) !important; }

/* Stats / text helpers */
.sec-hint, .upload-hint { color: #9BB8E0 !important; }
.chip-ok { color: #4ADE80 !important; }
.chip-info { color: #818CF8 !important; }
.chip-wait { color: #F59E0B !important; }

/* ═══ GLOBAL BACKGROUND & TRANSPARENCY — v3 ═══════════════════════ */
html, body {
  background-image: url('/static/bg.jpg') !important;
  background-size: cover !important;
  background-position: center center !important;
  background-attachment: fixed !important;
  background-repeat: no-repeat !important;
  background-color: #08101E !important;
}
/* Single very-light overlay — image stays visible */
body::before {
  content: '' !important;
  position: fixed !important;
  inset: 0 !important;
  z-index: 0 !important;
  background: rgba(3, 7, 18, 0.45) !important;
  pointer-events: none !important;
}
body > * { position: relative; z-index: 1; }

/* ── Top-bar: fully transparent glass, no black border ── */
.top-bar {
  background: rgba(6, 10, 28, 0.55) !important;
  backdrop-filter: blur(18px) !important;
  -webkit-backdrop-filter: blur(18px) !important;
  border-bottom: 1px solid rgba(255,255,255,0.08) !important;
  position: sticky !important;
  top: 0 !important;
  z-index: 200 !important;
}

/* ── Cards / Panels — glass, no solid black fill ── */
.card, .section-card, .panel, .sec-card, .stats-box, .action-card {
  background: rgba(8, 14, 38, 0.70) !important;
  border: 1px solid rgba(255,255,255,0.10) !important;
  backdrop-filter: blur(16px) !important;
  -webkit-backdrop-filter: blur(16px) !important;
  box-shadow: 0 4px 32px rgba(0,0,0,0.35) !important;
}
.panel-head, .card-header {
  background: rgba(10, 18, 52, 0.72) !important;
  border-bottom: 1px solid rgba(255,255,255,0.07) !important;
}
.panel-body { background: rgba(5, 10, 30, 0.60) !important; }

/* ── Upload zones ── */
.upload-zone, .dz {
  background: rgba(6, 12, 34, 0.55) !important;
  border: 2px dashed rgba(99,102,241,0.55) !important;
}

/* ── Typography: all white/light ── */
body, h1, h2, h3, h4, p, span, div, td, th, label, a {
  color: #EEF3FF !important;
}
.hero-title, .brand-title, .card-title, .admin-hero-title {
  color: #FFFFFF !important;
  text-shadow: 0 2px 16px rgba(0,0,0,0.7) !important;
  font-weight: 800 !important;
}
.hero-sub, .brand-sub, .admin-hero-sub, .sec-hint {
  color: rgba(210, 225, 255, 0.80) !important;
}
.top-mini-brand, .top-page-label, .top-brand .brand-text {
  color: #FFFFFF !important;
  font-weight: 700 !important;
}
.back-btn {
  color: #C5D5FF !important;
  background: rgba(99,102,241,0.18) !important;
  border: 1px solid rgba(99,102,241,0.35) !important;
}
.back-btn:hover { background: rgba(99,102,241,0.32) !important; }

/* ── Inputs: legible on transparent backgrounds ── */
input[type=text], input[type=number], input[type=date],
input[type=password], select, textarea,
.input-num, .action-select, .rr-input, .rr-select, .date-input {
  background: rgba(5, 9, 28, 0.80) !important;
  border: 1.5px solid rgba(99,102,241,0.40) !important;
  color: #EEF3FF !important;
  font-size: 14px !important;
}
input::placeholder, textarea::placeholder {
  color: rgba(180, 200, 240, 0.50) !important;
}

/* ── Muted / secondary text ── */
.muted, .sec-label, [style*="color:#475569"],
[style*="color:#64748B"], [style*="color:#374167"] {
  color: rgba(190, 210, 255, 0.70) !important;
}

/* ── Action + RR cards ── */
.action-card {
  background: rgba(12, 18, 48, 0.78) !important;
  border-color: rgba(245,158,11,0.55) !important;
}
.rr-panel {
  background: rgba(8, 14, 42, 0.80) !important;
  border-color: rgba(99,102,241,0.40) !important;
}

/* ── Global font sizes ── */
body { font-size: 14px !important; }
h1 { font-size: 26px !important; }
h2 { font-size: 22px !important; }
h3, .hero-title { font-size: 22px !important; }
h4 { font-size: 17px !important; }

/* ── Dev-credit fixed footer ── */
.apsg-footer {
  position: fixed !important; bottom: 0 !important;
  left: 0 !important; right: 0 !important; z-index: 9999 !important;
  text-align: center !important; padding: .28rem 1rem !important;
  background: rgba(4, 7, 20, 0.70) !important;
  backdrop-filter: blur(8px) !important;
  border-top: 1px solid rgba(255,255,255,0.07) !important;
  font-size: 11px !important; font-weight: 600 !important;
  color: rgba(200, 220, 255, 0.70) !important;
  letter-spacing: .06em !important; pointer-events: none !important;
  user-select: none !important;
}
</style>
</head>
<body>


<!-- Wake Screen -->
<div class="wake" id="wakeScreen">
  <div class="wake-spinner"></div>
  <div class="wake-label">APSG Report</div>
  <div class="wake-sub" id="wakeMsg">Loading&hellip;</div>
</div>

<!-- Top Bar -->
<div class="top-bar">
  <div class="top-brand">
    <div class="brand-mark">&#9889;</div>
    <span class="brand-text">APSG <span>Report</span></span>
  </div>
  <div class="top-spacer"></div>
  <div class="top-right">
    <a href="/admin" id="adminBtn" class="top-pill pill-admin" style="display:none;">&#9881; Admin</a>
    <span class="top-pill pill-user">&#128100; {{ name }}</span>
    <a href="#" onclick="showProfileModal();return false;" class="top-pill pill-user" style="border-color:rgba(99,102,241,.25);" title="Change password">🔑</a>
    <a href="/logout" class="top-pill pill-logout">Sign Out</a>
  </div>
</div>

<!-- Main -->
<div class="page">

  <!-- Hero -->
  <div class="hero">
    <div class="hero-eyebrow"><span class="live-dot"></span>Internal Reporting Platform</div>
    <div class="hero-title">APSG (Staging Ground) Report</div>
    <div class="hero-sub">Staging Ground Report System &middot; Phase 3</div>
    <div class="hero-user">Welcome back, {{ name }}</div>
    {% if last_login %}<div style="font-size:.68rem;color:rgba(148,163,184,.5);margin-top:.3rem;">Last login: {{ last_login }}</div>{% endif %}
  </div>

  <!-- Cards -->
  <div class="section-label">Select a report to get started</div>

  <div class="cards-grid">

    <a href="/app/daily-report" class="app-card card-blue" style="--delay:.05s">
      <div class="card-icon-wrap">
        <div class="card-icon">&#128202;</div>
        <span class="card-arrow">&#8594;</span>
      </div>
      <div class="card-title">Daily Report</div>
      <div class="card-desc">Generate daily staging ground reports with Online vs WB comparison, validation, and Excel download.</div>
      <span class="card-badge badge-active">&#10003; Active</span>
    </a>

    <a href="/app/hourly-report" class="app-card card-indigo" style="--delay:.10s">
      <div class="card-icon-wrap">
        <div class="card-icon">&#9201;</div>
        <span class="card-arrow">&#8594;</span>
      </div>
      <div class="card-title">Hourly Report</div>
      <div class="card-desc">Upload WB data to view hourly SC / GE breakdown, queue status, wait time, and copy-ready Slack summaries.</div>
      <span class="card-badge badge-active">&#10003; Active</span>
    </a>

    <a href="/app/cga-report" class="app-card card-teal" style="--delay:.15s">
      <div class="card-icon-wrap">
        <div class="card-icon">&#128666;</div>
        <span class="card-arrow">&#8594;</span>
      </div>
      <div class="card-title">CGA Report</div>
      <div class="card-desc">Generate hourly trucks pivot table for CGA site codes. Upload source Excel/CSV, select date, download as PNG image.</div>
      <span class="card-badge badge-active">&#10003; Active</span>
    </a>

    <a href="/ct/" class="app-card card-gray" style="--delay:.20s">
      <div class="card-icon-wrap">
        <div class="card-icon">⏱</div>
        <span class="card-arrow">&#8594;</span>
      </div>
      <div class="card-title">Cycle Time Report</div>
      <div class="card-desc">Upload CT and Online Data files to generate the 7-sheet Excel workbook, Hourly Track summary, and PowerPoint with anomaly detection.</div>
      <span class="card-badge badge-active">&#10003; Active</span>
    </a>

    <a href="/app/excel-rejection" class="app-card card-green" style="--delay:.25s">
      <div class="card-icon-wrap">
        <div class="card-icon">&#128203;</div>
        <span class="card-arrow">&#8594;</span>
      </div>
      <div class="card-title">Excel Rejection Report</div>
      <div class="card-desc">Convert Excel data into formatted PowerPoint rejection reports grouped by source site.</div>
      <span class="card-badge badge-active">&#10003; Active</span>
    </a>

    <a href="/app/photo-merge" class="app-card card-purple" style="--delay:.30s">
      <div class="card-icon-wrap">
        <div class="card-icon">&#128444;</div>
        <span class="card-arrow">&#8594;</span>
      </div>
      <div class="card-title">PPT Alignment</div>
      <div class="card-desc">Merge Top Photo and Front Photo PPT files with auto token matching and layout control.</div>
      <span class="card-badge badge-active">&#10003; Active</span>
    </a>

    <a href="/app/ppt-rejection" class="app-card card-amber" style="--delay:.35s">
      <div class="card-icon-wrap">
        <div class="card-icon">&#128193;</div>
        <span class="card-arrow">&#8594;</span>
      </div>
      <div class="card-title">Monthly Rejection Report Filter</div>
      <div class="card-desc">Upload bulk PPT files, filter by E-Token/date/reason, and generate grouped rejection reports with ZIP export.</div>
      <span class="card-badge badge-active">&#10003; Active</span>
    </a>

  </div>

  <!-- Footer -->
  <div class="page-footer">
    <div class="dev-credit" onclick="showMotivation()" title="Click for a message">
      <span style="font-size:1.1rem;line-height:1;">&#129489;</span>
      Developed by&nbsp;<strong>Karthi</strong>
      <span style="font-size:.7rem;color:rgba(99,102,241,.7);margin-left:4px;">✦</span>
    </div>
    <div class="footer-meta">APSG (Staging Ground) Report &middot; v2.0 &middot; Internal Use Only</div>
  </div>

</div><!-- /page -->

<!-- Change Password Modal -->
<div id="profileModal" style="display:none;position:fixed;inset:0;z-index:9997;
  align-items:center;justify-content:center;background:rgba(2,5,18,.65);backdrop-filter:blur(6px);">
  <div style="max-width:360px;width:90%;background:rgba(10,16,42,.96);
    border:1px solid rgba(99,102,241,.35);border-radius:18px;padding:1.8rem 2rem;
    box-shadow:0 24px 80px rgba(0,0,0,.6);">
    <div style="font-size:.95rem;font-weight:700;margin-bottom:1.2rem;text-align:center;">🔑 Change Password</div>
    <div id="profileErr" style="display:none;font-size:.78rem;color:#F87171;background:rgba(248,113,113,.08);
      border:1px solid rgba(248,113,113,.2);border-radius:8px;padding:.5rem .75rem;margin-bottom:.8rem;"></div>
    <div id="profileOk" style="display:none;font-size:.78rem;color:#4ADE80;background:rgba(74,222,128,.08);
      border:1px solid rgba(74,222,128,.2);border-radius:8px;padding:.5rem .75rem;margin-bottom:.8rem;"></div>
    <div style="margin-bottom:.75rem;">
      <label style="font-size:.67rem;font-weight:700;color:rgba(148,163,184,.7);letter-spacing:.1em;text-transform:uppercase;display:block;margin-bottom:.35rem;">Current Password</label>
      <input type="password" id="proCurPw" style="width:100%;padding:.65rem .9rem;background:rgba(5,9,28,.8);border:1.5px solid rgba(99,102,241,.35);border-radius:10px;color:#EEF3FF;font-size:14px;">
    </div>
    <div style="margin-bottom:.75rem;">
      <label style="font-size:.67rem;font-weight:700;color:rgba(148,163,184,.7);letter-spacing:.1em;text-transform:uppercase;display:block;margin-bottom:.35rem;">New Password</label>
      <input type="password" id="proNewPw" style="width:100%;padding:.65rem .9rem;background:rgba(5,9,28,.8);border:1.5px solid rgba(99,102,241,.35);border-radius:10px;color:#EEF3FF;font-size:14px;">
    </div>
    <div style="margin-bottom:1.2rem;">
      <label style="font-size:.67rem;font-weight:700;color:rgba(148,163,184,.7);letter-spacing:.1em;text-transform:uppercase;display:block;margin-bottom:.35rem;">Confirm New Password</label>
      <input type="password" id="proConPw" style="width:100%;padding:.65rem .9rem;background:rgba(5,9,28,.8);border:1.5px solid rgba(99,102,241,.35);border-radius:10px;color:#EEF3FF;font-size:14px;">
    </div>
    <div style="display:flex;gap:.6rem;">
      <button onclick="submitProfilePw()" style="flex:1;padding:.75rem;border-radius:10px;border:none;
        background:linear-gradient(135deg,#4338CA,#6366F1);color:#fff;font-size:.88rem;font-weight:700;cursor:pointer;">
        Update Password
      </button>
      <button onclick="document.getElementById('profileModal').style.display='none'" style="padding:.75rem 1rem;border-radius:10px;
        border:1px solid rgba(148,163,184,.2);background:rgba(255,255,255,.04);
        color:rgba(148,163,184,.7);font-size:.88rem;cursor:pointer;">Cancel</button>
    </div>
  </div>
</div>

<script>
function showProfileModal(){
  document.getElementById('profileErr').style.display='none';
  document.getElementById('profileOk').style.display='none';
  document.getElementById('proCurPw').value='';
  document.getElementById('proNewPw').value='';
  document.getElementById('proConPw').value='';
  document.getElementById('profileModal').style.display='flex';
  setTimeout(()=>document.getElementById('proCurPw').focus(),80);
}
async function submitProfilePw(){
  const err=document.getElementById('profileErr');
  const ok=document.getElementById('profileOk');
  err.style.display='none'; ok.style.display='none';
  const cur=document.getElementById('proCurPw').value;
  const nw=document.getElementById('proNewPw').value;
  const con=document.getElementById('proConPw').value;
  if(!cur||!nw||!con){err.textContent='All fields required.';err.style.display='block';return;}
  try{
    const r=await fetch('/api/user/change_password',{method:'POST',
      headers:{'Content-Type':'application/json'},
      body:JSON.stringify({current_password:cur,new_password:nw,confirm:con})});
    const d=await r.json();
    if(d.ok){ok.textContent='✅ Password updated successfully.';ok.style.display='block';
      setTimeout(()=>document.getElementById('profileModal').style.display='none',1800);}
    else{err.textContent=d.error||'Failed.';err.style.display='block';}
  }catch(e){err.textContent='Connection error.';err.style.display='block';}
}
document.getElementById('profileModal').addEventListener('click',function(e){
  if(e.target===this) this.style.display='none';
});
document.addEventListener('keydown',e=>{
  if(e.key==='Escape') document.getElementById('profileModal').style.display='none';
});
</script>
<div id="motivPopup" style="
  display:none;position:fixed;inset:0;z-index:9998;align-items:center;justify-content:center;
  background:rgba(2,5,18,.65);backdrop-filter:blur(6px);transition:opacity .22s;">
  <div id="motivInner" style="
    max-width:390px;width:92%;background:rgba(10,16,42,.96);
    border:1px solid rgba(99,102,241,.35);border-radius:20px;
    padding:2rem 2.2rem;text-align:center;position:relative;
    box-shadow:0 24px 80px rgba(0,0,0,.6),0 0 0 1px rgba(99,102,241,.15);
    transition:transform .22s,opacity .22s;">
    <button id="motivCloseBtn" style="
      position:absolute;top:.8rem;right:1rem;background:none;border:none;
      color:rgba(148,163,184,.7);font-size:1.3rem;cursor:pointer;line-height:1;
      width:28px;height:28px;display:flex;align-items:center;justify-content:center;
      border-radius:50%;transition:all .15s;">✕</button>
    <div style="font-size:2.2rem;margin-bottom:.6rem;">&#129489;</div>
    <div id="motivSource" style="font-size:.72rem;font-weight:700;letter-spacing:.1em;
      text-transform:uppercase;color:rgba(99,102,241,.85);margin-bottom:.75rem;min-height:1.1rem;"></div>
    <div id="motivQuote" style="font-size:.95rem;font-weight:600;color:#EEF3FF;
      line-height:1.68;margin-bottom:.5rem;min-height:3.2rem;"></div>
    <div style="font-size:.67rem;color:rgba(148,163,184,.4);margin-top:.55rem;">
      — Karthi &nbsp;&#10024; &nbsp;&middot;&nbsp; Click Next for another interesting thing
    </div>
    <div style="margin-top:1.4rem;display:flex;gap:.6rem;justify-content:center;">
      <button id="motivNextBtn" style="
        padding:.5rem 1.2rem;border-radius:8px;border:1px solid rgba(99,102,241,.4);
        background:rgba(99,102,241,.12);color:#818CF8;font-size:.78rem;font-weight:600;
        cursor:pointer;transition:all .2s;">Next &#8594;</button>
      <button id="motivCloseBtn2" style="
        padding:.5rem 1.2rem;border-radius:8px;border:1px solid rgba(148,163,184,.18);
        background:rgba(255,255,255,.04);color:rgba(148,163,184,.7);font-size:.78rem;font-weight:600;
        cursor:pointer;transition:all .2s;">Close</button>
    </div>
  </div>
</div>

<script>
const role="{{ role }}";
if(role==="admin") document.getElementById("adminBtn").style.display="inline-flex";

// ── Card entrance animation ────────────────────────────────────────────────────
document.querySelectorAll('.app-card').forEach(card => {
  const delay = card.style.getPropertyValue('--delay') || '0s';
  card.style.opacity = '0';
  card.style.transform = 'translateY(28px)';
  card.style.transition = `opacity .5s ease ${delay}, transform .5s ease ${delay}`;
  requestAnimationFrame(() => requestAnimationFrame(() => {
    card.style.opacity = '1';
    card.style.transform = 'translateY(0)';
  }));
});

// ── Wake screen ────────────────────────────────────────────────────────────────
window.addEventListener("load", () => {
  const ws = document.getElementById("wakeScreen");
  const wm = document.getElementById("wakeMsg");
  fetch("/api/health").then(r => {
    wm.textContent = r.ok ? "Ready!" : "Server warming up...";
    setTimeout(() => {
      ws.style.opacity = "0"; ws.style.transition = "opacity .5s";
      setTimeout(() => ws.style.display = "none", 520);
    }, r.ok ? 300 : 1800);
  }).catch(() => setTimeout(() => ws.style.display = "none", 1500));
});

// ── Motivational popup — fixed close + dynamic interesting things ─────────────
const _cards = [
  {cat:"💡 Productivity Tip",  text:"The 2-minute rule: if a task takes less than 2 minutes, do it now."},
  {cat:"🔧 Tech Fact",         text:"Excel's maximum row count — 1,048,576 — is exactly 2 to the power of 20."},
  {cat:"📜 Did You Know",      text:"The Post-it Note was invented accidentally in 1968 when the adhesive came out too weak."},
  {cat:"✨ Insight",           text:"You don't rise to the level of your goals — you fall to the level of your systems."},
  {cat:"📊 Work & Data",       text:"Bad data costs businesses an average of $12.9M per year. Clean data is the foundation."},
  {cat:"🌿 Nature Fact",       text:"Ants have been farming fungi for over 50 million years — long before humans discovered agriculture."},
  {cat:"💡 Productivity Tip",  text:"Writing down 3 priorities the night before reduces morning decision fatigue by 40%."},
  {cat:"🔧 Tech Fact",         text:"Python was named after Monty Python's Flying Circus — not the snake."},
  {cat:"✨ Insight",           text:"Most overnight successes took 10 years. The work was invisible; only the result was sudden."},
  {cat:"📊 Work & Data",       text:"A report that saves 1 hour per week pays back 2 hours of build time in just 2 weeks."},
];
let _qi = 0, _motivOpen = false;

function _motivClose() {
  if (!_motivOpen) return;
  _motivOpen = false;
  const p   = document.getElementById('motivPopup');
  const inn = document.getElementById('motivInner');
  // Set transition FIRST, THEN change values
  p.style.transition   = 'opacity .22s';
  inn.style.transition = 'opacity .22s, transform .22s';
  p.style.opacity      = '0';
  inn.style.opacity    = '0';
  inn.style.transform  = 'scale(.92)';
  document.body.style.overflow = '';
  setTimeout(() => {
    p.style.display    = 'none';
    p.style.opacity    = '';
    inn.style.opacity  = '';
    inn.style.transform = '';
  }, 230);
}

async function showMotivation() {
  if (_motivOpen) return;
  const p   = document.getElementById('motivPopup');
  const inn = document.getElementById('motivInner');
  if (!p || !inn) return;

  // Show a local card immediately — no blank flash
  _qi = Math.floor(Math.random() * _cards.length);
  _setCard(_cards[_qi].cat, _cards[_qi].text);

  // Fetch a fresh interesting thing from server (non-blocking)
  fetch('/api/quote').then(r => r.ok ? r.json() : null).then(d => {
    if (d && d.quote) _setCard(d.category || '', d.quote);
  }).catch(() => {});

  // Animate open
  _motivOpen = true;
  document.body.style.overflow = 'hidden';
  p.style.display      = 'flex';
  p.style.opacity      = '0';
  inn.style.opacity    = '0';
  inn.style.transform  = 'scale(.9) translateY(12px)';
  p.style.transition   = 'none';
  inn.style.transition = 'none';
  requestAnimationFrame(() => requestAnimationFrame(() => {
    p.style.transition   = 'opacity .25s';
    inn.style.transition = 'opacity .25s, transform .28s cubic-bezier(.34,1.56,.64,1)';
    p.style.opacity      = '1';
    inn.style.opacity    = '1';
    inn.style.transform  = 'scale(1) translateY(0)';
  }));
}

function _setCard(cat, text) {
  const el  = document.getElementById('motivQuote');
  const src = document.getElementById('motivSource');
  if (el)  el.textContent  = text;
  if (src) src.textContent = cat || '';
}

function nextQuote() {
  const el  = document.getElementById('motivQuote');
  const src = document.getElementById('motivSource');
  if (!el) return;
  el.style.transition = 'opacity .18s'; el.style.opacity = '0';
  if (src) { src.style.transition = 'opacity .18s'; src.style.opacity = '0'; }
  setTimeout(() => {
    fetch('/api/quote').then(r => r.ok ? r.json() : null).then(d => {
      if (d && d.quote) _setCard(d.category||'', d.quote);
      else { _qi=(_qi+1)%_cards.length; _setCard(_cards[_qi].cat, _cards[_qi].text); }
      el.style.opacity='1'; if(src) src.style.opacity='1';
    }).catch(() => {
      _qi=(_qi+1)%_cards.length; _setCard(_cards[_qi].cat, _cards[_qi].text);
      el.style.opacity='1'; if(src) src.style.opacity='1';
    });
  }, 160);
}

// Wire up close buttons and events
document.addEventListener('DOMContentLoaded', () => {
  document.getElementById('motivCloseBtn').addEventListener('click',  _motivClose);
  document.getElementById('motivCloseBtn2').addEventListener('click', _motivClose);
  document.getElementById('motivNextBtn').addEventListener('click',   nextQuote);
  // Click outside (on the overlay backdrop) closes it
  document.getElementById('motivPopup').addEventListener('click', e => {
    if (e.target === document.getElementById('motivPopup')) _motivClose();
  });
});
// ESC key
document.addEventListener('keydown', e => { if (e.key === 'Escape') _motivClose(); });
</script>
<div class="apsg-footer">✦ Internal Reporting Platform — APSG Staging Ground &nbsp;·&nbsp; Developed by Karthik</div>
</body>
</html>"""


PPT_REJECTION_HTML = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>PPT Alignment — APSG (Staging Ground) Report</title>
<style>

/* ═══ MODERN UI BASE (Blue/Purple/Cyan Theme) ═══ */
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&family=Poppins:wght@600;700;800&display=swap');
*, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }
:root {
  --bg:transparent; --card-bg: rgba(13,18,35,0.88);
  --indigo: #6366F1; --indigo-l: #818CF8; --cyan: #22D3EE;
  --purple: #A855F7; --green: #10B981; --amber: #F59E0B; --red: #F87171;
  --text: #E8EEF8; --muted: #64748B; --border: rgba(99,102,241,0.15);
}
body {
  font-family: 'Inter', system-ui, -apple-system, sans-serif;
  min-height: 100vh; font-size: 14px; line-height: 1.6;
  background:transparent; color: var(--text);
  position: relative;
}
body::before {
  content: ''; position: fixed; inset: 0; z-index: 0; pointer-events: none;
  background:
    radial-gradient(ellipse 70% 55% at 5% 10%, rgba(99,102,241,.09) 0%, transparent 65%),
    radial-gradient(ellipse 55% 45% at 95% 90%, rgba(168,85,247,.07) 0%, transparent 65%);
}
body::after {
  content: ''; position: fixed; inset: 0; z-index: 0; pointer-events: none; opacity: .016;
  background-image: linear-gradient(var(--indigo) 1px, transparent 1px),
    linear-gradient(90deg, var(--indigo) 1px, transparent 1px);
  background-size: 60px 60px;
}

/* ── Modern Top Bar ── */
.top-bar {
  position: sticky; top: 0; z-index: 200; height: 56px;
  display: flex; align-items: center; padding: 0 1.5rem; gap: .75rem;
  background: rgba(6,9,22,.96); backdrop-filter: blur(20px); -webkit-backdrop-filter: blur(20px);
  border-bottom: 1px solid var(--border);
}
.top-mini-brand {
  font-family: 'Poppins', sans-serif;
  font-size: .82rem; font-weight: 700; color: var(--text);
  letter-spacing: -.01em; white-space: nowrap; flex-shrink: 0;
  display: flex; align-items: center; gap: .45rem;
}
.top-mini-brand::before {
  content: ''; width: 8px; height: 8px; border-radius: 50%;
  background: linear-gradient(135deg, var(--indigo), var(--cyan)); flex-shrink: 0;
}
.top-sep { width: 1px; height: 18px; background: var(--border); flex-shrink: 0; }
.top-page-label { font-size: .75rem; font-weight: 600; color: var(--muted); white-space: nowrap; }
.top-brand-tag { font-size: .7rem; color: var(--muted); }
.top-spacer { flex: 1; }
.back-btn {
  background: rgba(99,102,241,.08); border: 1px solid rgba(99,102,241,.18);
  border-radius: 8px; padding: .3rem .9rem; font-size: .7rem; font-weight: 600;
  color: var(--indigo-l); text-decoration: none; transition: all .2s; white-space: nowrap;
}
.back-btn:hover { background: rgba(99,102,241,.18); transform: translateX(-2px); }

/* ── Cards & Containers ── */
.container, .page-content { position: relative; z-index: 1; }
.card, .section-card, .panel {
  background: var(--card-bg); border: 1px solid var(--border);
  border-radius: 16px; backdrop-filter: blur(16px); -webkit-backdrop-filter: blur(16px);
  box-shadow: 0 8px 32px rgba(0,0,0,.4), inset 0 1px 0 rgba(255,255,255,.03);
  transition: box-shadow .25s, border-color .25s;
}
.card:hover, .section-card:hover { border-color: rgba(99,102,241,.25); }

/* ── Upload Zone — Modern drag & drop ── */
.upload-zone, .dz {
  border: 2px dashed rgba(99,102,241,.3); border-radius: 14px;
  padding: 2rem; text-align: center; cursor: pointer;
  transition: all .22s; background: rgba(99,102,241,.03);
  position: relative;
}
.upload-zone:hover, .dz:hover, .upload-zone.drag-over, .dz.drag-over {
  border-color: var(--indigo); background: rgba(99,102,241,.08);
  box-shadow: 0 0 0 4px rgba(99,102,241,.12);
}
.upload-zone.ok, .dz.ok {
  border-color: var(--green); background: rgba(16,185,129,.06);
  border-style: solid;
}
.upload-zone.ok:hover, .dz.ok:hover {
  border-color: var(--green); background: rgba(16,185,129,.1);
  box-shadow: 0 0 0 4px rgba(16,185,129,.1);
}
.upload-icon { font-size: 2rem; margin-bottom: .5rem; display: block; }
.upload-label { font-size: .82rem; color: var(--muted); font-weight: 500; }
.upload-hint { font-size: .7rem; color: rgba(100,116,139,.6); margin-top: .25rem; }
.upload-filename { font-size: .78rem; color: var(--green); font-weight: 600; margin-top: .4rem; }

/* ── Modern Buttons ── */
.btn-primary, .btn-generate, .btn-teal, .modal-btn-primary {
  background:transparent;
  color: #fff; border: none; border-radius: 10px;
  padding: .7rem 1.4rem; font-size: .85rem; font-weight: 700;
  font-family: 'Inter', sans-serif; cursor: pointer; letter-spacing: .01em;
  box-shadow: 0 4px 18px rgba(99,102,241,.35);
  position: relative; overflow: hidden;
  transition: transform .2s, box-shadow .2s;
}
.btn-primary::before, .btn-generate::before, .btn-teal::before, .modal-btn-primary::before {
  content: ''; position: absolute; inset: 0;
  background: linear-gradient(135deg, transparent, rgba(255,255,255,.12), transparent);
  transform: translateX(-100%); transition: transform .45s;
}
.btn-primary:hover, .btn-generate:hover, .modal-btn-primary:hover {
  transform: translateY(-2px) scale(1.01);
  box-shadow: 0 8px 28px rgba(99,102,241,.55);
}
.btn-primary:hover::before, .btn-generate:hover::before, .modal-btn-primary:hover::before {
  transform: translateX(100%);
}
.btn-primary:active, .btn-generate:active { transform: translateY(0) scale(.98); }
.btn-primary:disabled, .btn-generate:disabled { opacity: .45; cursor: not-allowed; transform: none; }

.btn-teal {
  background:transparent;
  box-shadow: 0 4px 18px rgba(16,185,129,.3);
}
.btn-teal:hover { box-shadow: 0 8px 28px rgba(16,185,129,.5); }

.btn-download, .btn-dl, .dl-btn {
  background: linear-gradient(135deg, #0D7A5F, #10B981);
  color: #fff; border: none; border-radius: 10px;
  padding: .65rem 1.3rem; font-size: .82rem; font-weight: 700;
  cursor: pointer; transition: all .22s;
  box-shadow: 0 4px 16px rgba(16,185,129,.3);
}
.btn-download:hover, .btn-dl:hover, .dl-btn:hover {
  transform: translateY(-2px); box-shadow: 0 8px 26px rgba(16,185,129,.5);
}

.btn-secondary, .btn-gray {
  background: rgba(30,41,86,.6); color: var(--muted);
  border: 1px solid rgba(99,102,241,.2); border-radius: 10px;
  padding: .65rem 1.2rem; font-size: .82rem; font-weight: 600;
  cursor: pointer; transition: all .2s;
}
.btn-secondary:hover, .btn-gray:hover { background: rgba(99,102,241,.1); color: var(--indigo-l); }

.btn-del, .btn-danger {
  background: rgba(239,68,68,.08); color: var(--red);
  border: 1px solid rgba(239,68,68,.2); border-radius: 10px;
  padding: .6rem 1.2rem; font-size: .8rem; font-weight: 600;
  cursor: pointer; transition: all .2s;
}
.btn-del:hover, .btn-danger:hover { background: rgba(239,68,68,.16); }
.btn-del:hover { animation: shake .3s ease; }

@keyframes shake {
  0%,100% { transform: translateX(0); }
  25% { transform: translateX(-3px); }
  75% { transform: translateX(3px); }
}

/* ── Form Inputs ── */
input[type=text], input[type=number], input[type=date], input[type=password],
select, textarea {
  background: rgba(6,10,24,.8); border: 1.5px solid rgba(30,41,86,.8);
  border-radius: 10px; color: var(--text); padding: .65rem .9rem;
  font-size: .85rem; font-family: 'Inter', sans-serif;
  transition: all .2s; width: 100%;
}
input:focus, select:focus, textarea:focus {
  outline: none; border-color: var(--indigo);
  box-shadow: 0 0 0 3px rgba(99,102,241,.16);
  background: rgba(10,14,32,.9);
}
input::placeholder, textarea::placeholder { color: rgba(100,116,139,.45); font-weight: 300; }

/* ── Alerts / Status ── */
.alert-success, .alert.success, .msg.success {
  background: rgba(16,185,129,.08); border: 1px solid rgba(16,185,129,.2);
  color: #34D399; border-radius: 10px; padding: .65rem .9rem;
  font-size: .8rem; font-weight: 500;
}
.alert-error, .alert.error, .msg.error {
  background: rgba(239,68,68,.08); border: 1px solid rgba(239,68,68,.2);
  color: var(--red); border-radius: 10px; padding: .65rem .9rem;
  font-size: .8rem; font-weight: 500;
}
.alert-warn, .alert.warn { 
  background: rgba(245,158,11,.08); border: 1px solid rgba(245,158,11,.2);
  color: var(--amber); border-radius: 10px; padding: .65rem .9rem;
  font-size: .8rem; font-weight: 500;
}

/* ── Tables ── */
table { width: 100%; border-collapse: collapse; font-size: .8rem; }
thead th {
  background: rgba(99,102,241,.08); color: var(--muted);
  font-weight: 700; font-size: .68rem; letter-spacing: .06em;
  text-transform: uppercase; padding: .65rem .9rem; text-align: left;
  border-bottom: 1px solid var(--border);
}
tbody tr { border-bottom: 1px solid rgba(30,45,80,.3); transition: background .15s; }
tbody tr:hover { background: rgba(99,102,241,.04); }
tbody td { padding: .6rem .9rem; color: var(--text); }

/* ── Tabs ── */
.tab { position: relative; transition: all .2s; }
.tab.active { color: var(--indigo-l) !important; }
.tab.active::after {
  content: ''; position: absolute; bottom: -1px; left: 15%; right: 15%;
  height: 2px; border-radius: 2px; background: var(--indigo);
  animation: tabIn .2s ease;
}
@keyframes tabIn { from { left: 50%; right: 50%; } to { left: 15%; right: 15%; } }

/* ── Spinner ── */
@keyframes spin { to { transform: rotate(360deg); } }
.spinner { width: 20px; height: 20px; border: 2px solid rgba(99,102,241,.2);
  border-top-color: var(--indigo); border-radius: 50%; animation: spin 1s linear infinite; }



@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&display=swap');
*, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }
body {
  font-family: 'Inter', system-ui, -apple-system, sans-serif;
  min-height: 100vh; font-size: 14px; line-height: 1.6;
  background:transparent;
  color: #E8EEF8;
  position: relative;
}


/* ══ ANIMATED BUTTONS ══════════════════════════════════════════════════════════ */
/* Primary action button — pulse + shimmer */
.btn-primary, .btn-generate, .btn-teal, .modal-btn-primary {
  position: relative; overflow: hidden;
  transition: transform .2s, box-shadow .2s !important;
}
.btn-primary::after, .btn-generate::after, .btn-teal::after, .modal-btn-primary::after {
  content: ''; position: absolute; top: -50%; left: -75%;
  width: 50%; height: 200%; background: rgba(255,255,255,.15);
  transform: skewX(-20deg); transition: left .5s ease;
  pointer-events: none;
}
.btn-primary:hover::after, .btn-generate:hover::after,
.btn-teal:hover::after, .modal-btn-primary:hover::after {
  left: 150%;
}
.btn-primary:hover, .btn-generate:hover, .modal-btn-primary:hover {
  transform: translateY(-3px) scale(1.02) !important;
  box-shadow: 0 8px 28px rgba(99,102,241,.55) !important;
}
.btn-primary:active, .btn-generate:active, .modal-btn-primary:active {
  transform: translateY(0) scale(.98) !important;
}

/* Download button — bounce */
.btn-dl, .btn-download, .dl-btn {
  animation: gentlePulse 3s ease-in-out infinite;
  transition: transform .2s, box-shadow .2s !important;
}
.btn-dl:hover, .btn-download:hover, .dl-btn:hover {
  animation: none;
  transform: translateY(-3px) scale(1.02) !important;
  box-shadow: 0 8px 28px rgba(16,185,129,.5) !important;
}
@keyframes gentlePulse {
  0%,100% { box-shadow: 0 4px 16px rgba(16,185,129,.3); }
  50%      { box-shadow: 0 4px 24px rgba(16,185,129,.55); }
}

/* Danger/delete button — shake on hover */
.btn-del:hover {
  animation: shake .35s ease;
}
@keyframes shake {
  0%,100% { transform: translateX(0); }
  20%     { transform: translateX(-3px); }
  40%     { transform: translateX(3px); }
  60%     { transform: translateX(-2px); }
  80%     { transform: translateX(2px); }
}

/* Secondary / back buttons — slide arrow */
.back-btn, .btn-secondary, .btn-gray, .btn-refresh {
  transition: all .2s !important;
}
.back-btn:hover {
  padding-left: .55rem !important;
  letter-spacing: .02em;
}

/* Upload zone — glow pulse when empty */
.upload-zone:not(.ok) {
  animation: uploadGlow 4s ease-in-out infinite;
}
.upload-zone:hover, .dz:hover {
  animation: none;
}
@keyframes uploadGlow {
  0%,100% { border-color: rgba(99,102,241,.3); }
  50%      { border-color: rgba(99,102,241,.65); box-shadow: 0 0 18px rgba(99,102,241,.15); }
}

/* App cards on dashboard — float in on load */
.app-card {
  animation: cardFadeIn .5s ease both;
}
.app-card:nth-child(1) { animation-delay: .05s; }
.app-card:nth-child(2) { animation-delay: .10s; }
.app-card:nth-child(3) { animation-delay: .15s; }
.app-card:nth-child(4) { animation-delay: .20s; }
.app-card:nth-child(5) { animation-delay: .25s; }
@keyframes cardFadeIn {
  from { opacity: 0; transform: translateY(18px); }
  to   { opacity: 1; transform: translateY(0); }
}

/* Stat cards pop-in */
.stat-card, .stat {
  animation: popIn .4s cubic-bezier(.34,1.56,.64,1) both;
}
.stat-card:nth-child(1), .stat:nth-child(1) { animation-delay: .05s; }
.stat-card:nth-child(2), .stat:nth-child(2) { animation-delay: .10s; }
.stat-card:nth-child(3), .stat:nth-child(3) { animation-delay: .15s; }
.stat-card:nth-child(4), .stat:nth-child(4) { animation-delay: .20s; }
@keyframes popIn {
  from { opacity: 0; transform: scale(.88); }
  to   { opacity: 1; transform: scale(1); }
}

/* Login button — ripple */
.btn.btn-primary, button.btn-primary {
  transition: transform .15s, box-shadow .15s !important;
}

/* Tab buttons — active indicator slide */
.tab {
  transition: all .2s !important;
  position: relative;
}
.tab.active::after {
  content: '';
  position: absolute; bottom: -1px; left: 10%; right: 10%;
  height: 2px; border-radius: 2px;
  background: #6366F1;
  animation: tabSlide .25s ease;
}
@keyframes tabSlide {
  from { left: 50%; right: 50%; }
  to   { left: 10%; right: 10%; }
}

/* Generate/action buttons in forms — shimmer idle */
.btn-generate:not(:disabled), .btn-teal:not(:disabled) {
  background-size: 200% 100% !important;
  transition: background-position .4s, transform .2s, box-shadow .2s !important;
}

/* ── Compact top bar (40px) for all internal pages ── */
.top-bar {
  position: sticky; top: 0; z-index: 200;
  height: 40px; display: flex; align-items: center;
  padding: 0 1.2rem; gap: .75rem;
  background: rgba(5,8,20,.96); backdrop-filter: blur(16px);
  border-bottom: 1px solid rgba(99,102,241,.12);
}
.top-mini-brand {
  font-size: .75rem; font-weight: 800; color: #6366F1;
  letter-spacing: .04em; white-space: nowrap; flex-shrink: 0;
}
.top-brand-text {
  font-size: .78rem; font-weight: 800; color: #818CF8;
  letter-spacing: .02em; white-space: nowrap; flex-shrink: 0;
}
.top-sep { width: 1px; height: 18px; background: rgba(99,102,241,.18); flex-shrink: 0; }
.top-page-label {
  font-size: .78rem; font-weight: 600; color: #94A3B8;
  white-space: nowrap; overflow: hidden; text-overflow: ellipsis;
}
.top-spacer { flex: 1; }
.back-btn {
  background: rgba(99,102,241,.1); border: 1px solid rgba(99,102,241,.2);
  border-radius: 6px; padding: .22rem .7rem; font-size: .7rem; font-weight: 600;
  color: #818CF8; text-decoration: none; white-space: nowrap; transition: all .2s;
  flex-shrink: 0;
}
.back-btn:hover { background: rgba(99,102,241,.2); }

.container, .main, .page { padding: 1rem 1.2rem 2rem; max-width: 1400px; margin: 0 auto; position: relative; z-index: 1; }


/* PPT Rejection specific */
:root{--g:#00B050;--mid:#0F172A;--light:#E8EEF8;--muted:#6B7280;--danger:#ef4444;--warn:#f59e0b;--info:#3B82F6;}
.card{background:rgba(15,23,42,.8);border:1px solid rgba(99,102,241,.12);border-radius:12px;padding:1.2rem;margin-bottom:1rem;backdrop-filter:blur(10px);}
.card h2{font-size:.82rem;font-weight:700;color:#818CF8;margin-bottom:.85rem;display:flex;align-items:center;gap:.5rem;}
.upload-zone{border:2px dashed rgba(99,102,241,.3);border-radius:9px;padding:2rem;text-align:center;cursor:pointer;transition:.2s;position:relative;}
.upload-zone:hover,.upload-zone.dragover{border-color:#6366F1;background:rgba(99,102,241,.06);}
.upload-zone input{position:absolute;inset:0;opacity:0;cursor:pointer;width:100%;height:100%;}
.upload-zone .icon{font-size:2rem;margin-bottom:.4rem;}
.upload-zone p{color:var(--muted);font-size:.82rem;}
.file-list{margin-top:.6rem;display:flex;flex-wrap:wrap;gap:.35rem;}
.file-tag{background:rgba(99,102,241,.12);border:1px solid rgba(99,102,241,.25);color:#818CF8;padding:.2rem .6rem;border-radius:20px;font-size:.72rem;display:flex;align-items:center;gap:.3rem;}
.file-tag button{background:none;border:none;color:#F87171;cursor:pointer;font-size:.85rem;line-height:1;padding:0;}
.status-bar{border-radius:7px;padding:.6rem .85rem;font-size:.79rem;display:none;align-items:center;gap:.5rem;margin-top:.6rem;}
.status-bar.show{display:flex;}
.status-bar.ok{background:rgba(16,185,129,.1);border:1px solid rgba(16,185,129,.2);color:#10B981;}
.status-bar.error{background:rgba(239,68,68,.1);border:1px solid rgba(239,68,68,.3);color:#F87171;}
.status-bar.warn{background:rgba(245,158,11,.1);border:1px solid rgba(245,158,11,.3);color:#F59E0B;}
.prog-wrap{margin-top:.5rem;display:none;}.prog-wrap.show{display:block;}
.prog-bar{height:5px;border-radius:3px;background:rgba(255,255,255,.08);overflow:hidden;}
.prog-fill{height:100%;background:#6366F1;width:0;transition:width .3s;border-radius:3px;}
.prog-lbl{font-size:.72rem;color:var(--muted);margin-top:.25rem;}
.filter-grid{display:grid;grid-template-columns:repeat(auto-fit,minmax(160px,1fr));gap:.65rem;}
.field label{display:block;font-size:.72rem;color:var(--muted);margin-bottom:.3rem;font-weight:500;}
.field input,.field select{width:100%;padding:.5rem .75rem;background:rgba(8,15,40,.7);border:1px solid rgba(99,102,241,.15);border-radius:7px;color:var(--light);font-size:.82rem;transition:.2s;font-family:'Inter',sans-serif;}
.field input:focus,.field select:focus{outline:none;border-color:#6366F1;}
.field input::placeholder{color:var(--muted);}select option{background:#1a1a2e;}
.btn-row{display:flex;gap:.6rem;flex-wrap:wrap;margin-top:.6rem;align-items:center;}
.btn{padding:.5rem 1.1rem;border-radius:7px;border:none;cursor:pointer;font-size:.8rem;font-weight:600;display:flex;align-items:center;gap:.4rem;transition:.2s;font-family:'Inter',sans-serif;}
.btn-primary{background:#6366F1;color:#fff;}.btn-primary:hover{background:#4F46E5;transform:translateY(-1px);}
.btn-secondary{background:rgba(255,255,255,.06);color:var(--light);border:1px solid rgba(255,255,255,.1);}
.btn-secondary:hover{background:rgba(255,255,255,.1);}
.btn-dl{background:#3B82F6;color:#fff;}.btn-dl:hover{background:#2563EB;transform:translateY(-1px);}
.btn-xl{background:#7C3AED;color:#fff;}.btn-xl:hover{background:#6D28D9;transform:translateY(-1px);}
.btn-zip{background:#D97706;color:#fff;}.btn-zip:hover{background:#B45309;transform:translateY(-1px);}
.btn:disabled{opacity:.4;cursor:not-allowed;transform:none!important;}
.stats-row{display:flex;gap:.65rem;flex-wrap:wrap;margin-bottom:1rem;}
.stat{flex:1;min-width:80px;background:rgba(99,102,241,.08);border:1px solid rgba(99,102,241,.15);border-radius:9px;padding:.75rem;text-align:center;}
.stat .val{font-size:1.6rem;font-weight:800;color:#818CF8;}
.stat .lbl{font-size:.65rem;color:var(--muted);margin-top:.15rem;}
.result-meta{display:flex;align-items:center;gap:.65rem;margin-bottom:.5rem;flex-wrap:wrap;}
.count-badge{background:rgba(99,102,241,.12);border:1px solid rgba(99,102,241,.25);color:#818CF8;padding:.2rem .7rem;border-radius:20px;font-size:.76rem;font-weight:600;}
.tbl-wrap{overflow-x:auto;border-radius:7px;border:1px solid rgba(99,102,241,.1);}
table{width:100%;border-collapse:collapse;font-size:.76rem;}
thead{background:rgba(99,102,241,.1);}
thead th{padding:.55rem .7rem;text-align:center;font-weight:600;color:#818CF8;white-space:nowrap;border-bottom:1px solid rgba(99,102,241,.18);}
tbody tr{border-bottom:1px solid rgba(255,255,255,.03);}
tbody tr:hover{background:rgba(255,255,255,.03);}
tbody td{padding:.45rem .7rem;text-align:center;color:var(--light);}
tbody td.left{text-align:left;}
tr.group-header td{background:rgba(99,102,241,.1);color:#818CF8;font-weight:700;font-size:.72rem;text-align:left;padding:.4rem .9rem;}
.badge-no{display:inline-block;padding:.12rem .48rem;border-radius:10px;font-size:.68rem;font-weight:600;background:rgba(239,68,68,.15);color:#F87171;}
.badge-reason{display:inline-block;padding:.12rem .48rem;border-radius:10px;font-size:.68rem;font-weight:600;background:rgba(245,158,11,.12);color:#F59E0B;}
.no-data{text-align:center;padding:2rem;color:var(--muted);}

/* ══ ANIMATED BUTTONS ══════════════════════════════════════════════════════════ */
/* Primary action button — pulse + shimmer */
.btn-primary, .btn-generate, .btn-teal, .modal-btn-primary {
  position: relative; overflow: hidden;
  transition: transform .2s, box-shadow .2s !important;
}
.btn-primary::after, .btn-generate::after, .btn-teal::after, .modal-btn-primary::after {
  content: ''; position: absolute; top: -50%; left: -75%;
  width: 50%; height: 200%; background: rgba(255,255,255,.15);
  transform: skewX(-20deg); transition: left .5s ease;
  pointer-events: none;
}
.btn-primary:hover::after, .btn-generate:hover::after,
.btn-teal:hover::after, .modal-btn-primary:hover::after {
  left: 150%;
}
.btn-primary:hover, .btn-generate:hover, .modal-btn-primary:hover {
  transform: translateY(-3px) scale(1.02) !important;
  box-shadow: 0 8px 28px rgba(99,102,241,.55) !important;
}
.btn-primary:active, .btn-generate:active, .modal-btn-primary:active {
  transform: translateY(0) scale(.98) !important;
}

/* Download button — bounce */
.btn-dl, .btn-download, .dl-btn {
  animation: gentlePulse 3s ease-in-out infinite;
  transition: transform .2s, box-shadow .2s !important;
}
.btn-dl:hover, .btn-download:hover, .dl-btn:hover {
  animation: none;
  transform: translateY(-3px) scale(1.02) !important;
  box-shadow: 0 8px 28px rgba(16,185,129,.5) !important;
}
@keyframes gentlePulse {
  0%,100% { box-shadow: 0 4px 16px rgba(16,185,129,.3); }
  50%      { box-shadow: 0 4px 24px rgba(16,185,129,.55); }
}

/* Danger/delete button — shake on hover */
.btn-del:hover {
  animation: shake .35s ease;
}
@keyframes shake {
  0%,100% { transform: translateX(0); }
  20%     { transform: translateX(-3px); }
  40%     { transform: translateX(3px); }
  60%     { transform: translateX(-2px); }
  80%     { transform: translateX(2px); }
}

/* Secondary / back buttons — slide arrow */
.back-btn, .btn-secondary, .btn-gray, .btn-refresh {
  transition: all .2s !important;
}
.back-btn:hover {
  padding-left: .55rem !important;
  letter-spacing: .02em;
}

/* Upload zone — glow pulse when empty */
.upload-zone:not(.ok) {
  animation: uploadGlow 4s ease-in-out infinite;
}
.upload-zone:hover, .dz:hover {
  animation: none;
}
@keyframes uploadGlow {
  0%,100% { border-color: rgba(99,102,241,.3); }
  50%      { border-color: rgba(99,102,241,.65); box-shadow: 0 0 18px rgba(99,102,241,.15); }
}

/* App cards on dashboard — float in on load */
.app-card {
  animation: cardFadeIn .5s ease both;
}
.app-card:nth-child(1) { animation-delay: .05s; }
.app-card:nth-child(2) { animation-delay: .10s; }
.app-card:nth-child(3) { animation-delay: .15s; }
.app-card:nth-child(4) { animation-delay: .20s; }
.app-card:nth-child(5) { animation-delay: .25s; }
@keyframes cardFadeIn {
  from { opacity: 0; transform: translateY(18px); }
  to   { opacity: 1; transform: translateY(0); }
}

/* Stat cards pop-in */
.stat-card, .stat {
  animation: popIn .4s cubic-bezier(.34,1.56,.64,1) both;
}
.stat-card:nth-child(1), .stat:nth-child(1) { animation-delay: .05s; }
.stat-card:nth-child(2), .stat:nth-child(2) { animation-delay: .10s; }
.stat-card:nth-child(3), .stat:nth-child(3) { animation-delay: .15s; }
.stat-card:nth-child(4), .stat:nth-child(4) { animation-delay: .20s; }
@keyframes popIn {
  from { opacity: 0; transform: scale(.88); }
  to   { opacity: 1; transform: scale(1); }
}

/* Login button — ripple */
.btn.btn-primary, button.btn-primary {
  transition: transform .15s, box-shadow .15s !important;
}

/* Tab buttons — active indicator slide */
.tab {
  transition: all .2s !important;
  position: relative;
}
.tab.active::after {
  content: '';
  position: absolute; bottom: -1px; left: 10%; right: 10%;
  height: 2px; border-radius: 2px;
  background: #6366F1;
  animation: tabSlide .25s ease;
}
@keyframes tabSlide {
  from { left: 50%; right: 50%; }
  to   { left: 10%; right: 10%; }
}

/* Generate/action buttons in forms — shimmer idle */
.btn-generate:not(:disabled), .btn-teal:not(:disabled) {
  background-size: 200% 100% !important;
  transition: background-position .4s, transform .2s, box-shadow .2s !important;
}

.spinner{animation:spin 1s linear infinite;display:inline-block;}
@keyframes spin{to{transform:rotate(360deg);}}
/* ── Password management button styles ── */
.btn-pw{background:rgba(99,102,241,.12);color:#818CF8;border:1px solid rgba(99,102,241,.25);}
.btn-pw:hover{background:rgba(99,102,241,.22);}
.btn-reset{background:rgba(245,158,11,.1);color:#F59E0B;border:1px solid rgba(245,158,11,.22);}
.btn-reset:hover{background:rgba(245,158,11,.18);}
.btn-del{background:rgba(239,68,68,.08);color:#F87171;border:1px solid rgba(239,68,68,.18);}
.btn-del:hover{background:rgba(239,68,68,.16);}
/* ── Modal backdrop + card ── */
.modal-overlay{position:fixed;inset:0;z-index:9000;background:rgba(0,0,0,.65);
  display:none;align-items:center;justify-content:center;backdrop-filter:blur(4px);}
.modal-card{background:#0D1529;border:1px solid rgba(99,102,241,.25);border-radius:16px;
  padding:1.8rem 2rem;width:100%;max-width:420px;box-shadow:0 24px 60px rgba(0,0,0,.6);}
.modal-title{font-size:1rem;font-weight:800;color:#F1F5FF;margin-bottom:.25rem;
  display:flex;align-items:center;gap:.5rem;}
.modal-sub{font-size:.78rem;color:#475569;margin-bottom:1.3rem;}
.modal-field{margin-bottom:.9rem;}
.modal-field label{display:block;font-size:.68rem;font-weight:700;color:#475569;
  letter-spacing:.08em;text-transform:uppercase;margin-bottom:.35rem;}
.modal-field input{width:100%;padding:.72rem .9rem;background:rgba(8,15,40,.9);
  border:1.5px solid rgba(30,41,86,.9);border-radius:10px;color:#E8EEF8;
  font-size:.9rem;font-family:'Inter',sans-serif;transition:all .2s;}
.modal-field input:focus{outline:none;border-color:#6366F1;box-shadow:0 0 0 3px rgba(99,102,241,.18);}
.modal-field input::placeholder{color:#334155;}
.show-pw-row{display:flex;align-items:center;gap:.5rem;margin-bottom:1rem;
  font-size:.76rem;color:#475569;cursor:pointer;}
.show-pw-row input[type=checkbox]{width:14px;height:14px;accent-color:#6366F1;cursor:pointer;}
.modal-actions{display:flex;gap:.6rem;margin-top:.5rem;}
.modal-btn{flex:1;padding:.75rem;border-radius:10px;border:none;cursor:pointer;
  font-size:.88rem;font-weight:700;font-family:'Inter',sans-serif;transition:all .2s;}
.modal-btn-primary{background:linear-gradient(135deg,#4338CA,#6366F1);color:#fff;
  box-shadow:0 4px 16px rgba(99,102,241,.35);}
.modal-btn-primary:hover{transform:translateY(-1px);box-shadow:0 6px 20px rgba(99,102,241,.5);}
.modal-btn-cancel{background:rgba(255,255,255,.06);color:#94A3B8;
  border:1px solid rgba(255,255,255,.1);}
.modal-btn-cancel:hover{background:rgba(255,255,255,.1);}
.modal-result{padding:.55rem .8rem;border-radius:8px;font-size:.8rem;font-weight:500;
  margin-top:.75rem;display:none;}
.modal-result.result-ok{display:block;background:rgba(16,185,129,.1);
  border:1px solid rgba(16,185,129,.22);color:#34D399;}
.modal-result.result-err{display:block;background:rgba(239,68,68,.1);
  border:1px solid rgba(239,68,68,.22);color:#F87171;}
/* Temp PW display */
.temp-pw-box{background:rgba(8,15,40,.9);border:1.5px solid rgba(99,102,241,.3);
  border-radius:10px;padding:1rem 1.2rem;margin:.8rem 0;text-align:center;}
.temp-pw-label{font-size:.68rem;color:#475569;font-weight:600;letter-spacing:.08em;
  text-transform:uppercase;margin-bottom:.5rem;}
.temp-pw-value{font-size:1.6rem;font-weight:900;color:#818CF8;letter-spacing:.12em;
  font-family:'Courier New',monospace;}
.temp-pw-warn{font-size:.72rem;color:#F59E0B;margin-top:.6rem;line-height:1.5;}




/* Action cards */
.action-card { background: rgba(10,16,42,.85) !important; border-color: rgba(245,158,11,.5) !important; }
.rr-panel { background: rgba(6,12,38,.88) !important; border-color: rgba(99,102,241,.35) !important; }

/* Stats / text helpers */
.sec-hint, .upload-hint { color: #9BB8E0 !important; }
.chip-ok { color: #4ADE80 !important; }
.chip-info { color: #818CF8 !important; }
.chip-wait { color: #F59E0B !important; }

/* ── Developed by Karthik — fixed footer ── */
.dev-credit {
  position: fixed; bottom: 0; left: 0; right: 0; z-index: 9999;
  text-align: center; padding: .3rem 1rem;
  background: rgba(5, 8, 22, 0.75); backdrop-filter: blur(8px);
  border-top: 1px solid rgba(99,102,241,.2);
  font-size: 11px; font-weight: 600; color: rgba(160,180,220,.75);
  letter-spacing: .06em; font-family: 'Inter', system-ui, sans-serif;
  pointer-events: none; user-select: none;
}

/* ═══ GLOBAL BACKGROUND & TRANSPARENCY — v3 ═══════════════════════ */
html, body {
  background-image: url('/static/bg.jpg') !important;
  background-size: cover !important;
  background-position: center center !important;
  background-attachment: fixed !important;
  background-repeat: no-repeat !important;
  background-color: #08101E !important;
}
/* Single very-light overlay — image stays visible */
body::before {
  content: '' !important;
  position: fixed !important;
  inset: 0 !important;
  z-index: 0 !important;
  background: rgba(3, 7, 18, 0.45) !important;
  pointer-events: none !important;
}
body > * { position: relative; z-index: 1; }

/* ── Top-bar: fully transparent glass, no black border ── */
.top-bar {
  background: rgba(6, 10, 28, 0.55) !important;
  backdrop-filter: blur(18px) !important;
  -webkit-backdrop-filter: blur(18px) !important;
  border-bottom: 1px solid rgba(255,255,255,0.08) !important;
  position: sticky !important;
  top: 0 !important;
  z-index: 200 !important;
}

/* ── Cards / Panels — glass, no solid black fill ── */
.card, .section-card, .panel, .sec-card, .stats-box, .action-card {
  background: rgba(8, 14, 38, 0.70) !important;
  border: 1px solid rgba(255,255,255,0.10) !important;
  backdrop-filter: blur(16px) !important;
  -webkit-backdrop-filter: blur(16px) !important;
  box-shadow: 0 4px 32px rgba(0,0,0,0.35) !important;
}
.panel-head, .card-header {
  background: rgba(10, 18, 52, 0.72) !important;
  border-bottom: 1px solid rgba(255,255,255,0.07) !important;
}
.panel-body { background: rgba(5, 10, 30, 0.60) !important; }

/* ── Upload zones ── */
.upload-zone, .dz {
  background: rgba(6, 12, 34, 0.55) !important;
  border: 2px dashed rgba(99,102,241,0.55) !important;
}

/* ── Typography: all white/light ── */
body, h1, h2, h3, h4, p, span, div, td, th, label, a {
  color: #EEF3FF !important;
}
.hero-title, .brand-title, .card-title, .admin-hero-title {
  color: #FFFFFF !important;
  text-shadow: 0 2px 16px rgba(0,0,0,0.7) !important;
  font-weight: 800 !important;
}
.hero-sub, .brand-sub, .admin-hero-sub, .sec-hint {
  color: rgba(210, 225, 255, 0.80) !important;
}
.top-mini-brand, .top-page-label, .top-brand .brand-text {
  color: #FFFFFF !important;
  font-weight: 700 !important;
}
.back-btn {
  color: #C5D5FF !important;
  background: rgba(99,102,241,0.18) !important;
  border: 1px solid rgba(99,102,241,0.35) !important;
}
.back-btn:hover { background: rgba(99,102,241,0.32) !important; }

/* ── Inputs: legible on transparent backgrounds ── */
input[type=text], input[type=number], input[type=date],
input[type=password], select, textarea,
.input-num, .action-select, .rr-input, .rr-select, .date-input {
  background: rgba(5, 9, 28, 0.80) !important;
  border: 1.5px solid rgba(99,102,241,0.40) !important;
  color: #EEF3FF !important;
  font-size: 14px !important;
}
input::placeholder, textarea::placeholder {
  color: rgba(180, 200, 240, 0.50) !important;
}

/* ── Muted / secondary text ── */
.muted, .sec-label, [style*="color:#475569"],
[style*="color:#64748B"], [style*="color:#374167"] {
  color: rgba(190, 210, 255, 0.70) !important;
}

/* ── Action + RR cards ── */
.action-card {
  background: rgba(12, 18, 48, 0.78) !important;
  border-color: rgba(245,158,11,0.55) !important;
}
.rr-panel {
  background: rgba(8, 14, 42, 0.80) !important;
  border-color: rgba(99,102,241,0.40) !important;
}

/* ── Global font sizes ── */
body { font-size: 14px !important; }
h1 { font-size: 26px !important; }
h2 { font-size: 22px !important; }
h3, .hero-title { font-size: 22px !important; }
h4 { font-size: 17px !important; }

/* ── Dev-credit fixed footer ── */
.apsg-footer {
  position: fixed !important; bottom: 0 !important;
  left: 0 !important; right: 0 !important; z-index: 9999 !important;
  text-align: center !important; padding: .28rem 1rem !important;
  background: rgba(4, 7, 20, 0.70) !important;
  backdrop-filter: blur(8px) !important;
  border-top: 1px solid rgba(255,255,255,0.07) !important;
  font-size: 11px !important; font-weight: 600 !important;
  color: rgba(200, 220, 255, 0.70) !important;
  letter-spacing: .06em !important; pointer-events: none !important;
  user-select: none !important;
}
</style>
</head>
<body>
<div class="top-bar">
  <span class="top-mini-brand">APSG</span>
  <div class="top-sep"></div>
  <span class="top-page-label">📁 Monthly Rejection Report Filter</span>
  <div class="top-spacer"></div>
  <span style="font-size:.65rem;color:#6366F1;font-weight:700;white-space:nowrap;margin-right:.5rem;letter-spacing:.01em;">✦ Karthi</span>
  <a href="/" class="back-btn">← Dashboard</a>
</div>

<div class="container" style="padding-top:.75rem; position:relative; z-index:1; max-width:1400px; margin:0 auto;">
  <div class="card">
    <h2>📁 Upload PPT Files</h2>
    <div class="upload-zone" id="dropZone">
      <input type="file" id="fileInput" multiple accept=".ppt,.pptx" style="display:none">
      <div class="icon">📊</div>
      <strong>Drop files here or click to browse</strong>
      <p>Supports multiple .PPTX files — any quantity</p>
    </div>
    <div class="file-list" id="fileList"></div>
    <div class="prog-wrap" id="progWrap">
      <div class="prog-bar"><div class="prog-fill" id="progFill"></div></div>
      <div class="prog-lbl" id="progLbl"></div>
    </div>
    <div class="status-bar" id="uploadStatus"></div>
    <div class="btn-row">
      <button class="btn btn-primary" id="uploadBtn" onclick="uploadFiles()">⬆ Upload &amp; Extract</button>
      <button class="btn btn-secondary" onclick="clearAll()">🗑 Clear All</button>
    </div>
  </div>
  <div class="stats-row" id="statsRow" style="display:none">
    <div class="stat"><div class="val" id="sTotal">0</div><div class="lbl">Total Records</div></div>
    <div class="stat"><div class="val" id="sFiltered">0</div><div class="lbl">Filtered</div></div>
    <div class="stat"><div class="val" id="sVehicles">0</div><div class="lbl">Vehicles</div></div>
    <div class="stat"><div class="val" id="sSites">0</div><div class="lbl">Sites</div></div>
    <div class="stat"><div class="val" id="sFiles">0</div><div class="lbl">Files</div></div>
  </div>
  <div class="card">
    <h2>🔍 Search &amp; Filter</h2>
    <div class="filter-grid">
      <div class="field"><label>E-Token Search</label>
        <input type="text" id="searchKw" placeholder="Search E-Token…" oninput="onSearchInput()"></div>
      <div class="field"><label>E-Token</label>
        <select id="selEToken" onchange="onETokenChange()"><option value="">All</option></select></div>
      <div class="field"><label>Reject Reason</label>
        <select id="selReason" onchange="applyFilters()"><option value="">All</option></select></div>
      <div class="field"><label>Date From</label>
        <input type="date" id="dateFrom" onchange="applyFilters()" style="color-scheme:dark;"></div>
      <div class="field"><label>Date To</label>
        <input type="date" id="dateTo" onchange="applyFilters()" style="color-scheme:dark;"></div>
    </div>
    <div class="btn-row">
      <button class="btn btn-secondary" onclick="resetFilters()">↺ Reset</button>
      <button class="btn btn-primary" onclick="applyFilters()">🔎 Apply</button>
    </div>
  </div>
  <div class="card">
    <h2>📋 Preview &amp; Export</h2>
    <div class="result-meta" id="resultMeta" style="display:none">
      <span class="count-badge" id="vehicleCountBadge">0 Vehicles</span>
      <span class="count-badge" id="recordCountBadge">0 Records</span>
      <span class="count-badge" id="siteCountBadge">0 Sites</span>
    </div>
    <div class="tbl-wrap">
      <table>
        <thead><tr>
          <th>#</th><th>Ticket No</th><th>Veh No</th><th>Material</th>
          <th>Source Site</th><th>Date</th><th>Time</th>
          <th>E-Token</th><th>Accepted</th><th>Reject Reason</th>
        </tr></thead>
        <tbody id="previewBody">
          <tr><td colspan="10" class="no-data">Upload PPT files to see data</td></tr>
        </tbody>
      </table>
    </div>
    <div class="btn-row" style="margin-top:1rem">
      <button class="btn btn-dl" id="dlBtn" onclick="downloadReport()" disabled>⬇ Download PPT</button>
      <button class="btn btn-zip" id="zipBtn" onclick="downloadZip()" disabled>📦 Download ZIP</button>
      <button class="btn btn-xl" id="xlBtn" onclick="downloadExcel()" disabled>📊 Excel</button>
      <span id="dlStatus" style="font-size:.78rem;color:var(--muted);"></span>
    </div>
  </div>
</div>
<script>
let allFiles=[],allRecords=[],filtered=[],totalLoaded=0;
const fileInput=document.getElementById('fileInput');
const dropZone=document.getElementById('dropZone');

// Click to open file picker
dropZone.addEventListener('click', () => fileInput.click());
fileInput.addEventListener('change', () => addFiles(fileInput.files));

// Drag & drop events on the outer zone div
dropZone.addEventListener('dragover', e => { e.preventDefault(); e.stopPropagation(); dropZone.classList.add('dragover'); });
dropZone.addEventListener('dragleave', e => { if (!dropZone.contains(e.relatedTarget)) dropZone.classList.remove('dragover'); });
dropZone.addEventListener('drop', e => { e.preventDefault(); e.stopPropagation(); dropZone.classList.remove('dragover'); addFiles(e.dataTransfer.files); });
function addFiles(files){for(const f of files)if(f.name.match(/\.pptx?$/i)&&!allFiles.find(x=>x.name===f.name&&x.size===f.size))allFiles.push(f);renderFileList();}
function renderFileList(){document.getElementById('fileList').innerHTML=allFiles.map((f,i)=>`<div class="file-tag">📄 ${esc(f.name)} <small>(${(f.size/1024).toFixed(0)}KB)</small><button onclick="removeFile(${i})">×</button></div>`).join('');}
function removeFile(i){allFiles.splice(i,1);renderFileList();}
function clearAll(){allFiles=[];allRecords=[];filtered=[];totalLoaded=0;renderFileList();renderTable([]);setStats(0,0,0,0,0);document.getElementById('statsRow').style.display='none';document.getElementById('resultMeta').style.display='none';setStatus('','');['dlBtn','zipBtn','xlBtn'].forEach(id=>document.getElementById(id).disabled=true);}
async function uploadFiles(){
  if(!allFiles.length){setStatus('Please select at least one PPT file.','error');return;}
  const btn=document.getElementById('uploadBtn');
  btn.disabled=true;btn.innerHTML='<span class="spinner">⟳</span> Processing…';
  setStatus('','');allRecords=[];filtered=[];totalLoaded=0;
  const BATCH=5,errs=[];
  try{
    for(let i=0;i<allFiles.length;i+=BATCH){
      const batch=allFiles.slice(i,i+BATCH);
      showProg(true,Math.round((i/allFiles.length)*90),`Processing ${i+1}–${Math.min(i+BATCH,allFiles.length)} of ${allFiles.length}…`);
      const fd=new FormData();
      batch.forEach(f=>fd.append('files',f));
      let res,data;
      try{
        res=await fetch(`/api/ppt/upload?reset=${i===0}`,{method:'POST',body:fd});
        data=await res.json();
      }catch(ne){errs.push(`Network: ${ne.message}`);continue;}
      if(data.error&&!(data.records&&data.records.length)){errs.push(data.error);continue;}
      if(data.records)allRecords=data.records;
      if(data.files_loaded)totalLoaded+=data.files_loaded;
      if(data.errors&&data.errors.length)errs.push(...data.errors);
    }
    showProg(true,100,'Done!');
    filtered=[...allRecords];
    populateDropdowns(allRecords);renderTable(filtered);updateMeta(filtered);
    const v=countUnique(filtered,'veh_no'),s=countUnique(filtered,'source_site');
    setStats(allRecords.length,filtered.length,v,s,totalLoaded);
    document.getElementById('statsRow').style.display='flex';
    ['dlBtn','zipBtn','xlBtn'].forEach(id=>document.getElementById(id).disabled=allRecords.length===0);
    let msg=`✓ Extracted ${allRecords.length} records from ${totalLoaded} file(s).`;
    let type='ok';
    if(errs.length){msg+=` ⚠ ${errs.length} issue(s).`;type='warn';}
    setStatus(msg,type);
  }catch(e){setStatus('Upload error: '+e.message,'error');}
  finally{btn.disabled=false;btn.innerHTML='⬆ Upload &amp; Extract';setTimeout(()=>showProg(false),2000);}
}
function countUnique(recs,key){return new Set(recs.map(r=>r[key]).filter(Boolean)).size;}
function onSearchInput(){const kw=document.getElementById('searchKw').value.trim();const esel=document.getElementById('selEToken');if(kw.length>=1){const m=[...new Set(allRecords.filter(r=>r.e_token.toLowerCase().includes(kw.toLowerCase())).map(r=>r.e_token))];rebuildSel(esel,m);}else{rebuildSel(esel,[...new Set(allRecords.map(r=>r.e_token))].filter(Boolean));}onETokenChange();}
function onETokenChange(){const et=document.getElementById('selEToken').value;const reasons=[...new Set(allRecords.filter(r=>!et||r.e_token===et).map(r=>r.reject_reason))].filter(Boolean);rebuildSel(document.getElementById('selReason'),reasons);applyFilters();}
function applyFilters(){const kw=document.getElementById('searchKw').value.trim().toLowerCase();const et=document.getElementById('selEToken').value;const rs=document.getElementById('selReason').value;const df=pd(document.getElementById('dateFrom').value);const dt_=pd(document.getElementById('dateTo').value);filtered=allRecords.filter(r=>{if(kw.length>=3&&!r.e_token.toLowerCase().includes(kw))return false;if(et&&r.e_token!==et)return false;if(rs&&r.reject_reason!==rs)return false;if(df||dt_){const d=pd(r.date);if(!d)return false;if(df&&d<df)return false;if(dt_&&d>dt_)return false;}return true;});filtered.sort((a,b)=>a.source_site.localeCompare(b.source_site)||a.e_token.localeCompare(b.e_token));renderTable(filtered);updateMeta(filtered);const v=countUnique(filtered,'veh_no'),s=countUnique(filtered,'source_site');setStats(allRecords.length,filtered.length,v,s,totalLoaded);['dlBtn','zipBtn','xlBtn'].forEach(id=>document.getElementById(id).disabled=filtered.length===0);}
function resetFilters(){document.getElementById('searchKw').value='';document.getElementById('dateFrom').value='';document.getElementById('dateTo').value='';populateDropdowns(allRecords);filtered=[...allRecords];filtered.sort((a,b)=>a.source_site.localeCompare(b.source_site)||a.e_token.localeCompare(b.e_token));renderTable(filtered);updateMeta(filtered);const v=countUnique(filtered,'veh_no'),s=countUnique(filtered,'source_site');setStats(allRecords.length,filtered.length,v,s,totalLoaded);}
function populateDropdowns(recs){rebuildSel(document.getElementById('selEToken'),[...new Set(recs.map(r=>r.e_token))].filter(Boolean));rebuildSel(document.getElementById('selReason'),[...new Set(recs.map(r=>r.reject_reason))].filter(Boolean));}
function rebuildSel(sel,opts){const cur=sel.value;sel.innerHTML='<option value="">All</option>'+opts.map(o=>`<option value="${esc(o)}"${o===cur?' selected':''}>${esc(o)}</option>`).join('');}
function renderTable(recs){const tb=document.getElementById('previewBody');if(!recs.length){tb.innerHTML='<tr><td colspan="10" class="no-data">No records match current filters</td></tr>';return;}let html='',lastSite='',sn=0;for(const r of recs){if(r.source_site!==lastSite){const cnt=recs.filter(x=>x.source_site===r.source_site).length;html+=`<tr class="group-header"><td colspan="10">📍 ${esc(r.source_site)} | ${cnt} record(s)</td></tr>`;lastSite=r.source_site;sn=0;}sn++;html+=`<tr><td>${sn}</td><td>${esc(r.ticket_no)}</td><td>${esc(r.veh_no)}</td><td>${esc(r.material)}</td><td>${esc(r.source_site)}</td><td>${esc(r.date)}</td><td>${esc(r.time)}</td><td style="font-size:.7rem">${esc(r.e_token)}</td><td><span class="badge-no">${esc(r.accepted||'NO')}</span></td><td class="left"><span class="badge-reason">${esc(r.reject_reason)}</span></td></tr>`;}tb.innerHTML=html;}
function updateMeta(recs){const show=recs.length>0;document.getElementById('resultMeta').style.display=show?'flex':'none';if(show){const v=countUnique(recs,'veh_no'),s=countUnique(recs,'source_site');document.getElementById('vehicleCountBadge').textContent=`🚛 ${v} Vehicle${v!==1?'s':''}`;document.getElementById('recordCountBadge').textContent=`📋 ${recs.length} Record${recs.length!==1?'s':''}`;document.getElementById('siteCountBadge').textContent=`📍 ${s} Source Site${s!==1?'s':''}`; }}
async function downloadReport(){if(!filtered.length)return;await doDownload('/api/ppt/generate','Rejection_Report.pptx','application/vnd.openxmlformats-officedocument.presentationml.presentation','⬇ Download PPT');}
async function downloadExcel(){if(!filtered.length)return;await doDownload('/api/ppt/export_excel','Rejection_Report.xlsx','application/vnd.openxmlformats-officedocument.spreadsheetml.sheet','📊 Excel');}
async function downloadZip(){if(!filtered.length)return;await doDownload('/api/ppt/generate_zip','Rejection_Report.zip','application/zip','📦 Download ZIP');}
async function doDownload(url,defaultName,mime,btnLabel){const btn=url.includes('zip')?document.getElementById('zipBtn'):url.includes('generate')?document.getElementById('dlBtn'):document.getElementById('xlBtn');const st=document.getElementById('dlStatus');btn.disabled=true;btn.innerHTML=`<span class="spinner">⟳</span> Generating…`;st.textContent='Building…';try{const res=await fetch(url,{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({records:filtered})});if(!res.ok){const e=await res.json().catch(()=>({error:`HTTP ${res.status}`}));st.textContent='Error: '+(e.error||'failed');return;}const blob=await res.blob();const dl=URL.createObjectURL(blob);const a=document.createElement('a');a.href=dl;a.download=defaultName;a.click();URL.revokeObjectURL(dl);st.textContent=`✓ Downloaded (${filtered.length} records)`;} catch(e){st.textContent='Error: '+e.message;}finally{btn.disabled=filtered.length===0;btn.innerHTML=btnLabel;}}
function setStats(total,filt,v,s,files){document.getElementById('sTotal').textContent=total;document.getElementById('sFiltered').textContent=filt;document.getElementById('sVehicles').textContent=v;document.getElementById('sSites').textContent=s;document.getElementById('sFiles').textContent=files;}
function setStatus(msg,type){const el=document.getElementById('uploadStatus');el.textContent=msg;el.className='status-bar'+(msg?' show':'')+(type?' '+type:'');}
function showProg(show,pct,lbl){document.getElementById('progWrap').className='prog-wrap'+(show?' show':'');if(show){document.getElementById('progFill').style.width=(pct||0)+'%';document.getElementById('progLbl').textContent=lbl||'';}}
function pd(str){if(!str)return null;const iso=str.match(/^(\d{4})-(\d{2})-(\d{2})$/);if(iso)return new Date(+iso[1],+iso[2]-1,+iso[3]);const dmy=str.match(/^(\d{1,2})-(\d{1,2})-(\d{4})$/);return dmy?new Date(+dmy[3],+dmy[2]-1,+dmy[1]):null;}
function esc(s){return String(s||'').replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;').replace(/"/g,'&quot;');}
</script>
<div class="apsg-footer">✦ Internal Reporting Platform — APSG Staging Ground &nbsp;·&nbsp; Developed by Karthik</div>
</body>
</html>"""

EXCEL_REJECTION_HTML = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>Excel Rejection Report — APSG</title>
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800&family=Poppins:wght@700;800&display=swap" rel="stylesheet">
<style>
*{box-sizing:border-box;margin:0;padding:0;}
:root{--bg:#0A0F1E;--s1:#111827;--s2:#1C2333;--s3:#232D42;
  --bdr:rgba(255,255,255,.08);--t1:#F1F5F9;--t2:#94A3B8;--mu:#64748B;
  --acc:#2DD4BF;--ind:#6366F1;--ind-l:#818CF8;--grn:#10B981;--r:14px;}
html,body{background:url('/static/bg.jpg') center/cover fixed #08101E;}
body::before{content:'';position:fixed;inset:0;z-index:0;background:rgba(3,7,18,.52);pointer-events:none;}
body{font-family:'Inter',sans-serif;color:var(--t1);min-height:100vh;padding-bottom:4rem;}
body>*{position:relative;z-index:1;}
.topbar{position:sticky;top:0;z-index:200;height:44px;display:flex;align-items:center;
  padding:0 1.2rem;gap:.65rem;background:rgba(6,10,28,.65);backdrop-filter:blur(18px);
  border-bottom:1px solid var(--bdr);}
.topbar .brand{font-family:'Poppins',sans-serif;font-size:.8rem;font-weight:800;color:var(--ind-l);}
.topbar .sep{width:1px;height:18px;background:var(--bdr);}
.topbar .lbl{font-size:.76rem;font-weight:600;color:var(--t2);}
.topbar .sp{flex:1;}
.back{background:rgba(99,102,241,.18);border:1px solid rgba(99,102,241,.35);border-radius:7px;
  padding:.28rem .8rem;font-size:.7rem;font-weight:600;color:#C5D5FF;text-decoration:none;transition:.2s;}
.back:hover{background:rgba(99,102,241,.32);}
.wrap{max-width:820px;margin:0 auto;padding:1.6rem 1rem 4rem;}

/* hero */
.hero{text-align:center;padding:.5rem 0 1.4rem;}
.hero-icon{font-size:2.2rem;display:block;margin-bottom:.3rem;}
.hero-title{font-family:'Poppins',sans-serif;font-size:1.55rem;font-weight:800;color:#fff;
  text-shadow:0 2px 16px rgba(0,0,0,.7);}
.hero-sub{font-size:.8rem;color:rgba(210,225,255,.7);margin:.3rem 0;}

/* card */
.card{background:rgba(8,14,38,.75);border:1px solid var(--bdr);border-radius:var(--r);
  padding:1.4rem;margin-bottom:1rem;backdrop-filter:blur(14px);}
.card-title{font-size:.72rem;font-weight:700;color:var(--ind-l);text-transform:uppercase;
  letter-spacing:.1em;margin-bottom:.9rem;}

/* upload zones */
.dz{border:2px dashed rgba(99,102,241,.35);border-radius:11px;padding:1.5rem 1rem;
  text-align:center;cursor:pointer;transition:.22s;background:rgba(5,9,28,.6);user-select:none;}
.dz:hover,.dz.over{border-color:var(--ind);background:rgba(99,102,241,.09);}
.dz.ok{border-color:var(--grn);border-style:solid;background:rgba(16,185,129,.06);}
.dz .ico{font-size:1.8rem;display:block;margin-bottom:.3rem;}
.dz .hint{font-size:.72rem;color:var(--mu);margin-top:.25rem;}
.dz .fn{font-size:.78rem;font-weight:700;margin-top:.35rem;}
.dz.ok .fn{color:var(--grn);}

/* compare upload section — ppt-dz reused by swap feature */
.ppt-dz{border:2px dashed rgba(45,212,191,.3);border-radius:10px;padding:1.1rem;
  text-align:center;cursor:pointer;transition:.22s;background:rgba(45,212,191,.03);}
.ppt-dz:hover,.ppt-dz.over{border-color:var(--acc);background:rgba(45,212,191,.08);}
.ppt-dz.ok{border-color:var(--acc);border-style:solid;background:rgba(45,212,191,.06);}
.ppt-dz .fn{font-size:.76rem;color:var(--acc);font-weight:700;margin-top:.35rem;}

/* date selector */
.di{width:100%;padding:.65rem .9rem;background:rgba(5,9,28,.85);
  border:1.5px solid rgba(99,102,241,.4);border-radius:9px;color:var(--t1);
  font-size:.87rem;font-family:'Inter',sans-serif;color-scheme:dark;}
.di:focus{outline:none;border-color:var(--ind);}
.di-label{font-size:.72rem;font-weight:600;color:var(--t2);margin-bottom:.4rem;display:block;}

/* badge row */
.badges{display:flex;flex-wrap:wrap;gap:.35rem;margin:.6rem 0 .3rem;}
.badge{background:rgba(99,102,241,.1);border:1px solid rgba(99,102,241,.2);
  color:var(--ind-l);padding:.22rem .65rem;border-radius:6px;font-size:.72rem;font-weight:600;}

/* status */
.status-box{padding:.65rem 1rem;border-radius:9px;font-size:.82rem;margin-top:.65rem;display:none;}
.status-box.show{display:block;}
.status-box.ok{background:rgba(16,185,129,.09);border:1px solid rgba(16,185,129,.25);color:#34D399;}
.status-box.err{background:rgba(239,68,68,.09);border:1px solid rgba(239,68,68,.25);color:#FCA5A5;}
.status-box.info{background:rgba(99,102,241,.09);border:1px solid rgba(99,102,241,.25);color:var(--ind-l);}

/* generate button */
.btn-gen{width:100%;padding:.9rem;border:none;border-radius:11px;font-size:.95rem;font-weight:800;
  font-family:'Inter',sans-serif;cursor:pointer;transition:.22s;margin-top:.75rem;}
.btn-gen:disabled{opacity:.38;cursor:not-allowed;transform:none;}
.btn-gen.normal{background:linear-gradient(135deg,#1740C0,#2563EB 55%,#3B82F6);color:#fff;
  box-shadow:0 4px 20px rgba(37,99,235,.45);}
.btn-gen.normal:hover:not(:disabled){box-shadow:0 6px 28px rgba(37,99,235,.65);transform:translateY(-2px);}

/* download button */
.btn-dl{width:100%;padding:.85rem;background:rgba(5,9,28,.8);border:1px solid rgba(99,102,241,.35);
  border-radius:10px;color:var(--ind-l);font-size:.9rem;font-weight:700;cursor:pointer;
  transition:.2s;margin-top:.65rem;font-family:'Inter',sans-serif;display:none;}
.btn-dl:hover{background:rgba(99,102,241,.1);}

/* result modal */
.modal-bg{position:fixed;inset:0;z-index:9000;background:rgba(2,5,18,.78);
  backdrop-filter:blur(8px);display:none;align-items:center;justify-content:center;}
.modal-bg.show{display:flex;}
.modal{background:rgba(10,16,42,.97);border:1px solid rgba(99,102,241,.35);border-radius:20px;
  padding:2rem 2.2rem;width:90%;max-width:430px;
  box-shadow:0 24px 80px rgba(0,0,0,.7);text-align:center;
  animation:popIn .3s cubic-bezier(.34,1.56,.64,1);}
@keyframes popIn{from{opacity:0;transform:scale(.88) translateY(16px);}to{opacity:1;transform:none;}}
.modal-icon{font-size:3rem;display:block;margin-bottom:.6rem;}
.modal-title{font-family:'Poppins',sans-serif;font-size:1.1rem;font-weight:800;color:#fff;margin-bottom:.4rem;}
.modal-msg{font-size:.84rem;color:rgba(210,225,255,.8);line-height:1.6;margin-bottom:1.2rem;white-space:pre-line;}
.modal-stats{display:flex;gap:.6rem;justify-content:center;margin-bottom:1.3rem;flex-wrap:wrap;}
.modal-stat{background:rgba(99,102,241,.1);border:1px solid rgba(99,102,241,.25);
  border-radius:10px;padding:.65rem 1rem;min-width:90px;}
.modal-stat-n{font-size:1.5rem;font-weight:900;color:var(--ind-l);line-height:1;}
.modal-stat-l{font-size:.65rem;color:var(--mu);margin-top:.2rem;}
.modal-stat.grn{background:rgba(16,185,129,.1);border-color:rgba(16,185,129,.25);}
.modal-stat.grn .modal-stat-n{color:#34D399;}
.modal-btns{display:flex;gap:.65rem;justify-content:center;}
.mbtn{padding:.6rem 1.4rem;border-radius:9px;border:none;cursor:pointer;
  font-size:.85rem;font-weight:700;font-family:'Inter',sans-serif;transition:.2s;}
.mbtn-dl{background:linear-gradient(135deg,#065F38,#10B981);color:#fff;
  box-shadow:0 4px 16px rgba(16,185,129,.35);}
.mbtn-dl:hover{box-shadow:0 6px 24px rgba(16,185,129,.5);transform:translateY(-1px);}
.mbtn-close{background:rgba(255,255,255,.07);border:1px solid rgba(255,255,255,.12);
  color:rgba(148,163,184,.8);}
.mbtn-close:hover{background:rgba(255,255,255,.12);}
.footer-bar{position:fixed;bottom:0;left:0;right:0;z-index:9999;text-align:center;
  padding:.28rem 1rem;background:rgba(4,7,20,.72);backdrop-filter:blur(8px);
  border-top:1px solid rgba(255,255,255,.07);font-size:11px;font-weight:600;
  color:rgba(200,220,255,.7);letter-spacing:.06em;pointer-events:none;user-select:none;}
/* ── Image Slide Replacement (swap mode) ─────────────────────────────────── */
.swap-toggle{display:flex;align-items:flex-start;gap:.75rem;padding:.9rem 1rem;
  background:rgba(45,212,191,.07);border:1px solid rgba(45,212,191,.2);border-radius:10px;
  margin-top:.65rem;cursor:pointer;transition:.18s;}
.swap-toggle:hover{background:rgba(45,212,191,.12);}
.swap-toggle input[type=checkbox]{width:17px;height:17px;accent-color:var(--acc);
  cursor:pointer;flex-shrink:0;margin-top:2px;}
.swap-label{flex:1;}
.swap-label strong{font-size:.85rem;font-weight:700;color:var(--t1);display:block;}
.swap-label span{font-size:.73rem;color:var(--t2);margin-top:.2rem;display:block;line-height:1.5;}
.swap-body{border:1px solid rgba(45,212,191,.25);border-radius:11px;padding:1.1rem;
  margin-top:.65rem;background:rgba(45,212,191,.04);}
.swap-body-title{font-size:.74rem;font-weight:700;color:var(--acc);margin-bottom:.6rem;}
</style>
</head>
<body>

<div class="topbar">
  <span class="brand">APSG</span>
  <div class="sep"></div>
  <span class="lbl">📋 Excel Rejection Report</span>
  <div class="sp"></div>
  <span style="font-size:.65rem;color:var(--ind-l);font-weight:700;margin-right:.35rem;">✦ Karthi</span>
  <a href="/" class="back">← Dashboard</a>
</div>

<div class="wrap">
  <div class="hero">
    <span class="hero-icon">📊</span>
    <div class="hero-title">Rejection Report Generator</div>
    <div class="hero-sub">APSG Staging Ground · Convert Excel → PowerPoint</div>
  </div>

  <div class="card">
    <div class="card-title">📂 Step 1 — Upload Source Excel File</div>
    <div class="dz" id="excelZone">
      <input type="file" id="excelInput" accept=".xlsx,.xls,.csv" style="display:none">
      <span class="ico">📂</span>
      <strong style="font-size:.84rem;">Click to browse or drag &amp; drop</strong>
      <div class="hint">.xlsx · .xls · .csv accepted</div>
      <div class="fn" id="excelName"></div>
    </div>

    <!-- ── Image Slide Replacement Toggle (NEW — independent of Compare Mode) ── -->
    <label class="swap-toggle" id="swapToggleLabel" for="swapToggle" style="display:none;">
      <input type="checkbox" id="swapToggle" onchange="onSwapToggle()">
      <div class="swap-label">
        <strong>🔄 Use Existing PPT for Image Slide Replacement</strong>
        <span>Upload a Secondary PPT. For each image/detail slide where both
Ticket No. and E-Token match, that slide will be replaced by the
corresponding slide from the Secondary PPT. Summary slides are fully ignored.
Serial number continuity is strictly preserved.</span>
      </div>
    </label>

    <!-- Secondary PPT upload — shown when swap mode is ON -->
    <div id="swapBody" class="swap-body" style="display:none;">
      <div class="swap-body-title">🖼 Upload Secondary PPT (Image Slide Source)</div>
      <div class="ppt-dz" id="swapPptZone">
        <input type="file" id="swapPptInput" accept=".ppt,.pptx" style="display:none">
        <div style="font-size:1.4rem;margin-bottom:.2rem;">📊</div>
        <div style="font-size:.8rem;font-weight:600;color:var(--t2);">Drop Secondary PPT here or click to browse</div>
        <div style="font-size:.7rem;color:var(--mu);margin-top:.2rem;">Image/detail slides from this file will replace matched slides</div>
        <div class="fn" id="swapPptName"></div>
      </div>
    </div>
  </div>

  <!-- Step 2: Date (shown after file upload) -->
  <div class="card" id="dateCard" style="display:none;">
    <div class="card-title">📅 Step 2 — Select Report Date</div>
    <label class="di-label">Filter by Operational Date</label>
    <select class="di" id="dateSelect" onchange="onDateChange()">
      <option value="">All dates</option>
    </select>
    <div class="badges" id="badgeRow"></div>
  </div>

  <!-- Step 3: Generate -->
  <div class="card" id="genCard" style="display:none;">
    <div class="card-title" id="genCardTitle">🚀 Step 3 — Generate Report</div>
    <div class="status-box" id="statusBox"></div>
    <button class="btn-gen normal" id="btnGen" onclick="onGenerate()" disabled>
      ⚡ Generate Rejection Report
    </button>
    <button class="btn-dl" id="btnDl" onclick="downloadResult()">
      ⬇ Download PowerPoint
    </button>
  </div>
</div>

<!-- Result Modal -->
<div class="modal-bg" id="modalBg" onclick="if(event.target===this)closeModal()">
  <div class="modal">
    <span class="modal-icon" id="mIcon">✅</span>
    <div class="modal-title" id="mTitle"></div>
    <div class="modal-msg" id="mMsg"></div>
    <div class="modal-stats" id="mStats"></div>
    <div class="modal-btns">
      <button class="mbtn mbtn-dl" id="mDlBtn" onclick="downloadResult();closeModal();" style="display:none;">
        ⬇ Download PPT
      </button>
      <button class="mbtn mbtn-close" onclick="closeModal()">Close</button>
    </div>
  </div>
</div>

<div class="footer-bar">✦ Internal Reporting Platform — APSG Staging Ground &nbsp;·&nbsp; Developed by Karthik</div>

<script>
// ── State ──────────────────────────────────────────────────────────────────
var _excelFile = null, _selDate = '', _blob = null, _blobName = '';
var _swapMode  = false, _swapPptFile = null;
var esc = function(s){ return String(s||'').replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;'); };

// ── Excel upload (click + drag-drop) ──────────────────────────────────────
var excelZone = document.getElementById('excelZone');
var excelInput= document.getElementById('excelInput');
excelZone.addEventListener('click',  function(){ excelInput.click(); });
excelInput.addEventListener('change',function(){ if(excelInput.files[0]) applyExcel(excelInput.files[0]); });
excelZone.addEventListener('dragover', function(e){ e.preventDefault(); e.stopPropagation(); excelZone.classList.add('over'); });
excelZone.addEventListener('dragleave',function(e){ if(!excelZone.contains(e.relatedTarget)) excelZone.classList.remove('over'); });
excelZone.addEventListener('drop',     function(e){ e.preventDefault(); e.stopPropagation(); excelZone.classList.remove('over'); var f=e.dataTransfer.files[0]; if(f) applyExcel(f); });

function applyExcel(f) {
  _excelFile = f;
  document.getElementById('excelName').textContent = '✓ ' + f.name;
  excelZone.classList.add('ok');
  document.getElementById('swapToggleLabel').style.display = 'flex';
  loadDates();
}

// ── Swap zone (secondary PPT) ─────────────────────────────────────────────
var swapPptZone  = document.getElementById('swapPptZone');
var swapPptInput = document.getElementById('swapPptInput');
swapPptZone.addEventListener('click',  function(){ swapPptInput.click(); });
swapPptInput.addEventListener('change',function(){ if(swapPptInput.files[0]) applySwapPpt(swapPptInput.files[0]); });
swapPptZone.addEventListener('dragover', function(e){ e.preventDefault(); e.stopPropagation(); swapPptZone.classList.add('over'); });
swapPptZone.addEventListener('dragleave',function(e){ if(!swapPptZone.contains(e.relatedTarget)) swapPptZone.classList.remove('over'); });
swapPptZone.addEventListener('drop',     function(e){ e.preventDefault(); e.stopPropagation(); swapPptZone.classList.remove('over'); var f=e.dataTransfer.files[0]; if(f) applySwapPpt(f); });

function applySwapPpt(f) {
  if(!f.name.match(/\.pptx?$/i)){
    setStatus('Please upload a .pptx file for the secondary PPT.', 'err'); return;
  }
  _swapPptFile = f;
  document.getElementById('swapPptName').textContent = '✓ ' + f.name;
  swapPptZone.classList.add('ok');
  checkReady();
}

// ── Swap mode toggle — NEW ────────────────────────────────────────────────
function onSwapToggle() {
  _swapMode = document.getElementById('swapToggle').checked;
  document.getElementById('swapBody').style.display = _swapMode ? 'block' : 'none';
  if (!_swapMode) {
    _swapPptFile = null;
    swapPptZone.classList.remove('ok','over');
    document.getElementById('swapPptName').textContent = '';
    swapPptInput.value = '';
  }
  checkReady();
}

// ── Button state ───────────────────────────────────────────────────────────
function refreshBtn() {
  var btn = document.getElementById('btnGen');
  btn.className = 'btn-gen normal';
  btn.innerHTML = '⚡ Generate Rejection Report';
  document.getElementById('genCardTitle').textContent = '🚀 Step 3 — Generate Report';
}

function checkReady() {
  document.getElementById('btnGen').disabled = !_excelFile;
}

// ── Generate ───────────────────────────────────────────────────────────────
async function onGenerate() {
  clearStatus();
  _blob = null; _blobName = '';
  document.getElementById('btnDl').style.display = 'none';
  document.getElementById('btnGen').disabled = true;
  await doGenerate();
  document.getElementById('btnGen').disabled = false;
}
async function loadDates() {
  var fd = new FormData(); fd.append('file', _excelFile);
  try {
    var r = await fetch('/api/excel/dates', {method:'POST', body:fd});
    var d = await r.json();
    if (d.error) { setStatus('❌ ' + d.error, 'err'); return; }
    var sel = document.getElementById('dateSelect');
    sel.innerHTML = '<option value="">All dates</option>' +
      (d.dates||[]).map(function(x){ return '<option value="'+esc(x.value)+'">'+esc(x.label)+'</option>'; }).join('');
    if (d.dates && d.dates.length) { sel.value = d.dates[0].value; _selDate = d.dates[0].value; }
    document.getElementById('dateCard').style.display = 'block';
    document.getElementById('genCard').style.display  = 'block';
    await updateBadges();
    checkReady();
  } catch(ex) { setStatus('❌ ' + ex.message, 'err'); }
}

function onDateChange() {
  _selDate = document.getElementById('dateSelect').value;
  updateBadges();
}

async function updateBadges() {
  if (!_excelFile) return;
  try {
    var fd = new FormData(); fd.append('file', _excelFile); fd.append('date', _selDate);
    var r = await fetch('/api/excel/preview', {method:'POST', body:fd});
    var d = await r.json();
    if (d.error) return;
    var html = '';
    if (d.badges) d.badges.forEach(function(b){
      html += '<span class="badge">'+esc(b.site)+' <span style="opacity:.5">×'+b.count+'</span></span>';
    });
    document.getElementById('badgeRow').innerHTML = html;
  } catch(ex) {}
}

// ── Button state ───────────────────────────────────────────────────────────
function refreshBtn() {
  var btn = document.getElementById('btnGen');
  btn.className = 'btn-gen normal';
  btn.innerHTML = '⚡ Generate Rejection Report';
  document.getElementById('genCardTitle').textContent = '🚀 Step 3 — Generate Report';
}

function checkReady() {
  document.getElementById('btnGen').disabled = !_excelFile;
}

// ── Generate ───────────────────────────────────────────────────────────────
async function onGenerate() {
  clearStatus();
  _blob = null; _blobName = '';
  document.getElementById('btnDl').style.display = 'none';
  document.getElementById('btnGen').disabled = true;
  await doGenerate();
  document.getElementById('btnGen').disabled = false;
}

// Normal generation
async function doGenerate() {
  setStatus('⏳ Generating report…', 'info');
  try {
    var fd = new FormData();
    fd.append('file', _excelFile);
    fd.append('date', _selDate);
    // NEW — include secondary PPT for image slide replacement if swap mode is on
    if (_swapMode && _swapPptFile) {
      fd.append('secondary_ppt', _swapPptFile);
    }
    var r = await fetch('/api/excel/generate', {method:'POST', body:fd});
    if (!r.ok) { var j=await r.json().catch(function(){return{};}); setStatus('❌ '+(j.error||'Server error'), 'err'); return; }
    _blob = await r.blob();
    var cd = r.headers.get('Content-Disposition')||'';
    _blobName = (cd.match(/filename="([^"]+)"/)||[])[1] || 'Rejection_Report.pptx';
    var swapped = parseInt(r.headers.get('X-Replaced-Count')||'0', 10);
    var swapNote = (swapped > 0)
      ? '\n' + swapped + ' image slide'+(swapped===1?'':'s')+' replaced from Secondary PPT.'
      : (_swapMode && _swapPptFile ? '\nNo matching slides found in Secondary PPT.' : '');
    setStatus('✅ Ready — '+esc(_blobName)+swapNote, 'ok');
    document.getElementById('btnDl').style.display = 'block';
    showModal({icon:'✅', title:'Report Ready',
      msg:'Your rejection report has been generated.'+swapNote+'\nClick Download to save the file.',
      stats: swapped > 0 ? [{n:swapped, l:'Slides Replaced', cls:'grn'}] : [],
      showDl:true});
  } catch(ex) { setStatus('❌ ' + ex.message, 'err'); }
}

// ── Download ──────────────────────────────────────────────────────────────
function downloadResult() {
  if (!_blob) return;
  var url = URL.createObjectURL(_blob);
  var a   = document.createElement('a'); a.href = url; a.download = _blobName; a.click();
  setTimeout(function(){ URL.revokeObjectURL(url); }, 1500);
}

// ── Modal ─────────────────────────────────────────────────────────────────
function showModal(o) {
  document.getElementById('mIcon').textContent  = o.icon;
  document.getElementById('mTitle').textContent = o.title;
  document.getElementById('mMsg').textContent   = o.msg;
  var sEl = document.getElementById('mStats'); sEl.innerHTML = '';
  if (o.stats && o.stats.length) {
    o.stats.forEach(function(s){
      sEl.innerHTML += '<div class="modal-stat '+s.cls+'"><div class="modal-stat-n">'+s.n+'</div><div class="modal-stat-l">'+esc(s.l)+'</div></div>';
    });
  }
  document.getElementById('mDlBtn').style.display = o.showDl ? 'inline-flex' : 'none';
  document.getElementById('modalBg').classList.add('show');
}
function closeModal(){ document.getElementById('modalBg').classList.remove('show'); }
document.addEventListener('keydown', function(e){ if(e.key==='Escape') closeModal(); });

// ── Status helper ─────────────────────────────────────────────────────────
function setStatus(msg, type) {
  var el = document.getElementById('statusBox');
  el.className = 'status-box show ' + (type||'');
  el.textContent = msg;
}
function clearStatus() { document.getElementById('statusBox').className = 'status-box'; }
</script>
</body>
</html>"""







PHOTO_MERGE_HTML = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"/>
<meta name="viewport" content="width=device-width,initial-scale=1"/>
<title>APSG PowerPoint Merger</title>
<link href="https://fonts.googleapis.com/css2?family=Syne:wght@400;600;700;800&family=DM+Sans:ital,opsz,wght@0,9..40,300;0,9..40,400;0,9..40,500;0,9..40,600&display=swap" rel="stylesheet"/>
<style>
*,*::before,*::after{box-sizing:border-box;margin:0;padding:0;}
:root{
  --bg:#0A0E1A; --s1:#111827; --s2:#1A2234; --s3:#222E42;
  --b1:rgba(255,255,255,.07); --b2:rgba(255,255,255,.13);
  --teal:#00D4AA; --teal-d:rgba(0,212,170,.12); --teal-g:rgba(0,212,170,.28);
  --amber:#FFB340; --amber-d:rgba(255,179,64,.1);
  --indigo:#818CF8; --indigo-d:rgba(129,140,248,.12);
  --emerald:#34D399; --emerald-d:rgba(52,211,153,.12);
  --red:#FF5B5B; --text:#E8EDF5; --text2:#8A94A6; --text3:#3D4A5C;
  --r:14px; --rs:8px;
}
html{scroll-behavior:smooth;}
body{font-family:'DM Sans',sans-serif;background:transparent;color:var(--text);
  min-height:100vh;font-size:14px;line-height:1.6;overflow-x:hidden;}
body::before{content:'';position:fixed;inset:0;z-index:0;pointer-events:none;
  background:
    radial-gradient(ellipse 70% 50% at 5% 5%,rgba(0,212,170,.05) 0%,transparent 55%),
    radial-gradient(ellipse 50% 40% at 95% 85%,rgba(129,140,248,.05) 0%,transparent 50%),
    radial-gradient(ellipse 40% 60% at 50% 50%,rgba(255,179,64,.025) 0%,transparent 60%);}
.page{max-width:960px;margin:0 auto;padding:0 20px 60px;position:relative;z-index:1;}

/* ── HEADER ── */
header{display:flex;align-items:center;justify-content:space-between;
  padding:26px 0 28px;border-bottom:1px solid var(--b1);margin-bottom:32px;}
.logo{display:flex;align-items:center;gap:14px;}
.logo-mark{width:44px;height:44px;border-radius:12px;flex-shrink:0;
  background:linear-gradient(135deg,var(--teal),#00A882);font-size:20px;
  display:flex;align-items:center;justify-content:center;
  box-shadow:0 0 28px var(--teal-g);}
.logo-text h1{font-family:'Syne',sans-serif;font-size:19px;font-weight:800;
  letter-spacing:-.3px;}
.logo-text p{font-size:11.5px;color:var(--text2);margin-top:1px;}
.by-tag{font-size:11px;color:var(--text3);display:flex;align-items:center;gap:6px;}
.by-tag b{color:var(--teal);background:var(--teal-d);padding:3px 9px;border-radius:20px;
  border:1px solid rgba(0,212,170,.2);font-weight:600;}

/* ── CARDS ── */
.card{background:var(--s1);border:1px solid var(--b1);border-radius:var(--r);
  margin-bottom:18px;overflow:hidden;transition:border-color .2s;}
.card:hover{border-color:var(--b2);}
.card-head{display:flex;align-items:center;gap:12px;padding:16px 20px;
  border-bottom:1px solid var(--b1);}
.badge{width:34px;height:34px;border-radius:9px;flex-shrink:0;
  display:flex;align-items:center;justify-content:center;font-size:15px;}
.badge.teal{background:var(--teal-d);} .badge.amber{background:var(--amber-d);}
.badge.indigo{background:var(--indigo-d);}
.ch h3{font-family:'Syne',sans-serif;font-size:13.5px;font-weight:700;}
.ch p{font-size:11px;color:var(--text2);margin-top:2px;}
.card-body{padding:20px;}

/* ── UPLOAD ── */
.drop-grid{display:grid;grid-template-columns:1fr 1fr;gap:12px;}
@media(max-width:560px){.drop-grid{grid-template-columns:1fr;}}
.dz{border:1.5px dashed var(--b2);border-radius:var(--r);padding:26px 16px;
  text-align:center;cursor:pointer;transition:all .2s;background:var(--s2);
  display:flex;flex-direction:column;align-items:center;gap:5px;min-height:130px;justify-content:center;}
.dz:hover,.dz.drag{border-color:var(--teal);background:var(--teal-d);}
.dz.ok{border-color:var(--teal);border-style:solid;background:rgba(0,212,170,.07);}
.dz-icon{font-size:24px;} .dz-label{font-family:'Syne',sans-serif;font-size:13px;font-weight:700;}
.dz-hint{font-size:11px;color:var(--text3);}
.dz-name{font-size:11px;color:var(--teal);font-weight:600;word-break:break-all;max-width:180px;}
input[type=file]{display:none;}

/* ── PREVIEW CARD ── */
.prev-card{background:var(--s1);border:1.5px solid rgba(0,212,170,.22);
  border-radius:var(--r);margin-bottom:18px;overflow:hidden;
  box-shadow:0 0 48px rgba(0,212,170,.05);}
.prev-head{display:flex;align-items:center;gap:12px;padding:16px 20px;
  background:linear-gradient(90deg,rgba(0,212,170,.07),transparent);
  border-bottom:1px solid rgba(0,212,170,.13);}

/* Slide canvas */
.slide-wrap{background:var(--s2);border-radius:10px;overflow:hidden;
  margin-bottom:20px;border:1px solid var(--b1);}
.slide-bar{background:var(--s1);padding:7px 14px;font-size:10px;font-weight:600;
  color:var(--text2);letter-spacing:.8px;text-transform:uppercase;
  border-bottom:1px solid var(--b1);display:flex;align-items:center;gap:7px;}
.live-dot{width:6px;height:6px;border-radius:50%;background:var(--teal);
  box-shadow:0 0 6px var(--teal);animation:pulse 2s ease-in-out infinite;}
@keyframes pulse{0%,100%{opacity:1;}50%{opacity:.3;}}
.slide-canvas{width:100%;padding:10px;background:#F8FAFC;}
.slide-inner{position:relative;width:100%;}

/* ── CONTROL GROUPS ── */
.ctrl-grid{display:grid;grid-template-columns:1fr 1fr;gap:16px;}
@media(max-width:640px){.ctrl-grid{grid-template-columns:1fr;}}

.ctrl-group{background:var(--s2);border-radius:var(--r);padding:18px;
  border:1px solid var(--b1);transition:border-color .2s;}
.ctrl-group:hover{border-color:var(--b2);}
.ctrl-group.top-g{border-top:3px solid var(--indigo);}
.ctrl-group.frt-g{border-top:3px solid var(--emerald);}
.ctrl-group.gap-g{border-top:3px solid var(--amber);grid-column:1/-1;}

.g-title{font-family:'Syne',sans-serif;font-size:11.5px;font-weight:700;
  color:var(--text2);letter-spacing:.5px;text-transform:uppercase;
  margin-bottom:14px;display:flex;align-items:center;gap:7px;}
.g-dot{width:8px;height:8px;border-radius:50%;flex-shrink:0;}
.top-g .g-dot{background:var(--indigo);} .frt-g .g-dot{background:var(--emerald);}
.gap-g .g-dot{background:var(--amber);}

/* independence badge */
.ind-badge{margin-left:auto;font-size:9.5px;padding:2px 7px;border-radius:20px;
  font-weight:600;letter-spacing:.3px;}
.top-g .ind-badge{background:rgba(129,140,248,.15);color:var(--indigo);border:1px solid rgba(129,140,248,.2);}
.frt-g .ind-badge{background:rgba(52,211,153,.15);color:var(--emerald);border:1px solid rgba(52,211,153,.2);}

.field-2col{display:grid;grid-template-columns:1fr 1fr;gap:10px;margin-bottom:10px;}
.field{display:flex;flex-direction:column;gap:4px;}
.field label{font-size:11px;font-weight:500;color:var(--text2);}
.sr{display:flex;align-items:center;gap:7px;}
.sr input[type=range]{flex:1;height:4px;appearance:none;background:var(--b2);
  border-radius:99px;outline:none;cursor:pointer;}
.top-g .sr input[type=range]{accent-color:var(--indigo);}
.frt-g .sr input[type=range]{accent-color:var(--emerald);}
.gap-g .sr input[type=range]{accent-color:var(--amber);}
input[type=range]::-webkit-slider-thumb{appearance:none;width:15px;height:15px;
  border-radius:50%;cursor:pointer;box-shadow:0 0 6px rgba(0,0,0,.4);}
.top-g input[type=range]::-webkit-slider-thumb{background:var(--indigo);}
.frt-g input[type=range]::-webkit-slider-thumb{background:var(--emerald);}
.gap-g input[type=range]::-webkit-slider-thumb{background:var(--amber);}
.vbadge{min-width:42px;padding:2px 6px;border-radius:5px;font-size:12px;font-weight:700;
  text-align:center;background:var(--s1);border:1px solid var(--b2);}
.top-g .vbadge{color:var(--indigo);} .frt-g .vbadge{color:var(--emerald);}
.gap-g .vbadge{color:var(--amber);}
.num-in{width:100%;padding:7px 10px;background:var(--s1);border:1px solid var(--b2);
  border-radius:var(--rs);font-size:13px;font-weight:500;color:var(--text);
  outline:none;transition:border-color .2s;}
.num-in:focus{border-color:var(--teal);}
.fhint{font-size:10px;color:var(--text3);}

/* gap row */
.gap-row{display:grid;grid-template-columns:1fr auto;gap:12px;align-items:end;}

/* save default row */
.save-row{display:flex;gap:8px;align-items:center;margin-top:14px;
  padding-top:12px;border-top:1px solid var(--b1);}
.save-btn{padding:7px 14px;background:var(--teal-d);border:1px solid rgba(0,212,170,.25);
  border-radius:var(--rs);color:var(--teal);font-size:12px;font-weight:600;cursor:pointer;
  transition:all .2s;}
.save-btn:hover{background:rgba(0,212,170,.2);}
.reset-btn{padding:7px 14px;background:transparent;border:1px solid var(--b2);
  border-radius:var(--rs);color:var(--text2);font-size:12px;font-weight:600;cursor:pointer;
  transition:all .2s;}
.reset-btn:hover{border-color:var(--text2);color:var(--text);}
.save-status{font-size:11px;color:var(--teal);opacity:0;transition:opacity .4s;}

/* ── ADVANCED ── */
.adv-grid{display:grid;grid-template-columns:1fr 1fr 1fr;gap:12px;}
@media(max-width:600px){.adv-grid{grid-template-columns:1fr;}}
.text-in,.sel-in{width:100%;padding:8px 11px;background:var(--s2);
  border:1px solid var(--b1);border-radius:var(--rs);font-size:13px;
  color:var(--text);outline:none;transition:border-color .2s;}
.text-in:focus,.sel-in:focus{border-color:var(--teal);}

/* ── RUN BUTTON ── */
.run-btn{width:100%;padding:13px;background:linear-gradient(135deg,var(--teal),#00A882);
  color:#041A14;border:none;border-radius:var(--r);
  font-family:'Syne',sans-serif;font-size:13px;font-weight:700;cursor:pointer;
  display:flex;align-items:center;justify-content:center;gap:8px;
  transition:all .2s;box-shadow:0 4px 28px var(--teal-g);margin-bottom:18px;letter-spacing:.2px;}
.run-btn:hover:not(:disabled){transform:translateY(-1px);box-shadow:0 6px 36px var(--teal-g);}
.run-btn:disabled{opacity:.35;cursor:not-allowed;transform:none;}

/* ── PROGRESS ── */
.prog-card{display:none;background:var(--s1);border:1px solid var(--b1);
  border-radius:var(--r);padding:20px;margin-bottom:18px;}
.prog-head{display:flex;align-items:center;gap:10px;margin-bottom:12px;}
.spinner{width:20px;height:20px;border-radius:50%;flex-shrink:0;
  border:2px solid rgba(0,212,170,.2);border-top-color:var(--teal);
  animation:spin .85s linear infinite;}
@keyframes spin{to{transform:rotate(360deg);}}
.prog-title{font-family:'Syne',sans-serif;font-size:14px;font-weight:700;}
.prog-track{background:var(--s2);border-radius:99px;height:5px;overflow:hidden;margin-bottom:8px;}
.prog-bar{height:100%;width:0%;border-radius:99px;transition:width .4s;
  background:linear-gradient(90deg,var(--teal),#00F5D4);}
.prog-st{font-size:11.5px;color:var(--text2);margin-bottom:10px;}
.log-term{background:#050A12;border-radius:var(--rs);padding:12px;
  font-family:'Courier New',monospace;font-size:11px;line-height:1.8;
  max-height:180px;overflow-y:auto;border:1px solid var(--b1);}
.log-ok{color:var(--teal);} .log-err{color:var(--red);} .log-li{color:#3D5A80;}

/* ── RESULT ── */
.res-card{display:none;background:var(--s1);border:1px solid rgba(0,212,170,.18);
  border-radius:var(--r);padding:20px;margin-bottom:18px;
  box-shadow:0 0 36px rgba(0,212,170,.05);}
.res-title{font-family:'Syne',sans-serif;font-size:16px;font-weight:800;
  display:flex;align-items:center;gap:10px;margin-bottom:6px;}
.res-sub{font-size:12px;color:var(--text2);margin-bottom:14px;}
.stats-row{display:flex;flex-wrap:wrap;gap:8px;margin-bottom:16px;}
.stat{background:var(--s2);border:1px solid var(--b1);border-radius:10px;
  padding:9px 14px;text-align:center;min-width:72px;}
.stat strong{display:block;font-size:21px;font-weight:800;line-height:1.1;}
.stat span{font-size:10px;color:var(--text2);text-transform:uppercase;letter-spacing:.5px;}
.dl-btn{display:inline-flex;align-items:center;gap:8px;
  background:linear-gradient(135deg,var(--teal),#00A882);color:#041A14;
  padding:10px 22px;border-radius:var(--rs);border:none;
  font-family:'Syne',sans-serif;font-size:13px;font-weight:800;cursor:pointer;
  margin-right:10px;box-shadow:0 4px 20px var(--teal-g);transition:all .2s;}
.dl-btn:hover{transform:translateY(-1px);}
.new-btn{display:inline-flex;align-items:center;gap:6px;
  background:transparent;border:1px solid var(--b2);color:var(--text2);
  padding:10px 18px;border-radius:var(--rs);font-size:13px;font-weight:600;cursor:pointer;
  transition:all .2s;}
.new-btn:hover{border-color:var(--text2);color:var(--text);}



/* Action cards */
.action-card { background: rgba(10,16,42,.85) !important; border-color: rgba(245,158,11,.5) !important; }
.rr-panel { background: rgba(6,12,38,.88) !important; border-color: rgba(99,102,241,.35) !important; }

/* Stats / text helpers */
.sec-hint, .upload-hint { color: #9BB8E0 !important; }
.chip-ok { color: #4ADE80 !important; }
.chip-info { color: #818CF8 !important; }
.chip-wait { color: #F59E0B !important; }

/* ── Developed by Karthik — fixed footer ── */
.dev-credit {
  position: fixed; bottom: 0; left: 0; right: 0; z-index: 9999;
  text-align: center; padding: .3rem 1rem;
  background: rgba(5, 8, 22, 0.75); backdrop-filter: blur(8px);
  border-top: 1px solid rgba(99,102,241,.2);
  font-size: 11px; font-weight: 600; color: rgba(160,180,220,.75);
  letter-spacing: .06em; font-family: 'Inter', system-ui, sans-serif;
  pointer-events: none; user-select: none;
}

/* ═══ GLOBAL BACKGROUND & TRANSPARENCY — v3 ═══════════════════════ */
html, body {
  background-image: url('/static/bg.jpg') !important;
  background-size: cover !important;
  background-position: center center !important;
  background-attachment: fixed !important;
  background-repeat: no-repeat !important;
  background-color: #08101E !important;
}
/* Single very-light overlay — image stays visible */
body::before {
  content: '' !important;
  position: fixed !important;
  inset: 0 !important;
  z-index: 0 !important;
  background: rgba(3, 7, 18, 0.45) !important;
  pointer-events: none !important;
}
body > * { position: relative; z-index: 1; }

/* ── Top-bar: fully transparent glass, no black border ── */
.top-bar {
  background: rgba(6, 10, 28, 0.55) !important;
  backdrop-filter: blur(18px) !important;
  -webkit-backdrop-filter: blur(18px) !important;
  border-bottom: 1px solid rgba(255,255,255,0.08) !important;
  position: sticky !important;
  top: 0 !important;
  z-index: 200 !important;
}

/* ── Cards / Panels — glass, no solid black fill ── */
.card, .section-card, .panel, .sec-card, .stats-box, .action-card {
  background: rgba(8, 14, 38, 0.70) !important;
  border: 1px solid rgba(255,255,255,0.10) !important;
  backdrop-filter: blur(16px) !important;
  -webkit-backdrop-filter: blur(16px) !important;
  box-shadow: 0 4px 32px rgba(0,0,0,0.35) !important;
}
.panel-head, .card-header {
  background: rgba(10, 18, 52, 0.72) !important;
  border-bottom: 1px solid rgba(255,255,255,0.07) !important;
}
.panel-body { background: rgba(5, 10, 30, 0.60) !important; }

/* ── Upload zones ── */
.upload-zone, .dz {
  background: rgba(6, 12, 34, 0.55) !important;
  border: 2px dashed rgba(99,102,241,0.55) !important;
}

/* ── Typography: all white/light ── */
body, h1, h2, h3, h4, p, span, div, td, th, label, a {
  color: #EEF3FF !important;
}
.hero-title, .brand-title, .card-title, .admin-hero-title {
  color: #FFFFFF !important;
  text-shadow: 0 2px 16px rgba(0,0,0,0.7) !important;
  font-weight: 800 !important;
}
.hero-sub, .brand-sub, .admin-hero-sub, .sec-hint {
  color: rgba(210, 225, 255, 0.80) !important;
}
.top-mini-brand, .top-page-label, .top-brand .brand-text {
  color: #FFFFFF !important;
  font-weight: 700 !important;
}
.back-btn {
  color: #C5D5FF !important;
  background: rgba(99,102,241,0.18) !important;
  border: 1px solid rgba(99,102,241,0.35) !important;
}
.back-btn:hover { background: rgba(99,102,241,0.32) !important; }

/* ── Inputs: legible on transparent backgrounds ── */
input[type=text], input[type=number], input[type=date],
input[type=password], select, textarea,
.input-num, .action-select, .rr-input, .rr-select, .date-input {
  background: rgba(5, 9, 28, 0.80) !important;
  border: 1.5px solid rgba(99,102,241,0.40) !important;
  color: #EEF3FF !important;
  font-size: 14px !important;
}
input::placeholder, textarea::placeholder {
  color: rgba(180, 200, 240, 0.50) !important;
}

/* ── Muted / secondary text ── */
.muted, .sec-label, [style*="color:#475569"],
[style*="color:#64748B"], [style*="color:#374167"] {
  color: rgba(190, 210, 255, 0.70) !important;
}

/* ── Action + RR cards ── */
.action-card {
  background: rgba(12, 18, 48, 0.78) !important;
  border-color: rgba(245,158,11,0.55) !important;
}
.rr-panel {
  background: rgba(8, 14, 42, 0.80) !important;
  border-color: rgba(99,102,241,0.40) !important;
}

/* ── Global font sizes ── */
body { font-size: 14px !important; }
h1 { font-size: 26px !important; }
h2 { font-size: 22px !important; }
h3, .hero-title { font-size: 22px !important; }
h4 { font-size: 17px !important; }

/* ── Dev-credit fixed footer ── */
.apsg-footer {
  position: fixed !important; bottom: 0 !important;
  left: 0 !important; right: 0 !important; z-index: 9999 !important;
  text-align: center !important; padding: .28rem 1rem !important;
  background: rgba(4, 7, 20, 0.70) !important;
  backdrop-filter: blur(8px) !important;
  border-top: 1px solid rgba(255,255,255,0.07) !important;
  font-size: 11px !important; font-weight: 600 !important;
  color: rgba(200, 220, 255, 0.70) !important;
  letter-spacing: .06em !important; pointer-events: none !important;
  user-select: none !important;
}
</style>
</head>
<body>

<div class="top-bar">
  <span class="top-mini-brand">APSG</span>
  <div class="top-sep"></div>
  <span class="top-page-label">🖼 PPT Alignment — Top/Front Photo Merge</span>
  <div class="top-spacer"></div>
  <span style="font-size:.65rem;color:#6366F1;font-weight:700;white-space:nowrap;margin-right:.5rem;letter-spacing:.01em;">✦ Karthi</span>
  <a href="/" class="back-btn">← Dashboard</a>
</div>

<div class="page">

<!-- HEADER -->
<header>
  <div class="logo">
    <div class="logo-mark">📊</div>
    <div class="logo-text">
      <h1>APSG PowerPoint Merger</h1>
      <p>Independent Photo Layout · v6.0</p>
    </div>
  </div>
  <div class="by-tag">Created by <b>Karthi</b></div>
</header>

<!-- UPLOAD -->
<div class="card">
  <div class="card-head">
    <div class="badge teal">📁</div>
    <div class="ch"><h3>Upload PPTX Files</h3><p>Drop both files — reference auto-detected</p></div>
  </div>
  <div class="card-body">
    <div class="drop-grid">
      <div>
        <div class="dz" id="dzA" onclick="document.getElementById('iA').click()"
             ondragover="ev(event,'A',true)" ondragleave="ev(event,'A',false)" ondrop="drp(event,'A')">
          <div class="dz-icon">🪨</div><div class="dz-label">Top Photo File</div>
          <div class="dz-hint">Click or drag .pptx</div><div class="dz-name" id="nA"></div>
        </div>
        <input type="file" accept=".pptx" id="iA" onchange="onFile('A',this)"/>
      </div>
      <div>
        <div class="dz" id="dzB" onclick="document.getElementById('iB').click()"
             ondragover="ev(event,'B',true)" ondragleave="ev(event,'B',false)" ondrop="drp(event,'B')">
          <div class="dz-icon">🚛</div><div class="dz-label">Front Photo File</div>
          <div class="dz-hint">Click or drag .pptx</div><div class="dz-name" id="nB"></div>
        </div>
        <input type="file" accept=".pptx" id="iB" onchange="onFile('B',this)"/>
      </div>
    </div>
  </div>
</div>

<!-- PREVIEW + CONTROLS -->
<div class="prev-card">
  <div class="prev-head">
    <div class="badge amber">🖼️</div>
    <div class="ch">
      <h3>Photo Layout — Live Preview &amp; Independent Controls</h3>
      <p>Top Photo and Front Photo are completely independent — adjusting one never moves the other</p>
    </div>
  </div>
  <div class="card-body">

    <!-- Live preview -->
    <div class="slide-wrap">
      <div class="slide-bar"><div class="live-dot"></div>Live Slide Preview</div>
      <div class="slide-canvas" id="slideCanvas">
        <div class="slide-inner" id="slideInner"></div>
      </div>
    </div>

    <!-- Controls -->
    <div class="ctrl-grid">

      <!-- TOP PHOTO — independent -->
      <div class="ctrl-group top-g">
        <div class="g-title">
          <div class="g-dot"></div>🪨 Top Photo (Soil/Material)
          <span class="ind-badge">INDEPENDENT</span>
        </div>
        <div class="field-2col">
          <div class="field">
            <label>Height (cm)</label>
            <div class="sr">
              <input type="range" id="sTopH" min="4" max="14" step="0.01" value="9.11"
                     oninput="set('TopH',this.value)"/>
              <span class="vbadge" id="vTopH">9.11</span>
            </div>
            <input type="number" class="num-in" id="nTopH" min="4" max="14" step="0.01" value="9.11"
                   oninput="set('TopH',this.value)"/>
          </div>
          <div class="field">
            <label>Width (cm)</label>
            <div class="sr">
              <input type="range" id="sTopW" min="4" max="22" step="0.01" value="15.28"
                     oninput="set('TopW',this.value)"/>
              <span class="vbadge" id="vTopW">15.28</span>
            </div>
            <input type="number" class="num-in" id="nTopW" min="4" max="22" step="0.01" value="15.28"
                   oninput="set('TopW',this.value)"/>
          </div>
        </div>
        <div class="field">
          <label>Left Position (cm) — independent absolute position</label>
          <div class="sr">
            <input type="range" id="sTopL" min="0" max="20" step="0.01" value="0.95"
                   oninput="set('TopL',this.value)"/>
            <span class="vbadge" id="vTopL">0.95</span>
          </div>
          <input type="number" class="num-in" id="nTopL" min="0" max="20" step="0.01" value="0.95"
                 oninput="set('TopL',this.value)"/>
          <span class="fhint">Move left edge — Front Photo stays where it is</span>
        </div>
        <div class="save-row">
          <button class="save-btn" onclick="saveDefault('top')">💾 Save as Default</button>
          <button class="reset-btn" onclick="resetTop()">↺ Reset Top</button>
          <span class="save-status" id="topSaved">✓ Saved!</span>
        </div>
      </div>

      <!-- FRONT PHOTO — independent -->
      <div class="ctrl-group frt-g">
        <div class="g-title">
          <div class="g-dot"></div>🚛 Front Photo (Truck Plate)
          <span class="ind-badge">INDEPENDENT</span>
        </div>
        <div class="field-2col">
          <div class="field">
            <label>Height (cm)</label>
            <div class="sr">
              <input type="range" id="sFrtH" min="4" max="14" step="0.01" value="9.11"
                     oninput="set('FrtH',this.value)"/>
              <span class="vbadge" id="vFrtH">9.11</span>
            </div>
            <input type="number" class="num-in" id="nFrtH" min="4" max="14" step="0.01" value="9.11"
                   oninput="set('FrtH',this.value)"/>
          </div>
          <div class="field">
            <label>Width (cm)</label>
            <div class="sr">
              <input type="range" id="sFrtW" min="4" max="22" step="0.01" value="15.51"
                     oninput="set('FrtW',this.value)"/>
              <span class="vbadge" id="vFrtW">15.51</span>
            </div>
            <input type="number" class="num-in" id="nFrtW" min="4" max="22" step="0.01" value="15.51"
                   oninput="set('FrtW',this.value)"/>
          </div>
        </div>
        <div class="field">
          <label>Left Position (cm) — independent absolute position</label>
          <div class="sr">
            <input type="range" id="sFrtL" min="0" max="32" step="0.01" value="18.73"
                   oninput="set('FrtL',this.value)"/>
            <span class="vbadge" id="vFrtL">18.73</span>
          </div>
          <input type="number" class="num-in" id="nFrtL" min="0" max="32" step="0.01" value="18.73"
                 oninput="set('FrtL',this.value)"/>
          <span class="fhint">Move left edge — Top Photo stays where it is</span>
        </div>
        <div class="save-row">
          <button class="save-btn" onclick="saveDefault('frt')">💾 Save as Default</button>
          <button class="reset-btn" onclick="resetFrt()">↺ Reset Front</button>
          <span class="save-status" id="frtSaved">✓ Saved!</span>
        </div>
      </div>

      <!-- GAP (visual-only helper) -->
      <div class="ctrl-group gap-g">
        <div class="g-title"><div class="g-dot"></div>⬌ Visual Gap Helper</div>
        <div class="gap-row">
          <div class="field">
            <label>Desired gap between photos (cm) — sets Front Photo left automatically</label>
            <div class="sr">
              <input type="range" id="sGap" min="0" max="6" step="0.1" value="2.5"
                     oninput="applyGap(this.value)"/>
              <span class="vbadge" id="vGap">2.5</span>
            </div>
            <input type="number" class="num-in" id="nGap" min="0" max="6" step="0.1" value="2.5"
                   oninput="applyGap(this.value)"/>
            <span class="fhint">Shortcut: sets Front Photo left = Top left + Top width + gap. After applying, both photos remain independent.</span>
          </div>
          <div style="display:flex;flex-direction:column;gap:8px;padding-bottom:24px;">
            <button class="save-btn" onclick="applyGapToFrt()">Apply Gap →</button>
            <button class="reset-btn" onclick="resetAll()">↺ Reset All</button>
          </div>
        </div>
      </div>

    </div>
  </div>
</div>

<!-- ADVANCED -->
<div class="card">
  <div class="card-head">
    <div class="badge indigo">⚙️</div>
    <div class="ch"><h3>Advanced Settings</h3><p>Only change if your template uses different shape names</p></div>
  </div>
  <div class="card-body">
    <div class="adv-grid">
      <div class="field"><label>Top Photo shape name</label>
        <input type="text" class="text-in" id="topPh" value="Rectangle 14"/></div>
      <div class="field"><label>Front Photo shape name</label>
        <input type="text" class="text-in" id="frontPh" value="Rectangle 15"/></div>
      <div class="field"><label>Log level</label>
        <select class="sel-in" id="verbSel">
          <option value="0">Standard</option>
          <option value="1">Verbose (debug)</option>
        </select></div>
    </div>
  </div>
</div>

<!-- RUN -->
<button class="run-btn" id="runBtn" onclick="runMerge()" disabled>
  ⚡ Generate Report
</button>

<!-- PROGRESS -->
<div class="prog-card" id="progCard">
  <div class="prog-head">
    <div class="spinner"></div>
    <div class="prog-title" id="progTitle">Processing…</div>
  </div>
  <div class="prog-track"><div class="prog-bar" id="progBar"></div></div>
  <div class="prog-st" id="progSt"></div>
  <div class="log-term" id="logBox"></div>
</div>

<!-- RESULT -->
<div class="res-card" id="resCard">
  <div class="res-title">✅ Report Generated Successfully</div>
  <div class="res-sub" id="resMsg"></div>
  <div class="stats-row" id="statsRow"></div>
  <button class="dl-btn" id="dlBtn" style="display:none">⬇ Download PPTX</button>
  <button class="new-btn" onclick="doReset()">↺ New Merge</button>
</div>

</div><!-- /page -->
<script>
// ════════════════════════════════════════════════════════════════════════════
//  STATE — each photo fully independent, no shared calculations
// ════════════════════════════════════════════════════════════════════════════
const FACTORY = { TopH:9.11, TopW:15.28, TopL:0.95, FrtH:9.11, FrtW:15.51, FrtL:18.73, Gap:2.5 };

// Load saved defaults from localStorage (or fall back to factory)
function loadDefaults() {
  const saved = JSON.parse(localStorage.getItem('apsg_defaults') || '{}');
  return Object.assign({}, FACTORY, saved);
}

let v = loadDefaults();

const files = { A: null, B: null };
const SLIDE_W = 33.867, SLIDE_H = 19.05;
const LABEL_TOP = 6.73, LABEL_H = 0.9, IMG_TOP = LABEL_TOP + LABEL_H + 0.1;

// ── File upload ──────────────────────────────────────────────────────────────
function onFile(tag, inp) {
  const f = inp.files[0];
  if (!f) return;
  if (!f.name.toLowerCase().endsWith('.pptx')) { alert('Please upload a .pptx file.'); return; }
  files[tag] = f;
  document.getElementById('n'+tag).textContent = f.name;
  document.getElementById('dz'+tag).classList.add('ok');
  checkReady();
}
function ev(e,tag,on){ e.preventDefault(); document.getElementById('dz'+tag).classList.toggle('drag',on); }
function drp(e,tag) {
  e.preventDefault(); document.getElementById('dz'+tag).classList.remove('drag');
  const f = e.dataTransfer.files[0];
  if (f) { files[tag]=f; document.getElementById('n'+tag).textContent=f.name; document.getElementById('dz'+tag).classList.add('ok'); checkReady(); }
}
function checkReady() { document.getElementById('runBtn').disabled = !(files.A && files.B); }

// ── CRITICAL: set() updates ONLY the named field — zero effect on others ────
function set(key, rawVal) {
  const val = parseFloat(rawVal);
  if (isNaN(val)) return;

  // Update ONLY this field's own state
  v[key] = val;

  const dec = key === 'Gap' ? 1 : 2;
  const disp = val.toFixed(dec);

  // Update ONLY this field's own UI elements
  const s = document.getElementById('s'+key);
  const n = document.getElementById('n'+key);
  const b = document.getElementById('v'+key);
  if (s && String(s.value) !== String(val)) s.value = val;
  if (n && n.value !== disp) n.value = disp;
  if (b) b.textContent = disp;

  // Redraw preview from current v[] state
  renderPreview();
}

// ── Gap helper: applies gap by moving FRONT LEFT only ─────────────────────
function applyGap(raw) {
  const g = parseFloat(raw);
  if (isNaN(g)) return;
  v.Gap = g;
  document.getElementById('sGap').value = g;
  document.getElementById('nGap').value = g.toFixed(1);
  document.getElementById('vGap').textContent = g.toFixed(1);
  // Only updates the gap display — does NOT auto-apply to front left
  renderPreview();
}
function applyGapToFrt() {
  // One-shot: set Front Left = Top Left + Top Width + Gap
  // After this, both remain independent — Gap slider no longer linked
  const newFrtL = v.TopL + v.TopW + v.Gap;
  set('FrtL', parseFloat(newFrtL.toFixed(2)));
}

// ── Save defaults ─────────────────────────────────────────────────────────
function saveDefault(who) {
  const saved = JSON.parse(localStorage.getItem('apsg_defaults') || '{}');
  if (who === 'top') {
    saved.TopH = v.TopH; saved.TopW = v.TopW; saved.TopL = v.TopL;
    flashSaved('topSaved');
  } else {
    saved.FrtH = v.FrtH; saved.FrtW = v.FrtW; saved.FrtL = v.FrtL;
    flashSaved('frtSaved');
  }
  localStorage.setItem('apsg_defaults', JSON.stringify(saved));
}
function flashSaved(id) {
  const el = document.getElementById(id);
  el.style.opacity = '1';
  setTimeout(() => { el.style.opacity = '0'; }, 2000);
}

// ── Reset helpers ─────────────────────────────────────────────────────────
function resetTop()  { ['TopH','TopW','TopL'].forEach(k => set(k, FACTORY[k])); }
function resetFrt()  { ['FrtH','FrtW','FrtL'].forEach(k => set(k, FACTORY[k])); }
function resetAll()  { Object.keys(FACTORY).forEach(k => set(k, FACTORY[k])); }

// ── Preview renderer — pure function of v[], no side effects ──────────────
function renderPreview() {
  const canvas = document.getElementById('slideCanvas');
  const inner  = document.getElementById('slideInner');
  const CW = canvas.offsetWidth - 20;
  const CH = CW * (SLIDE_H / SLIDE_W);
  inner.style.width  = CW + 'px';
  inner.style.height = CH + 'px';
  canvas.style.height = (CH + 20) + 'px';

  const sc = CW / SLIDE_W;
  const cm = n => n * sc;

  // Top photo uses ONLY v.TopL, v.TopW, v.TopH — no front values
  const tL = v.TopL, tW = v.TopW, tH = v.TopH;
  // Front photo uses ONLY v.FrtL, v.FrtW, v.FrtH — no top values
  const fL = v.FrtL, fW = v.FrtW, fH = v.FrtH;
  // Gap shown for visual reference only
  const gap = fL - (tL + tW);

  let html = `<div style="position:absolute;inset:0;background:#F8FAFC;border-radius:6px;"></div>`;

  // Title
  html += `<div style="position:absolute;left:${cm(.86)}px;top:${cm(.63)}px;
    width:${cm(30)}px;height:${cm(1.35)}px;background:#E5E7EB;border-radius:3px;
    display:flex;align-items:center;padding:0 ${cm(.28)}px;
    font-size:${cm(.31)}px;font-weight:700;color:#374151;white-space:nowrap;overflow:hidden;">
    1. Rejection Due To Stone Found at Material Platform</div>`;

  // Table
  html += `<div style="position:absolute;left:${cm(.95)}px;top:${cm(2.67)}px;
    width:${cm(31.9)}px;height:${cm(2.38)}px;background:#F3F4F6;
    border:1px solid #D1D5DB;border-radius:3px;
    display:flex;align-items:center;justify-content:center;
    font-size:${cm(.24)}px;color:#9CA3AF;">Ticket Data Table</div>`;

  // TOP PHOTO label — uses only top values
  html += `<div style="position:absolute;left:${cm(tL)}px;top:${cm(LABEL_TOP)}px;
    width:${cm(tW)}px;height:${cm(LABEL_H)}px;
    background:#EEF2FF;border-radius:3px;
    display:flex;align-items:center;justify-content:center;
    font-size:${cm(.27)}px;font-weight:700;color:#4338CA;letter-spacing:.8px;">TOP PHOTO</div>`;

  // FRONT PHOTO label — uses only front values (independent position)
  html += `<div style="position:absolute;left:${cm(fL)}px;top:${cm(LABEL_TOP)}px;
    width:${cm(fW)}px;height:${cm(LABEL_H)}px;
    background:#ECFDF5;border-radius:3px;
    display:flex;align-items:center;justify-content:center;
    font-size:${cm(.27)}px;font-weight:700;color:#065F46;letter-spacing:.8px;">FRONT PHOTO</div>`;

  // TOP PHOTO image — ONLY tL, tW, tH used
  html += `<div style="position:absolute;
    left:${cm(tL)}px;top:${cm(IMG_TOP)}px;
    width:${cm(tW)}px;height:${cm(tH)}px;
    background:linear-gradient(135deg,#EEF2FF,#C7D2FE);
    border:2px solid #818CF8;border-radius:4px;overflow:hidden;
    display:flex;flex-direction:column;align-items:center;justify-content:center;gap:2px;">
    <div style="font-size:${cm(.5)}px;">🪨</div>
    <div style="font-size:${cm(.21)}px;font-weight:700;color:#3730A3;">Top Photo</div>
    <div style="font-size:${cm(.175)}px;color:#818CF8;background:rgba(129,140,248,.12);
      padding:1px ${cm(.08)}px;border-radius:3px;margin-top:2px;">
      ${tW.toFixed(2)} × ${tH.toFixed(2)} cm &nbsp;|&nbsp; left: ${tL.toFixed(2)} cm</div>
  </div>`;

  // Gap visual (amber) — purely decorative, non-linked
  if (gap > 0.05) {
    const gL = tL + tW, gW = gap;
    const gH = Math.max(tH, fH);
    html += `<div style="position:absolute;
      left:${cm(gL)}px;top:${cm(IMG_TOP)}px;
      width:${cm(gW)}px;height:${cm(gH)}px;
      background:repeating-linear-gradient(45deg,
        rgba(255,179,64,.07) 0,rgba(255,179,64,.07) 3px,transparent 3px,transparent 7px);
      border-left:1.5px dashed #F59E0B;border-right:1.5px dashed #F59E0B;
      display:flex;align-items:center;justify-content:center;overflow:hidden;">
      <div style="font-size:${cm(.17)}px;font-weight:700;color:#D97706;
        writing-mode:vertical-lr;transform:rotate(180deg);white-space:nowrap;">
        ${gW.toFixed(2)} cm</div></div>`;
  }

  // FRONT PHOTO image — ONLY fL, fW, fH used — no dependency on top
  html += `<div style="position:absolute;
    left:${cm(fL)}px;top:${cm(IMG_TOP)}px;
    width:${cm(fW)}px;height:${cm(fH)}px;
    background:linear-gradient(135deg,#ECFDF5,#A7F3D0);
    border:2px solid #34D399;border-radius:4px;overflow:hidden;
    display:flex;flex-direction:column;align-items:center;justify-content:center;gap:2px;">
    <div style="font-size:${cm(.5)}px;">🚛</div>
    <div style="font-size:${cm(.21)}px;font-weight:700;color:#065F46;">Front Photo</div>
    <div style="font-size:${cm(.175)}px;color:#34D399;background:rgba(52,211,153,.12);
      padding:1px ${cm(.08)}px;border-radius:3px;margin-top:2px;">
      ${fW.toFixed(2)} × ${fH.toFixed(2)} cm &nbsp;|&nbsp; left: ${fL.toFixed(2)} cm</div>
  </div>`;

  // Footer
  html += `<div style="position:absolute;left:0;bottom:${cm(.25)}px;width:100%;
    height:${cm(1.15)}px;background:#F3F4F6;border-top:1px solid #E5E7EB;
    display:flex;align-items:center;padding:0 ${cm(.35)}px;
    font-size:${cm(.21)}px;color:#9CA3AF;font-weight:500;">
    TOA-SAMSUNG C&amp;T JOINT VENTURE</div>`;

  inner.innerHTML = html;
}

// ── Merge ─────────────────────────────────────────────────────────────────
let curJob=null, pollT=null;

async function runMerge() {
  document.getElementById('runBtn').disabled = true;
  document.getElementById('progCard').style.display = 'block';
  document.getElementById('resCard').style.display  = 'none';
  document.getElementById('progBar').style.width    = '0%';
  document.getElementById('logBox').innerHTML        = '';
  document.getElementById('progTitle').textContent  = 'Processing…';
  document.getElementById('progSt').textContent     = 'Uploading files…';

  logLine(`Top: left=${v.TopL.toFixed(2)} w=${v.TopW.toFixed(2)} h=${v.TopH.toFixed(2)} cm`,'ok');
  logLine(`Front: left=${v.FrtL.toFixed(2)} w=${v.FrtW.toFixed(2)} h=${v.FrtH.toFixed(2)} cm`,'ok');
  logLine(`(Both fully independent — no shared values)`,'li');

  const fd = new FormData();
  fd.append('file_a', files.A); fd.append('file_b', files.B);
  fd.append('top_placeholder',  document.getElementById('topPh').value.trim()||'Rectangle 14');
  fd.append('front_placeholder',document.getElementById('frontPh').value.trim()||'Rectangle 15');
  fd.append('verbose', document.getElementById('verbSel').value==='1'?'1':'0');
  // Each photo sends its own independent values
  fd.append('top_h_cm',     v.TopH);
  fd.append('top_w_cm',     v.TopW);
  fd.append('top_left_cm',  v.TopL);
  fd.append('front_h_cm',   v.FrtH);
  fd.append('front_w_cm',   v.FrtW);
  fd.append('front_left_cm',v.FrtL);
  fd.append('center_gap_cm',v.Gap);

  try {
    const r = await fetch('/api/merge',{method:'POST',body:fd});
    const j = await r.json();
    if (j.error){showErr(j.error);return;}
    curJob=j.job_id; pollT=setInterval(poll,700);
  } catch(e){showErr(String(e));}
}

async function poll() {
  try {
    const r=await fetch('/api/job/'+curJob); const j=await r.json();
    document.getElementById('progBar').style.width=(j.progress||0)+'%';
    document.getElementById('progSt').textContent=j.status||'';
    if(j.log_lines) j.log_lines.slice(-5).forEach(([lv,msg])=>logLine(msg,lv==='ERROR'?'err':'li'));
    if(j.status==='complete'){clearInterval(pollT);showResult(j);}
    else if(j.status==='error'){clearInterval(pollT);showErr(j.error||'Error');}
  }catch(e){clearInterval(pollT);showErr(String(e));}
}

function showResult(j) {
  document.getElementById('progCard').style.display='none';
  document.getElementById('resCard').style.display='block';
  const st=j.stats||{};
  document.getElementById('resMsg').textContent=`Output: ${j.output_name||''}  ·  Reference: ${j.reference_file||'?'}`;
  const sr=document.getElementById('statsRow'); sr.innerHTML='';
  [['Total',st.total||0],['Merged',st.merged||0],['Both ✓',st.both||0],
   ['Top only',st.top_only||0],['Front only',st.front_only||0],['None',st.none_found||0]]
  .forEach(([l,n])=>{const d=document.createElement('div');d.className='stat';
    d.innerHTML=`<strong>${n}</strong><span>${l}</span>`;sr.appendChild(d);});
  const db=document.getElementById('dlBtn');
  db.style.display=j.has_result?'inline-flex':'none';
  db.onclick=()=>{window.location='/api/download/'+curJob;};
}
function showErr(msg){
  document.getElementById('progTitle').textContent='❌ Error';
  document.getElementById('progSt').textContent=msg;
  document.getElementById('runBtn').disabled=false;
}
function doReset(){
  files.A=null;files.B=null;curJob=null;
  ['A','B'].forEach(t=>{document.getElementById('n'+t).textContent='';
    document.getElementById('dz'+t).classList.remove('ok','drag');});
  document.getElementById('progCard').style.display='none';
  document.getElementById('resCard').style.display='none';
  document.getElementById('logBox').innerHTML='';
  document.getElementById('runBtn').disabled=true;
}
function logLine(msg,cls){
  const b=document.getElementById('logBox');
  const d=document.createElement('div');d.className='log-'+cls;
  d.textContent=msg;b.appendChild(d);b.scrollTop=b.scrollHeight;
}

// ── Init: load saved defaults into UI ────────────────────────────────────
window.addEventListener('load',()=>{
  Object.keys(v).forEach(k=>set(k,v[k]));
  renderPreview();
});
window.addEventListener('resize', renderPreview);
</script>
<div class="apsg-footer">✦ Internal Reporting Platform — APSG Staging Ground &nbsp;·&nbsp; Developed by Karthik</div>
</body>
</html>
"""

DAILY_REPORT_HTML = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>Daily Report — APSG Report</title>
<style>

/* ═══ MODERN UI BASE (Blue/Purple/Cyan Theme) ═══ */
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&family=Poppins:wght@600;700;800&display=swap');
*, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }
:root {
  --bg:transparent; --card-bg: rgba(13,18,35,0.88);
  --indigo: #6366F1; --indigo-l: #818CF8; --cyan: #22D3EE;
  --purple: #A855F7; --green: #10B981; --amber: #F59E0B; --red: #F87171;
  --text: #E8EEF8; --muted: #64748B; --border: rgba(99,102,241,0.15);
}
body {
  font-family: 'Inter', system-ui, -apple-system, sans-serif;
  min-height: 100vh; font-size: 14px; line-height: 1.6;
  background:transparent; color: var(--text);
  position: relative;
}
body::before {
  content: ''; position: fixed; inset: 0; z-index: 0; pointer-events: none;
  background:
    radial-gradient(ellipse 70% 55% at 5% 10%, rgba(99,102,241,.09) 0%, transparent 65%),
    radial-gradient(ellipse 55% 45% at 95% 90%, rgba(168,85,247,.07) 0%, transparent 65%);
}
body::after {
  content: ''; position: fixed; inset: 0; z-index: 0; pointer-events: none; opacity: .016;
  background-image: linear-gradient(var(--indigo) 1px, transparent 1px),
    linear-gradient(90deg, var(--indigo) 1px, transparent 1px);
  background-size: 60px 60px;
}

/* ── Modern Top Bar ── */
.top-bar {
  position: sticky; top: 0; z-index: 200; height: 56px;
  display: flex; align-items: center; padding: 0 1.5rem; gap: .75rem;
  background: rgba(6,9,22,.96); backdrop-filter: blur(20px); -webkit-backdrop-filter: blur(20px);
  border-bottom: 1px solid var(--border);
}
.top-mini-brand {
  font-family: 'Poppins', sans-serif;
  font-size: .82rem; font-weight: 700; color: var(--text);
  letter-spacing: -.01em; white-space: nowrap; flex-shrink: 0;
  display: flex; align-items: center; gap: .45rem;
}
.top-mini-brand::before {
  content: ''; width: 8px; height: 8px; border-radius: 50%;
  background: linear-gradient(135deg, var(--indigo), var(--cyan)); flex-shrink: 0;
}
.top-sep { width: 1px; height: 18px; background: var(--border); flex-shrink: 0; }
.top-page-label { font-size: .75rem; font-weight: 600; color: var(--muted); white-space: nowrap; }
.top-brand-tag { font-size: .7rem; color: var(--muted); }
.top-spacer { flex: 1; }
.back-btn {
  background: rgba(99,102,241,.08); border: 1px solid rgba(99,102,241,.18);
  border-radius: 8px; padding: .3rem .9rem; font-size: .7rem; font-weight: 600;
  color: var(--indigo-l); text-decoration: none; transition: all .2s; white-space: nowrap;
}
.back-btn:hover { background: rgba(99,102,241,.18); transform: translateX(-2px); }
.main { position: relative; z-index: 1; max-width: 1400px; margin: 0 auto; padding: .7rem 1rem 2rem; }

/* ── Cards & Containers ── */
.container, .page-content { position: relative; z-index: 1; }
.card, .section-card, .panel {
  background: var(--card-bg); border: 1px solid var(--border);
  border-radius: 16px; backdrop-filter: blur(16px); -webkit-backdrop-filter: blur(16px);
  box-shadow: 0 8px 32px rgba(0,0,0,.4), inset 0 1px 0 rgba(255,255,255,.03);
  transition: box-shadow .25s, border-color .25s;
}
.card:hover, .section-card:hover { border-color: rgba(99,102,241,.25); }

/* ── Upload Zone — Modern drag & drop ── */
.upload-zone, .dz {
  border: 2px dashed rgba(99,102,241,.3); border-radius: 14px;
  padding: 2rem; text-align: center; cursor: pointer;
  transition: all .22s; background: rgba(99,102,241,.03);
  position: relative;
}
.upload-zone:hover, .dz:hover, .upload-zone.drag-over, .dz.drag-over {
  border-color: var(--indigo); background: rgba(99,102,241,.08);
  box-shadow: 0 0 0 4px rgba(99,102,241,.12);
}
.upload-zone.ok, .dz.ok {
  border-color: var(--green); background: rgba(16,185,129,.06);
  border-style: solid;
}
.upload-zone.ok:hover, .dz.ok:hover {
  border-color: var(--green); background: rgba(16,185,129,.1);
  box-shadow: 0 0 0 4px rgba(16,185,129,.1);
}
.upload-icon { font-size: 2rem; margin-bottom: .5rem; display: block; }
.upload-label { font-size: .82rem; color: var(--muted); font-weight: 500; }
.upload-hint { font-size: .7rem; color: rgba(100,116,139,.6); margin-top: .25rem; }
.upload-filename { font-size: .78rem; color: var(--green); font-weight: 600; margin-top: .4rem; }

/* ── Modern Buttons ── */
.btn-primary, .btn-generate, .btn-teal, .modal-btn-primary {
  background:transparent;
  color: #fff; border: none; border-radius: 10px;
  padding: .7rem 1.4rem; font-size: .85rem; font-weight: 700;
  font-family: 'Inter', sans-serif; cursor: pointer; letter-spacing: .01em;
  box-shadow: 0 4px 18px rgba(99,102,241,.35);
  position: relative; overflow: hidden;
  transition: transform .2s, box-shadow .2s;
}
.btn-primary::before, .btn-generate::before, .btn-teal::before, .modal-btn-primary::before {
  content: ''; position: absolute; inset: 0;
  background: linear-gradient(135deg, transparent, rgba(255,255,255,.12), transparent);
  transform: translateX(-100%); transition: transform .45s;
}
.btn-primary:hover, .btn-generate:hover, .modal-btn-primary:hover {
  transform: translateY(-2px) scale(1.01);
  box-shadow: 0 8px 28px rgba(99,102,241,.55);
}
.btn-primary:hover::before, .btn-generate:hover::before, .modal-btn-primary:hover::before {
  transform: translateX(100%);
}
.btn-primary:active, .btn-generate:active { transform: translateY(0) scale(.98); }
.btn-primary:disabled, .btn-generate:disabled { opacity: .45; cursor: not-allowed; transform: none; }

.btn-teal {
  background:transparent;
  box-shadow: 0 4px 18px rgba(16,185,129,.3);
}
.btn-teal:hover { box-shadow: 0 8px 28px rgba(16,185,129,.5); }

.btn-download, .btn-dl, .dl-btn {
  background: linear-gradient(135deg, #0D7A5F, #10B981);
  color: #fff; border: none; border-radius: 10px;
  padding: .65rem 1.3rem; font-size: .82rem; font-weight: 700;
  cursor: pointer; transition: all .22s;
  box-shadow: 0 4px 16px rgba(16,185,129,.3);
}
.btn-download:hover, .btn-dl:hover, .dl-btn:hover {
  transform: translateY(-2px); box-shadow: 0 8px 26px rgba(16,185,129,.5);
}

.btn-secondary, .btn-gray {
  background: rgba(30,41,86,.6); color: var(--muted);
  border: 1px solid rgba(99,102,241,.2); border-radius: 10px;
  padding: .65rem 1.2rem; font-size: .82rem; font-weight: 600;
  cursor: pointer; transition: all .2s;
}
.btn-secondary:hover, .btn-gray:hover { background: rgba(99,102,241,.1); color: var(--indigo-l); }

.btn-del, .btn-danger {
  background: rgba(239,68,68,.08); color: var(--red);
  border: 1px solid rgba(239,68,68,.2); border-radius: 10px;
  padding: .6rem 1.2rem; font-size: .8rem; font-weight: 600;
  cursor: pointer; transition: all .2s;
}
.btn-del:hover, .btn-danger:hover { background: rgba(239,68,68,.16); }
.btn-del:hover { animation: shake .3s ease; }

@keyframes shake {
  0%,100% { transform: translateX(0); }
  25% { transform: translateX(-3px); }
  75% { transform: translateX(3px); }
}

/* ── Form Inputs ── */
input[type=text], input[type=number], input[type=date], input[type=password],
select, textarea {
  background: rgba(6,10,24,.8); border: 1.5px solid rgba(30,41,86,.8);
  border-radius: 10px; color: var(--text); padding: .65rem .9rem;
  font-size: .85rem; font-family: 'Inter', sans-serif;
  transition: all .2s; width: 100%;
}
input:focus, select:focus, textarea:focus {
  outline: none; border-color: var(--indigo);
  box-shadow: 0 0 0 3px rgba(99,102,241,.16);
  background: rgba(10,14,32,.9);
}
input::placeholder, textarea::placeholder { color: rgba(100,116,139,.45); font-weight: 300; }

/* ── Alerts / Status ── */
.alert-success, .alert.success, .msg.success {
  background: rgba(16,185,129,.08); border: 1px solid rgba(16,185,129,.2);
  color: #34D399; border-radius: 10px; padding: .65rem .9rem;
  font-size: .8rem; font-weight: 500;
}
.alert-error, .alert.error, .msg.error {
  background: rgba(239,68,68,.08); border: 1px solid rgba(239,68,68,.2);
  color: var(--red); border-radius: 10px; padding: .65rem .9rem;
  font-size: .8rem; font-weight: 500;
}
.alert-warn, .alert.warn { 
  background: rgba(245,158,11,.08); border: 1px solid rgba(245,158,11,.2);
  color: var(--amber); border-radius: 10px; padding: .65rem .9rem;
  font-size: .8rem; font-weight: 500;
}

/* ── Tables ── */
table { width: 100%; border-collapse: collapse; font-size: .8rem; }
thead th {
  background: rgba(99,102,241,.08); color: var(--muted);
  font-weight: 700; font-size: .68rem; letter-spacing: .06em;
  text-transform: uppercase; padding: .65rem .9rem; text-align: left;
  border-bottom: 1px solid var(--border);
}
tbody tr { border-bottom: 1px solid rgba(30,45,80,.3); transition: background .15s; }
tbody tr:hover { background: rgba(99,102,241,.04); }
tbody td { padding: .6rem .9rem; color: var(--text); }

/* ── Tabs ── */
.tab { position: relative; transition: all .2s; }
.tab.active { color: var(--indigo-l) !important; }
.tab.active::after {
  content: ''; position: absolute; bottom: -1px; left: 15%; right: 15%;
  height: 2px; border-radius: 2px; background: var(--indigo);
  animation: tabIn .2s ease;
}
@keyframes tabIn { from { left: 50%; right: 50%; } to { left: 15%; right: 15%; } }

/* ── Spinner ── */
@keyframes spin { to { transform: rotate(360deg); } }
.spinner { width: 20px; height: 20px; border: 2px solid rgba(99,102,241,.2);
  border-top-color: var(--indigo); border-radius: 50%; animation: spin 1s linear infinite; }

/* ══ ANIMATED BUTTONS — Daily Report ══════════════════════════════════════════ */
.btn-primary::after{content:'';position:absolute;top:-50%;left:-75%;width:50%;height:200%;
  background:rgba(255,255,255,.15);transform:skewX(-20deg);transition:left .5s;pointer-events:none;}
.btn-primary:hover::after{left:150%;}
.btn-dl{animation:gentlePulse 3s ease-in-out infinite;transition:transform .2s,box-shadow .2s !important;}
.btn-dl:hover{animation:none;transform:translateY(-3px) !important;box-shadow:0 8px 28px rgba(16,185,129,.5) !important;}
@keyframes gentlePulse{0%,100%{box-shadow:0 4px 16px rgba(16,185,129,.3);}50%{box-shadow:0 4px 24px rgba(16,185,129,.55);}}
.upload-zone:not(.ok){animation:uploadGlow 4s ease-in-out infinite;}
.upload-zone:hover{animation:none;}
@keyframes uploadGlow{0%,100%{border-color:rgba(99,102,241,.3);}50%{border-color:rgba(99,102,241,.65);box-shadow:0 0 18px rgba(99,102,241,.12);}}

.hero{text-align:center;padding:.8rem 1rem .9rem;}
.hero-icon{font-size:2rem;display:block;margin-bottom:.25rem;}
.hero-title{font-size:1.5rem;font-weight:900;color:#F1F5FF;letter-spacing:-.03em;}
.hero-sub{font-size:.82rem;color:#374167;margin:.4rem 0;}
.hero-pill{display:inline-block;padding:.22rem .9rem;background:rgba(99,102,241,.1);border:1px solid rgba(99,102,241,.3);border-radius:20px;font-size:.7rem;color:#818CF8;font-style:italic;}
.hbar{height:1px;margin:.3rem 0 .9rem;background:linear-gradient(90deg,transparent,#1E2456 40%,#1E2456 60%,transparent);}
.two-col{display:grid;grid-template-columns:1fr 1fr;gap:1.5rem;}
@media(max-width:900px){.two-col{grid-template-columns:1fr;}}
.panel{border-radius:14px;overflow:hidden;margin-bottom:1rem;}
.panel-head{padding:.75rem 1.1rem;font-size:1rem;font-weight:800;letter-spacing:.015em;}
.panel-online .panel-head{background:linear-gradient(90deg,#0D1B40,#0A1530);color:#818CF8;border-left:4px solid #6366F1;}
.panel-wb .panel-head{background:linear-gradient(90deg,#0B1F18,#081610);color:#34D399;border-left:4px solid #10B981;}
.panel-body{background:#060C1C;border:1px solid #111827;border-top:none;padding:1.2rem;border-radius:0 0 14px 14px;}
.sec-card{background:#080F24;border:1px solid #162040;border-radius:12px;padding:.85rem 1.1rem;margin:.4rem 0 .85rem;}
.sec-label{font-size:.62rem;font-weight:700;letter-spacing:.12em;text-transform:uppercase;color:#6366F1;margin-bottom:.3rem;}
.sec-label-wb{color:#10B981;}
.sec-title{font-size:.95rem;font-weight:800;color:#E8EEF8;margin-bottom:.15rem;}
.sec-hint{font-size:.75rem;color:#2D3F60;line-height:1.5;}
.upload-zone{border:2px dashed rgba(99,102,241,.35);border-radius:12px;padding:1.8rem;text-align:center;cursor:pointer;position:relative;transition:.2s;background:rgba(8,15,40,.6);}
.upload-zone:hover{border-color:#6366F1;background:rgba(99,102,241,.06);}
.upload-zone input{position:absolute;inset:0;opacity:0;cursor:pointer;width:100%;height:100%;}
.upload-zone-wb{border-color:rgba(16,185,129,.35);}
.upload-zone-wb:hover{border-color:#10B981;background:rgba(16,185,129,.05);}
.chip-ok{display:inline-block;padding:.25rem .8rem;background:rgba(16,185,129,.12);color:#10B981;border:1px solid rgba(16,185,129,.3);border-radius:7px;font-size:.75rem;font-weight:700;}
.chip-wait{display:inline-block;padding:.25rem .8rem;background:rgba(251,191,36,.08);color:#F59E0B;border:1px solid rgba(251,191,36,.22);border-radius:7px;font-size:.75rem;font-weight:700;}
.chip-info{display:inline-block;padding:.25rem .8rem;background:rgba(99,102,241,.1);color:#818CF8;border:1px solid rgba(99,102,241,.25);border-radius:7px;font-size:.75rem;font-weight:700;}
.date-input{width:100%;padding:.6rem .9rem;background:#070D1E;border:1.5px solid #1E2456;border-radius:9px;color:#C8D8F8;font-size:.88rem;font-family:'Inter',sans-serif;}
.date-input:focus{outline:none;border-color:#6366F1;}
.stats-box{background:#060C1C;border:1px solid #111827;border-radius:12px;padding:.85rem 1rem;margin:.5rem 0 .85rem;}
.stats-box-wb{border-color:#0D2A1A;}
.stats-title{font-size:.65rem;font-weight:700;letter-spacing:.12em;text-transform:uppercase;margin-bottom:.55rem;}
.stats-title-online{color:#6366F1;}
.stats-title-wb{color:#10B981;}
.stats-inner table{width:100%;border-collapse:collapse;font-size:.82rem;}
.stats-inner td{padding:.28rem 0;color:#6B7280;}
.stats-inner td:last-child{text-align:right;font-weight:700;}
.val-ok{background:#031A10;border:1px solid #065F38;border-radius:11px;padding:.7rem 1rem;margin:.5rem 0;color:#10B981;font-size:.82rem;}
.val-err{background:#1A0508;border:1px solid #7F1D1D;border-radius:11px;padding:.7rem 1rem;margin:.5rem 0;color:#F87171;font-size:.82rem;}
.val-warn{background:#150E02;border:1px solid #78450A;border-radius:11px;padding:.7rem 1rem;margin:.5rem 0;color:#F59E0B;font-size:.82rem;}
.cmp-table{width:100%;border-collapse:collapse;font-size:.82rem;margin:.5rem 0 .8rem;}
.cmp-table th{background:#080F24;color:#818CF8;font-weight:700;padding:.5rem .8rem;border:1px solid #1E2456;text-align:center;}
.cmp-table td{padding:.42rem .8rem;border:1px solid #111827;color:#C8D8F8;text-align:center;}
.cmp-match{color:#10B981;font-weight:900;}
.cmp-miss{color:#F87171;font-weight:900;}
.pivot-outer{background:#060C1C;border:1px solid #111827;border-radius:13px;padding:1rem;margin:.8rem 0;overflow-x:auto;}
.pivot-title{font-size:.88rem;font-weight:800;color:#10B981;letter-spacing:.06em;text-align:center;padding:.4rem 1rem;background:rgba(16,185,129,.07);border-radius:8px;border:1px solid rgba(16,185,129,.18);text-transform:uppercase;margin-bottom:.8rem;}
.pivot-tbl{width:100%;border-collapse:collapse;font-size:.8rem;min-width:500px;}
.pivot-tbl th{background:#080F24;color:#818CF8;font-weight:700;padding:.5rem .7rem;border:1px solid #1E2456;text-align:center;}
.pivot-tbl td{padding:.42rem .7rem;border:1px solid #111827;color:#C8D8F8;text-align:center;}
.pivot-mat td{background:#07101E;color:#818CF8;font-weight:800;text-align:left;font-size:.84rem;border-top:2px solid #1E2456;}
.pivot-grand td{background:#08122A;color:#10B981;font-weight:900;border-top:2px solid #1E2456;}
.btn{width:100%;padding:.82rem 1rem;border-radius:12px;border:none;cursor:pointer;font-size:.95rem;font-weight:700;font-family:'Inter',sans-serif;transition:all .18s;margin-top:.5rem;}
.btn-primary{background:linear-gradient(135deg,#4338CA,#6366F1 55%,#818CF8);color:#fff;box-shadow:0 4px 20px rgba(99,102,241,.4);}
.btn-primary:hover{transform:translateY(-2px);box-shadow:0 6px 28px rgba(99,102,241,.55);}
.btn-primary:disabled{opacity:.4;cursor:not-allowed;transform:none;}
.btn-dl{background:linear-gradient(135deg,#065F38,#059669 55%,#10B981);color:#fff;box-shadow:0 4px 20px rgba(5,150,105,.4);}
.btn-dl:hover{transform:translateY(-2px);box-shadow:0 6px 28px rgba(16,185,129,.5);}
.btn-dl:disabled{opacity:.4;cursor:not-allowed;transform:none;}
.action-card{background:#0A1020;border:1.5px solid #F59E0B;border-radius:12px;padding:1rem;margin:.6rem 0;}
.action-title{font-size:.82rem;font-weight:800;color:#F59E0B;margin-bottom:.5rem;}
.action-select,.input-num{width:100%;padding:.55rem .8rem;background:#070D1E;border:1.5px solid #1E2456;border-radius:8px;color:#C8D8F8;font-size:.85rem;font-family:'Inter',sans-serif;margin-bottom:.3rem;}
.rr-panel{margin-top:.8rem;padding:.75rem;background:rgba(99,102,241,.07);border:1.5px solid rgba(99,102,241,.3);border-radius:10px;}
.rr-panel-title{font-size:.76rem;font-weight:800;color:#818CF8;margin-bottom:.5rem;letter-spacing:.03em;}
.rr-input{width:100%;padding:.48rem .75rem;background:#070D1E;border:1.5px solid #1E2456;border-radius:8px;color:#C8D8F8;font-size:.83rem;font-family:'Inter',sans-serif;margin-bottom:.4rem;}
.rr-select{width:100%;padding:.48rem .75rem;background:#070D1E;border:1.5px solid #1E2456;border-radius:8px;color:#C8D8F8;font-size:.83rem;font-family:'Inter',sans-serif;margin-bottom:.4rem;}
.rr-btn{display:inline-flex;align-items:center;gap:.4rem;padding:.45rem 1.1rem;background:linear-gradient(135deg,#4F46E5,#7C3AED);color:#fff;border:none;border-radius:8px;font-size:.8rem;font-weight:700;cursor:pointer;transition:opacity .2s;}
.rr-btn:hover{opacity:.85;}
.rr-btn:disabled{opacity:.4;cursor:not-allowed;}
.rr-status{font-size:.74rem;margin-top:.3rem;min-height:1.2rem;}
.step-row{display:flex;gap:.4rem;align-items:center;margin:.2rem 0 1rem;flex-wrap:wrap;}
.step-pip{display:inline-flex;align-items:center;gap:.3rem;padding:.2rem .7rem;border-radius:20px;font-size:.7rem;font-weight:700;}
.step-done{background:rgba(16,185,129,.15);color:#10B981;border:1px solid rgba(16,185,129,.35);}
.step-wait{background:rgba(245,158,11,.08);color:#F59E0B;border:1px solid rgba(245,158,11,.22);}
.step-arr{color:#1E2456;font-size:.7rem;}
hr{border:none;border-top:1px solid #0C1325;margin:1rem 0;}
.spinner{animation:spin 1s linear infinite;display:inline-block;}
@keyframes spin{to{transform:rotate(360deg);}}



/* Action cards */
.action-card { background: rgba(10,16,42,.85) !important; border-color: rgba(245,158,11,.5) !important; }
.rr-panel { background: rgba(6,12,38,.88) !important; border-color: rgba(99,102,241,.35) !important; }

/* Stats / text helpers */
.sec-hint, .upload-hint { color: #9BB8E0 !important; }
.chip-ok { color: #4ADE80 !important; }
.chip-info { color: #818CF8 !important; }
.chip-wait { color: #F59E0B !important; }

/* ── Developed by Karthik — fixed footer ── */
.dev-credit {
  position: fixed; bottom: 0; left: 0; right: 0; z-index: 9999;
  text-align: center; padding: .3rem 1rem;
  background: rgba(5, 8, 22, 0.75); backdrop-filter: blur(8px);
  border-top: 1px solid rgba(99,102,241,.2);
  font-size: 11px; font-weight: 600; color: rgba(160,180,220,.75);
  letter-spacing: .06em; font-family: 'Inter', system-ui, sans-serif;
  pointer-events: none; user-select: none;
}

/* ═══ GLOBAL BACKGROUND & TRANSPARENCY — v3 ═══════════════════════ */
html, body {
  background-image: url('/static/bg.jpg') !important;
  background-size: cover !important;
  background-position: center center !important;
  background-attachment: fixed !important;
  background-repeat: no-repeat !important;
  background-color: #08101E !important;
}
/* Single very-light overlay — image stays visible */
body::before {
  content: '' !important;
  position: fixed !important;
  inset: 0 !important;
  z-index: 0 !important;
  background: rgba(3, 7, 18, 0.45) !important;
  pointer-events: none !important;
}
body > * { position: relative; z-index: 1; }

/* ── Top-bar: fully transparent glass, no black border ── */
.top-bar {
  background: rgba(6, 10, 28, 0.55) !important;
  backdrop-filter: blur(18px) !important;
  -webkit-backdrop-filter: blur(18px) !important;
  border-bottom: 1px solid rgba(255,255,255,0.08) !important;
  position: sticky !important;
  top: 0 !important;
  z-index: 200 !important;
}

/* ── Cards / Panels — glass, no solid black fill ── */
.card, .section-card, .panel, .sec-card, .stats-box, .action-card {
  background: rgba(8, 14, 38, 0.70) !important;
  border: 1px solid rgba(255,255,255,0.10) !important;
  backdrop-filter: blur(16px) !important;
  -webkit-backdrop-filter: blur(16px) !important;
  box-shadow: 0 4px 32px rgba(0,0,0,0.35) !important;
}
.panel-head, .card-header {
  background: rgba(10, 18, 52, 0.72) !important;
  border-bottom: 1px solid rgba(255,255,255,0.07) !important;
}
.panel-body { background: rgba(5, 10, 30, 0.60) !important; }

/* ── Upload zones ── */
.upload-zone, .dz {
  background: rgba(6, 12, 34, 0.55) !important;
  border: 2px dashed rgba(99,102,241,0.55) !important;
}

/* ── Typography: all white/light ── */
body, h1, h2, h3, h4, p, span, div, td, th, label, a {
  color: #EEF3FF !important;
}
.hero-title, .brand-title, .card-title, .admin-hero-title {
  color: #FFFFFF !important;
  text-shadow: 0 2px 16px rgba(0,0,0,0.7) !important;
  font-weight: 800 !important;
}
.hero-sub, .brand-sub, .admin-hero-sub, .sec-hint {
  color: rgba(210, 225, 255, 0.80) !important;
}
.top-mini-brand, .top-page-label, .top-brand .brand-text {
  color: #FFFFFF !important;
  font-weight: 700 !important;
}
.back-btn {
  color: #C5D5FF !important;
  background: rgba(99,102,241,0.18) !important;
  border: 1px solid rgba(99,102,241,0.35) !important;
}
.back-btn:hover { background: rgba(99,102,241,0.32) !important; }

/* ── Inputs: legible on transparent backgrounds ── */
input[type=text], input[type=number], input[type=date],
input[type=password], select, textarea,
.input-num, .action-select, .rr-input, .rr-select, .date-input {
  background: rgba(5, 9, 28, 0.80) !important;
  border: 1.5px solid rgba(99,102,241,0.40) !important;
  color: #EEF3FF !important;
  font-size: 14px !important;
}
input::placeholder, textarea::placeholder {
  color: rgba(180, 200, 240, 0.50) !important;
}

/* ── Muted / secondary text ── */
.muted, .sec-label, [style*="color:#475569"],
[style*="color:#64748B"], [style*="color:#374167"] {
  color: rgba(190, 210, 255, 0.70) !important;
}

/* ── Action + RR cards ── */
.action-card {
  background: rgba(12, 18, 48, 0.78) !important;
  border-color: rgba(245,158,11,0.55) !important;
}
.rr-panel {
  background: rgba(8, 14, 42, 0.80) !important;
  border-color: rgba(99,102,241,0.40) !important;
}

/* ── Global font sizes ── */
body { font-size: 14px !important; }
h1 { font-size: 26px !important; }
h2 { font-size: 22px !important; }
h3, .hero-title { font-size: 22px !important; }
h4 { font-size: 17px !important; }

/* ── Dev-credit fixed footer ── */
.apsg-footer {
  position: fixed !important; bottom: 0 !important;
  left: 0 !important; right: 0 !important; z-index: 9999 !important;
  text-align: center !important; padding: .28rem 1rem !important;
  background: rgba(4, 7, 20, 0.70) !important;
  backdrop-filter: blur(8px) !important;
  border-top: 1px solid rgba(255,255,255,0.07) !important;
  font-size: 11px !important; font-weight: 600 !important;
  color: rgba(200, 220, 255, 0.70) !important;
  letter-spacing: .06em !important; pointer-events: none !important;
  user-select: none !important;
}
</style>
</head>
<body>
<div class="top-bar">
  <span class="top-mini-brand">APSG</span>
  <div class="top-sep"></div>
  <span class="top-page-label">📊 Daily Report Generator</span>
  <div class="top-spacer"></div>
  <span style="font-size:.65rem;color:#6366F1;font-weight:700;white-space:nowrap;margin-right:.5rem;letter-spacing:.01em;">✦ Karthi</span>
  <a href="/" class="back-btn">← Dashboard</a>
</div>

<div class="main">
  <div class="hero">
    <span class="hero-icon">📊</span>
    <div class="hero-title">Daily Report Generator</div>
    <div class="hero-sub">Staging Ground Report System · Phase 3</div>
    <span class="hero-pill">✦ WB Summary Contract Report · WB Server Comparison</span>
  </div>
  <div class="hbar"></div>
  <div class="two-col">
    <div>
      <div class="panel panel-online">
        <div class="panel-head">🌐 ONLINE SUMMARY / ONLINE DATA</div>
        <div class="panel-body">
          <div class="sec-card">
            <div class="sec-label">Step 1 — Upload</div>
            <div class="sec-title">📂 Upload Online Export</div>
            <div class="sec-hint">Drag & drop or click — .xlsx, .xls, .csv accepted</div>
          </div>
          <div class="upload-zone" onclick="document.getElementById('onlineFile').click()">
            <input type="file" id="onlineFile" accept=".xlsx,.xls,.csv" onchange="onOnlineFileChange()" onclick="event.stopPropagation()">
            <div style="font-size:2rem;margin-bottom:.4rem">📂</div>
            <strong style="font-size:.9rem">Drop or click to upload Online export</strong>
            <p style="font-size:.76rem;color:#2D3F60;margin-top:.3rem">.xlsx · .xls · .csv</p>
          </div>
          <div id="onlineFileStatus" style="margin:.5rem 0"></div>
          <hr>
          <div id="onlineDateSection" style="display:none">
            <div class="sec-card">
              <div class="sec-label">Step 2 — Filter Date</div>
              <div class="sec-title">📅 Filter Date (Online)</div>
              <div class="sec-hint">Filtered on <strong style="color:#60A5FA">WB In Time</strong></div>
            </div>
            <input type="date" class="date-input" id="onlineDate" onchange="onOnlineDateChange()" style="color-scheme:dark;">
            <div id="onlineDateStatus" style="margin:.5rem 0"></div>
            <hr>
          </div>
          <div id="onlineValidationSection" style="display:none">
            <div id="onlineActionCards"></div>
            <div id="onlineValResult"></div>
            <div id="onlineStats" style="display:none"></div>
            <hr>
            <div class="step-row" id="stepTracker"></div>
            <button class="btn btn-primary" id="generateBtn" onclick="generateReport()" disabled>⚡ Generate Report</button>
            <div id="generateStatus" style="margin:.5rem 0;font-size:.82rem;color:#818CF8"></div>
            <button class="btn btn-dl" id="downloadBtn" onclick="downloadReport()" disabled style="display:none">⬇️ Download Report</button>
          </div>
        </div>
      </div>
    </div>
    <div>
      <div class="panel panel-wb">
        <div class="panel-head">⚖️ WB DATA <span style="font-size:.72rem;font-weight:600;opacity:.65">WB Server · Comparison View</span></div>
        <div class="panel-body">
          <div class="sec-card">
            <div class="sec-label sec-label-wb">Step 1 — Upload</div>
            <div class="sec-title">📂 Upload WB Transaction List</div>
            <div class="sec-hint">.xlsx / .xls / .csv — Weighbridge server export</div>
          </div>
          <div class="upload-zone upload-zone-wb" onclick="document.getElementById('wbFile').click()">
            <input type="file" id="wbFile" accept=".xlsx,.xls,.csv" onchange="onWbFileChange()" onclick="event.stopPropagation()">
            <div style="font-size:2rem;margin-bottom:.4rem">⚖️</div>
            <strong style="font-size:.9rem">Drop or click to upload WB file</strong>
            <p style="font-size:.76rem;color:#2D3F60;margin-top:.3rem">.xlsx · .xls · .csv</p>
          </div>
          <div id="wbFileStatus" style="margin:.5rem 0"></div>
          <hr>
          <div id="wbDateSection" style="display:none">
            <div class="sec-card">
              <div class="sec-label sec-label-wb">Step 2 — Filter Date</div>
              <div class="sec-title">📅 Filter Date (WB)</div>
              <div class="sec-hint">Aligns with Online for comparison</div>
            </div>
            <input type="date" class="date-input" id="wbDate" onchange="onWbDateChange()" style="color-scheme:dark;border-color:#1E3A2A;">
            <hr>
          </div>
          <div id="wbResultSection" style="display:none">
            <div id="wbStats"></div>
            <div id="comparisonTable"></div>
            <div id="pivotSection"></div>
            <button class="btn btn-dl" id="pivotDlBtn" onclick="downloadPivot()" style="display:none;margin-top:.5rem">⬇️ Download Pivot (Excel)</button>
          </div>
        </div>
      </div>
    </div>
  </div>
</div>
<script>
// ── State ─────────────────────────────────────────────────────────────────────
let onlineTmpPath='',onlineDate='',wbTmpPath='',wbDate='';
let corrections={},wbDecisions={},reportBlob=null,reportName='',pivotTmpPath='',pivotFname='';
let onlineStats=null,wbStats=null;
const _incompleteItems={};   // idx → item metadata (token, in_weight)
const _rrSerials={};         // idx → saved RR serial string
const _rrReasons={};         // idx → saved RR reason string
let _validationTimer=null;   // debounce handle
let _wbPivotTimer=null;

// ═══════════════════════════════════════════════════════════════════════════════
// CORE FIX: scheduleValidation NEVER clears cards while user is typing.
// It only fires AFTER a 900 ms idle gap, and _syncActionCards NEVER destroys
// a card that still has an active correction or a focused input.
// ═══════════════════════════════════════════════════════════════════════════════

function _userIsTyping(){
  const active=document.activeElement;
  if(!active)return false;
  const tag=active.tagName.toLowerCase();
  if(tag==='input'||tag==='textarea')return true;
  return false;
}

function scheduleValidation(){
  clearTimeout(_validationTimer);
  _validationTimer=setTimeout(function(){
    if(_userIsTyping()){scheduleValidation();return;} // re-defer while typing
    runOnlineValidation();
  },900);
}

// ── File upload ───────────────────────────────────────────────────────────────
async function onOnlineFileChange(){
  const f=document.getElementById('onlineFile').files[0];if(!f)return;
  setEl('onlineFileStatus','<span class="chip-info">⏳ Uploading…</span>');
  const fd=new FormData();fd.append('file',f);
  try{
    const res=await fetch('/api/daily/upload',{method:'POST',body:fd});
    const d=await res.json();
    if(d.error){setEl('onlineFileStatus',`<span class="chip-wait">❌ ${esc(d.error)}</span>`);return;}
    onlineTmpPath=d.tmp_path;
    setEl('onlineFileStatus',`<span class="chip-ok">✓ ${esc(f.name)} — ${d.rows.toLocaleString()} rows loaded</span>`);
    const di=document.getElementById('onlineDate');
    di.min=d.min_date;di.max=d.max_date;di.value=d.max_date;onlineDate=d.max_date;
    show('onlineDateSection');show('onlineValidationSection');
    corrections={};await runOnlineValidation();
  }catch(e){setEl('onlineFileStatus',`<span class="chip-wait">❌ ${e.message}</span>`);}
}

async function onOnlineDateChange(){
  onlineDate=document.getElementById('onlineDate').value;corrections={};
  clearEl('onlineActionCards');clearEl('onlineValResult');clearEl('onlineStats');
  hide('downloadBtn');document.getElementById('generateBtn').disabled=true;
  await runOnlineValidation();
}

// ── Validation ────────────────────────────────────────────────────────────────
async function runOnlineValidation(){
  if(!onlineTmpPath||!onlineDate)return;
  setEl('onlineValResult','<div class="chip-info">🔍 Running validation…</div>');
  try{
    const res=await fetch('/api/daily/validate',{method:'POST',headers:{'Content-Type':'application/json'},
      body:JSON.stringify({tmp_path:onlineTmpPath,filter_date:onlineDate,corrections})});
    const d=await res.json();
    if(d.error){setEl('onlineValResult',`<div class="val-err">❌ ${esc(d.error)}</div>`);return;}
    if(d.rows===0){setEl('onlineValResult','<div class="val-warn">⚠️ No records for selected date.</div>');return;}
    _syncActionCards(d.incomplete);
    if(d.all_resolved||d.incomplete.length===0){
      if(d.errors&&d.errors.length>0){
        setEl('onlineValResult',d.errors.map(e=>`<div class="val-err">❌ ${esc(e)}</div>`).join(''));
      }else{
        setEl('onlineValResult','<div class="val-ok">✅ All validations completed successfully</div>');
      }
      onlineStats=d.stats;renderOnlineStats(d.stats);show('onlineStats');
      renderComparisonTable();
    }else{clearEl('onlineValResult');}
    updateStepTracker();updateGenerateBtn();
  }catch(e){setEl('onlineValResult',`<div class="val-err">❌ ${e.message}</div>`);}
}

// ── Smart card sync ───────────────────────────────────────────────────────────
// KEY RULES:
//  1. Never destroy a card that has a focused input inside it.
//  2. Never destroy a card that has a pending correction (user has made a choice).
//  3. Only add cards for new incomplete items; skip existing ones entirely.
function _syncActionCards(incomplete){
  const container=document.getElementById('onlineActionCards');
  if(!container)return;

  // Store metadata for any new items
  if(incomplete&&incomplete.length){
    for(const item of incomplete)_incompleteItems[item.idx]=item;
  }

  // Build set of idx values the server says are still incomplete
  const serverIncomplete=new Set((incomplete||[]).map(i=>i.idx));

  // Remove cards ONLY if: server no longer lists them AND no correction AND not focused
  container.querySelectorAll('[id^="acard-"]').forEach(el=>{
    const idx=parseInt(el.id.replace('acard-',''));
    const hasFocus=el.contains(document.activeElement);
    const hasCorrection=!!corrections[idx];
    if(!serverIncomplete.has(idx)&&!hasFocus&&!hasCorrection){
      el.remove();
    }
  });

  // Add cards for items not yet in the DOM
  for(const item of (incomplete||[])){
    if(document.getElementById(`acard-${item.idx}`))continue;
    const corr=corrections[item.idx]||{};
    const sel=corr.Accepted||'';
    const card=document.createElement('div');
    card.className='action-card';card.id=`acard-${item.idx}`;
    card.innerHTML=`
      <div class="action-title">⚠️ Action Required — E-Token: <span style="color:#F59E0B;font-weight:900">${esc(item.token)}</span></div>
      <div style="font-size:12px;color:#9BB4D4;margin-bottom:.4rem;">Missing Date Out / Time Out — Select action to resolve:</div>
      <select class="action-select" id="dec-${item.idx}">
        <option value="" ${sel===''?'selected':''}>— select —</option>
        <option value="YES" ${sel==='YES'?'selected':''}>YES — Accept</option>
        <option value="NO" ${sel==='NO'?'selected':''}>NO — Reject</option>
      </select>
      <div id="dec-detail-${item.idx}"></div>
      <div id="rr-panel-${item.idx}"></div>`;
    card.querySelector('select').addEventListener('change',function(){
      onDecisionChange(item.idx,item.in_weight,this.value);
    });
    container.appendChild(card);
    if(corrections[item.idx]){
      renderDecisionDetail(item.idx,item.in_weight,corrections[item.idx].Accepted,true);
    }
  }
}

// ── Decision handlers ─────────────────────────────────────────────────────────
function onDecisionChange(idx,inW,decision){
  renderDecisionDetail(idx,inW,decision,false);
}

function renderDecisionDetail(idx,inW,decision,restoring){
  const el=document.getElementById(`dec-detail-${idx}`);if(!el)return;

  if(decision==='NO'){
    corrections[idx]={Accepted:'NO','Out Weight':inW,'Net Weight':0};
    el.innerHTML='<span class="chip-ok">✓ Marked Rejected</span>';
    _ensureRRPanel(idx,inW,'NO',0,0);
    if(!restoring)scheduleValidation();

  }else if(decision==='YES'){
    const prev=(corrections[idx]&&corrections[idx]['Out Weight'])||0;
    // Only build the input DOM once — never replace it
    if(!document.getElementById(`outw-${idx}`)){
      el.innerHTML=`
        <div style="font-size:12px;color:#9BB4D4;font-weight:700;letter-spacing:.06em;text-transform:uppercase;margin-bottom:.4rem;">📥 Enter Net Weight → Out Weight</div>
        <div style="display:flex;align-items:center;gap:.5rem;margin-bottom:.3rem;">
          <label style="font-size:13px;color:#C8D8F8;font-weight:600;white-space:nowrap;min-width:90px;">Out Weight (T):</label>
          <input type="number" class="input-num" id="outw-${idx}"
            placeholder="e.g. 18.250" min="0" max="${inW}" step="0.001"
            value="${prev||''}" autocomplete="off"
            style="font-size:14px;padding:.6rem .9rem;border:1.5px solid rgba(99,102,241,.5);background:rgba(6,10,28,.9);color:#E8F0FF;border-radius:9px;flex:1;">
        </div>
        <div style="font-size:12px;color:#6B7EC8;margin-bottom:.2rem;">In Weight: <strong style="color:#10B981">${inW} T</strong></div>
        <div id="outw-s-${idx}" style="margin-top:.35rem;font-size:13px;font-weight:600;min-height:1.4rem;"></div>`;
      const inp=document.getElementById(`outw-${idx}`);
      // Use 'input' event only — fires on each keystroke, updates status label only
      inp.addEventListener('input',function(){
        _applyOutWValue(idx,inW,parseFloat(this.value)||0,false);
      });
      // Only schedule validation on blur (when user leaves the field)
      inp.addEventListener('blur',function(){
        const v=parseFloat(this.value)||0;
        if(v>0)scheduleValidation();
      });
      // Stop any key event from bubbling out of this input
      inp.addEventListener('keydown',function(e){e.stopPropagation();});
    }
    if(prev>0)_applyOutWValue(idx,inW,prev,true);
    else _ensureRRPanel(idx,inW,'YES',0,0);

  }else{
    delete corrections[idx];
    el.innerHTML='';
    const rp=document.getElementById(`rr-panel-${idx}`);if(rp)rp.innerHTML='';
  }
}

function onOutW(idx,inW,outW){_applyOutWValue(idx,inW,outW,false);}

function _applyOutWValue(idx,inW,outW,restoring){
  const statusEl=document.getElementById(`outw-s-${idx}`);if(!statusEl)return;
  if(outW<=0){
    statusEl.innerHTML='<span style="color:#F59E0B">Enter Out Weight to proceed</span>';
    delete corrections[idx];
    _updateRRPanelWeightLabel(idx,0);
    return;
  }
  const net=Math.round((inW-outW)*1000)/1000;
  corrections[idx]={Accepted:'YES','Out Weight':outW,'Net Weight':net};
  statusEl.innerHTML=`<span class="chip-ok">✓ Out: ${outW}T, Net: ${net}T</span>`;
  // Show RR panel & update weight label — never destroys the serial input
  _ensureRRPanel(idx,inW,'YES',outW,net);
  _updateRRPanelWeightLabel(idx,outW);
  // NOTE: do NOT call scheduleValidation() here — the blur handler does that
  // so typing never triggers a re-render that kills focus
}

// ── RR Panel ──────────────────────────────────────────────────────────────────
function _ensureRRPanel(idx,inW,decision,outW,net){
  const panel=document.getElementById(`rr-panel-${idx}`);
  if(!panel)return;

  // Already rendered — just update weight label, preserve all inputs
  if(panel.querySelector('.rr-panel')){
    _updateRRPanelWeightLabel(idx,outW);
    return;
  }

  const savedSerial=_rrSerials[idx]||'';

  // All three options always visible; default pre-selection based on decision
  const reasonOptions=`
    <option value="Accepted Towing Vehicle" ${decision==='YES'?'selected':''}>Option A — Accepted Towing Vehicle</option>
    <option value="Rejected Towing Vehicle" ${decision==='NO'?'selected':''}>Option B — Rejected Towing Vehicle</option>
    <option value="Late Time / Breakdown">Option C — Late Time / Breakdown</option>`;

  const wrapper=document.createElement('div');
  wrapper.className='rr-panel';
  wrapper.innerHTML=`
    <div class="rr-panel-title">📄 Rectification Report — ${decision==='YES'?'✅ Acceptance Format':'❌ Rejection Format'}</div>

    <div style="margin:.6rem 0 .25rem;font-size:11px;color:#9BB4D4;font-weight:700;letter-spacing:.07em;text-transform:uppercase;">📋 Reason / Narrative Type</div>
    <select class="rr-select" id="rr-reason-${idx}"
      style="font-size:13px;padding:.55rem .9rem;margin-bottom:.6rem;border:1.5px solid rgba(99,102,241,.5);background:rgba(6,10,28,.9);color:#E8F0FF;border-radius:9px;width:100%;">
      ${reasonOptions}
    </select>

    <div style="margin:.3rem 0 .25rem;font-size:11px;color:#9BB4D4;font-weight:700;letter-spacing:.07em;text-transform:uppercase;">🔢 Serial Number</div>
    <input class="rr-input" id="rr-serial-${idx}" type="text"
      placeholder="Enter RR Serial No. (e.g. 0290)" maxlength="10"
      value="${esc(savedSerial)}" autocomplete="off" spellcheck="false"
      style="font-size:14px;padding:.6rem .9rem;border:1.5px solid rgba(99,102,241,.5);background:rgba(6,10,28,.9);color:#E8F0FF;border-radius:9px;width:100%;margin-bottom:.6rem;">



    <button class="rr-btn" id="rr-btn-${idx}" type="button"
      style="width:100%;padding:.7rem;margin-top:.4rem;font-size:13px;font-weight:700;letter-spacing:.02em;border-radius:9px;">
      ⬇ Download Rectification Report
    </button>
    <div class="rr-status" id="rr-status-${idx}" style="font-size:13px;margin-top:.45rem;min-height:1.4rem;"></div>`;

  panel.appendChild(wrapper);

  // Attach events AFTER DOM insertion
  const serialEl=panel.querySelector(`#rr-serial-${idx}`);
  if(serialEl){
    serialEl.addEventListener('input',function(){_rrSerials[idx]=this.value;});
    serialEl.addEventListener('keydown',function(e){e.stopPropagation();});
    serialEl.addEventListener('click',function(e){e.stopPropagation();});
    serialEl.addEventListener('mousedown',function(e){e.stopPropagation();});
  }
  const reasonDropEl=panel.querySelector(`#rr-reason-${idx}`);
  if(reasonDropEl){
    reasonDropEl.addEventListener('keydown',function(e){e.stopPropagation();});
    reasonDropEl.addEventListener('click',function(e){e.stopPropagation();});
  }
  const btn=panel.querySelector(`#rr-btn-${idx}`);
  if(btn)btn.addEventListener('click',function(e){e.stopPropagation();generateRR(idx);});
}
function _updateRRPanelWeightLabel(idx,outW){
  // Weight reference UI removed — computed server-side from reason
}

// ── Generate Rectification Report ─────────────────────────────────────────────
async function generateRR(idx){
  const item=_incompleteItems[idx]||{};
  const token=item.token||'';
  const corr=corrections[idx]||{};
  const decision=corr.Accepted||'';
  const outW=corr['Out Weight']||0;
  const net=corr['Net Weight']||0;
  const serialEl=document.getElementById(`rr-serial-${idx}`);
  const serial=(serialEl?serialEl.value:'').trim();
  // Read reason dropdown (weight_label computed server-side)  // Read reason dropdown
  const reasonEl=document.getElementById(`rr-reason-${idx}`);
  const reason=reasonEl?reasonEl.value:(decision==='YES'?'Accepted Towing Vehicle':'Rejected Towing Vehicle');

  const statusEl=document.getElementById(`rr-status-${idx}`);
  const btn=document.getElementById(`rr-btn-${idx}`);

  if(!serial){if(statusEl)statusEl.innerHTML='<span style="color:#F59E0B">⚠ Enter RR Serial Number first</span>';return;}
  if(!decision){if(statusEl)statusEl.innerHTML='<span style="color:#F59E0B">⚠ Select YES or NO first</span>';return;}
  if(!onlineTmpPath){if(statusEl)statusEl.innerHTML='<span style="color:#F87171">❌ Online file not loaded</span>';return;}

  if(btn){btn.disabled=true;btn.textContent='⏳ Generating…';}
  if(statusEl)statusEl.innerHTML='<span style="color:#818CF8">⏳ Building report…</span>';
  try{
    const res=await fetch('/api/daily/generate_rr',{
      method:'POST',headers:{'Content-Type':'application/json'},
      body:JSON.stringify({
        tmp_path:onlineTmpPath,token,rr_serial:serial,accepted:decision,
        out_weight:outW,net_weight:net,
        reason:reason,filter_date:onlineDate
      })
    });
    if(!res.ok){
      const e=await res.json().catch(()=>({error:`HTTP ${res.status}`}));
      if(statusEl)statusEl.innerHTML=`<span style="color:#F87171">❌ ${esc(e.error||'Failed')}</span>`;
      return;
    }
    const blob=await res.blob();
    const cd=res.headers.get('Content-Disposition')||'';
    // Try quoted filename first, then unquoted
    const fnMatch=cd.match(/filename="([^"]+)"/) || cd.match(/filename=([^;\s]+)/);
    const fname=fnMatch?decodeURIComponent(fnMatch[1]):`RR-B-${serial}-${token}.docx`;
    const url=URL.createObjectURL(blob);
    const a=document.createElement('a');a.href=url;a.download=fname;a.click();
    URL.revokeObjectURL(url);
    if(statusEl)statusEl.innerHTML=`<span style="color:#4ADE80">✅ Downloaded: ${esc(fname)}</span>`;
  }catch(ex){
    if(statusEl)statusEl.innerHTML=`<span style="color:#F87171">❌ ${esc(ex.message)}</span>`;
  }finally{
    if(btn){btn.disabled=false;btn.textContent='⬇ Download Rectification Report';}
  }
}
// ── Stats / tracker ───────────────────────────────────────────────────────────
function renderOnlineStats(s){
  setEl('onlineStats',`<div class="stats-box"><div class="stats-title stats-title-online">Online Data — Summary</div>
    <div class="stats-inner"><table>
      <tr><td>Accepted Loads</td><td style="color:#818CF8">${s.accepted}</td></tr>
      <tr><td>Rejected Loads</td><td style="color:#F87171">${s.rejected}</td></tr>
      <tr><td>Total Weight In (T)</td><td style="color:#10B981">${s.wi.toFixed(2)}</td></tr>
      <tr><td>Total Weight Out (T)</td><td style="color:#10B981">${s.wo.toFixed(2)}</td></tr>
      <tr><td>Total Net Weight (T)</td><td style="color:#10B981;font-weight:800">${s.nw.toFixed(2)}</td></tr>
    </table></div></div>`);show('onlineStats');
}

function updateStepTracker(){
  const s=(ok,l)=>`<span class="step-pip ${ok?'step-done':'step-wait'}">${ok?'✔':'○'} ${l}</span>`;
  setEl('stepTracker',`${s(!!onlineStats,'Online Validated')}<span class="step-arr">→</span>${s(!!wbStats,'WB Compared')}<span class="step-arr">→</span>${s(!!onlineStats&&!!wbStats,'Ready to Download')}`);
}
function updateGenerateBtn(){
  const btn=document.getElementById('generateBtn');
  if(btn)btn.disabled=!(onlineTmpPath&&onlineDate&&onlineStats);
}

// ── WB Side ───────────────────────────────────────────────────────────────────
async function onWbFileChange(){
  const f=document.getElementById('wbFile').files[0];if(!f)return;
  setEl('wbFileStatus','<span class="chip-info">⏳ Uploading…</span>');
  const fd=new FormData();fd.append('file',f);
  try{
    const res=await fetch('/api/daily/wb_upload',{method:'POST',body:fd});
    const d=await res.json();
    if(d.error){setEl('wbFileStatus',`<span class="chip-wait">❌ ${esc(d.error)}</span>`);return;}
    wbTmpPath=d.tmp_path;wbDecisions={};
    setEl('wbFileStatus',`<span class="chip-ok">✓ ${esc(f.name)} — ${d.rows.toLocaleString()} rows loaded</span>`);
    const di=document.getElementById('wbDate');
    di.min=d.min_date;di.max=d.max_date;di.value=d.max_date;wbDate=d.max_date;
    show('wbDateSection');show('wbResultSection');await runWbPivot();
  }catch(e){setEl('wbFileStatus',`<span class="chip-wait">❌ ${e.message}</span>`);}
}

async function onWbDateChange(){wbDate=document.getElementById('wbDate').value;wbDecisions={};await runWbPivot();}

function scheduleWbPivot(){
  clearTimeout(_wbPivotTimer);
  _wbPivotTimer=setTimeout(function(){
    if(_userIsTyping()){scheduleWbPivot();return;}
    runWbPivot();
  },900);
}

async function runWbPivot(){
  if(!wbTmpPath||!wbDate)return;
  setEl('wbStats','<div class="chip-info">⏳ Processing WB data…</div>');
  try{
    const res=await fetch('/api/daily/wb_pivot',{method:'POST',headers:{'Content-Type':'application/json'},
      body:JSON.stringify({
        wb_tmp_path:wbTmpPath,filter_date:wbDate,wb_decisions:wbDecisions,
        online_tmp_path:onlineTmpPath,online_date:onlineDate,corrections
      })});
    const d=await res.json();
    if(d.error){setEl('wbStats',`<div class="val-err">❌ ${esc(d.error)}</div>`);return;}
    if(d.incomplete&&d.incomplete.length>0){renderWbActionCards(d.incomplete);return;}
    wbStats=d.wb_stats;pivotTmpPath=d.pivot_tmp;pivotFname=d.pivot_fname;
    if(d.errors&&d.errors.length>0){
      setEl('wbStats',d.errors.map(e=>`<div class="val-err">❌ ${esc(e)}</div>`).join('')+renderWbStatsHtml(d.wb_stats));
    }else{
      setEl('wbStats','<div class="val-ok">✅ All validations completed successfully</div>'+renderWbStatsHtml(d.wb_stats));
    }
    renderComparisonTable();renderPivot(d.rows);
    show('pivotDlBtn');updateStepTracker();
  }catch(e){setEl('wbStats',`<div class="val-err">❌ ${e.message}</div>`);}
}

function renderWbStatsHtml(s){
  if(!s||typeof s.accepted==='undefined')return '';
  return `<div class="stats-box stats-box-wb"><div class="stats-title stats-title-wb">Weighbridge Data — Summary</div>
    <div class="stats-inner"><table>
      <tr><td>Accepted Loads</td><td style="color:#34D399">${s.accepted}</td></tr>
      <tr><td>Rejected Loads</td><td style="color:#F87171">${s.rejected}</td></tr>
      <tr><td>Total Weight In (T)</td><td style="color:#10B981">${(s.wi||0).toFixed(2)}</td></tr>
      <tr><td>Total Weight Out (T)</td><td style="color:#10B981">${(s.wo||0).toFixed(2)}</td></tr>
      <tr><td>Total Net Weight (T)</td><td style="color:#10B981;font-weight:800">${(s.nw||0).toFixed(2)}</td></tr>
    </table></div></div>`;
}

function renderWbActionCards(incomplete){
  if(!incomplete||!incomplete.length)return;
  let html='<div class="val-warn">⚠️ Action Required — WB rows with missing Date Out / Time Out</div>';
  for(const item of incomplete){
    const dec=wbDecisions[item.idx]||{};
    html+=`<div class="action-card" id="wb-acard-${item.idx}">
      <div class="action-title">⚠️ E-Token [${esc(item.etoken)}] – Missing Date Out / Time Out</div>
      <select class="action-select" id="wb-dec-${item.idx}" onchange="onWbDecisionChange(${item.idx},${item.in_weight})">
        <option value="" ${!dec.decision&&dec.decision!==0?'selected':''}>— Select —</option>
        <option value="0" ${dec.decision===0?'selected':''}>0 — Rejected (remove row)</option>
        <option value="1" ${dec.decision===1?'selected':''}>1 — Accepted (enter Out Weight)</option>
      </select>
      <div id="wb-dec-detail-${item.idx}"></div>
    </div>`;
  }
  setEl('wbStats',html);
  for(const item of incomplete){if(wbDecisions[item.idx])renderWbDecisionDetail(item.idx,item.in_weight,wbDecisions[item.idx].decision);}
}

function onWbDecisionChange(idx,inW){
  const v=document.getElementById(`wb-dec-${idx}`).value;
  renderWbDecisionDetail(idx,inW,v===''?null:parseInt(v));
}

function renderWbDecisionDetail(idx,inW,decision){
  const el=document.getElementById(`wb-dec-detail-${idx}`);if(!el)return;
  if(decision===0||decision==='0'){
    wbDecisions[idx]={decision:0};
    el.innerHTML='<span class="chip-ok">✓ Marked Rejected — row will be removed.</span>';
    scheduleWbPivot();
  }else if(decision===1||decision==='1'){
    if(!document.getElementById(`wb-outw-${idx}`)){
      const prev=(wbDecisions[idx]&&wbDecisions[idx].out_weight)||0;
      el.innerHTML=`<input type="number" class="input-num" id="wb-outw-${idx}"
        placeholder="Out Weight (T) [In Weight=${inW}T]" min="0" max="${inW||9999}"
        step="0.001" value="${prev||''}" autocomplete="off">
        <div id="wb-outw-s-${idx}" style="margin-top:.3rem;font-size:.78rem"></div>`;
      const inp=document.getElementById(`wb-outw-${idx}`);
      inp.addEventListener('input',function(){onWbOutW(idx,inW,parseFloat(this.value)||0);});
      inp.addEventListener('blur',function(){
        const v=parseFloat(this.value)||0;
        if(v>0)scheduleWbPivot();
      });
      inp.addEventListener('keydown',function(e){e.stopPropagation();});
      if(prev>0)_applyWbOutWValue(idx,inW,prev,true);
    }
  }else{delete wbDecisions[idx];el.innerHTML='';}
}

function onWbOutW(idx,inW,outW){_applyWbOutWValue(idx,inW,outW,false);}
function _applyWbOutWValue(idx,inW,outW,restoring){
  const el=document.getElementById(`wb-outw-s-${idx}`);if(!el)return;
  if(outW<=0){el.innerHTML='<span style="color:#F59E0B">Enter Out Weight to proceed</span>';delete wbDecisions[idx];return;}
  const net=Math.round((inW-outW)*1000)/1000;
  wbDecisions[idx]={decision:1,out_weight:outW};
  el.innerHTML=`<span class="chip-ok">✓ Out: ${outW}T, Net: ${net}T</span>`;
  // validation only on blur, not here
}

function renderWbStats(s){setEl('wbStats',renderWbStatsHtml(s));}

function renderComparisonTable(){
  if(!onlineStats||!wbStats)return;
  const on=onlineStats,wb=wbStats;
  const m=(a,b,wt)=>Math.abs(a-b)<(wt?.01:.5)?`<td class="cmp-match">✔ Validated</td>`:`<td class="cmp-miss">✘ Mismatch</td>`;
  setEl('comparisonTable',`<table class="cmp-table">
    <thead><tr><th style="text-align:left">Metric</th><th>Online</th><th>Weighbridge</th><th>Remark</th></tr></thead>
    <tbody>
      <tr><td style="text-align:left;color:#C8D8F8">Accepted Loads</td><td style="color:#818CF8;font-weight:700">${on.accepted}</td><td style="color:#34D399;font-weight:700">${wb.accepted}</td>${m(on.accepted,wb.accepted,false)}</tr>
      <tr><td style="text-align:left;color:#C8D8F8">Rejected Loads</td><td style="color:#818CF8;font-weight:700">${on.rejected}</td><td style="color:#34D399;font-weight:700">${wb.rejected}</td>${m(on.rejected,wb.rejected,false)}</tr>
      <tr><td style="text-align:left;color:#C8D8F8">Weight In (T)</td><td style="color:#818CF8;font-weight:700">${on.wi.toFixed(2)}</td><td style="color:#34D399;font-weight:700">${wb.wi.toFixed(2)}</td>${m(on.wi,wb.wi,true)}</tr>
      <tr><td style="text-align:left;color:#C8D8F8">Weight Out (T)</td><td style="color:#818CF8;font-weight:700">${on.wo.toFixed(2)}</td><td style="color:#34D399;font-weight:700">${wb.wo.toFixed(2)}</td>${m(on.wo,wb.wo,true)}</tr>
      <tr><td style="text-align:left;color:#C8D8F8">Net Weight (T)</td><td style="color:#818CF8;font-weight:700;font-size:.95rem">${on.nw.toFixed(2)}</td><td style="color:#34D399;font-weight:700;font-size:.95rem">${wb.nw.toFixed(2)}</td>${m(on.nw,wb.nw,true)}</tr>
    </tbody></table>`);show('comparisonTable');
}

function renderPivot(rows){
  if(!rows||!rows.length){clearEl('pivotSection');return;}
  let html=`<div class="pivot-outer"><div class="pivot-title">📊 WB Pivot — Net Weight by Material &amp; Site</div><table class="pivot-tbl"><thead><tr>`;
  html+=`<th>LABEL</th><th>LOADS</th><th>WEIGHT IN (T)</th><th>WEIGHT OUT (T)</th><th>NET WEIGHT (T)</th>`;
  html+=`</tr></thead><tbody>`;
  rows.forEach(r=>{
    const isMat  = r['type']==='mat_header';
    const isGrand= r['label']==='GRAND TOTAL';
    const cls    = isGrand ? 'pivot-grand' : isMat ? 'pivot-mat' : '';
    html+=`<tr class="${cls}">`;
    if(isMat){
      // Material header row: span all data columns, show only label
      html+=`<td colspan="5" style="text-align:left;font-weight:800;letter-spacing:.04em;">${esc(r['label'])}</td>`;
    } else {
      html+=`<td>${esc(r['label'])}</td>`;
      html+=`<td>${r['loads']||0}</td>`;
      html+=`<td>${(+(r['wi']||0)).toFixed(2)}</td>`;
      html+=`<td>${(+(r['wo']||0)).toFixed(2)}</td>`;
      html+=`<td>${(+(r['nw']||0)).toFixed(2)}</td>`;
    }
    html+=`</tr>`;
  });
  html+=`</tbody></table></div>`;
  setEl('pivotSection',html);
}

async function generateReport(){
  const btn=document.getElementById('generateBtn');
  if(btn){btn.disabled=true;btn.textContent='⏳ Generating…';}
  setEl('generateStatus','<span style="color:#818CF8">⏳ Generating report…</span>');
  try{
    const res=await fetch('/api/daily/generate',{method:'POST',headers:{'Content-Type':'application/json'},
      body:JSON.stringify({tmp_path:onlineTmpPath,filter_date:onlineDate,corrections})});
    if(!res.ok){const e=await res.json().catch(()=>({error:`HTTP ${res.status}`}));setEl('generateStatus',`<span style="color:#F87171">❌ ${esc(e.error||'Failed')}</span>`);return;}
    const blob=await res.blob();
    const cd=res.headers.get('Content-Disposition')||'';
    const fnMatch=cd.match(/filename="([^"]+)"/);
    reportName=fnMatch?fnMatch[1]:'daily_report.xlsx';
    reportBlob=blob;
    setEl('generateStatus',`<span style="color:#4ADE80">✅ Ready: ${esc(reportName)}</span>`);
    show('downloadBtn');document.getElementById('downloadBtn').disabled=false;
  }catch(e){setEl('generateStatus',`<span style="color:#F87171">❌ ${e.message}</span>`);}
  finally{if(btn){btn.disabled=false;btn.textContent='⚡ Generate Report';}}
}

function downloadReport(){
  if(!reportBlob)return;
  const url=URL.createObjectURL(reportBlob);
  const a=document.createElement('a');a.href=url;a.download=reportName;a.click();
  URL.revokeObjectURL(url);
}

function downloadPivot(){if(!pivotTmpPath)return;window.location.href=`/api/daily/wb_pivot_download?path=${encodeURIComponent(pivotTmpPath)}&fname=${encodeURIComponent(pivotFname)}`;}

function setEl(id,html){const el=document.getElementById(id);if(el)el.innerHTML=html;}
function clearEl(id){const el=document.getElementById(id);if(el)el.innerHTML='';}
function show(id){const el=document.getElementById(id);if(el)el.style.display='block';}
function hide(id){const el=document.getElementById(id);if(el)el.style.display='none';}
function esc(s){return String(s||'').replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;').replace(/"/g,'&quot;');}
</script>
<div class="apsg-footer">✦ Internal Reporting Platform — APSG Staging Ground &nbsp;·&nbsp; Developed by Karthik</div>
</body>
</html>"""


ADMIN_HTML = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>Admin Panel — APSG Report</title>
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&display=swap" rel="stylesheet">
<style>

/* ═══ MODERN UI BASE (Blue/Purple/Cyan Theme) ═══ */
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&family=Poppins:wght@600;700;800&display=swap');
*, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }
:root {
  --bg:transparent; --card-bg: rgba(13,18,35,0.88);
  --indigo: #6366F1; --indigo-l: #818CF8; --cyan: #22D3EE;
  --purple: #A855F7; --green: #10B981; --amber: #F59E0B; --red: #F87171;
  --text: #E8EEF8; --muted: #64748B; --border: rgba(99,102,241,0.15);
}
body {
  font-family: 'Inter', system-ui, -apple-system, sans-serif;
  min-height: 100vh; font-size: 14px; line-height: 1.6;
  background:transparent; color: var(--text);
  position: relative;
}
body::before {
  content: ''; position: fixed; inset: 0; z-index: 0; pointer-events: none;
  background:
    radial-gradient(ellipse 70% 55% at 5% 10%, rgba(99,102,241,.09) 0%, transparent 65%),
    radial-gradient(ellipse 55% 45% at 95% 90%, rgba(168,85,247,.07) 0%, transparent 65%);
}
body::after {
  content: ''; position: fixed; inset: 0; z-index: 0; pointer-events: none; opacity: .016;
  background-image: linear-gradient(var(--indigo) 1px, transparent 1px),
    linear-gradient(90deg, var(--indigo) 1px, transparent 1px);
  background-size: 60px 60px;
}

/* ── Modern Top Bar ── */
.top-bar {
  position: sticky; top: 0; z-index: 200; height: 56px;
  display: flex; align-items: center; padding: 0 1.5rem; gap: .75rem;
  background: rgba(6,9,22,.96); backdrop-filter: blur(20px); -webkit-backdrop-filter: blur(20px);
  border-bottom: 1px solid var(--border);
}
.top-mini-brand {
  font-family: 'Poppins', sans-serif;
  font-size: .82rem; font-weight: 700; color: var(--text);
  letter-spacing: -.01em; white-space: nowrap; flex-shrink: 0;
  display: flex; align-items: center; gap: .45rem;
}
.top-mini-brand::before {
  content: ''; width: 8px; height: 8px; border-radius: 50%;
  background: linear-gradient(135deg, var(--indigo), var(--cyan)); flex-shrink: 0;
}
.top-sep { width: 1px; height: 18px; background: var(--border); flex-shrink: 0; }
.top-page-label { font-size: .75rem; font-weight: 600; color: var(--muted); white-space: nowrap; }
.top-brand-tag { font-size: .7rem; color: var(--muted); }
.top-spacer { flex: 1; }
.back-btn {
  background: rgba(99,102,241,.08); border: 1px solid rgba(99,102,241,.18);
  border-radius: 8px; padding: .3rem .9rem; font-size: .7rem; font-weight: 600;
  color: var(--indigo-l); text-decoration: none; transition: all .2s; white-space: nowrap;
}
.back-btn:hover { background: rgba(99,102,241,.18); transform: translateX(-2px); }
/* ── Cards & Containers ── */
.container, .page-content { position: relative; z-index: 1; }
.card, .section-card, .panel {
  background: var(--card-bg); border: 1px solid var(--border);
  border-radius: 16px; backdrop-filter: blur(16px); -webkit-backdrop-filter: blur(16px);
  box-shadow: 0 8px 32px rgba(0,0,0,.4), inset 0 1px 0 rgba(255,255,255,.03);
  transition: box-shadow .25s, border-color .25s;
}
.card:hover, .section-card:hover { border-color: rgba(99,102,241,.25); }

/* ── Upload Zone — Modern drag & drop ── */
.upload-zone, .dz {
  border: 2px dashed rgba(99,102,241,.3); border-radius: 14px;
  padding: 2rem; text-align: center; cursor: pointer;
  transition: all .22s; background: rgba(99,102,241,.03);
  position: relative;
}
.upload-zone:hover, .dz:hover, .upload-zone.drag-over, .dz.drag-over {
  border-color: var(--indigo); background: rgba(99,102,241,.08);
  box-shadow: 0 0 0 4px rgba(99,102,241,.12);
}
.upload-zone.ok, .dz.ok {
  border-color: var(--green); background: rgba(16,185,129,.06);
  border-style: solid;
}
.upload-zone.ok:hover, .dz.ok:hover {
  border-color: var(--green); background: rgba(16,185,129,.1);
  box-shadow: 0 0 0 4px rgba(16,185,129,.1);
}
.upload-icon { font-size: 2rem; margin-bottom: .5rem; display: block; }
.upload-label { font-size: .82rem; color: var(--muted); font-weight: 500; }
.upload-hint { font-size: .7rem; color: rgba(100,116,139,.6); margin-top: .25rem; }
.upload-filename { font-size: .78rem; color: var(--green); font-weight: 600; margin-top: .4rem; }

/* ── Modern Buttons ── */
.btn-primary, .btn-generate, .btn-teal, .modal-btn-primary {
  background:transparent;
  color: #fff; border: none; border-radius: 10px;
  padding: .7rem 1.4rem; font-size: .85rem; font-weight: 700;
  font-family: 'Inter', sans-serif; cursor: pointer; letter-spacing: .01em;
  box-shadow: 0 4px 18px rgba(99,102,241,.35);
  position: relative; overflow: hidden;
  transition: transform .2s, box-shadow .2s;
}
.btn-primary::before, .btn-generate::before, .btn-teal::before, .modal-btn-primary::before {
  content: ''; position: absolute; inset: 0;
  background: linear-gradient(135deg, transparent, rgba(255,255,255,.12), transparent);
  transform: translateX(-100%); transition: transform .45s;
}
.btn-primary:hover, .btn-generate:hover, .modal-btn-primary:hover {
  transform: translateY(-2px) scale(1.01);
  box-shadow: 0 8px 28px rgba(99,102,241,.55);
}
.btn-primary:hover::before, .btn-generate:hover::before, .modal-btn-primary:hover::before {
  transform: translateX(100%);
}
.btn-primary:active, .btn-generate:active { transform: translateY(0) scale(.98); }
.btn-primary:disabled, .btn-generate:disabled { opacity: .45; cursor: not-allowed; transform: none; }

.btn-teal {
  background:transparent;
  box-shadow: 0 4px 18px rgba(16,185,129,.3);
}
.btn-teal:hover { box-shadow: 0 8px 28px rgba(16,185,129,.5); }

.btn-download, .btn-dl, .dl-btn {
  background: linear-gradient(135deg, #0D7A5F, #10B981);
  color: #fff; border: none; border-radius: 10px;
  padding: .65rem 1.3rem; font-size: .82rem; font-weight: 700;
  cursor: pointer; transition: all .22s;
  box-shadow: 0 4px 16px rgba(16,185,129,.3);
}
.btn-download:hover, .btn-dl:hover, .dl-btn:hover {
  transform: translateY(-2px); box-shadow: 0 8px 26px rgba(16,185,129,.5);
}

.btn-secondary, .btn-gray {
  background: rgba(30,41,86,.6); color: var(--muted);
  border: 1px solid rgba(99,102,241,.2); border-radius: 10px;
  padding: .65rem 1.2rem; font-size: .82rem; font-weight: 600;
  cursor: pointer; transition: all .2s;
}
.btn-secondary:hover, .btn-gray:hover { background: rgba(99,102,241,.1); color: var(--indigo-l); }

.btn-del, .btn-danger {
  background: rgba(239,68,68,.08); color: var(--red);
  border: 1px solid rgba(239,68,68,.2); border-radius: 10px;
  padding: .6rem 1.2rem; font-size: .8rem; font-weight: 600;
  cursor: pointer; transition: all .2s;
}
.btn-del:hover, .btn-danger:hover { background: rgba(239,68,68,.16); }
.btn-del:hover { animation: shake .3s ease; }

@keyframes shake {
  0%,100% { transform: translateX(0); }
  25% { transform: translateX(-3px); }
  75% { transform: translateX(3px); }
}

/* ── Form Inputs ── */
input[type=text], input[type=number], input[type=date], input[type=password],
select, textarea {
  background: rgba(6,10,24,.8); border: 1.5px solid rgba(30,41,86,.8);
  border-radius: 10px; color: var(--text); padding: .65rem .9rem;
  font-size: .85rem; font-family: 'Inter', sans-serif;
  transition: all .2s; width: 100%;
}
input:focus, select:focus, textarea:focus {
  outline: none; border-color: var(--indigo);
  box-shadow: 0 0 0 3px rgba(99,102,241,.16);
  background: rgba(10,14,32,.9);
}
input::placeholder, textarea::placeholder { color: rgba(100,116,139,.45); font-weight: 300; }

/* ── Alerts / Status ── */
.alert-success, .alert.success, .msg.success {
  background: rgba(16,185,129,.08); border: 1px solid rgba(16,185,129,.2);
  color: #34D399; border-radius: 10px; padding: .65rem .9rem;
  font-size: .8rem; font-weight: 500;
}
.alert-error, .alert.error, .msg.error {
  background: rgba(239,68,68,.08); border: 1px solid rgba(239,68,68,.2);
  color: var(--red); border-radius: 10px; padding: .65rem .9rem;
  font-size: .8rem; font-weight: 500;
}
.alert-warn, .alert.warn { 
  background: rgba(245,158,11,.08); border: 1px solid rgba(245,158,11,.2);
  color: var(--amber); border-radius: 10px; padding: .65rem .9rem;
  font-size: .8rem; font-weight: 500;
}

/* ── Tables ── */
table { width: 100%; border-collapse: collapse; font-size: .8rem; }
thead th {
  background: rgba(99,102,241,.08); color: var(--muted);
  font-weight: 700; font-size: .68rem; letter-spacing: .06em;
  text-transform: uppercase; padding: .65rem .9rem; text-align: left;
  border-bottom: 1px solid var(--border);
}
tbody tr { border-bottom: 1px solid rgba(30,45,80,.3); transition: background .15s; }
tbody tr:hover { background: rgba(99,102,241,.04); }
tbody td { padding: .6rem .9rem; color: var(--text); }

/* ── Tabs ── */
.tab { position: relative; transition: all .2s; }
.tab.active { color: var(--indigo-l) !important; }
.tab.active::after {
  content: ''; position: absolute; bottom: -1px; left: 15%; right: 15%;
  height: 2px; border-radius: 2px; background: var(--indigo);
  animation: tabIn .2s ease;
}
@keyframes tabIn { from { left: 50%; right: 50%; } to { left: 15%; right: 15%; } }

/* ── Spinner ── */
@keyframes spin { to { transform: rotate(360deg); } }
.spinner { width: 20px; height: 20px; border: 2px solid rgba(99,102,241,.2);
  border-top-color: var(--indigo); border-radius: 50%; animation: spin 1s linear infinite; }



@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&display=swap');
*, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }
body {
  font-family: 'Inter', system-ui, -apple-system, sans-serif;
  min-height: 100vh; font-size: 14px; line-height: 1.6;
  background:transparent;
  color: #E8EEF8;
}


/* ══ ANIMATED BUTTONS ══════════════════════════════════════════════════════════ */
/* Primary action button — pulse + shimmer */
.btn-primary, .btn-generate, .btn-teal, .modal-btn-primary {
  position: relative; overflow: hidden;
  transition: transform .2s, box-shadow .2s !important;
}
.btn-primary::after, .btn-generate::after, .btn-teal::after, .modal-btn-primary::after {
  content: ''; position: absolute; top: -50%; left: -75%;
  width: 50%; height: 200%; background: rgba(255,255,255,.15);
  transform: skewX(-20deg); transition: left .5s ease;
  pointer-events: none;
}
.btn-primary:hover::after, .btn-generate:hover::after,
.btn-teal:hover::after, .modal-btn-primary:hover::after {
  left: 150%;
}
.btn-primary:hover, .btn-generate:hover, .modal-btn-primary:hover {
  transform: translateY(-3px) scale(1.02) !important;
  box-shadow: 0 8px 28px rgba(99,102,241,.55) !important;
}
.btn-primary:active, .btn-generate:active, .modal-btn-primary:active {
  transform: translateY(0) scale(.98) !important;
}

/* Download button — bounce */
.btn-dl, .btn-download, .dl-btn {
  animation: gentlePulse 3s ease-in-out infinite;
  transition: transform .2s, box-shadow .2s !important;
}
.btn-dl:hover, .btn-download:hover, .dl-btn:hover {
  animation: none;
  transform: translateY(-3px) scale(1.02) !important;
  box-shadow: 0 8px 28px rgba(16,185,129,.5) !important;
}
@keyframes gentlePulse {
  0%,100% { box-shadow: 0 4px 16px rgba(16,185,129,.3); }
  50%      { box-shadow: 0 4px 24px rgba(16,185,129,.55); }
}

/* Danger/delete button — shake on hover */
.btn-del:hover {
  animation: shake .35s ease;
}
@keyframes shake {
  0%,100% { transform: translateX(0); }
  20%     { transform: translateX(-3px); }
  40%     { transform: translateX(3px); }
  60%     { transform: translateX(-2px); }
  80%     { transform: translateX(2px); }
}

/* Secondary / back buttons — slide arrow */
.back-btn, .btn-secondary, .btn-gray, .btn-refresh {
  transition: all .2s !important;
}
.back-btn:hover {
  padding-left: .55rem !important;
  letter-spacing: .02em;
}

/* Upload zone — glow pulse when empty */
.upload-zone:not(.ok) {
  animation: uploadGlow 4s ease-in-out infinite;
}
.upload-zone:hover, .dz:hover {
  animation: none;
}
@keyframes uploadGlow {
  0%,100% { border-color: rgba(99,102,241,.3); }
  50%      { border-color: rgba(99,102,241,.65); box-shadow: 0 0 18px rgba(99,102,241,.15); }
}

/* App cards on dashboard — float in on load */
.app-card {
  animation: cardFadeIn .5s ease both;
}
.app-card:nth-child(1) { animation-delay: .05s; }
.app-card:nth-child(2) { animation-delay: .10s; }
.app-card:nth-child(3) { animation-delay: .15s; }
.app-card:nth-child(4) { animation-delay: .20s; }
.app-card:nth-child(5) { animation-delay: .25s; }
@keyframes cardFadeIn {
  from { opacity: 0; transform: translateY(18px); }
  to   { opacity: 1; transform: translateY(0); }
}

/* Stat cards pop-in */
.stat-card, .stat {
  animation: popIn .4s cubic-bezier(.34,1.56,.64,1) both;
}
.stat-card:nth-child(1), .stat:nth-child(1) { animation-delay: .05s; }
.stat-card:nth-child(2), .stat:nth-child(2) { animation-delay: .10s; }
.stat-card:nth-child(3), .stat:nth-child(3) { animation-delay: .15s; }
.stat-card:nth-child(4), .stat:nth-child(4) { animation-delay: .20s; }
@keyframes popIn {
  from { opacity: 0; transform: scale(.88); }
  to   { opacity: 1; transform: scale(1); }
}

/* Login button — ripple */
.btn.btn-primary, button.btn-primary {
  transition: transform .15s, box-shadow .15s !important;
}

/* Tab buttons — active indicator slide */
.tab {
  transition: all .2s !important;
  position: relative;
}
.tab.active::after {
  content: '';
  position: absolute; bottom: -1px; left: 10%; right: 10%;
  height: 2px; border-radius: 2px;
  background: #6366F1;
  animation: tabSlide .25s ease;
}
@keyframes tabSlide {
  from { left: 50%; right: 50%; }
  to   { left: 10%; right: 10%; }
}

/* Generate/action buttons in forms — shimmer idle */
.btn-generate:not(:disabled), .btn-teal:not(:disabled) {
  background-size: 200% 100% !important;
  transition: background-position .4s, transform .2s, box-shadow .2s !important;
}

/* ── Compact top bar (40px) for all internal pages ── */
.top-bar {
  position: sticky; top: 0; z-index: 200;
  height: 40px; display: flex; align-items: center;
  padding: 0 1.2rem; gap: .75rem;
  background: rgba(5,8,20,.96); backdrop-filter: blur(16px);
  border-bottom: 1px solid rgba(99,102,241,.12);
}
.top-mini-brand {
  font-size: .75rem; font-weight: 800; color: #6366F1;
  letter-spacing: .04em; white-space: nowrap; flex-shrink: 0;
}
.top-brand-text {
  font-size: .78rem; font-weight: 800; color: #818CF8;
  letter-spacing: .02em; white-space: nowrap; flex-shrink: 0;
}
.top-sep { width: 1px; height: 18px; background: rgba(99,102,241,.18); flex-shrink: 0; }
.top-page-label {
  font-size: .78rem; font-weight: 600; color: #94A3B8;
  white-space: nowrap; overflow: hidden; text-overflow: ellipsis;
}
.top-spacer { flex: 1; }
.back-btn {
  background: rgba(99,102,241,.1); border: 1px solid rgba(99,102,241,.2);
  border-radius: 6px; padding: .22rem .7rem; font-size: .7rem; font-weight: 600;
  color: #818CF8; text-decoration: none; white-space: nowrap; transition: all .2s;
  flex-shrink: 0;
}
.back-btn:hover { background: rgba(99,102,241,.2); }

.container, .main, .page { padding: 1rem 1.2rem 2rem; max-width: 1400px; margin: 0 auto; }


.page-content{padding:.75rem 1.2rem 2.5rem;max-width:1200px;margin:0 auto;}
.admin-hero{display:flex;align-items:center;gap:.8rem;margin-bottom:1.3rem;padding:.9rem 1.2rem;
  background:rgba(245,158,11,.06);border:1px solid rgba(245,158,11,.13);border-radius:12px;}
.admin-hero-icon{font-size:1.6rem;}
.admin-hero-title{font-size:1rem;font-weight:800;color:#F1F5FF;}
.admin-hero-sub{font-size:.74rem;color:#475569;margin-top:.1rem;}
.tabs{display:flex;gap:.4rem;margin-bottom:1.2rem;flex-wrap:wrap;}
.tab{padding:.42rem 1rem;border-radius:8px;border:1px solid rgba(99,102,241,.18);background:transparent;
  color:#475569;cursor:pointer;font-size:.78rem;font-weight:600;font-family:'Inter',sans-serif;transition:all .2s;}
.tab:hover{background:rgba(99,102,241,.06);color:#818CF8;}
.tab.active{background:rgba(99,102,241,.14);border-color:rgba(99,102,241,.35);color:#818CF8;}
.panel{display:none;}.panel.active{display:block;}
.stats-grid{display:grid;grid-template-columns:repeat(auto-fit,minmax(140px,1fr));gap:.85rem;margin-bottom:1.2rem;}
.stat-card{background:rgba(13,19,40,.8);border:1px solid rgba(99,102,241,.14);border-radius:11px;padding:.9rem;text-align:center;}
.stat-num{font-size:2rem;font-weight:900;color:#6366F1;line-height:1.1;}
.stat-num.green{color:#10B981;}.stat-num.amber{color:#F59E0B;}
.stat-lbl{font-size:.64rem;color:#475569;text-transform:uppercase;letter-spacing:.08em;margin-top:.25rem;font-weight:500;}
.card{background:rgba(13,19,40,.8);border:1px solid rgba(99,102,241,.12);border-radius:12px;padding:1.1rem;margin-bottom:.9rem;}
.card-header{display:flex;align-items:center;justify-content:space-between;margin-bottom:.85rem;flex-wrap:wrap;gap:.4rem;}
.card-title{font-size:.84rem;font-weight:700;color:#818CF8;}
.btn-action{padding:.38rem .8rem;border-radius:7px;border:none;cursor:pointer;font-size:.72rem;font-weight:600;font-family:'Inter',sans-serif;transition:.2s;text-decoration:none;display:inline-flex;align-items:center;gap:.3rem;}
.btn-dl{background:rgba(16,185,129,.1);color:#10B981;border:1px solid rgba(16,185,129,.2);}
.btn-dl:hover{background:rgba(16,185,129,.18);}
.btn-refresh{background:rgba(99,102,241,.1);color:#818CF8;border:1px solid rgba(99,102,241,.2);}
.btn-refresh:hover{background:rgba(99,102,241,.18);}
.tbl-wrap{overflow-x:auto;border-radius:8px;border:1px solid rgba(99,102,241,.1);}
table{width:100%;border-collapse:collapse;font-size:.76rem;}
thead{background:rgba(99,102,241,.1);}
thead th{padding:.55rem .8rem;text-align:left;font-weight:700;color:#818CF8;white-space:nowrap;border-bottom:1px solid rgba(99,102,241,.18);font-size:.72rem;letter-spacing:.02em;}
tbody tr{border-bottom:1px solid rgba(255,255,255,.04);transition:background .15s;}
tbody tr:hover{background:rgba(99,102,241,.04);}
tbody td{padding:.48rem .8rem;color:#C8D8F8;}
.role-admin{color:#F59E0B;font-weight:700;font-size:.7rem;background:rgba(245,158,11,.1);padding:.12rem .45rem;border-radius:5px;border:1px solid rgba(245,158,11,.2);}
.role-user{color:#60A5FA;font-weight:600;font-size:.7rem;background:rgba(59,130,246,.1);padding:.12rem .45rem;border-radius:5px;border:1px solid rgba(59,130,246,.2);}
.action-chip{display:inline-block;padding:.12rem .5rem;border-radius:5px;font-size:.68rem;font-weight:600;background:rgba(99,102,241,.12);color:#818CF8;border:1px solid rgba(99,102,241,.2);}
.chip-login{background:rgba(16,185,129,.1);color:#10B981;border-color:rgba(16,185,129,.2);}
.chip-logout{background:rgba(107,114,128,.1);color:#6B7280;border-color:rgba(107,114,128,.18);}
.chip-open{background:rgba(59,130,246,.1);color:#60A5FA;border-color:rgba(59,130,246,.2);}
.no-data{text-align:center;padding:2rem;color:#334155;}
.online-dot{width:7px;height:7px;border-radius:50%;background:#10B981;display:inline-block;box-shadow:0 0 5px #10B981;}
.offline-dot{width:7px;height:7px;border-radius:50%;background:#475569;display:inline-block;}
.analytics-grid{display:grid;grid-template-columns:1fr 1fr;gap:.85rem;margin-bottom:.85rem;}
@media(max-width:700px){.analytics-grid{grid-template-columns:1fr;}}
.mini-bar-row{display:flex;align-items:center;gap:.5rem;margin-bottom:.35rem;font-size:.72rem;}
.mini-bar-label{width:100px;white-space:nowrap;overflow:hidden;text-overflow:ellipsis;color:#C8D8F8;flex-shrink:0;}
.mini-bar-track{flex:1;height:5px;background:rgba(99,102,241,.12);border-radius:3px;overflow:hidden;}
.mini-bar-fill{height:100%;background:linear-gradient(90deg,#4338CA,#6366F1);border-radius:3px;transition:width .6s;}
.mini-bar-val{width:28px;text-align:right;color:#818CF8;font-weight:600;flex-shrink:0;}
.spinner{animation:spin 1s linear infinite;display:inline-block;}
@keyframes spin{to{transform:rotate(360deg);}}




/* Action cards */
.action-card { background: rgba(10,16,42,.85) !important; border-color: rgba(245,158,11,.5) !important; }
.rr-panel { background: rgba(6,12,38,.88) !important; border-color: rgba(99,102,241,.35) !important; }

/* Stats / text helpers */
.sec-hint, .upload-hint { color: #9BB8E0 !important; }
.chip-ok { color: #4ADE80 !important; }
.chip-info { color: #818CF8 !important; }
.chip-wait { color: #F59E0B !important; }

/* ── Developed by Karthik — fixed footer ── */
.dev-credit {
  position: fixed; bottom: 0; left: 0; right: 0; z-index: 9999;
  text-align: center; padding: .3rem 1rem;
  background: rgba(5, 8, 22, 0.75); backdrop-filter: blur(8px);
  border-top: 1px solid rgba(99,102,241,.2);
  font-size: 11px; font-weight: 600; color: rgba(160,180,220,.75);
  letter-spacing: .06em; font-family: 'Inter', system-ui, sans-serif;
  pointer-events: none; user-select: none;
}

/* ═══ GLOBAL BACKGROUND & TRANSPARENCY — v3 ═══════════════════════ */
html, body {
  background-image: url('/static/bg.jpg') !important;
  background-size: cover !important;
  background-position: center center !important;
  background-attachment: fixed !important;
  background-repeat: no-repeat !important;
  background-color: #08101E !important;
}
/* Single very-light overlay — image stays visible */
body::before {
  content: '' !important;
  position: fixed !important;
  inset: 0 !important;
  z-index: 0 !important;
  background: rgba(3, 7, 18, 0.45) !important;
  pointer-events: none !important;
}
body > * { position: relative; z-index: 1; }

/* ── Top-bar: fully transparent glass, no black border ── */
.top-bar {
  background: rgba(6, 10, 28, 0.55) !important;
  backdrop-filter: blur(18px) !important;
  -webkit-backdrop-filter: blur(18px) !important;
  border-bottom: 1px solid rgba(255,255,255,0.08) !important;
  position: sticky !important;
  top: 0 !important;
  z-index: 200 !important;
}

/* ── Cards / Panels — glass, no solid black fill ── */
.card, .section-card, .panel, .sec-card, .stats-box, .action-card {
  background: rgba(8, 14, 38, 0.70) !important;
  border: 1px solid rgba(255,255,255,0.10) !important;
  backdrop-filter: blur(16px) !important;
  -webkit-backdrop-filter: blur(16px) !important;
  box-shadow: 0 4px 32px rgba(0,0,0,0.35) !important;
}
.panel-head, .card-header {
  background: rgba(10, 18, 52, 0.72) !important;
  border-bottom: 1px solid rgba(255,255,255,0.07) !important;
}
.panel-body { background: rgba(5, 10, 30, 0.60) !important; }

/* ── Upload zones ── */
.upload-zone, .dz {
  background: rgba(6, 12, 34, 0.55) !important;
  border: 2px dashed rgba(99,102,241,0.55) !important;
}

/* ── Typography: all white/light ── */
body, h1, h2, h3, h4, p, span, div, td, th, label, a {
  color: #EEF3FF !important;
}
.hero-title, .brand-title, .card-title, .admin-hero-title {
  color: #FFFFFF !important;
  text-shadow: 0 2px 16px rgba(0,0,0,0.7) !important;
  font-weight: 800 !important;
}
.hero-sub, .brand-sub, .admin-hero-sub, .sec-hint {
  color: rgba(210, 225, 255, 0.80) !important;
}
.top-mini-brand, .top-page-label, .top-brand .brand-text {
  color: #FFFFFF !important;
  font-weight: 700 !important;
}
.back-btn {
  color: #C5D5FF !important;
  background: rgba(99,102,241,0.18) !important;
  border: 1px solid rgba(99,102,241,0.35) !important;
}
.back-btn:hover { background: rgba(99,102,241,0.32) !important; }

/* ── Inputs: legible on transparent backgrounds ── */
input[type=text], input[type=number], input[type=date],
input[type=password], select, textarea,
.input-num, .action-select, .rr-input, .rr-select, .date-input {
  background: rgba(5, 9, 28, 0.80) !important;
  border: 1.5px solid rgba(99,102,241,0.40) !important;
  color: #EEF3FF !important;
  font-size: 14px !important;
}
input::placeholder, textarea::placeholder {
  color: rgba(180, 200, 240, 0.50) !important;
}

/* ── Muted / secondary text ── */
.muted, .sec-label, [style*="color:#475569"],
[style*="color:#64748B"], [style*="color:#374167"] {
  color: rgba(190, 210, 255, 0.70) !important;
}

/* ── Action + RR cards ── */
.action-card {
  background: rgba(12, 18, 48, 0.78) !important;
  border-color: rgba(245,158,11,0.55) !important;
}
.rr-panel {
  background: rgba(8, 14, 42, 0.80) !important;
  border-color: rgba(99,102,241,0.40) !important;
}

/* ── Global font sizes ── */
body { font-size: 14px !important; }
h1 { font-size: 26px !important; }
h2 { font-size: 22px !important; }
h3, .hero-title { font-size: 22px !important; }
h4 { font-size: 17px !important; }

/* ── Dev-credit fixed footer ── */
.apsg-footer {
  position: fixed !important; bottom: 0 !important;
  left: 0 !important; right: 0 !important; z-index: 9999 !important;
  text-align: center !important; padding: .28rem 1rem !important;
  background: rgba(4, 7, 20, 0.70) !important;
  backdrop-filter: blur(8px) !important;
  border-top: 1px solid rgba(255,255,255,0.07) !important;
  font-size: 11px !important; font-weight: 600 !important;
  color: rgba(200, 220, 255, 0.70) !important;
  letter-spacing: .06em !important; pointer-events: none !important;
  user-select: none !important;
}
</style>
</head>
<body>
<div class="top-bar">
  <span class="top-mini-brand">APSG · Admin</span>
  <div class="top-sep"></div>
  <span class="top-page-label">⚙ Admin Panel</span>
  <div class="top-spacer"></div>
  <span style="font-size:.65rem;color:#6366F1;font-weight:700;white-space:nowrap;margin-right:.5rem;letter-spacing:.01em;">✦ Karthi</span>
  <a href="/" class="back-btn">← Dashboard</a>
</div>
<div class="page-content">
  <div class="admin-hero">
    <div class="admin-hero-icon">⚙</div>
    <div>
      <div class="admin-hero-title">Admin Panel</div>
      <div class="admin-hero-sub">Manage users, view activity logs, and analyse system usage.</div>
    </div>
  </div>
  <div class="tabs">
    <button class="tab active" onclick="showTab('analytics')">📊 Analytics Dashboard</button>
    <button class="tab" onclick="showTab('users')">👥 Users</button>
    <button class="tab" onclick="showTab('activity')">📋 Activity Log</button>
  </div>
  <!-- ── Analytics Panel ── -->
  <div class="panel active" id="tab-analytics">
    <div class="stats-grid" id="analyticsStats">
      <div class="stat-card"><div class="stat-num spinner">⟳</div><div class="stat-lbl">Loading…</div></div>
    </div>
    <div class="analytics-grid">
      <div class="card">
        <div class="card-header"><div class="card-title">👤 User Login Activity</div></div>
        <div id="loginBarChart" class="mini-bar-wrap"><div class="no-data"><span class="spinner">⟳</span></div></div>
      </div>
      <div class="card">
        <div class="card-header"><div class="card-title">🛠 Module Usage</div></div>
        <div id="moduleBarChart" class="mini-bar-wrap"><div class="no-data"><span class="spinner">⟳</span></div></div>
      </div>
    </div>
    <div class="card">
      <div class="card-header">
        <div class="card-title">👥 User Activity Summary</div>
        <div style="display:flex;gap:.5rem">
          <button class="btn-action btn-refresh" onclick="loadAnalytics()">↺ Refresh</button>
          <a href="/api/admin/analytics/download" class="btn-action btn-dl">⬇ Download Excel</a>
        </div>
      </div>
      <div class="tbl-wrap">
        <table>
          <thead><tr><th>#</th><th>Username</th><th>Full Name</th><th>Role</th><th>Total Actions</th><th>First Seen</th><th>Last Active</th><th>Status</th></tr></thead>
          <tbody id="analyticsTbl"><tr><td colspan="8" class="no-data"><span class="spinner">⟳</span> Loading…</td></tr></tbody>
        </table>
      </div>
    </div>
  </div>
  <!-- ── Users Panel ── -->
  <div class="panel" id="tab-users">
    <div class="card">
      <div class="card-header">
        <div class="card-title">Registered Users</div>
        <button class="btn-action btn-refresh" onclick="loadUsers()">↺ Refresh</button>
      </div>
      <div class="tbl-wrap">
        <table>
          <thead><tr><th>#</th><th>Username</th><th>Password</th><th>Full Name</th><th>Email</th><th>Verified</th><th>Role</th><th>Registered</th><th>Actions</th></tr></thead>
          <tbody id="usersTbl"><tr><td colspan="9" class="no-data"><span class="spinner">⟳</span></td></tr></tbody>
        </table>
      </div>
    </div>
  </div>
  <!-- ── Activity Panel ── -->
  <div class="panel" id="tab-activity">
    <div class="card">
      <div class="card-header">
        <div class="card-title">Activity Log (latest 200)</div>
        <div style="display:flex;gap:.4rem;margin-left:auto;">
          <button class="btn-action btn-refresh" onclick="loadActivity('all')" id="actBtnAll" style="background:rgba(99,102,241,.25)">All</button>
          <button class="btn-action btn-refresh" onclick="loadActivity('login')" id="actBtnLogin">🔐 Logins Only</button>
        </div>
        <div style="display:flex;gap:.5rem">
          <button class="btn-action btn-refresh" onclick="loadActivity()">↺ Refresh</button>
          <a href="/api/admin/activity/download" class="btn-action btn-dl">⬇ Download</a>
        </div>
      </div>
      <div class="tbl-wrap">
        <table>
          <thead><tr><th>Timestamp</th><th>Username</th><th>Action</th><th>Detail</th><th>IP</th></tr></thead>
          <tbody id="activityTbl"><tr><td colspan="5" class="no-data">Click tab to load</td></tr></tbody>
        </table>
      </div>
    </div>
  </div>
</div>
<script>
function showTab(t){
  document.querySelectorAll('.tab').forEach((el,i)=>{
    el.classList.toggle('active',['analytics','users','activity'][i]===t);
  });
  document.querySelectorAll('.panel').forEach(p=>p.classList.remove('active'));
  document.getElementById(`tab-${t}`).classList.add('active');
  if(t==='analytics')loadAnalytics();
  else if(t==='users')loadUsers();
  else loadActivity();
}

async function loadAnalytics(){
  try{
    const res=await fetch('/api/admin/analytics');
    const d=await res.json();
    // Stats cards
    const now=new Date();
    const recentCutoff=new Date(now-24*60*60*1000);
    document.getElementById('analyticsStats').innerHTML=`
      <div class="stat-card"><div class="stat-num">${d.total_users}</div><div class="stat-lbl">Total Users</div></div>
      <div class="stat-card"><div class="stat-num green">${d.total_logins}</div><div class="stat-lbl">Total Logins</div></div>
      <div class="stat-card"><div class="stat-num amber">${d.total_actions}</div><div class="stat-lbl">Total Actions</div></div>
      <div class="stat-card"><div class="stat-num">${d.users.filter(u=>u.role==='admin').length}</div><div class="stat-lbl">Admin Users</div></div>
    `;
    // Login bar chart
    if(d.logins&&d.logins.length>0){
      const maxL=Math.max(...d.logins.map(l=>l.cnt));
      document.getElementById('loginBarChart').innerHTML=d.logins.slice(0,10).map(l=>`
        <div class="mini-bar-row">
          <div class="mini-bar-label" title="${esc(l.username)}">${esc(l.username)}</div>
          <div class="mini-bar-track"><div class="mini-bar-fill" style="width:${Math.round(l.cnt/maxL*100)}%"></div></div>
          <div class="mini-bar-val">${l.cnt}</div>
        </div>`).join('');
    }else{document.getElementById('loginBarChart').innerHTML='<div class="no-data">No login data yet</div>';}
    // Module usage chart
    const moduleMap={'OPEN_APP':'Open App','PPT_UPLOAD':'PPT Upload','PPT_GENERATE':'PPT Generate','EXCEL_REJECTION':'Excel Rejection','PHOTO_MERGE_START':'Photo Merge','DAILY_REPORT':'Daily Report','PPT_EXCEL':'PPT Excel','PPT_ZIP':'PPT ZIP'};
    const modCounts={};
    // Count from activity summary
    for(const u of d.activity_summary){
      if(u.actions_used){u.actions_used.split(',').forEach(a=>{
        const key=moduleMap[a.trim()]||a.trim();
        if(a.trim()!=='LOGIN'&&a.trim()!=='LOGOUT'&&a.trim()!=='REGISTER')
          modCounts[key]=(modCounts[key]||0)+1;
      });}
    }
    const modEntries=Object.entries(modCounts).sort((a,b)=>b[1]-a[1]).slice(0,8);
    if(modEntries.length>0){
      const maxM=Math.max(...modEntries.map(e=>e[1]));
      document.getElementById('moduleBarChart').innerHTML=modEntries.map(([k,v])=>`
        <div class="mini-bar-row">
          <div class="mini-bar-label" title="${esc(k)}">${esc(k)}</div>
          <div class="mini-bar-track"><div class="mini-bar-fill" style="width:${Math.round(v/maxM*100)}%;background:linear-gradient(90deg,#065F38,#10B981)"></div></div>
          <div class="mini-bar-val">${v}</div>
        </div>`).join('');
    }else{document.getElementById('moduleBarChart').innerHTML='<div class="no-data">No module usage data yet</div>';}
    // Activity summary table
    const actMap={};
    for(const a of d.activity_summary){actMap[a.username]=a;}
    document.getElementById('analyticsTbl').innerHTML=d.users.length?d.users.map((u,i)=>{
      const act=actMap[u.username]||{};
      const lastTs=act.last_seen?new Date(act.last_seen):null;
      const isRecent=lastTs&&(new Date()-lastTs)<24*60*60*1000;
      const status=isRecent?'<span class="online-dot"></span> Active (24h)':'<span class="offline-dot"></span> Inactive';
      return`<tr>
        <td style="color:#475569">${i+1}</td>
        <td><strong>${esc(u.username)}</strong></td>
        <td>${esc(u.name)}</td>
        <td><span class="${u.role==='admin'?'role-admin':'role-user'}">${u.role}</span></td>
        <td style="color:#818CF8;font-weight:600">${act.total_actions||0}</td>
        <td style="font-size:.72rem;color:#475569">${act.first_seen?act.first_seen.slice(0,16):'—'}</td>
        <td style="font-size:.72rem;color:#475569">${act.last_seen?act.last_seen.slice(0,16):'—'}</td>
        <td style="font-size:.72rem">${status}</td>
      </tr>`;
    }).join(''):'<tr><td colspan="8" class="no-data">No users found</td></tr>';
  }catch(e){console.error(e);}
}

async function loadUsers(){
  const tb=document.getElementById('usersTbl');
  tb.innerHTML='<tr><td colspan="9" class="no-data"><span class="spinner">⟳</span></td></tr>';
  try{
    const res=await fetch('/api/admin/users');const d=await res.json();
    tb.innerHTML=d.length?d.map((u,i)=>{
      const verified = u.email_verified
        ? '<span style="color:#10B981;font-weight:700;font-size:.78rem;">✓ Verified</span>'
        : '<span style="color:#F59E0B;font-weight:600;font-size:.78rem;">⏳ Unverified</span>';
      const isAdmin = u.role==='admin';
      return `<tr>
        <td style="color:#475569;font-size:.72rem">${i+1}</td>
        <td><strong>${esc(u.username)}</strong></td>
        <td>
          <span style="display:flex;align-items:center;gap:.4rem;">
            <span id="pw-${esc(u.username)}" style="font-family:monospace;font-size:.8rem;color:#C8D8F8;letter-spacing:.08em;">••••••••</span>
            <button onclick="togglePwView('${esc(u.username)}')" title="Show/Hide"
              style="background:none;border:none;cursor:pointer;font-size:.85rem;padding:.1rem .2rem;color:#818CF8;">
              <span id="pw-eye-${esc(u.username)}">👁</span>
            </button>
          </span>
        </td>
        <td>${esc(u.name)}</td>
        <td style="font-size:.78rem;color:#A5B4FC">${esc(u.email||'—')}</td>
        <td>${verified}</td>
        <td><span class="${isAdmin?'role-admin':'role-user'}">${u.role}</span></td>
        <td style="font-size:.7rem;color:#475569">${(u.created_at||'—').slice(0,16)}</td>
        <td>
          <div style="display:flex;gap:.35rem;flex-wrap:wrap">
            <button class="btn-action" style="background:rgba(99,102,241,.15);color:#818CF8;border-color:rgba(99,102,241,.3);"
              onclick="openMasterResetPw('${esc(u.username)}','${esc(u.name)}')" title="Reset Password">
              🔑 Reset PW
            </button>
            ${!isAdmin?`<button class="btn-action btn-del"
              onclick="confirmDeleteUser('${esc(u.username)}','${esc(u.name)}')" title="Delete user">
              🗑 Delete
            </button>`:'<span style="font-size:.7rem;color:#475569;">Protected</span>'}
          </div>
        </td>
      </tr>`;
    }).join(''):'<tr><td colspan="9" class="no-data">No users found</td></tr>';
  }catch(e){tb.innerHTML=`<tr><td colspan="9" class="no-data">Error: ${e.message}</td></tr>`;}
}

// ── Master Password Reset ─────────────────────────────────────────────────────
let _masterResetUsername = '';
let _masterToken = '';

function openMasterResetPw(username, name) {
  _masterResetUsername = username;
  _masterToken = '';
  document.getElementById('masterResetUser').textContent  = `${name} (@${username})`;
  document.getElementById('masterResetUser2').textContent = `${name} (@${username})`;
  document.getElementById('masterPwInput').value  = '';
  document.getElementById('newPwInput').value     = '';
  document.getElementById('newPwConfirm').value   = '';
  document.getElementById('masterStep1Result').innerHTML = '';
  document.getElementById('masterStep2Result').innerHTML = '';
  document.getElementById('masterStep1').style.display = 'block';
  document.getElementById('masterStep2').style.display = 'none';
  document.getElementById('masterResetModal').style.display = 'flex';
  setTimeout(() => document.getElementById('masterPwInput').focus(), 80);
}
function closeMasterReset() {
  document.getElementById('masterResetModal').style.display = 'none';
  _masterToken = '';
}

async function submitMasterVerify() {
  const pw = document.getElementById('masterPwInput').value;
  const res_el = document.getElementById('masterStep1Result');
  res_el.innerHTML = '';
  if (!pw) { res_el.innerHTML = '<span style="color:#F87171">Please enter your admin password.</span>'; return; }
  try {
    const r = await fetch('/api/admin/verify_master', {method:'POST',
      headers:{'Content-Type':'application/json'}, body: JSON.stringify({password: pw})});
    const d = await r.json();
    if (d.ok) {
      _masterToken = d.token;
      document.getElementById('masterStep1').style.display = 'none';
      document.getElementById('masterStep2').style.display = 'block';
      setTimeout(() => document.getElementById('newPwInput').focus(), 80);
    } else {
      res_el.innerHTML = `<span style="color:#F87171">❌ ${esc(d.error||'Incorrect password')}</span>`;
    }
  } catch(e) { res_el.innerHTML = `<span style="color:#F87171">Connection error</span>`; }
}

async function submitNewPassword() {
  const newPw   = document.getElementById('newPwInput').value.trim();
  const confirm = document.getElementById('newPwConfirm').value.trim();
  const res_el  = document.getElementById('masterStep2Result');
  res_el.innerHTML = '';
  if (!newPw || newPw.length < 6) { res_el.innerHTML = '<span style="color:#F87171">Password must be at least 6 characters.</span>'; return; }
  if (newPw !== confirm) { res_el.innerHTML = '<span style="color:#F87171">Passwords do not match.</span>'; return; }
  try {
    const r = await fetch(`/api/admin/user/${encodeURIComponent(_masterResetUsername)}/reset_password`, {
      method:'POST', headers:{'Content-Type':'application/json'},
      body: JSON.stringify({master_token: _masterToken, new_password: newPw, confirm})});
    const d = await r.json();
    if (d.ok) {
      res_el.innerHTML = `<span style="color:#4ADE80">✅ Password reset for @${esc(_masterResetUsername)}</span>`;
      setTimeout(closeMasterReset, 1800);
    } else {
      res_el.innerHTML = `<span style="color:#F87171">❌ ${esc(d.error||'Failed')}</span>`;
    }
  } catch(e) { res_el.innerHTML = `<span style="color:#F87171">Connection error</span>`; }
}

// ── Delete User ───────────────────────────────────────────────────────────────
let _deleteUsername = '';

function confirmDeleteUser(username, name) {
  _deleteUsername = username;
  document.getElementById('deleteModalUser').textContent = `${name} (@${username})`;
  document.getElementById('deleteModalResult').innerHTML = '';
  document.getElementById('confirmDeleteBtn').disabled = false;
  document.getElementById('deleteModal').style.display = 'flex';
}
function closeDeleteModal() {
  document.getElementById('deleteModal').style.display = 'none';
  _deleteUsername = '';
}

async function executeDelete() {
  if (!_deleteUsername) return;
  const btn = document.getElementById('confirmDeleteBtn');
  btn.disabled = true; btn.textContent = '⏳ Deleting…';
  try {
    const r = await fetch(`/api/admin/user/${encodeURIComponent(_deleteUsername)}/delete`, {method:'DELETE'});
    const ct = r.headers.get('Content-Type')||'';
    if (!ct.includes('application/json')) {
      document.getElementById('deleteModalResult').innerHTML = '<span style="color:#F87171">Session expired. Please reload.</span>';
      btn.disabled=false;btn.textContent='🗑 Yes, Delete Account';return;
    }
    const d = await r.json();
    if (d.ok) {
      document.getElementById('deleteModalResult').innerHTML = `<span style="color:#4ADE80">✅ Account deleted.</span>`;
      setTimeout(() => { closeDeleteModal(); loadUsers(); }, 1200);
    } else {
      document.getElementById('deleteModalResult').innerHTML = `<span style="color:#F87171">❌ ${esc(d.error||'Failed')}</span>`;
      btn.disabled=false;btn.textContent='🗑 Yes, Delete Account';
    }
  } catch(e) {
    document.getElementById('deleteModalResult').innerHTML = `<span style="color:#F87171">Connection error</span>`;
    btn.disabled=false;btn.textContent='🗑 Yes, Delete Account';
  }
}

const _pwCache = {};
async function togglePwView(username){
  const el=document.getElementById(`pw-${username}`);
  const eyeEl=document.getElementById(`pw-eye-${username}`);
  if(!el)return;
  if(el.dataset.shown==='1'){
    el.textContent='••••••••'; el.dataset.shown='0';
    eyeEl.textContent='👁';
    return;
  }
  // Already fetched?
  if(_pwCache[username]){
    el.textContent=_pwCache[username]; el.dataset.shown='1';
    eyeEl.textContent='🙈';
    return;
  }
  el.textContent='⏳';
  try{
    const res=await fetch(`/api/admin/user/${encodeURIComponent(username)}/view_password`);
    const ct=res.headers.get('Content-Type')||'';
    if(!ct.includes('application/json')){el.textContent='Login required';return;}
    const d=await res.json();
    if(d.ok){
      _pwCache[username]=d.password;
      el.textContent=d.password; el.dataset.shown='1';
      eyeEl.textContent='🙈';
      if(d.source==='generated'){
        showModalResult(`⚠️ A new temp password was generated for @${username}: <strong>${d.password}</strong><br>Share it with the user.`);
      }
    }else{
      el.textContent='Error'; el.dataset.shown='0';
    }
  }catch(e){el.textContent='Error';}
}

function openChangePw(username, name){
  document.getElementById('pwModalUser').textContent=`${name} (@${username})`;
  document.getElementById('pwModalUsername').value=username;
  document.getElementById('pwModalNewPw').value='';
  document.getElementById('pwModalConfirm').value='';
  document.getElementById('pwModalResult').innerHTML='';
  document.getElementById('pwShowToggle').checked=false;
  document.getElementById('pwModalNewPw').type='password';
  document.getElementById('pwModalConfirm').type='password';
  document.getElementById('pwModal').style.display='flex';
}
function closePwModal(){document.getElementById('pwModal').style.display='none';}

async function submitChangePw(){
  const username=document.getElementById('pwModalUsername').value;
  const newPw=document.getElementById('pwModalNewPw').value.trim();
  const confirm=document.getElementById('pwModalConfirm').value.trim();
  if(!newPw||newPw.length<6){showModalResult('Password must be at least 6 characters.','error');return;}
  if(newPw!==confirm){showModalResult('Passwords do not match.','error');return;}
  try{
    const res=await fetch(`/api/admin/user/${encodeURIComponent(username)}/password`,{
      method:'POST',headers:{'Content-Type':'application/json'},
      body:JSON.stringify({password:newPw})});
    // Check content-type before parsing — server may return HTML on session expiry
    const ct=res.headers.get('Content-Type')||'';
    if(!ct.includes('application/json')){
      showModalResult('Session expired. Please <a href="/login" style="color:#818CF8">log in again</a>.','error');
      return;
    }
    const d=await res.json();
    if(d.ok){showModalResult(`✅ Password updated for @${username}`,'ok');setTimeout(closePwModal,1800);}
    else if(d.redirect){showModalResult('Session expired. <a href="/login" style="color:#818CF8">Log in again</a>.','error');}
    else{showModalResult(d.error||'Failed to update password.','error');}
  }catch(e){showModalResult('Connection error: '+e.message,'error');}
}

// ── Reset Password (generates temp password) ──────────────────────────────────
async function resetPw(username){
  if(!confirm(`Generate a new temporary password for @${username}?`))return;
  try{
    const res=await fetch(`/api/admin/user/${encodeURIComponent(username)}/view_password`);
    const d=await res.json();
    if(d.ok){
      document.getElementById('tempPwUser').textContent=`${d.name} (@${d.username})`;
      document.getElementById('tempPwValue').textContent=d.temp_password;
      document.getElementById('tempPwModal').style.display='flex';
    }else{alert(d.error||'Failed');}
  }catch(e){alert(e.message);}
}
function closeTempPwModal(){document.getElementById('tempPwModal').style.display='none';}
function copyTempPw(){
  const txt=document.getElementById('tempPwValue').textContent;
  navigator.clipboard.writeText(txt).then(()=>{
    const btn=document.getElementById('copyTempBtn');
    btn.textContent='✅ Copied!';setTimeout(()=>btn.textContent='📋 Copy',1500);
  });
}

// ── Delete User ────────────────────────────────────────────────────────────────
async function deleteUser(username,name){
  if(!confirm(`Delete user "${name}" (@${username})?\nThis action cannot be undone.`))return;
  try{
    const res=await fetch(`/api/admin/user/${encodeURIComponent(username)}/delete`,{method:'DELETE'});
    const d=await res.json();
    if(d.ok){loadUsers();loadAnalytics();}
    else{alert(d.error||'Delete failed');}
  }catch(e){alert(e.message);}
}

function showModalResult(msg,type){
  const el=document.getElementById('pwModalResult');
  el.textContent=msg;
  el.className='modal-result '+(type==='ok'?'result-ok':'result-err');
}
function togglePwShow(checked){
  const t=checked?'text':'password';
  document.getElementById('pwModalNewPw').type=t;
  document.getElementById('pwModalConfirm').type=t;
}

async function loadActivity(filter='all'){
  document.getElementById('actBtnAll').style.background=filter==='all'?'rgba(99,102,241,.25)':''; 
  const loginBtn=document.getElementById('actBtnLogin'); if(loginBtn) loginBtn.style.background=filter==='login'?'rgba(16,185,129,.25)':'';
  const tb=document.getElementById('activityTbl');
  tb.innerHTML='<tr><td colspan="5" class="no-data"><span class="spinner">⟳</span></td></tr>';
  try{
    const actUrl=filter==='login'?'/api/admin/activity?limit=200&action=LOGIN':'/api/admin/activity?limit=200';
    const res=await fetch(actUrl);const d=await res.json();
    const chipClass=(a)=>a==='LOGIN'?'chip-login':a==='LOGOUT'?'chip-logout':a.startsWith('OPEN')?'chip-open':'';
    tb.innerHTML=d.length?d.map(r=>`<tr>
      <td style="font-size:.72rem;color:#475569;white-space:nowrap">${r.ts||''}</td>
      <td><strong>${esc(r.username)}</strong></td>
      <td><span class="action-chip ${chipClass(r.action)}">${esc(r.action)}</span></td>
      <td style="font-size:.75rem;color:#6B7280">${esc(r.detail||'')}</td>
      <td style="font-size:.72rem;color:#475569">${r.ip||''}</td>
    </tr>`).join(''):'<tr><td colspan="5" class="no-data">No activity yet</td></tr>';
  }catch(e){tb.innerHTML=`<tr><td colspan="5" class="no-data">Error: ${e.message}</td></tr>`;}
}

function esc(s){return String(s||'').replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;').replace(/"/g,'&quot;');}
loadAnalytics();
</script>

<!-- ── Master Password Reset Modal (2-step: verify admin pw → set new pw) ── -->
<div class="modal-overlay" id="masterResetModal" onclick="if(event.target===this)closeMasterReset()">
  <div class="modal-card" style="max-width:400px">

    <!-- Step 1: Verify master password -->
    <div id="masterStep1">
      <div class="modal-title">🛡 Admin Authentication</div>
      <div class="modal-sub" id="masterResetUser"></div>
      <div style="font-size:.78rem;color:#94A3B8;margin-bottom:1rem;line-height:1.5;
        background:rgba(245,158,11,.07);border:1px solid rgba(245,158,11,.2);
        border-radius:8px;padding:.6rem .85rem;">
        ⚠️ Enter your <strong style="color:#F59E0B">admin password</strong> to authorise this password reset.
      </div>
      <div class="modal-field">
        <label>Your Admin Password</label>
        <input type="password" id="masterPwInput" placeholder="Enter your admin password">
      </div>
      <div class="modal-result" id="masterStep1Result"></div>
      <div class="modal-actions">
        <button class="modal-btn modal-btn-cancel" onclick="closeMasterReset()">✕ Cancel</button>
        <button class="modal-btn modal-btn-primary" onclick="submitMasterVerify()">🔓 Authenticate →</button>
      </div>
    </div>

    <!-- Step 2: Set new password (shown after auth) -->
    <div id="masterStep2" style="display:none">
      <div class="modal-title">🔑 Reset Password</div>
      <div class="modal-sub" id="masterResetUser2"></div>
      <div style="font-size:.75rem;color:#34D399;margin-bottom:1rem;
        background:rgba(16,185,129,.07);border:1px solid rgba(16,185,129,.2);
        border-radius:8px;padding:.5rem .85rem;">✅ Admin authenticated — set new password below</div>
      <div class="modal-field">
        <label>New Password</label>
        <input type="password" id="newPwInput" placeholder="Minimum 6 characters">
      </div>
      <div class="modal-field">
        <label>Confirm New Password</label>
        <input type="password" id="newPwConfirm" placeholder="Re-enter new password">
      </div>
      <div class="modal-result" id="masterStep2Result"></div>
      <div class="modal-actions">
        <button class="modal-btn modal-btn-cancel" onclick="closeMasterReset()">✕ Cancel</button>
        <button class="modal-btn modal-btn-primary" onclick="submitNewPassword()">✔ Save Password</button>
      </div>
    </div>

  </div>
</div>

<!-- ── Delete Confirmation Modal ── -->
<div class="modal-overlay" id="deleteModal" onclick="if(event.target===this)closeDeleteModal()">
  <div class="modal-card" style="max-width:380px;text-align:center">
    <div style="font-size:2.5rem;margin-bottom:.5rem;">🗑</div>
    <div class="modal-title" style="color:#F87171">Delete User Account</div>
    <div class="modal-sub" id="deleteModalUser"></div>
    <div style="font-size:.82rem;color:#94A3B8;margin:1rem 0;line-height:1.55;
      background:rgba(248,113,113,.06);border:1px solid rgba(248,113,113,.18);
      border-radius:8px;padding:.7rem .9rem;">
      ⚠️ This will <strong style="color:#F87171">permanently delete</strong> this account and all associated data.
      This action <strong>cannot be undone</strong>.
    </div>
    <div class="modal-result" id="deleteModalResult"></div>
    <div class="modal-actions">
      <button class="modal-btn modal-btn-cancel" onclick="closeDeleteModal()">✕ Cancel</button>
      <button class="modal-btn" id="confirmDeleteBtn"
        style="background:rgba(248,113,113,.15);color:#F87171;border-color:rgba(248,113,113,.4);"
        onclick="executeDelete()">🗑 Yes, Delete Account</button>
    </div>
  </div>
</div>

<div class="apsg-footer">✦ Internal Reporting Platform — APSG Staging Ground &nbsp;·&nbsp; Developed by Karthik</div>
</body>
</html>"""


# ═══════════════════════════════════════════════════════════════════════════════
#  STARTUP
# ═══════════════════════════════════════════════════════════════════════════════

# Initialize DB (must be after all function definitions)
try:
    init_db()
    print("✓ Database initialized")
except Exception as _init_e:
    print(f"DB init warning: {_init_e}")

if __name__ == "__main__":
    port=int(os.environ.get("PORT",5000))
    app.run(host="0.0.0.0",port=port,debug=False)
